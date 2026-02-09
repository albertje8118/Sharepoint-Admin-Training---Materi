"""
Generate Course Introduction / Opening PPTX
Modern SharePoint Online for Administrators — 3-Day Course
Engaging design matching Modules 1-4 style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG     = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEC, 0xF9)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT    = RGBColor(0x24, 0x24, 0x24)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT  = RGBColor(0xD1, 0x34, 0x38)
YELLOW_WARN = RGBColor(0xF2, 0xC8, 0x11)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ═══════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════
def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6), icon="▸", font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon} {bullet}"; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.name = font_name; p.space_after = spacing
    return txBox

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)

def add_footer_bar(slide, slide_num, total):
    bar_top = SLIDE_HEIGHT - Inches(0.45)
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_WIDTH, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), bar_top + Inches(0.05), Inches(7), Inches(0.35),
                 "Modern SharePoint Online for Administrators  |  Course Introduction",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), bar_top + Inches(0.05),
                 Inches(1.2), Inches(0.35), f"{slide_num} / {total}",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)

def section_divider(slide, section_title, subtitle="", icon=""):
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_TEAL)
    if icon:
        add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                     icon, font_size=60, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                 section_title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

TOTAL_SLIDES = 23
slide_counter = [0]

def new_slide(notes=""):
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if notes: add_speaker_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome slide. Introduce yourself and the course. Let participants settle in. "
              "This is a modern replacement for legacy M55238B, fully aligned with 2026 Microsoft 365 administration.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Modern SharePoint Online\nfor Administrators",
             font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.0), Inches(3.3), Inches(3.3), Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(3.7), Inches(11), Inches(0.6),
             "3-Day Instructor-Led Training  ·  2026 Aligned",
             font_size=22, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(2), Inches(4.6), Inches(9), Inches(0.5),
             "Microsoft 365  ·  SharePoint Online  ·  Entra ID  ·  Purview  ·  PowerShell",
             font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(5.4), Inches(2.3), Inches(0.04), ACCENT_PURPLE)

add_text_box(s, Inches(2), Inches(5.8), Inches(9), Inches(0.8),
             "Replaces legacy M55238B  ·  MOC-style structure\nModules + Topics + Hands-on Labs",
             font_size=14, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Welcome & Housekeeping
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome participants. Cover logistics: schedule, breaks, Wi-Fi, restrooms, how to ask questions. "
              "Get a quick show of hands on experience level.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "👋  Welcome!", font_size=36, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Let's get settled before we begin",
             font_size=16, color=MID_GRAY)

items = [
    ("🕘", "Schedule", "3 full days: 09:00 – 17:00\nBreaks every 90 minutes\nLunch: 12:00 – 13:00", ACCENT_BLUE),
    ("📶", "Connectivity", "Wi-Fi credentials provided\nLab tenant credentials on your desk\nBring your own device or use provided", ACCENT_TEAL),
    ("❓", "Questions", "Ask anytime — we encourage it!\nParking lot for deeper topics\nAll materials provided digitally", ACCENT_PURPLE),
    ("🎯", "Expectations", "Hands-on: ~40% theory, ~60% labs\nReal-world scenarios throughout\nNo prior SharePoint admin required", ORANGE),
]
for i, (icon, title, desc, color) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 2.5)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(2.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(2.2), color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(1.0), y + Inches(0.15), Inches(4.5), Inches(0.5),
                 title, font_size=20, bold=True, color=color)
    add_text_box(s, x + Inches(1.0), y + Inches(0.7), Inches(4.5), Inches(1.3),
                 desc, font_size=14, color=DARK_TEXT)

add_footer_bar(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Meet Your Trainer
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Introduce yourself. Share your background, experience, and what makes you excited about this course.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "👨\u200d🏫  Meet Your Trainer", font_size=36, bold=True, color=DARK_BG)

# Photo placeholder (left side)
photo_box = add_rounded_rect(s, Inches(1.2), Inches(1.8), Inches(3.8), Inches(4.2), WHITE)
photo_box.shadow.inherit = False
add_shape_rect(s, Inches(1.2), Inches(1.8), Inches(3.8), Inches(0.08), ACCENT_PURPLE)
# Inner photo area
photo_area = add_rounded_rect(s, Inches(1.7), Inches(2.3), Inches(2.8), Inches(2.8), LIGHT_GRAY)
add_text_box(s, Inches(1.7), Inches(3.2), Inches(2.8), Inches(1.0),
             "📷\nInsert Photo Here", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1.5), Inches(5.3), Inches(3.2), Inches(0.5),
             "[Replace with trainer photo]", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Info fields (right side)
fields = [
    ("👤", "Name", "[Your Full Name]", ACCENT_BLUE),
    ("💼", "Title", "[Your Job Title / Role]", ACCENT_TEAL),
    ("🏢", "Company", "[Your Company / Organization]", ACCENT_PURPLE),
    ("📧", "Contact", "[Email or LinkedIn — optional]", ORANGE),
]
for i, (icon, label, placeholder, color) in enumerate(fields):
    y = Inches(1.8 + i * 1.15)
    card = add_rounded_rect(s, Inches(5.8), y, Inches(6.5), Inches(0.95), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, Inches(5.8), y, Inches(0.08), Inches(0.95), color)
    add_text_box(s, Inches(6.1), y + Inches(0.05), Inches(0.5), Inches(0.45),
                 icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(6.7), y + Inches(0.05), Inches(2.0), Inches(0.4),
                 label, font_size=14, bold=True, color=color)
    add_text_box(s, Inches(6.7), y + Inches(0.45), Inches(5.3), Inches(0.4),
                 placeholder, font_size=18, bold=True, color=DARK_TEXT)

# Bio section
bio_box = add_rounded_rect(s, Inches(5.8), Inches(6.4), Inches(6.5), Inches(0.65), LIGHT_BLUE)
add_text_box(s, Inches(6.1), Inches(6.45), Inches(5.9), Inches(0.55),
             "💬  [Add a brief bio: experience, certifications, areas of expertise, fun fact]",
             font_size=13, color=DARK_TEXT)

add_footer_bar(s, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — About This Course
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Position the course: modern replacement for M55238B, covers SPO/M365/Entra/Purview/PowerShell. "
              "Emphasize it is 2026-aligned with current admin center UIs.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "About This Course", font_size=36, bold=True, color=DARK_BG)

# Info cards
info = [
    ("📅", "Duration", "3 Days"),
    ("📊", "Level", "Intermediate"),
    ("🎓", "Style", "MOC-Format"),
    ("🧪", "Labs", "12 Hands-on"),
]
for i, (icon, label, value) in enumerate(info):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(s, x, Inches(1.3), Inches(2.8), Inches(1.6), WHITE)
    card.shadow.inherit = False
    add_text_box(s, x + Inches(0.3), Inches(1.4), Inches(2.2), Inches(0.5),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(1.9), Inches(2.2), Inches(0.4),
                 label, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.3), Inches(2.2), Inches(0.5),
                 value, font_size=24, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Audience box
add_rounded_rect(s, Inches(0.8), Inches(3.2), Inches(5.8), Inches(2.4), WHITE).shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(3.2), Inches(5.8), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(1.0), Inches(3.4), Inches(5.4), Inches(0.5),
             "👥  Target Audience", font_size=20, bold=True, color=ACCENT_TEAL)
audience = [
    "SharePoint Online Administrators",
    "Microsoft 365 Administrators",
    "IT Professionals managing collaboration",
    "Helpdesk leads supporting SharePoint",
]
add_bullet_frame(s, Inches(1.0), Inches(3.95), Inches(5.0), Inches(1.5),
                 audience, font_size=14, color=DARK_TEXT, icon="▸")

# Prerequisites box
add_rounded_rect(s, Inches(6.9), Inches(3.2), Inches(5.8), Inches(2.4), WHITE).shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(3.2), Inches(5.8), Inches(0.08), ORANGE)
add_text_box(s, Inches(7.1), Inches(3.4), Inches(5.4), Inches(0.5),
             "📋  Prerequisites", font_size=20, bold=True, color=ORANGE)
prereqs = [
    "Basic Microsoft 365 administration knowledge",
    "Familiarity with PowerShell fundamentals",
    "Understanding of identity concepts (users, groups)",
    "Web browser and modern device access",
]
add_bullet_frame(s, Inches(7.1), Inches(3.95), Inches(5.0), Inches(1.5),
                 prereqs, font_size=14, color=DARK_TEXT, icon="▸")

# Certification alignment at bottom
cert_box = add_rounded_rect(s, Inches(2.5), Inches(5.9), Inches(8.3), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(2.8), Inches(5.95), Inches(7.7), Inches(0.55),
             "🏅 Certification Alignment:  MS-102 (Microsoft 365 Administrator)  ·  SC-300 (Identity & Access Administrator)",
             font_size=14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Course Completion Outcomes
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("These are the 5 key outcomes. By end of 3 days, every learner should be able to do each of these.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🎯  What You'll Be Able to Do", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "By the end of this 3-day course, you will confidently:",
             font_size=16, color=MID_GRAY)

outcomes = [
    ("🏗️", "Administer SharePoint Online", "Using modern tools: admin centers, PowerShell, Graph API", ACCENT_BLUE),
    ("🔒", "Secure Collaboration", "Using Entra ID identities, Purview compliance, and sharing controls", ACCENT_TEAL),
    ("🗂️", "Design Architecture", "Scalable site structures, metadata, and information architecture", ACCENT_PURPLE),
    ("🔍", "Manage Discovery", "Configure Microsoft Search, managed properties, bookmarks & verticals", ORANGE),
    ("⚙️", "Automate Administration", "PowerShell scripting, bulk operations, monitoring & auditing", GREEN),
]
for i, (icon, title, desc, color) in enumerate(outcomes):
    y = Inches(1.7 + i * 1.05)
    card = add_rounded_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.85), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, Inches(0.8), y, Inches(0.08), Inches(0.85), color)
    add_text_box(s, Inches(1.1), y + Inches(0.1), Inches(0.6), Inches(0.55),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.7), y + Inches(0.1), Inches(4.0), Inches(0.5),
                 title, font_size=18, bold=True, color=color)
    add_text_box(s, Inches(5.8), y + Inches(0.1), Inches(6.5), Inches(0.6),
                 desc, font_size=14, color=DARK_TEXT)

add_footer_bar(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Section Divider: The 3-Day Journey
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: overview of the full 3-day schedule.")
section_divider(s, "Your 3-Day Journey", "12 Modules  ·  12 Labs  ·  From Foundations to Automation", "🗺️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — 3-Day Overview (Visual Map)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("High-level 3-day map. Each day has a clear theme. "
              "Day 1 = foundations, Day 2 = information architecture, Day 3 = governance.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "3-Day Course Overview", font_size=36, bold=True, color=DARK_BG)

days = [
    ("DAY 1", "Tenant Foundations\n& Site Management", "Modules 1 – 3", [
        "M365 & SharePoint intro",
        "Identity & external sharing",
        "Site collections & storage",
    ], ACCENT_BLUE, "🏗️"),
    ("DAY 2", "Information Architecture\nSearch & Customization", "Modules 4 – 7", [
        "Permissions & collaboration",
        "Metadata & Term Store",
        "Search & Microsoft Search",
        "Apps & customization",
    ], ACCENT_TEAL, "🗂️"),
    ("DAY 3", "Governance, Compliance\n& Automation", "Modules 8 – 12", [
        "Purview & compliance",
        "OneDrive administration",
        "PowerShell automation",
        "Monitoring & auditing",
        "Power Platform (optional)",
    ], ACCENT_PURPLE, "⚙️"),
]
for i, (day, theme, modules, topics, color, icon) in enumerate(days):
    x = Inches(0.6 + i * 4.2)
    card = add_rounded_rect(s, x, Inches(1.4), Inches(3.9), Inches(5.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.4), Inches(3.9), Inches(0.1), color)
    # Day label
    add_text_box(s, x + Inches(0.3), Inches(1.65), Inches(3.3), Inches(0.5),
                 f"{icon}  {day}", font_size=26, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Theme
    add_text_box(s, x + Inches(0.3), Inches(2.2), Inches(3.3), Inches(0.8),
                 theme, font_size=15, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # Modules badge
    mbadge = add_rounded_rect(s, x + Inches(0.9), Inches(3.05), Inches(2.1), Inches(0.35), color)
    tf = mbadge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = modules; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    # Topics
    add_bullet_frame(s, x + Inches(0.3), Inches(3.6), Inches(3.3), Inches(2.8),
                     topics, font_size=13, color=DARK_TEXT, icon="•")

add_footer_bar(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Day 1 Detail
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Detailed view of Day 1 modules, topics, and labs.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_BLUE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📅  Day 1 — Tenant Foundations & Site Management", font_size=30, bold=True, color=ACCENT_BLUE)

modules_d1 = [
    ("Module 1", "Introduction to Microsoft 365\nand SharePoint Online", [
        "M365 service architecture",
        "SharePoint admin centers",
        "Service limits & quotas",
    ], "Lab: Explore the M365\nEnvironment", ACCENT_BLUE),
    ("Module 2", "Identity, Access, and\nExternal Sharing", [
        "Entra ID fundamentals",
        "Guest access & B2B",
        "External sharing policies",
    ], "Lab: Configure Secure\nAccess", ACCENT_TEAL),
    ("Module 3", "Working with\nSite Collections", [
        "Team & Communication sites",
        "M365 Groups integration",
        "Storage & lifecycle",
    ], "Lab: Manage Site\nCollections", ACCENT_PURPLE),
]
for i, (mod, title, topics, lab, color) in enumerate(modules_d1):
    x = Inches(0.6 + i * 4.2)
    card = add_rounded_rect(s, x, Inches(1.3), Inches(3.9), Inches(5.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.3), Inches(3.9), Inches(0.08), color)
    # Module number
    badge = add_rounded_rect(s, x + Inches(1.2), Inches(1.5), Inches(1.5), Inches(0.4), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = mod; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    # Title
    add_text_box(s, x + Inches(0.2), Inches(2.1), Inches(3.5), Inches(0.8),
                 title, font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # Topics
    add_shape_rect(s, x + Inches(0.5), Inches(2.95), Inches(2.9), Inches(0.03), LIGHT_GRAY)
    add_text_box(s, x + Inches(0.2), Inches(3.1), Inches(3.5), Inches(0.4),
                 "Key Topics:", font_size=12, bold=True, color=MID_GRAY)
    add_bullet_frame(s, x + Inches(0.2), Inches(3.45), Inches(3.3), Inches(1.5),
                     topics, font_size=12, color=DARK_TEXT, icon="•")
    # Lab box
    lab_box = add_rounded_rect(s, x + Inches(0.3), Inches(5.1), Inches(3.3), Inches(1.0), LIGHT_BLUE)
    add_text_box(s, x + Inches(0.4), Inches(5.15), Inches(3.1), Inches(0.3),
                 "🧪 Hands-on Lab:", font_size=11, bold=True, color=ACCENT_BLUE)
    add_text_box(s, x + Inches(0.4), Inches(5.45), Inches(3.1), Inches(0.55),
                 lab, font_size=12, color=DARK_TEXT)

add_footer_bar(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Day 2 Detail
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Detailed view of Day 2 modules: Permissions, Metadata, Search, Apps.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_TEAL)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📅  Day 2 — Information Architecture, Search & Customization",
             font_size=28, bold=True, color=ACCENT_TEAL)

modules_d2 = [
    ("Module 4", "Permissions &\nCollaboration", [
        "Permission inheritance",
        "Sharing links & scopes",
        "SP vs M365 groups",
    ], "Lab: Design a\nPermission Model", ACCENT_BLUE),
    ("Module 5", "Metadata &\nTerm Store", [
        "Information architecture",
        "Managed metadata",
        "Term Store hierarchy",
    ], "Lab: Create & Manage\nMetadata", ACCENT_TEAL),
    ("Module 6", "Search &\nMicrosoft Search", [
        "Search architecture",
        "Bookmarks & Q&A",
        "Search verticals",
    ], "Lab: Configure Search\nExperience", ACCENT_PURPLE),
    ("Module 7", "Apps &\nCustomization", [
        "SPFx overview",
        "App Catalog",
        "App governance",
    ], "Lab: Managing\nApps", ORANGE),
]
for i, (mod, title, topics, lab, color) in enumerate(modules_d2):
    x = Inches(0.5 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.3), Inches(2.95), Inches(5.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.3), Inches(2.95), Inches(0.08), color)
    badge = add_rounded_rect(s, x + Inches(0.5), Inches(1.5), Inches(1.9), Inches(0.38), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = mod; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    add_text_box(s, x + Inches(0.15), Inches(2.1), Inches(2.65), Inches(0.7),
                 title, font_size=15, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.3), Inches(2.8), Inches(2.35), Inches(0.03), LIGHT_GRAY)
    add_bullet_frame(s, x + Inches(0.15), Inches(2.95), Inches(2.65), Inches(1.5),
                     topics, font_size=12, color=DARK_TEXT, icon="•")
    lab_box = add_rounded_rect(s, x + Inches(0.15), Inches(4.95), Inches(2.65), Inches(1.1), LIGHT_BLUE)
    add_text_box(s, x + Inches(0.25), Inches(5.0), Inches(2.45), Inches(0.3),
                 "🧪 Lab:", font_size=11, bold=True, color=ACCENT_BLUE)
    add_text_box(s, x + Inches(0.25), Inches(5.3), Inches(2.45), Inches(0.6),
                 lab, font_size=11, color=DARK_TEXT)

add_footer_bar(s, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Day 3 Detail
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Detailed view of Day 3 modules: Compliance, OneDrive, PowerShell, Monitoring, Power Platform.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📅  Day 3 — Governance, Compliance & Automation",
             font_size=30, bold=True, color=ACCENT_PURPLE)

modules_d3 = [
    ("Module 8", "Purview &\nCompliance", [
        "Retention policies",
        "Sensitivity labels",
        "eDiscovery & DLP",
    ], "Lab: Compliance\nControls", ACCENT_BLUE),
    ("Module 9", "OneDrive\nAdmin", [
        "Sharing & sync",
        "Storage policies",
        "Device access",
    ], "Lab: OneDrive\nSettings", ACCENT_TEAL),
    ("Module 10", "PowerShell\nAutomation", [
        "SPO Management Shell",
        "Graph PowerShell",
        "Bulk operations",
    ], "Lab: Automate\nAdmin Tasks", ACCENT_PURPLE),
    ("Module 11", "Monitoring\n& Auditing", [
        "Audit logs",
        "Usage analytics",
        "Governance practices",
    ], "Lab: Operational\nReview", ORANGE),
    ("Module 12", "Power Platform\n(Optional)", [
        "Power Automate",
        "Power Apps basics",
        "Workflow governance",
    ], "Lab: Request\nWorkflow", MID_GRAY),
]
for i, (mod, title, topics, lab, color) in enumerate(modules_d3):
    x = Inches(0.3 + i * 2.6)
    card = add_rounded_rect(s, x, Inches(1.3), Inches(2.4), Inches(5.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.3), Inches(2.4), Inches(0.08), color)
    badge = add_rounded_rect(s, x + Inches(0.2), Inches(1.5), Inches(2.0), Inches(0.35), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = mod; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    add_text_box(s, x + Inches(0.1), Inches(2.05), Inches(2.2), Inches(0.65),
                 title, font_size=14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.2), Inches(2.7), Inches(2.0), Inches(0.03), LIGHT_GRAY)
    add_bullet_frame(s, x + Inches(0.1), Inches(2.85), Inches(2.2), Inches(1.5),
                     topics, font_size=11, color=DARK_TEXT, icon="•")
    lab_box = add_rounded_rect(s, x + Inches(0.1), Inches(5.05), Inches(2.2), Inches(1.05), LIGHT_BLUE)
    add_text_box(s, x + Inches(0.2), Inches(5.1), Inches(2.0), Inches(0.25),
                 "🧪 Lab:", font_size=10, bold=True, color=ACCENT_BLUE)
    add_text_box(s, x + Inches(0.2), Inches(5.35), Inches(2.0), Inches(0.6),
                 lab, font_size=11, color=DARK_TEXT)

add_footer_bar(s, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Section Divider: The Microsoft 365 Ecosystem
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: Before we start, let's set the context of where SharePoint fits in the M365 ecosystem.")
section_divider(s, "The Microsoft 365 Ecosystem", "Where SharePoint Online fits in the bigger picture", "☁️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Microsoft 365 Platform Overview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Show the M365 platform: Azure AD/Entra at the base, SharePoint/OneDrive/Teams/Exchange as services, "
              "with Purview, Search, and Power Platform as cross-cutting capabilities.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Microsoft 365 Platform at a Glance", font_size=32, bold=True, color=DARK_BG)

# Layer 1: Foundation
add_shape_rect(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.2), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.4),
             "🔐  Foundation: Microsoft Entra ID (Identity + Access)",
             font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.4),
             "Authentication  ·  Authorization  ·  Conditional Access  ·  B2B/B2C  ·  Zero Trust",
             font_size=13, color=RGBColor(0xDD, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)

# Layer 2: Core Services (4 cards)
services = [
    ("📡", "SharePoint\nOnline", "Sites · Libraries\nContent Services", ACCENT_BLUE),
    ("☁️", "OneDrive for\nBusiness", "Personal files\nSync & share", ACCENT_TEAL),
    ("💬", "Microsoft\nTeams", "Chat · Meetings\nChannel files", ACCENT_PURPLE),
    ("📧", "Exchange\nOnline", "Mail · Calendar\nContacts", ORANGE),
]
for i, (icon, name, desc, color) in enumerate(services):
    x = Inches(0.8 + i * 3.05)
    card = add_rounded_rect(s, x, Inches(2.8), Inches(2.8), Inches(2.3), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(2.8), Inches(2.8), Inches(0.08), color)
    add_text_box(s, x + Inches(0.2), Inches(2.95), Inches(2.4), Inches(0.5),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.4), Inches(2.4), Inches(0.7),
                 name, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(4.1), Inches(2.4), Inches(0.7),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Layer 3: Cross-cutting (top banner)
add_shape_rect(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.2), LIGHT_BLUE)
add_text_box(s, Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.4),
             "Cross-Cutting Capabilities", font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
caps = ["🔍 Microsoft Search", "🛡️ Microsoft Purview", "⚡ Power Platform", "📊 Graph API"]
for i, cap in enumerate(caps):
    x = Inches(0.8 + i * 2.93)
    add_text_box(s, x, Inches(1.8), Inches(2.93), Inches(0.4),
                 cap, font_size=14, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SharePoint + OneDrive + Teams Triangle
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("The collaboration triangle. SharePoint provides content services, OneDrive provides personal storage, "
              "Teams provides the communication layer. They share a common content infrastructure.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "The Collaboration Triangle", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Three services, one content platform — understanding the relationships",
             font_size=16, color=MID_GRAY)

# SharePoint (center-left)
sp_card = add_rounded_rect(s, Inches(0.8), Inches(1.8), Inches(3.8), Inches(3.8), WHITE)
sp_card.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.8), Inches(3.8), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(2.05), Inches(3.4), Inches(0.5),
             "📡  SharePoint Online", font_size=22, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
sp_points = [
    "Team & Communication sites",
    "Document libraries & lists",
    "Intranet & portals",
    "Content services backbone",
    "Metadata & search integration",
]
add_bullet_frame(s, Inches(1.0), Inches(2.6), Inches(3.4), Inches(2.5),
                 sp_points, font_size=13, color=DARK_TEXT, icon="•")

# OneDrive (center)
od_card = add_rounded_rect(s, Inches(4.9), Inches(1.8), Inches(3.5), Inches(3.8), WHITE)
od_card.shadow.inherit = False
add_shape_rect(s, Inches(4.9), Inches(1.8), Inches(3.5), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(5.1), Inches(2.05), Inches(3.1), Inches(0.5),
             "☁️  OneDrive for Business", font_size=20, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
od_points = [
    "Personal file storage",
    "Sync to desktop/mobile",
    "Share files externally",
    "Built ON SharePoint",
    "Admin controls inherited",
]
add_bullet_frame(s, Inches(5.1), Inches(2.6), Inches(3.1), Inches(2.5),
                 od_points, font_size=13, color=DARK_TEXT, icon="•")

# Teams (right)
tm_card = add_rounded_rect(s, Inches(8.7), Inches(1.8), Inches(3.8), Inches(3.8), WHITE)
tm_card.shadow.inherit = False
add_shape_rect(s, Inches(8.7), Inches(1.8), Inches(3.8), Inches(0.08), ACCENT_PURPLE)
add_text_box(s, Inches(8.9), Inches(2.05), Inches(3.4), Inches(0.5),
             "💬  Microsoft Teams", font_size=22, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
tm_points = [
    "Chat & meetings hub",
    "Channel-based collaboration",
    "Files tab = SharePoint library",
    "M365 Group drives permissions",
    "Extensible with apps & bots",
]
add_bullet_frame(s, Inches(8.9), Inches(2.6), Inches(3.4), Inches(2.5),
                 tm_points, font_size=13, color=DARK_TEXT, icon="•")

# Connection arrows / labels
add_text_box(s, Inches(4.2), Inches(3.4), Inches(1.0), Inches(0.4),
             "↔", font_size=28, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(8.1), Inches(3.4), Inches(1.0), Inches(0.4),
             "↔", font_size=28, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

# Key insight at bottom
insight = add_rounded_rect(s, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.8), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(5.95), Inches(9.7), Inches(0.7),
             "💡 Key Insight:  Every Teams channel stores files in a SharePoint document library. "
             "Every OneDrive is technically a personal SharePoint site collection. "
             "As a SharePoint admin, you manage all three.",
             font_size=13, color=DARK_TEXT)

add_footer_bar(s, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Admin's Toolkit
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Overview of tools admins will use throughout this course: admin centers, PowerShell, Graph API.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🧰  Your Admin Toolkit", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Tools you'll use across all 3 days",
             font_size=16, color=MID_GRAY)

tools = [
    ("🌐", "Microsoft 365\nAdmin Center", "Tenant settings, users, groups,\nlicenses, service health",
     "admin.microsoft.com", ACCENT_BLUE),
    ("📡", "SharePoint\nAdmin Center", "Sites, sharing policies, storage,\nterm store, migration",
     "admin.sharepoint.com", ACCENT_TEAL),
    ("🔐", "Entra Admin\nCenter", "Identity, Conditional Access,\napp registrations, B2B",
     "entra.microsoft.com", ACCENT_PURPLE),
    ("🛡️", "Microsoft Purview\nPortal", "Compliance, retention, DLP,\nsensitivity labels, eDiscovery",
     "purview.microsoft.com", ORANGE),
    ("⚡", "PowerShell &\nGraph API", "SPO Management Shell, Graph PS,\nbulk ops, reporting, automation",
     "Shell + Graph Explorer", GREEN),
]
for i, (icon, name, desc, url, color) in enumerate(tools):
    if i < 4:
        row = 0; col = i
        x = Inches(0.5 + col * 3.15); y = Inches(1.7)
    else:
        x = Inches(3.8); y = Inches(4.2)
    w = Inches(2.95) if i < 4 else Inches(5.7)
    h = Inches(2.2) if i < 4 else Inches(2.0)
    card = add_rounded_rect(s, x, y, w, h, WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, w, Inches(0.08), color)
    add_text_box(s, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.85), y + Inches(0.15), w - Inches(1.1), Inches(0.65),
                 name, font_size=15, bold=True, color=color)
    add_text_box(s, x + Inches(0.85), y + Inches(0.75), w - Inches(1.1), Inches(0.8),
                 desc, font_size=12, color=DARK_TEXT)
    add_text_box(s, x + Inches(0.85), y + h - Inches(0.45), w - Inches(1.1), Inches(0.3),
                 url, font_size=10, color=MID_GRAY, font_name="Cascadia Code")

add_footer_bar(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Section Divider: Lab Environment
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: let's talk about the lab environment setup.")
section_divider(s, "Lab Environment", "Your sandbox for hands-on learning", "🧪")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Lab Environment Setup
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Explain the shared demo tenant, participant accounts (P01-P10), and the Northwind scenario. "
              "Stress: no tenant-wide changes by participants.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_TEAL)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🧪  Lab Environment", font_size=32, bold=True, color=DARK_BG)

# Environment cards
env_info = [
    ("🏢", "Shared Demo Tenant", "Microsoft 365 E3/E5 demo tenant\npre-configured with training users\nand Northwind scenario data", ACCENT_BLUE),
    ("👤", "Your Account", "You'll receive credentials:\nP01@tenant through P10@tenant\nSharePoint Admin role assigned", ACCENT_TEAL),
    ("🌐", "Your Practice Site", "NW-Pxx-ProjectSite (per participant)\nAll labs scoped to YOUR site\nNo tenant-wide changes!", ACCENT_PURPLE),
    ("📁", "Sample Content", "Pre-loaded document templates:\nContracts, metadata worksheets,\ntest files for each lab exercise", ORANGE),
]
for i, (icon, title, desc, color) in enumerate(env_info):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.3 + row * 2.4)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(2.1), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(2.1), color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(1.0), y + Inches(0.1), Inches(4.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color)
    add_text_box(s, x + Inches(1.0), y + Inches(0.65), Inches(4.5), Inches(1.3),
                 desc, font_size=14, color=DARK_TEXT)

# Safety warning
warn = add_rounded_rect(s, Inches(2.0), Inches(6.2), Inches(9.3), Inches(0.65), RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, Inches(2.3), Inches(6.25), Inches(8.7), Inches(0.55),
             "⚠️ Shared Tenant Rule:  Work ONLY inside your assigned NW-Pxx site. "
             "Do NOT modify tenant-wide settings unless instructed by the trainer.",
             font_size=13, bold=True, color=ORANGE)

add_footer_bar(s, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — The Northwind Scenario
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Introduce the Northwind scenario: a fictional company whose SharePoint environment you'll administer. "
              "Labs build on each other using this consistent scenario.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📋  The Northwind Scenario", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Your training company — all labs are set in the Northwind universe",
             font_size=16, color=MID_GRAY)

# Scenario description
add_rounded_rect(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.5), WHITE).shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(1.9), Inches(11.0), Inches(1.1),
             "Northwind Traders is a mid-sized company with 500 employees across 3 offices. They've recently migrated "
             "to Microsoft 365 and need a SharePoint administrator to set up their collaboration environment. "
             "You've been hired as that admin. Over the next 3 days, you'll configure their tenant, build sites, "
             "manage permissions, implement compliance, and automate operations.",
             font_size=15, color=DARK_TEXT)

# Lab progression
add_text_box(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
             "Lab Progression — Each lab builds on the previous one:",
             font_size=16, bold=True, color=DARK_TEXT)

lab_items = [
    ("Lab 1–3", "Set up tenant, configure identity & sharing, create sites", "Day 1", ACCENT_BLUE),
    ("Lab 4–5", "Design permissions, build metadata & information architecture", "Day 2a", ACCENT_TEAL),
    ("Lab 6–7", "Configure search, deploy apps & customization", "Day 2b", ACCENT_PURPLE),
    ("Lab 8–9", "Implement compliance, configure OneDrive policies", "Day 3a", ORANGE),
    ("Lab 10–12", "Automate with PowerShell, audit & monitor, build workflows", "Day 3b", GREEN),
]
for i, (labs, desc, timing, color) in enumerate(lab_items):
    y = Inches(4.1 + i * 0.55)
    add_shape_rect(s, Inches(0.8), y, Inches(0.08), Inches(0.45), color)
    add_text_box(s, Inches(1.1), y + Inches(0.02), Inches(1.5), Inches(0.4),
                 labs, font_size=14, bold=True, color=color)
    add_text_box(s, Inches(2.8), y + Inches(0.02), Inches(7.5), Inches(0.4),
                 desc, font_size=14, color=DARK_TEXT)
    add_text_box(s, Inches(10.5), y + Inches(0.02), Inches(2.0), Inches(0.4),
                 timing, font_size=12, bold=True, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

add_footer_bar(s, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Section Divider: Key Concepts Preview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: a quick primer on the key concepts we'll cover.")
section_divider(s, "Key Concepts Preview", "Setting the foundation before we dive in", "🧠")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — SharePoint Online in 2026 — What's Current
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Position SPO in 2026: Cloud-only, Entra-integrated, Purview-compliant, AI-ready. "
              "No more on-prem. Modern admin center. Graph-based APIs.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "SharePoint Online in 2026", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "What makes the modern platform different — and why this course exists",
             font_size=16, color=MID_GRAY)

evolutions = [
    ("☁️", "Cloud-Native", "No servers to manage. Microsoft handles infrastructure, patching, and scaling. "
     "You focus on configuration and governance.", ACCENT_BLUE),
    ("🔐", "Entra-Integrated", "Identity powered by Microsoft Entra ID. Zero Trust, Conditional Access, "
     "B2B collaboration built in from day one.", ACCENT_TEAL),
    ("🛡️", "Purview-Protected", "Compliance isn't an add-on. Retention, sensitivity labels, DLP, and eDiscovery "
     "are native to the platform.", ACCENT_PURPLE),
    ("📡", "Graph-Powered", "Microsoft Graph API is the unified endpoint. PowerShell modules, admin centers, "
     "and apps all use Graph under the hood.", ORANGE),
    ("🔍", "Search-Unified", "Microsoft Search spans SharePoint, OneDrive, Teams, and beyond. "
     "One search experience, admin-managed.", GREEN),
    ("⚡", "AI-Ready", "Copilot for Microsoft 365 relies on SharePoint content. "
     "Good admin practices = better AI results.", RGBColor(0x88, 0x00, 0xCC)),
]
for i, (icon, title, desc, color) in enumerate(evolutions):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.7 + row * 1.65)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(1.4), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(1.4), color)
    add_text_box(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=26, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.85), y + Inches(0.1), Inches(2.0), Inches(0.45),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s, x + Inches(0.85), y + Inches(0.55), Inches(4.7), Inches(0.7),
                 desc, font_size=12, color=DARK_TEXT)

add_footer_bar(s, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — SharePoint vs SharePoint Server (Quick Comparison)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("For participants coming from on-prem SharePoint, clarify the key differences. "
              "This course is 100% SharePoint Online / cloud-focused.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "SharePoint Online vs SharePoint Server", font_size=30, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
             "This course focuses exclusively on SharePoint Online (cloud)",
             font_size=15, color=MID_GRAY)

# Table header
header_y = Inches(1.6)
cols = [Inches(0.8), Inches(4.6), Inches(9.0)]
col_widths = [Inches(3.8), Inches(4.4), Inches(3.5)]
headers = ["Aspect", "SharePoint Online", "SharePoint Server"]
add_shape_rect(s, Inches(0.8), header_y, Inches(11.7), Inches(0.5), ACCENT_BLUE)
for j, h in enumerate(headers):
    add_text_box(s, cols[j], header_y + Inches(0.03), col_widths[j], Inches(0.4),
                 h, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

rows = [
    ("Infrastructure", "Microsoft-managed cloud", "Your servers & farms"),
    ("Updates", "Continuous (automatic)", "Manual patching cycles"),
    ("Identity", "Microsoft Entra ID", "Active Directory on-prem"),
    ("Administration", "Modern admin centers + PowerShell", "Central Admin + PS"),
    ("Customization", "SPFx, Power Platform", "Full-trust solutions, SPFx"),
    ("Storage", "Pooled tenant storage", "SQL Server databases"),
    ("Compliance", "Purview-native", "Separate configuration"),
    ("Scale", "Multi-tenant, global CDN", "Capacity planning required"),
    ("Cost Model", "Per-user licensing", "Server licensing + hardware"),
]
for i, (aspect, online, server) in enumerate(rows):
    y = header_y + Inches(0.5 + i * 0.53)
    bg_color = WHITE if i % 2 == 0 else NEAR_WHITE
    add_shape_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.53), bg_color, LIGHT_GRAY)
    vals = [aspect, online, server]
    for j, v in enumerate(vals):
        weight = j == 0
        c = DARK_TEXT if j != 1 else ACCENT_BLUE
        add_text_box(s, cols[j], y + Inches(0.07), col_widths[j], Inches(0.38),
                     v, font_size=13, bold=weight, color=c, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Learning Tips
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Practical tips for getting the most out of this course.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, GREEN)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "💡  Tips for Getting the Most Out of This Course", font_size=28, bold=True, color=DARK_BG)

tips = [
    ("🤔", "Ask Questions", "There are no dumb questions. If you're confused, someone else probably is too. "
     "Speak up anytime or use the parking lot."),
    ("🧪", "Do Every Lab", "Labs are where learning sticks. Follow the steps carefully, but also experiment. "
     "Your practice site is your sandbox."),
    ("📝", "Take Notes", "Jot down real-world connections: 'This would solve X problem at my organization.' "
     "These notes will be gold after the course."),
    ("🤝", "Help Each Other", "Pair up for labs. Explain concepts to your neighbor. Teaching is the best way "
     "to learn and verify your understanding."),
    ("📖", "Use References", "Every module links to official Microsoft documentation. Bookmark them! "
     "They'll be your go-to resource after the course ends."),
]
for i, (icon, title, desc) in enumerate(tips):
    y = Inches(1.3 + i * 1.15)
    card = add_rounded_rect(s, Inches(0.8), y, Inches(11.7), Inches(1.0), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, Inches(0.8), y, Inches(0.08), Inches(1.0), GREEN)
    add_text_box(s, Inches(1.1), y + Inches(0.15), Inches(0.6), Inches(0.5),
                 icon, font_size=26, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.7), y + Inches(0.1), Inches(2.5), Inches(0.45),
                 title, font_size=18, bold=True, color=GREEN)
    add_text_box(s, Inches(4.3), y + Inches(0.1), Inches(8.0), Inches(0.8),
                 desc, font_size=13, color=DARK_TEXT)

add_footer_bar(s, 21, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Icebreaker / Introductions
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Use this slide for participant introductions. Each person shares: name, role, "
              "experience with SharePoint, and what they hope to learn.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🙋  Your Turn — Introductions", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Let's get to know each other! Take 1 minute each.",
             font_size=16, color=MID_GRAY)

questions = [
    ("👤", "Your Name & Role", "What's your job title and\nwhat do you do day-to-day?", ACCENT_BLUE),
    ("📊", "Your SharePoint\nExperience", "None / Basic User /\nAdmin / Power User?", ACCENT_TEAL),
    ("🎯", "What You Hope\nto Learn", "One skill or topic you\nwant to master by Day 3", ACCENT_PURPLE),
    ("⚡", "One Fun Fact", "Something interesting\nabout you — keep it light!", ORANGE),
]
for i, (icon, title, desc, color) in enumerate(questions):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.5), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(2.05), Inches(2.3), Inches(0.5),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(2.65), Inches(2.5), Inches(0.8),
                 title, font_size=17, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.5), Inches(3.45), Inches(1.9), Inches(0.03), LIGHT_GRAY)
    add_text_box(s, x + Inches(0.2), Inches(3.6), Inches(2.5), Inches(1.0),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Note
add_text_box(s, Inches(2), Inches(5.8), Inches(9), Inches(0.6),
             "🎤 Trainer will go first to break the ice!",
             font_size=18, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Let's Begin!
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Closing slide of the introduction. Transition to Module 1.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Ready? Let's Begin! 🚀", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.0), Inches(3.3), Inches(3.3), Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
             "Next → Module 1: Introduction to Microsoft 365 and SharePoint Online",
             font_size=20, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(2), Inches(5.0), Inches(9), Inches(0.8),
             "3 Days  ·  12 Modules  ·  12 Labs  ·  1 Goal:\n"
             "Make you a confident, modern SharePoint Admin",
             font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(6.1), Inches(2.3), Inches(0.04), ACCENT_PURPLE)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "00-Course-Introduction.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {slide_counter[0]}")
