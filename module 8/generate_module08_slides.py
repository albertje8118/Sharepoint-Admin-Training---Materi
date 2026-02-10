#!/usr/bin/env python3
"""
Module 8 – Content Governance and Compliance with Microsoft Purview
Generates a 28-slide PPTX with the established design system.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design tokens ──────────────────────────────────────────────
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
DARK_BG   = RGBColor(0x1B, 0x1B, 0x2F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE   = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL   = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
ORANGE        = RGBColor(0xFF, 0x8C, 0x00)
GREEN         = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT    = RGBColor(0xD1, 0x34, 0x38)

FONT_BODY = "Segoe UI"
FONT_CODE = "Cascadia Code"
FOOTER_TEXT = "Module 8 | Content Governance & Compliance with Microsoft Purview"

# ── Helper functions ───────────────────────────────────────────
def add_solid_bg(slide, color):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font_name; p.alignment = alignment
    return tb

def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=16, color=DARK_GRAY, spacing=Pt(10), icon=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = f"{icon} " if icon else "• "
        p.text = prefix + item
        p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = FONT_BODY; p.space_after = spacing
    return tb

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE):
    add_shape_rect(slide, 0, 0, SLIDE_W, Inches(0.06), color)

def add_footer_bar(slide, label=FOOTER_TEXT):
    add_shape_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), Inches(7.08), Inches(10), Inches(0.35),
                 label, font_size=10, color=WHITE)

def section_divider(prs, title, subtitle="", accent=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(1), Inches(3.1), Inches(1.2), Inches(0.06), accent)
    add_text_box(slide, Inches(1), Inches(3.35), Inches(11), Inches(1),
                 title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.3), Inches(10), Inches(0.7),
                     subtitle, font_size=18, color=MID_GRAY)
    add_footer_bar(slide)
    return slide

def new_slide(prs, title, accent=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    add_top_bar(slide, accent)
    add_shape_rect(slide, Inches(0.6), Inches(0.45), Inches(0.08), Inches(0.55), accent)
    add_text_box(slide, Inches(0.85), Inches(0.4), Inches(11), Inches(0.65),
                 title, font_size=28, color=DARK_GRAY, bold=True)
    add_footer_bar(slide)
    return slide

# ── Presentation setup ─────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ================================================================
# SLIDE 1 – Title
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl, DARK_BG)
add_shape_rect(sl, 0, Inches(2.6), SLIDE_W, Inches(2.6), RGBColor(0x22, 0x22, 0x3A))
add_shape_rect(sl, Inches(1), Inches(3.55), Inches(1.5), Inches(0.07), ACCENT_PURPLE)
add_text_box(sl, Inches(1), Inches(2.75), Inches(11), Inches(0.7),
             "Module 8", font_size=22, color=ACCENT_TEAL, bold=True)
add_text_box(sl, Inches(1), Inches(3.7), Inches(11), Inches(1),
             "Content Governance & Compliance\nwith Microsoft Purview",
             font_size=40, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(5.0), Inches(10), Inches(0.5),
             "Scenario: Project Northwind Intranet Modernization  ·  Day 3 of 3",
             font_size=16, color=MID_GRAY)
add_footer_bar(sl)
add_speaker_notes(sl,
    "Welcome to Day 3. Module 8 introduces Microsoft Purview governance features "
    "that SharePoint admins need to understand: retention, sensitivity labels, "
    "eDiscovery, and DLP. Many of these settings are tenant-wide — we'll stress safety.")

# ================================================================
# SLIDE 2 – Why SharePoint admins care
# ================================================================
sl = new_slide(prs, "Why SharePoint Admins Care About Compliance")
cards = [
    ("🌐", "Tenant-Wide\nImpact", "Compliance policies often\napply across the entire\ntenant — one mistake\naffects everyone."),
    ("⚖️", "Regulatory\nRequirements", "Retention, classification,\nand eDiscovery are required\nby law in many industries."),
    ("🛡️", "Data\nProtection", "Sensitivity labels and DLP\nprevent accidental sharing\nof confidential content."),
    ("⚠️", "Shared Tenant\nSafety", "In training, the trainer\nleads policy creation.\nParticipants observe first."),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(0.6 + i * 3.1)
    add_rounded_rect(sl, x, Inches(1.6), Inches(2.8), Inches(4.6), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.3), Inches(1.8), Inches(2.2), Inches(0.7),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.7),
                 title, font_size=17, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(3.3), Inches(2.4), Inches(2.5),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Set the tone: compliance features are powerful but potentially disruptive. "
    "In shared training tenants, the trainer leads policy creation while participants observe. "
    "In production, always pilot before rolling out broadly.")

# ================================================================
# SLIDE 3 – Learning outcomes
# ================================================================
sl = new_slide(prs, "Learning Outcomes")
outcomes = [
    "Describe what Microsoft Purview is used for in Microsoft 365",
    "Explain retention labels, policies, and where they apply",
    "Explain sensitivity labels and label publishing policies",
    "Describe the modern eDiscovery workflow (cases, holds, searches)",
    "Explain why DLP is deployed using simulation and pilots",
]
for i, item in enumerate(outcomes):
    y = Inches(1.5 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.55), Inches(0.55), ACCENT_BLUE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1)
    p.font.size = Pt(22); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.55),
                 item, font_size=19, color=DARK_GRAY)
add_speaker_notes(sl,
    "Five learning outcomes covering the four Purview pillars for this module: "
    "retention, sensitivity labels, eDiscovery, and DLP.")

# ================================================================
# SLIDE 4 – Purview building blocks
# ================================================================
sl = new_slide(prs, "Microsoft Purview: The Four Pillars (Module 8)")
pillars = [
    ("🗄️", "Retention", "Data lifecycle &\nrecords management", ACCENT_BLUE),
    ("🏷️", "Sensitivity\nLabels", "Classification +\nprotection", ACCENT_PURPLE),
    ("🔍", "eDiscovery", "Cases, searches,\nholds, export", ACCENT_TEAL),
    ("🛡️", "DLP", "Data loss prevention\nsimulation-first", ORANGE),
]
for i, (icon, title, desc, color) in enumerate(pillars):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(sl, x, Inches(1.6), Inches(2.8), Inches(4.4), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(1.6), Inches(2.8), Inches(0.07), color)
    add_text_box(sl, x + Inches(0.3), Inches(1.9), Inches(2.2), Inches(0.7),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.7),
                 title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.6), Inches(2.2), Inches(1.5),
                 desc, font_size=15, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Microsoft Purview is the unified governance and compliance portal. "
    "In Module 8 we focus on these four pillars. Each has admin implications "
    "for SharePoint and OneDrive content management.")

# ================================================================
# SLIDE 5 – Section Divider: Retention
# ================================================================
section_divider(prs, "Section 1", "Retention: Policies and Labels", accent=ACCENT_BLUE)

# ================================================================
# SLIDE 6 – Retention: Policy vs Label
# ================================================================
sl = new_slide(prs, "Retention Policies vs Retention Labels")
# Headers
add_shape_rect(sl, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.65), ACCENT_BLUE)
add_text_box(sl, Inches(0.8), Inches(1.55), Inches(5.4), Inches(0.55),
             "Retention Policies (Broad)", font_size=18, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_shape_rect(sl, Inches(6.8), Inches(1.5), Inches(5.9), Inches(0.65), ACCENT_TEAL)
add_text_box(sl, Inches(7.0), Inches(1.55), Inches(5.5), Inches(0.55),
             "Retention Labels (Item-Level)", font_size=18, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
# Policy items
p_items = [
    "Apply to locations (SPO, OneDrive, Exchange…)",
    "Retain and/or delete content at scale",
    "Used for baseline governance",
    "Applied automatically — no user interaction",
    "Cannot mark items as records",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.5), Inches(5.4), Inches(3.5),
                 p_items, font_size=15, icon="•")
# Label items
l_items = [
    "Applied at the document/email level",
    "Only ONE retention label per item",
    "Must be published via label policies first",
    "Users can apply manually",
    "Can mark items as records",
]
add_bullet_frame(sl, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.5),
                 l_items, font_size=15, icon="•")
add_speaker_notes(sl,
    "Key distinction: policies target locations broadly, labels target individual items. "
    "From Microsoft docs: 'An item can have only one retention label applied at a time.' "
    "Both can be used together — they complement each other for a comprehensive strategy.")

# ================================================================
# SLIDE 7 – Capabilities comparison table
# ================================================================
sl = new_slide(prs, "Retention: Capabilities Comparison")
headers = ["Capability", "Policy", "Label"]
rows = [
    ["Retain and/or delete", "✅", "✅"],
    ["Applied automatically", "✅", "✅"],
    ["User can apply manually", "❌", "✅"],
    ["Mark item as a record", "❌", "✅"],
    ["Start period from event", "❌", "✅"],
    ["Disposition review", "❌", "✅"],
    ["Persists if content moved", "❌", "✅ (within M365)"],
]
# Header
col_widths = [Inches(5.5), Inches(2.5), Inches(4.2)]
for j, h in enumerate(headers):
    x = Inches(0.6) + sum(col_widths[:j])
    add_shape_rect(sl, x, Inches(1.4), col_widths[j], Inches(0.55), ACCENT_BLUE)
    add_text_box(sl, x + Inches(0.15), Inches(1.42), col_widths[j] - Inches(0.3),
                 Inches(0.5), h, font_size=14, color=WHITE, bold=True)
# Rows
for i, row in enumerate(rows):
    y = Inches(2.0 + i * 0.7)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        x = Inches(0.6) + sum(col_widths[:j])
        add_shape_rect(sl, x, y, col_widths[j], Inches(0.65), bg, MID_GRAY)
        add_text_box(sl, x + Inches(0.15), y + Inches(0.08), col_widths[j] - Inches(0.3),
                     Inches(0.5), cell, font_size=13, color=DARK_GRAY)
add_speaker_notes(sl,
    "Table sourced from Microsoft Learn: retention policies vs labels capabilities comparison. "
    "Key differences: labels support records management, disposition review, and event-based retention. "
    "Policies are simpler and best for broad baseline governance.")

# ================================================================
# SLIDE 8 – Publishing retention labels
# ================================================================
sl = new_slide(prs, "Publishing Retention Labels: The Path")
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Creating a label ≠ making it available. You must publish via a label policy.",
             font_size=18, color=DARK_GRAY, bold=True)
# Flow boxes
flow = [
    ("1. Create Label", "Purview portal →\nSolutions → Records Mgmt\nor Data Lifecycle Mgmt\n→ Labels", ACCENT_BLUE),
    ("2. Create Label Policy", "Solutions → Records Mgmt\nor Data Lifecycle Mgmt\n→ Policies → Label policies", ACCENT_PURPLE),
    ("3. Choose Locations", "SharePoint sites,\nOneDrive accounts,\nExchange mailboxes,\nM365 Groups", ACCENT_TEAL),
    ("4. Wait for Replication", "SPO / OneDrive:\ntypically within 1 day\n(allow up to 7 days)", ORANGE),
]
for i, (title, desc, color) in enumerate(flow):
    x = Inches(0.5 + i * 3.15)
    card = add_rounded_rect(sl, x, Inches(2.4), Inches(2.9), Inches(3.6), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(2.4), Inches(2.9), Inches(0.06), color)
    add_text_box(sl, x + Inches(0.2), Inches(2.6), Inches(2.5), Inches(0.5),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(2.2),
                 desc, font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(sl, x + Inches(2.9), Inches(3.8), Inches(0.4), Inches(0.5),
                     "→", font_size=20, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Four-step publishing flow. Emphasize step 4: replication delay. "
    "From Microsoft docs: 'When retention labels are published to SharePoint or OneDrive, "
    "those labels typically appear for users to select within one day. However, allow up to seven days.' "
    "In a training environment, pre-publish labels the day before if possible.")

# ================================================================
# SLIDE 9 – Retention timing warning
# ================================================================
sl = new_slide(prs, "Retention Label Timing: Plan for It", ORANGE)
add_rounded_rect(sl, Inches(0.6), Inches(1.5), Inches(12), Inches(1.4),
                 RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text_box(sl, Inches(0.9), Inches(1.65), Inches(11.4), Inches(1.1),
             "⏱️  Published retention labels to SharePoint / OneDrive typically appear within 1 day.\n"
             "Allow up to 7 days for full replication. Exchange labels run on a 7-day process cycle.",
             font_size=17, color=DARK_GRAY)
add_text_box(sl, Inches(0.8), Inches(3.3), Inches(11), Inches(0.5),
             "What If Labels Don't Appear?", font_size=20, color=ORANGE, bold=True)
troubleshoot = [
    "Check label policy Status in the Purview portal (look for 'Error')",
    "Use Set-RetentionCompliancePolicy -RetryDistribution (PowerShell)",
    "Verify the label policy includes the correct locations",
    "Confirm the user has permissions to the target library/site",
]
add_bullet_frame(sl, Inches(0.8), Inches(4.0), Inches(11), Inches(2.5),
                 troubleshoot, font_size=16, icon="→")
add_speaker_notes(sl,
    "Replication timing is a common source of frustration. "
    "From Microsoft docs: if labels don't appear after 7 days, check the policy status. "
    "If you see '(Error)', run Set-RetentionCompliancePolicy -RetryDistribution via PowerShell.")

# ================================================================
# SLIDE 10 – Section Divider: Sensitivity Labels
# ================================================================
section_divider(prs, "Section 2", "Sensitivity Labels: Classification + Protection",
                accent=ACCENT_PURPLE)

# ================================================================
# SLIDE 11 – Sensitivity labels overview
# ================================================================
sl = new_slide(prs, "Sensitivity Labels: What They Do", ACCENT_PURPLE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Classify and optionally protect content across Microsoft 365",
             font_size=18, color=DARK_GRAY)
features = [
    ("🏷️", "Classification", "Visual markings:\nheaders, footers,\nwatermarks"),
    ("🔐", "Encryption", "Restrict who can access\nand what actions they\ncan perform"),
    ("📊", "Scope", "Files, emails, meetings,\nTeams, Groups,\nSharePoint sites"),
    ("📤", "Persistence", "Label travels with the\ncontent — even when\nshared externally"),
]
for i, (icon, title, desc) in enumerate(features):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(sl, x, Inches(2.4), Inches(2.8), Inches(3.5), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.3), Inches(2.6), Inches(2.2), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.2), Inches(2.2), Inches(0.5),
                 title, font_size=17, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.9), Inches(2.2), Inches(1.5),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Sensitivity labels are different from retention labels: they classify and protect content. "
    "Labels support visual markings, encryption, and scope across files/emails/meetings/sites. "
    "Key difference from retention: sensitivity labels persist with content when shared externally.")

# ================================================================
# SLIDE 12 – Sensitivity vs Retention: key differences
# ================================================================
sl = new_slide(prs, "Sensitivity vs Retention Labels: Key Differences")
headers = ["Aspect", "Retention Labels", "Sensitivity Labels"]
rows = [
    ["Purpose", "Keep / delete content", "Classify / protect content"],
    ["Published to…", "Locations (sites, mailboxes)", "Users and groups"],
    ["Persists externally?", "No (within M365 only)", "Yes (travels with content)"],
    ["Can apply encryption?", "No", "Yes"],
    ["Items per document", "1 retention label", "1 sensitivity label"],
    ["Visual markings?", "No", "Yes (header/footer/watermark)"],
]
col_widths = [Inches(3.5), Inches(4.3), Inches(4.4)]
for j, h in enumerate(headers):
    x = Inches(0.6) + sum(col_widths[:j])
    add_shape_rect(sl, x, Inches(1.4), col_widths[j], Inches(0.55), ACCENT_PURPLE)
    add_text_box(sl, x + Inches(0.15), Inches(1.42), col_widths[j] - Inches(0.3),
                 Inches(0.5), h, font_size=14, color=WHITE, bold=True)
for i, row in enumerate(rows):
    y = Inches(2.0 + i * 0.8)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        x = Inches(0.6) + sum(col_widths[:j])
        add_shape_rect(sl, x, y, col_widths[j], Inches(0.75), bg, MID_GRAY)
        add_text_box(sl, x + Inches(0.15), y + Inches(0.1), col_widths[j] - Inches(0.3),
                     Inches(0.55), cell, font_size=13, color=DARK_GRAY)
add_speaker_notes(sl,
    "Critical comparison. Participants often confuse retention and sensitivity labels. "
    "Retention = lifecycle (keep/delete). Sensitivity = classification + protection. "
    "A document can have BOTH one retention label AND one sensitivity label simultaneously.")

# ================================================================
# SLIDE 13 – Create and publish sensitivity labels
# ================================================================
sl = new_slide(prs, "Sensitivity Labels: Create and Publish", ACCENT_PURPLE)
steps = [
    ("1", "Create the label", "Purview portal → Solutions →\nInformation Protection → Sensitivity labels → + Create a label"),
    ("2", "Define scope", "Files & data assets, Emails, Meetings,\nGroups & sites — choose what applies"),
    ("3", "Configure settings", "Visual markings (header/footer/watermark),\nencryption, access controls"),
    ("4", "Create publishing policy", "Solutions → Information Protection →\nPublishing policies → select labels + users/groups"),
    ("5", "Wait & validate", "Allow replication time (~1 hour minimum).\nPilot with test users first."),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.35 + i * 1.15)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.5), Inches(0.5), ACCENT_PURPLE)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(18); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.6), y, Inches(3.5), Inches(0.5),
                 title, font_size=16, color=ACCENT_PURPLE, bold=True)
    add_text_box(sl, Inches(5.3), y, Inches(7.3), Inches(0.8),
                 desc, font_size=13, color=DARK_GRAY)
add_speaker_notes(sl,
    "Five-step process. Note that sensitivity labels are published to USERS and GROUPS, "
    "not locations like retention labels. From Microsoft docs: 'Publish new labels to just a few "
    "test users first, wait for at least one hour, then verify the label behavior on SharePoint and OneDrive.' "
    "For training, avoid encryption complexity unless the tenant is pre-configured.")

# ================================================================
# SLIDE 14 – Sensitivity timing and pilot
# ================================================================
sl = new_slide(prs, "Sensitivity Labels: Timing and Pilot-First Approach", ORANGE)
# Timeline
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Recommended Rollout Strategy:", font_size=20, color=ORANGE, bold=True)
phases = [
    ("Phase 1", "Pilot (few users)", "Publish to test group\nValidate in SPO/OneDrive\nWait ≥ 1 hour", GREEN),
    ("Phase 2", "Expand", "Add more users to\npublishing policy\nMonitor audit logs", ACCENT_BLUE),
    ("Phase 3", "Broad rollout", "Make available to\nall standard users\nLabels fully synced", ACCENT_PURPLE),
]
for i, (phase, title, desc, color) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    card = add_rounded_rect(sl, x, Inches(2.3), Inches(3.8), Inches(3.8), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(2.3), Inches(3.8), Inches(0.06), color)
    badge = add_rounded_rect(sl, x + Inches(0.3), Inches(2.5), Inches(1.2), Inches(0.5), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = phase
    p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, x + Inches(0.3), Inches(3.2), Inches(3.2), Inches(0.5),
                 title, font_size=17, color=color, bold=True)
    add_text_box(sl, x + Inches(0.3), Inches(3.8), Inches(3.2), Inches(1.8),
                 desc, font_size=14, color=DARK_GRAY)
    if i < 2:
        add_text_box(sl, x + Inches(3.85), Inches(3.8), Inches(0.5), Inches(0.5),
                     "→", font_size=22, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "From Microsoft docs: 'Publish new labels to just a few test users first, wait for at least one hour, "
    "then verify the label behavior on SharePoint and OneDrive. Wait at least a day before making the label "
    "available to more users.' Always pilot before broad rollout.")

# ================================================================
# SLIDE 15 – Section Divider: eDiscovery
# ================================================================
section_divider(prs, "Section 3", "eDiscovery: Modern Experience", accent=ACCENT_TEAL)

# ================================================================
# SLIDE 16 – eDiscovery overview
# ================================================================
sl = new_slide(prs, "eDiscovery: What It Is and Why Admins Care", ACCENT_TEAL)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "eDiscovery helps organizations identify, preserve, and export content\n"
             "as evidence for legal and regulatory matters.",
             font_size=18, color=DARK_GRAY)
blocks = [
    ("📁", "Case", "Container for an\ninvestigation —\nscopes all activities"),
    ("🔒", "Hold", "Preserves content so\nit can't be permanently\ndeleted"),
    ("🔎", "Search", "Finds relevant content\nacross SharePoint,\nOneDrive, Exchange"),
    ("📤", "Export", "Produces deliverables\nfor legal / compliance\nworkflows"),
]
for i, (icon, title, desc) in enumerate(blocks):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(sl, x, Inches(2.6), Inches(2.8), Inches(3.4), LIGHT_GRAY, ACCENT_TEAL)
    add_text_box(sl, x + Inches(0.3), Inches(2.8), Inches(2.2), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.4), Inches(2.2), Inches(0.5),
                 title, font_size=18, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(4.0), Inches(2.2), Inches(1.5),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Four building blocks of the eDiscovery workflow. Case is the container, Hold preserves content, "
    "Search finds it, and Export produces it for legal teams. "
    "Classic eDiscovery experiences have been retired — use the modern experience in the Purview portal.")

# ================================================================
# SLIDE 17 – eDiscovery workflow
# ================================================================
sl = new_slide(prs, "eDiscovery Workflow (Modern Experience)", ACCENT_TEAL)
workflow = [
    ("1", "Create a Case", "Define investigation scope and name", ACCENT_TEAL),
    ("2", "Add Members", "Assign roles: eDiscovery Manager\nor eDiscovery Administrator", ACCENT_BLUE),
    ("3", "Place Content on Hold", "Preserve data in SharePoint,\nOneDrive, Exchange (optional)", ORANGE),
    ("4", "Run Search", "Use keywords, date ranges,\nlocation filters to find content", ACCENT_PURPLE),
    ("5", "Review and Export", "Review results, then export\nfor legal / compliance teams", GREEN),
]
for i, (num, title, desc, color) in enumerate(workflow):
    y = Inches(1.35 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.5), Inches(0.5), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(18); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.6), y, Inches(4), Inches(0.5),
                 title, font_size=17, color=color, bold=True)
    add_text_box(sl, Inches(6), y, Inches(6.5), Inches(0.7),
                 desc, font_size=14, color=DARK_GRAY)
add_speaker_notes(sl,
    "Five-step workflow. Note that Hold is optional but recommended when preservation is critical. "
    "Classic eDiscovery has been retired per Microsoft guidance. "
    "eDiscovery permissions are separate — eDiscovery Manager role is needed.")

# ================================================================
# SLIDE 18 – eDiscovery: 2026 alignment + shared tenant safety
# ================================================================
sl = new_slide(prs, "eDiscovery: 2026 Alignment & Shared Tenant Safety", RED_ACCENT)
# 2026 notice
add_rounded_rect(sl, Inches(0.6), Inches(1.5), Inches(12), Inches(1.2),
                 RGBColor(0xFD, 0xE8, 0xE8), RED_ACCENT)
add_text_box(sl, Inches(0.9), Inches(1.6), Inches(11.4), Inches(0.9),
             "⚠️  Classic eDiscovery experiences have been retired (per Microsoft guidance).\n"
             "Use the modern eDiscovery experience in the Microsoft Purview portal.",
             font_size=16, color=DARK_GRAY)
add_text_box(sl, Inches(0.8), Inches(3.1), Inches(11), Inches(0.5),
             "Shared Tenant Lab Rules:", font_size=20, color=RED_ACCENT, bold=True)
rules = [
    "Default: trainer-led demonstration (participants observe)",
    "Hands-on only if eDiscovery permissions are explicitly assigned",
    "Always scope searches and holds to NW-Pxx locations only",
    "Never place holds on other participants' content",
    "Treat eDiscovery as a privileged governance activity",
]
add_bullet_frame(sl, Inches(0.8), Inches(3.8), Inches(11), Inches(2.8),
                 rules, font_size=16, icon="🔒")
add_speaker_notes(sl,
    "Very important safety slide. eDiscovery is powerful — participants must not search or hold "
    "content outside their own NW-Pxx sites. In a shared training tenant, default to trainer-led demo. "
    "Hands-on only if the trainer explicitly assigns eDiscovery Manager permissions.")

# ================================================================
# SLIDE 19 – Section Divider: DLP
# ================================================================
section_divider(prs, "Section 4", "Data Loss Prevention (DLP)", accent=ORANGE)

# ================================================================
# SLIDE 20 – DLP overview
# ================================================================
sl = new_slide(prs, "DLP: Data Loss Prevention Overview", ORANGE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "DLP policies detect and optionally prevent risky actions — such as sharing\n"
             "sensitive content externally (credit cards, SSNs, health records, etc.).",
             font_size=18, color=DARK_GRAY)
add_text_box(sl, Inches(0.8), Inches(2.4), Inches(11), Inches(0.5),
             "The #1 Admin Rule: Don't \"turn on and pray\"", font_size=22,
             color=ORANGE, bold=True)
# Four-phase deployment
phases = [
    ("1", "Define Intent", "What data? From whom?\nWhich locations?", ACCENT_BLUE),
    ("2", "Simulation Mode", "Run policy in simulation\nto understand impact", ORANGE),
    ("3", "Simulation + Tips", "Show policy tips to\nusers — observe reactions", ACCENT_PURPLE),
    ("4", "Enforcement", "Enable blocking/restriction\nafter successful pilot", GREEN),
]
for i, (num, title, desc, color) in enumerate(phases):
    x = Inches(0.5 + i * 3.15)
    card = add_rounded_rect(sl, x, Inches(3.3), Inches(2.9), Inches(3.0), LIGHT_GRAY, color)
    badge = add_rounded_rect(sl, x + Inches(0.3), Inches(3.5), Inches(0.5), Inches(0.5), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(16); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, x + Inches(1.0), Inches(3.55), Inches(1.6), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(sl, x + Inches(0.3), Inches(4.2), Inches(2.3), Inches(1.5),
                 desc, font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Four-phase DLP deployment. From Microsoft docs: 'Make sure you understand the data you're "
    "protecting and the goals you want to achieve. Take time to design a policy before you implement it.' "
    "Simulation mode is essential — it lets you see what WOULD be flagged without blocking users.")

# ================================================================
# SLIDE 21 – DLP simulation mode
# ================================================================
sl = new_slide(prs, "DLP Simulation Mode: How It Works", ORANGE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Simulation mode runs the policy like a WhatIf — you see results without blocking users.",
             font_size=18, color=DARK_GRAY)
sim_items = [
    ("📊", "See Matched Content", "View which files/emails would\ntrigger the policy rules"),
    ("🔄", "Refine Rules", "Adjust conditions to reduce\nfalse positives before going live"),
    ("📈", "Scale Gradually", "Start with one SPO site,\nthen expand to more locations"),
    ("⏱️", "Plan Timing", "Simulation can take up to\n12 hours to complete"),
]
for i, (icon, title, desc) in enumerate(sim_items):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(sl, x, Inches(2.5), Inches(2.8), Inches(3.3), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.6),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.3), Inches(2.2), Inches(0.5),
                 title, font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.9), Inches(2.2), Inches(1.4),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "From Microsoft docs: 'The simulated deployment runs like the WhatIf parameter for PowerShell, "
    "for a specific point in time.' Simulation can take up to 12 hours. "
    "Use it to iteratively refine rules and reduce false positives before enforcement.")

# ================================================================
# SLIDE 22 – Admin governance mindset summary
# ================================================================
sl = new_slide(prs, "Admin Governance Mindset: Purview Summary")
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Three principles for every Purview feature:", font_size=20,
             color=ACCENT_BLUE, bold=True)
principles = [
    ("🎯", "Pilot First", "Always test with a small group before\ntenant-wide rollout. Monitor and refine.",
     GREEN),
    ("⏱️", "Plan for Delays", "Replication takes time. Publish labels/policies\nwell before you need them.",
     ORANGE),
    ("🔒", "Least Privilege", "Grant only the permissions needed.\neDiscovery roles are powerful — use carefully.",
     RED_ACCENT),
]
for i, (icon, title, desc, color) in enumerate(principles):
    y = Inches(2.3 + i * 1.6)
    add_shape_rect(sl, Inches(0.8), y, Inches(0.12), Inches(1.2), color)
    add_text_box(sl, Inches(1.2), y, Inches(0.8), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, Inches(2.0), y, Inches(3), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    add_text_box(sl, Inches(5.5), y, Inches(7), Inches(1.2),
                 desc, font_size=15, color=DARK_GRAY)
add_speaker_notes(sl,
    "Three overarching principles applying to all Purview features: "
    "pilot first, plan for replication delays, and use least privilege. "
    "These apply to retention, sensitivity labels, eDiscovery, and DLP equally.")

# ================================================================
# SLIDE 23 – Section Divider: Lab
# ================================================================
section_divider(prs, "Section 5", "Lab 8: Implementing Compliance Controls", accent=GREEN)

# ================================================================
# SLIDE 24 – Lab preview
# ================================================================
sl = new_slide(prs, "Lab 8: Hands-On Exercises", GREEN)
exercises = [
    ("Task 1", "Upload FAKE content", "Upload sample docs to your NW-Pxx site for labeling"),
    ("Task 2", "Apply sensitivity label", "Apply a published sensitivity label to a document (if available)"),
    ("Task 3", "Apply retention label", "Apply a published retention label to a document (if available)"),
    ("Task 4", "eDiscovery demo", "Trainer-led: walk through Case → Search → Hold workflow"),
    ("Task 5", "DLP awareness", "Trainer-led: review DLP policy creation in simulation mode"),
]
for i, (task, title, desc) in enumerate(exercises):
    y = Inches(1.4 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(1.2), Inches(0.5), GREEN)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = task
    p.font.size = Pt(13); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(2.3), y, Inches(3.5), Inches(0.5),
                 title, font_size=17, color=DARK_GRAY, bold=True)
    add_text_box(sl, Inches(6.2), y, Inches(6.5), Inches(0.5),
                 desc, font_size=14, color=DARK_GRAY)
# Safety note
add_rounded_rect(sl, Inches(0.6), Inches(6.2), Inches(12), Inches(0.9),
                 RGBColor(0xFD, 0xE8, 0xE8), RED_ACCENT)
add_text_box(sl, Inches(0.9), Inches(6.3), Inches(11.4), Inches(0.7),
             "⚠️  Label availability depends on pre-published policies. Tasks 4-5 are trainer-led by default.",
             font_size=14, color=DARK_GRAY, bold=True)
add_speaker_notes(sl,
    "Tasks 1-3 are hands-on if labels have been pre-published. Tasks 4-5 are trainer-led demos. "
    "Remind participants to use only FAKE sample content — never upload real sensitive data. "
    "If labels are not yet visible, explain the replication timing and move to the demo tasks.")

# ================================================================
# SLIDE 25 – Validation checklist
# ================================================================
sl = new_slide(prs, "Lab 8: Validation Checklist", GREEN)
checks = [
    "FAKE sample documents uploaded to NW-Pxx document library",
    "Sensitivity label applied to at least one document (if published)",
    "Retention label applied to at least one document (if published)",
    "Can explain the difference between retention and sensitivity labels",
    "Can describe the eDiscovery workflow: Case → Hold → Search → Export",
    "Can explain why DLP should use simulation mode before enforcement",
]
for i, item in enumerate(checks):
    y = Inches(1.4 + i * 0.95)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.45), Inches(0.45), GREEN)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = "✓"
    p.font.size = Pt(16); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.5), y, Inches(11), Inches(0.55),
                 item, font_size=15, color=DARK_GRAY)
add_speaker_notes(sl,
    "Six validation checkpoints. The first three depend on label availability. "
    "The last three are knowledge checks — every participant should be able to answer these. "
    "If labels weren't available, focus discussion on the conceptual differences.")

# ================================================================
# SLIDE 26 – Key takeaways
# ================================================================
sl = new_slide(prs, "Key Takeaways")
takeaways = [
    "Microsoft Purview is the unified governance + compliance portal for M365",
    "Retention policies target locations broadly; labels target individual items",
    "Retention labels must be published via label policies — plan for replication delays",
    "Sensitivity labels classify + protect; published to users/groups (not locations)",
    "A document can have BOTH one retention label AND one sensitivity label",
    "Classic eDiscovery is retired — use the modern experience in the Purview portal",
    "DLP: always start with simulation mode → pilot → then enforcement",
]
for i, item in enumerate(takeaways):
    y = Inches(1.4 + i * 0.8)
    add_shape_rect(sl, Inches(0.8), y + Inches(0.15), Inches(0.08), Inches(0.35), ACCENT_PURPLE)
    add_text_box(sl, Inches(1.1), y, Inches(11.5), Inches(0.65),
                 item, font_size=16, color=DARK_GRAY)
add_speaker_notes(sl,
    "Seven key takeaways covering all four Purview pillars. "
    "Emphasize: retention ≠ sensitivity labels (different purpose, different publish targets). "
    "DLP is the most impactful if mis-configured — simulation mode is not optional, it's essential.")

# ================================================================
# SLIDE 27 – Knowledge check
# ================================================================
sl = new_slide(prs, "Knowledge Check")
questions = [
    ("Q1", "What is the key difference between a retention policy and a retention label?"),
    ("Q2", "How long can it take for published retention labels to appear in SharePoint?"),
    ("Q3", "Are sensitivity labels published to locations or to users/groups?"),
    ("Q4", "What are the four building blocks of the modern eDiscovery workflow?"),
    ("Q5", "Why should DLP policies start in simulation mode?"),
]
colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ORANGE, GREEN]
for i, (qnum, text) in enumerate(questions):
    y = Inches(1.4 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.7), Inches(0.55), colors[i])
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = qnum
    p.font.size = Pt(16); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.8), y + Inches(0.05), Inches(10.5), Inches(0.55),
                 text, font_size=17, color=DARK_GRAY)
add_speaker_notes(sl,
    "Answers: Q1 — Policies target locations broadly; labels target individual items and support records. "
    "Q2 — Typically within 1 day, but allow up to 7 days. "
    "Q3 — To users and groups (retention labels are published to locations). "
    "Q4 — Case, Hold, Search, Export. "
    "Q5 — To understand impact and reduce false positives before blocking users.")

# ================================================================
# SLIDE 28 – Thank you / Next Module
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl, DARK_BG)
add_shape_rect(sl, Inches(1), Inches(3.1), Inches(1.2), Inches(0.06), ACCENT_TEAL)
add_text_box(sl, Inches(1), Inches(2.2), Inches(11), Inches(0.7),
             "End of Module 8", font_size=20, color=ACCENT_TEAL)
add_text_box(sl, Inches(1), Inches(3.35), Inches(11), Inches(1),
             "Content Governance & Compliance\nwith Microsoft Purview",
             font_size=36, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(4.8), Inches(10), Inches(0.6),
             "Up Next  →  Module 9: OneDrive Administration & Operational Controls",
             font_size=18, color=MID_GRAY)
add_footer_bar(sl)
add_speaker_notes(sl,
    "Module 8 complete. Next is Module 9: OneDrive for Business administration. "
    "Remind participants that compliance settings take effect over time — "
    "encourage them to revisit their labels after the next module.")

# ── Save ───────────────────────────────────────────────────────
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "Module-08-Slides.pptx")
prs.save(out_path)
print(f"✅ Saved {len(prs.slides)}-slide presentation → {out_path}")
