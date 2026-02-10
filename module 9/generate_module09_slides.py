#!/usr/bin/env python3
"""
Module 9 – OneDrive for Business Administration
Generates a 28-slide PPTX with the established design system.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design tokens ──────────────────────────────────────────────
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
DARK_BG   = RGBColor(0x1B, 0x1B, 0x2F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE   = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL   = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
ORANGE        = RGBColor(0xFF, 0x8C, 0x00)
GREEN         = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT    = RGBColor(0xD1, 0x34, 0x38)

FONT_BODY = "Segoe UI"
FONT_CODE = "Cascadia Code"
FOOTER_TEXT = "Module 9 | OneDrive for Business Administration"

# ── Helper functions ───────────────────────────────────────────
def add_solid_bg(slide, color):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font_name; p.alignment = alignment
    return tb

def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=16, color=DARK_GRAY, spacing=Pt(10), icon=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = f"{icon} " if icon else "• "
        p.text = prefix + item
        p.font.size = Pt(font_size); p.font.color.rgb = color
        p.font.name = FONT_BODY; p.space_after = spacing
    return tb

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE):
    add_shape_rect(slide, 0, 0, SLIDE_W, Inches(0.06), color)

def add_footer_bar(slide, label=FOOTER_TEXT):
    add_shape_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), Inches(7.08), Inches(10), Inches(0.35),
                 label, font_size=10, color=WHITE)

def section_divider(prs, title, subtitle="", accent=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(1), Inches(3.1), Inches(1.2), Inches(0.06), accent)
    add_text_box(slide, Inches(1), Inches(3.35), Inches(11), Inches(1),
                 title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.3), Inches(10), Inches(0.7),
                     subtitle, font_size=18, color=MID_GRAY)
    add_footer_bar(slide)
    return slide

def new_slide(prs, title, accent=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    add_top_bar(slide, accent)
    add_shape_rect(slide, Inches(0.6), Inches(0.45), Inches(0.08), Inches(0.55), accent)
    add_text_box(slide, Inches(0.85), Inches(0.4), Inches(11), Inches(0.65),
                 title, font_size=28, color=DARK_GRAY, bold=True)
    add_footer_bar(slide)
    return slide

# ── Presentation setup ─────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ================================================================
# SLIDE 1 – Title
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl, DARK_BG)
add_shape_rect(sl, 0, Inches(2.6), SLIDE_W, Inches(2.6), RGBColor(0x22, 0x22, 0x3A))
add_shape_rect(sl, Inches(1), Inches(3.55), Inches(1.5), Inches(0.07), ACCENT_BLUE)
add_text_box(sl, Inches(1), Inches(2.75), Inches(11), Inches(0.7),
             "Module 9", font_size=22, color=ACCENT_TEAL, bold=True)
add_text_box(sl, Inches(1), Inches(3.7), Inches(11), Inches(1),
             "OneDrive for Business\nAdministration",
             font_size=40, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(5.0), Inches(10), Inches(0.5),
             "Scenario: Project Northwind Intranet Modernization  ·  Day 3 of 3",
             font_size=16, color=MID_GRAY)
add_footer_bar(sl)
add_speaker_notes(sl,
    "Module 9 covers OneDrive for Business from a SharePoint admin perspective. "
    "OneDrive is built on SharePoint Online, so many controls live in the SharePoint admin center. "
    "We'll cover sharing, sync, storage, access control, and the deleted-user lifecycle.")

# ================================================================
# SLIDE 2 – Why admins care
# ================================================================
sl = new_slide(prs, "Why Is OneDrive \"SharePoint Admin Work\"?")
cards = [
    ("🏗️", "Same Platform", "OneDrive = personal site\non SharePoint Online\n(-my.sharepoint.com)"),
    ("⚙️", "Same Admin Center", "Sharing, sync, storage,\naccess control — all\nin SharePoint admin center"),
    ("🌐", "Tenant-Wide", "Policy changes affect\nevery user's OneDrive\nimmediately or within hours"),
    ("⚠️", "Shared Tenant", "In training: observe\nfirst, trainer changes\npolicies only"),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(0.6 + i * 3.1)
    add_rounded_rect(sl, x, Inches(1.6), Inches(2.8), Inches(4.6), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.3), Inches(1.8), Inches(2.2), Inches(0.7),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.6),
                 title, font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(3.3), Inches(2.4), Inches(2.5),
                 desc, font_size=15, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Key insight: OneDrive for Business is a personal site collection hosted on SharePoint Online. "
    "Each user's OneDrive lives under {tenant}-my.sharepoint.com. "
    "SharePoint admins manage OneDrive policies via the same admin center used for all SPO settings.")

# ================================================================
# SLIDE 3 – Learning outcomes
# ================================================================
sl = new_slide(prs, "Learning Outcomes")
outcomes = [
    "Explain how OneDrive for Business relates to SharePoint Online",
    "Locate OneDrive admin settings in the SharePoint admin center",
    "Describe key policy areas: sharing, sync, storage, retention, access control",
    "Explain the deleted-user OneDrive lifecycle and restore options",
]
for i, item in enumerate(outcomes):
    y = Inches(1.6 + i * 1.3)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.55), Inches(0.55), ACCENT_BLUE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1)
    p.font.size = Pt(22); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.55),
                 item, font_size=20, color=DARK_GRAY)
add_speaker_notes(sl,
    "Four learning outcomes — architecture relationship, where settings live, "
    "five policy areas, and the user lifecycle for departures.")

# ================================================================
# SLIDE 4 – OneDrive architecture
# ================================================================
sl = new_slide(prs, "OneDrive Architecture: Admin Mental Model")
# Diagram: Tenant → SPO → -my.sharepoint.com → Personal Sites
boxes = [
    ("Microsoft 365\nTenant", ACCENT_BLUE, Inches(0.6), Inches(2.5)),
    ("SharePoint\nOnline", ACCENT_PURPLE, Inches(3.5), Inches(2.5)),
    ("{tenant}-my\n.sharepoint.com", ACCENT_TEAL, Inches(6.4), Inches(2.0)),
    ("User A\nOneDrive", ORANGE, Inches(9.8), Inches(1.6)),
    ("User B\nOneDrive", ORANGE, Inches(9.8), Inches(3.2)),
    ("User C\nOneDrive", ORANGE, Inches(9.8), Inches(4.8)),
]
for label, color, x, y in boxes:
    w = Inches(2.3) if "User" not in label else Inches(2.5)
    h = Inches(1.2)
    box = add_rounded_rect(sl, x, y, w, h, color)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label
    p.font.size = Pt(13); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
# Arrows
for ax, ay, txt in [(Inches(2.9), Inches(2.9), "→"), (Inches(5.8), Inches(2.9), "→"),
                     (Inches(8.7), Inches(2.0), "→"), (Inches(8.7), Inches(3.5), "→"),
                     (Inches(8.7), Inches(5.1), "→")]:
    add_text_box(sl, ax, ay, Inches(0.8), Inches(0.5), txt, font_size=18, color=DARK_GRAY,
                 bold=True, alignment=PP_ALIGN.CENTER)
# Key point
add_rounded_rect(sl, Inches(0.6), Inches(5.5), Inches(8), Inches(1),
                 RGBColor(0xE3, 0xF2, 0xFD), ACCENT_BLUE)
add_text_box(sl, Inches(0.9), Inches(5.6), Inches(7.4), Inches(0.8),
             "💡 Each user's OneDrive is a SharePoint site collection.\n"
             "SharePoint policies, retention, and eDiscovery all apply.",
             font_size=14, color=DARK_GRAY)
add_speaker_notes(sl,
    "Each OneDrive is a personal site collection under {tenant}-my.sharepoint.com. "
    "This means SharePoint permission model, retention, eDiscovery, sensitivity labels, "
    "and DLP all work on OneDrive content — it's the same platform.")

# ================================================================
# SLIDE 5 – Section Divider: Where Settings Live
# ================================================================
section_divider(prs, "Section 1", "Where OneDrive Settings Live", accent=ACCENT_BLUE)

# ================================================================
# SLIDE 6 – Settings map
# ================================================================
sl = new_slide(prs, "OneDrive Settings in the SharePoint Admin Center")
settings = [
    ("🔗", "Sharing", "Organization-wide sharing\nlevel and default link type", "Sharing page"),
    ("🔄", "Sync", "Tenant sync controls,\nhide Sync button option", "Settings → Sync"),
    ("💾", "Storage Limit", "Default quota for all\nnew and existing users", "Settings → Storage limit"),
    ("🗄️", "Retention", "Days to keep deleted\nuser's OneDrive (30–3650)", "Settings → Retention"),
    ("🔒", "Access Control", "Unmanaged devices,\nnetwork location policy", "Access control page"),
    ("🔔", "Notifications", "External sharing\nnotification settings", "Settings → Notifications"),
]
for i, (icon, title, desc, path) in enumerate(settings):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.5 + row * 2.8)
    card = add_rounded_rect(sl, x, y, Inches(3.9), Inches(2.5), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.6),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.9), y + Inches(0.15), Inches(2.7), Inches(0.5),
                 title, font_size=17, color=ACCENT_BLUE, bold=True)
    add_text_box(sl, x + Inches(0.2), y + Inches(0.8), Inches(3.5), Inches(1),
                 desc, font_size=13, color=DARK_GRAY)
    add_text_box(sl, x + Inches(0.2), y + Inches(1.9), Inches(3.5), Inches(0.4),
                 f"📍 {path}", font_size=11, color=ACCENT_PURPLE)
add_speaker_notes(sl,
    "Six key setting areas in the SharePoint admin center. From Microsoft docs: "
    "'Many organizations use OneDrive without changing any of the options. "
    "To change these settings, use the SharePoint admin center.' "
    "Each setting area has a specific location in the admin center.")

# ================================================================
# SLIDE 7 – Section Divider: Sharing Controls
# ================================================================
section_divider(prs, "Section 2", "Sharing & Sync Controls", accent=ACCENT_TEAL)

# ================================================================
# SLIDE 8 – Sharing controls
# ================================================================
sl = new_slide(prs, "Sharing Controls: Policy vs User Action", ACCENT_TEAL)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Tenant policy sets the boundaries → Users share within those boundaries",
             font_size=18, color=DARK_GRAY, bold=True)
# Two columns
col_data = [
    ("Admin Policy (Boundary)", ACCENT_TEAL,
     ["Set organization-wide sharing level",
      "Configure default sharing link type",
      "Control 'Anyone' link expiration",
      "Restrict external sharing by domain",
      "OneDrive sharing ≤ SharePoint level"]),
    ("User Action (Within Boundary)", ACCENT_BLUE,
     ["Share files and folders with others",
      "Choose link type (within policy limits)",
      "Grant Edit or View permissions",
      "Share internally or externally (if allowed)",
      "Manage shared file access"]),
]
for i, (title, color, items) in enumerate(col_data):
    x = Inches(0.6 + i * 6.3)
    add_rounded_rect(sl, x, Inches(2.3), Inches(5.8), Inches(4.0), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(2.3), Inches(5.8), Inches(0.06), color)
    add_text_box(sl, x + Inches(0.3), Inches(2.5), Inches(5.2), Inches(0.5),
                 title, font_size=17, color=color, bold=True)
    add_bullet_frame(sl, x + Inches(0.3), Inches(3.2), Inches(5.2), Inches(2.8),
                     items, font_size=15, icon="•")
add_speaker_notes(sl,
    "Sharing controls mirror what we covered in Module 2, but here we emphasize the OneDrive angle. "
    "Note: OneDrive sharing level can never be MORE permissive than the SharePoint organization level. "
    "For training, keep sharing internal-only unless the trainer authorizes external test.")

# ================================================================
# SLIDE 9 – Sync: Shortcuts vs Sync
# ================================================================
sl = new_slide(prs, "Sync Controls: Shortcuts vs Sync Button", ACCENT_TEAL)
# Two columns
add_shape_rect(sl, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.65), GREEN)
add_text_box(sl, Inches(0.8), Inches(1.55), Inches(5.4), Inches(0.55),
             "✅  Shortcuts (Recommended)", font_size=18, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_shape_rect(sl, Inches(6.8), Inches(1.5), Inches(5.9), Inches(0.65), ORANGE)
add_text_box(sl, Inches(7.0), Inches(1.55), Inches(5.5), Inches(0.55),
             "⚠️  Sync Button (Legacy Approach)", font_size=18, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
# Shortcuts
s_items = [
    "Add shortcut to OneDrive",
    "Linked to user account (follows across devices)",
    "Only the specific folder is synced",
    "More performant than full library sync",
    "Microsoft-recommended approach",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.5), Inches(5.4), Inches(3),
                 s_items, font_size=15, icon="→")
# Sync button
b_items = [
    "Syncs entire library to the device",
    "Device-bound (doesn't follow the user)",
    "Can cause large sync queues",
    "Existing syncs not affected if hidden",
    "Admin can hide via Set-SPOTenant",
]
add_bullet_frame(sl, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3),
                 b_items, font_size=15, icon="→")
# PowerShell
add_rounded_rect(sl, Inches(0.6), Inches(5.8), Inches(12), Inches(1),
                 RGBColor(0x1E, 0x1E, 0x2E))
add_text_box(sl, Inches(0.9), Inches(5.9), Inches(11.4), Inches(0.8),
             "Set-SPOTenant -HideSyncButtonOnTeamSite $true",
             font_size=15, color=RGBColor(0xCE, 0xD4, 0xDA), font_name=FONT_CODE)
add_speaker_notes(sl,
    "From Microsoft docs: 'It's recommended to use shortcuts instead of using the Sync button. "
    "Shortcuts are more performant because rather than syncing the entire library, only the specific "
    "folder is synced. Additionally, because shortcuts are added to a user's OneDrive rather than "
    "to the device, it's easier to access content across all devices.' "
    "Admin command: Set-SPOTenant -HideSyncButtonOnTeamSite $true to hide the Sync button.")

# ================================================================
# SLIDE 10 – Section Divider: Storage & Access
# ================================================================
section_divider(prs, "Section 3", "Storage Policies & Access Control", accent=ACCENT_PURPLE)

# ================================================================
# SLIDE 11 – Storage policies
# ================================================================
sl = new_slide(prs, "Storage Policies: Quotas and Risks", ACCENT_PURPLE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Storage management has three layers:", font_size=18, color=DARK_GRAY, bold=True)
layers = [
    ("1", "License-Based Maximum", "Determined by the user's\nlicense plan (e.g., 1 TB or 5 TB)", ACCENT_BLUE),
    ("2", "Tenant Default", "Admin sets a default storage\nlimit for all new/existing users", ACCENT_PURPLE),
    ("3", "Per-User Override", "Admin can set a custom quota\nfor specific individual users", ACCENT_TEAL),
]
for i, (num, title, desc, color) in enumerate(layers):
    x = Inches(0.5 + i * 4.15)
    card = add_rounded_rect(sl, x, Inches(2.4), Inches(3.9), Inches(2.5), LIGHT_GRAY, color)
    badge = add_rounded_rect(sl, x + Inches(0.3), Inches(2.6), Inches(0.5), Inches(0.5), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = num
    p.font.size = Pt(18); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, x + Inches(1.0), Inches(2.65), Inches(2.6), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    add_text_box(sl, x + Inches(0.3), Inches(3.3), Inches(3.3), Inches(1.2),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
# Warning
add_rounded_rect(sl, Inches(0.6), Inches(5.4), Inches(12), Inches(1.2),
                 RGBColor(0xFD, 0xE8, 0xE8), RED_ACCENT)
add_text_box(sl, Inches(0.9), Inches(5.5), Inches(11.4), Inches(1),
             "⚠️  Risk: If you reduce the default storage limit below a user's current usage,\n"
             "their OneDrive becomes read-only until they reduce their storage or you increase the limit.",
             font_size=15, color=DARK_GRAY, bold=True)
add_speaker_notes(sl,
    "Three-layer storage model. From Microsoft docs: if you reduce storage below current usage, "
    "the OneDrive may become read-only. Set default storage via SharePoint admin center → Settings → Storage limit. "
    "Per-user overrides are useful for executives or special projects.")

# ================================================================
# SLIDE 12 – Device access controls
# ================================================================
sl = new_slide(prs, "Device Access Control: Unmanaged Devices", ACCENT_PURPLE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "SharePoint admin center → Access control → Unmanaged devices",
             font_size=16, color=ACCENT_PURPLE, bold=True)
# Three options
options = [
    ("Allow Full Access", "Users can access from\nany device, any app.\nNo restrictions.", GREEN, "🔓"),
    ("Allow Limited\n(Web-Only)", "Browser only — no download,\nprint, or sync.\nEditing can be restricted.", ORANGE, "🌐"),
    ("Block Access", "No access from\nunmanaged devices.\nMost restrictive.", RED_ACCENT, "🚫"),
]
for i, (title, desc, color, icon) in enumerate(options):
    x = Inches(0.5 + i * 4.15)
    card = add_rounded_rect(sl, x, Inches(2.3), Inches(3.9), Inches(3.0), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(2.3), Inches(3.9), Inches(0.06), color)
    add_text_box(sl, x + Inches(0.3), Inches(2.5), Inches(3.3), Inches(0.6),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.1), Inches(3.3), Inches(0.7),
                 title, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.9), Inches(3.3), Inches(1.2),
                 desc, font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
# Operational note
add_rounded_rect(sl, Inches(0.6), Inches(5.6), Inches(12), Inches(1.1),
                 RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text_box(sl, Inches(0.9), Inches(5.7), Inches(11.4), Inches(0.9),
             "⏱️  Changes can take up to 24 hours to take effect.\n"
             "Doesn't impact users already signed in. Uses Entra Conditional Access under the hood.",
             font_size=14, color=DARK_GRAY)
add_speaker_notes(sl,
    "Three access tiers for unmanaged devices. From Microsoft docs: 'If you revert back to Allow Full Access, "
    "it could take up to 24 hours for the changes to take effect.' "
    "These controls use Microsoft Entra Conditional Access policies. "
    "Recommendation: also block apps that don't use modern authentication to prevent bypass.")

# ================================================================
# SLIDE 13 – Access control details
# ================================================================
sl = new_slide(prs, "Limited Access: Advanced Configurations", ACCENT_PURPLE)
headers = ["Parameter", "Effect", "Use When"]
rows = [
    ["-AllowEditing $false", "Prevents editing Office files\nin browser", "Strict view-only needed"],
    ["-ReadOnlyForUnmanagedDevices\n$true", "Entire site read-only for\nimpacted users", "Full protection required"],
    ["-LimitedAccessFileType\nOfficeOnlineFilesOnly", "Preview only Office files;\nother files blocked", "Maximum security"],
    ["-LimitedAccessFileType\nWebPreviewableFiles", "Preview all files the browser\ncan render (default)", "Balance security/usability"],
]
col_widths = [Inches(4.2), Inches(4.2), Inches(3.8)]
# Header
for j, h in enumerate(headers):
    x = Inches(0.6) + sum(col_widths[:j])
    add_shape_rect(sl, x, Inches(1.4), col_widths[j], Inches(0.55), ACCENT_PURPLE)
    add_text_box(sl, x + Inches(0.15), Inches(1.42), col_widths[j] - Inches(0.3),
                 Inches(0.5), h, font_size=14, color=WHITE, bold=True)
# Rows
for i, row in enumerate(rows):
    y = Inches(2.0 + i * 1.2)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        x = Inches(0.6) + sum(col_widths[:j])
        add_shape_rect(sl, x, y, col_widths[j], Inches(1.15), bg, MID_GRAY)
        fs = 12 if j == 0 else 13
        fn = FONT_CODE if j == 0 else FONT_BODY
        add_text_box(sl, x + Inches(0.1), y + Inches(0.1), col_widths[j] - Inches(0.2),
                     Inches(0.95), cell, font_size=fs, color=DARK_GRAY, font_name=fn)
add_speaker_notes(sl,
    "PowerShell parameters for fine-tuning limited access. From Microsoft docs: "
    "'-LimitedAccessFileType WebPreviewableFiles (default) allows users to preview Office files. "
    "Warning: this option is known to cause problems with PDF and image file types.' "
    "OfficeOnlineFilesOnly is the most secure but may block legitimate file types.")

# ================================================================
# SLIDE 14 – Section Divider: User Lifecycle
# ================================================================
section_divider(prs, "Section 4", "User Lifecycle: Departures & Restore", accent=ORANGE)

# ================================================================
# SLIDE 15 – Deleted user lifecycle
# ================================================================
sl = new_slide(prs, "Deleted User OneDrive Lifecycle", ORANGE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "What happens when a user is deleted from Microsoft 365?",
             font_size=18, color=DARK_GRAY, bold=True)
# Timeline
timeline = [
    ("Day 0", "User deleted from\nM365 admin center", "Account deletion\nsynced to SharePoint", ACCENT_BLUE),
    ("Day 1–30*", "Retention period\n(default 30 days)", "Manager gets access;\nshared content\nstill accessible", GREEN),
    ("Day 23*", "7-day warning email\nsent to manager or\nsecondary owner", "Reminder before\ndeletion", ORANGE),
    ("After retention", "OneDrive enters\ndeleted state\n(93 days)", "Only SharePoint\nAdmin can restore", RED_ACCENT),
]
for i, (day, event, note, color) in enumerate(timeline):
    x = Inches(0.4 + i * 3.2)
    add_shape_rect(sl, x + Inches(1.4), Inches(2.4), Inches(0.06), Inches(0.5), color)
    badge = add_rounded_rect(sl, x + Inches(0.6), Inches(2.9), Inches(1.7), Inches(0.55), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = day
    p.font.size = Pt(13); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(1.3),
                 event, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), Inches(5.0), Inches(2.5), Inches(1.2),
                 note, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
             "* Configurable: 30–3,650 days (Settings → Retention in SharePoint admin center)",
             font_size=12, color=ACCENT_PURPLE)
add_speaker_notes(sl,
    "From Microsoft docs: 'The default retention period for OneDrive is 30 days, but you can change "
    "this in the SharePoint admin center (30 to 3650 days).' Manager gets automatic access by default. "
    "After the retention period, OneDrive remains in a deleted state for 93 days — "
    "only a SharePoint Administrator can restore it during that window.")

# ================================================================
# SLIDE 16 – Access delegation
# ================================================================
sl = new_slide(prs, "Access Delegation: Manager & Secondary Owner", ORANGE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "When a user is deleted, who gets automatic access to their OneDrive?",
             font_size=18, color=DARK_GRAY, bold=True)
# Two paths
paths = [
    ("Manager (Primary)", GREEN,
     ["Manager specified in Entra ID profile",
      "Gets email notification on user deletion",
      "Automatic access to the user's OneDrive",
      "Gets 7-day reminder before deletion"]),
    ("Secondary Owner (Fallback)", ACCENT_BLUE,
     ["Configured in SharePoint admin center",
      "More features → User profiles → My Sites",
      "Used when no manager is set",
      "Also receives email notification"]),
]
for i, (title, color, items) in enumerate(paths):
    x = Inches(0.6 + i * 6.3)
    add_rounded_rect(sl, x, Inches(2.5), Inches(5.8), Inches(3.5), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(2.5), Inches(5.8), Inches(0.06), color)
    add_text_box(sl, x + Inches(0.3), Inches(2.7), Inches(5.2), Inches(0.5),
                 title, font_size=17, color=color, bold=True)
    add_bullet_frame(sl, x + Inches(0.3), Inches(3.4), Inches(5.2), Inches(2.4),
                     items, font_size=15, icon="→")
# Key callout
add_rounded_rect(sl, Inches(0.6), Inches(6.3), Inches(12), Inches(0.7),
                 RGBColor(0xFD, 0xE8, 0xE8), RED_ACCENT)
add_text_box(sl, Inches(0.9), Inches(6.35), Inches(11.4), Inches(0.55),
             "⚠️  If neither manager nor secondary owner is set, no one has automatic access.",
             font_size=14, color=DARK_GRAY, bold=True)
add_speaker_notes(sl,
    "From Microsoft docs: 'By default, when a user is deleted, the user's manager is automatically "
    "given access to the user's OneDrive.' If no manager is set, the secondary owner configured in "
    "My Site Settings is used. If neither is configured, no one gets automatic access and the "
    "OneDrive will eventually be deleted without anyone being notified.")

# ================================================================
# SLIDE 17 – Restore options
# ================================================================
sl = new_slide(prs, "Restoring a Deleted OneDrive", ORANGE)
restore_options = [
    ("During Retention Period\n(default 30 days)", GREEN,
     ["OneDrive still accessible",
      "Manager/secondary owner has access",
      "Files can be downloaded/moved",
      "Re-create user account restores access"]),
    ("After Retention\n(93-day deleted state)", ORANGE,
     ["Only SharePoint Admin can restore",
      "Use Restore-DeletedSite or admin center",
      "Shared content no longer accessible",
      "Last chance before permanent deletion"]),
    ("After 93-Day Window", RED_ACCENT,
     ["Permanently deleted",
      "Cannot be recovered by admin",
      "Exception: Purview retention policies\nor eDiscovery holds may override",
      "Plan ahead!"]),
]
for i, (title, color, items) in enumerate(restore_options):
    x = Inches(0.4 + i * 4.2)
    card = add_rounded_rect(sl, x, Inches(1.5), Inches(3.9), Inches(5.0), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(1.5), Inches(3.9), Inches(0.06), color)
    add_text_box(sl, x + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.9),
                 title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(sl, x + Inches(0.2), Inches(2.8), Inches(3.5), Inches(3.2),
                     items, font_size=13, icon="→")
add_speaker_notes(sl,
    "Three restore windows. During the retention period, it's easy — files are still accessible. "
    "After retention, SPO admin can restore within the 93-day deleted state. "
    "After that, it's gone — unless Purview retention policies or eDiscovery holds were in place. "
    "From Microsoft docs: 'The OneDrive remains in a deleted state for 93 days and can only "
    "be restored by a SharePoint Administrator.'")

# ================================================================
# SLIDE 18 – Retention + eDiscovery interaction
# ================================================================
sl = new_slide(prs, "OneDrive + Purview: Retention & eDiscovery Holds")
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "Purview retention policies and eDiscovery holds can override\n"
             "the standard OneDrive deletion timeline.",
             font_size=18, color=DARK_GRAY)
interaction = [
    ("Purview Retention Policy", "If a retention policy covers OneDrive,\ncontent is retained for the full policy period\n— even if the user is deleted.",
     ACCENT_PURPLE),
    ("eDiscovery Hold", "If an eDiscovery hold is placed on a\nuser's OneDrive, content is preserved\nregardless of deletion or retention settings.",
     ACCENT_TEAL),
    ("Unlicensed Account Archive", "OneDrive accounts without a valid\nlicense are automatically archived on\nday 93. Holds are still honored.",
     ORANGE),
]
for i, (title, desc, color) in enumerate(interaction):
    y = Inches(2.5 + i * 1.5)
    add_shape_rect(sl, Inches(0.8), y, Inches(0.12), Inches(1.2), color)
    add_text_box(sl, Inches(1.2), y, Inches(4), Inches(0.5),
                 title, font_size=17, color=color, bold=True)
    add_text_box(sl, Inches(5.5), y, Inches(7), Inches(1.2),
                 desc, font_size=14, color=DARK_GRAY)
add_speaker_notes(sl,
    "Important for compliance-focused organizations. Purview retention policies keep content "
    "beyond the OneDrive retention setting. eDiscovery holds are the strongest — they override "
    "everything. From Microsoft docs: 'All OneDrive accounts that don't have a valid OneDrive "
    "license are automatically archived on their 93rd unlicensed day. Retention settings, "
    "retention policies, eDiscovery, and all holds are still honored.'")

# ================================================================
# SLIDE 19 – Section Divider: Common Support Scenarios
# ================================================================
section_divider(prs, "Section 5", "Common Support Scenarios", accent=RED_ACCENT)

# ================================================================
# SLIDE 20 – FAQ: support scenarios
# ================================================================
sl = new_slide(prs, "Common Admin Support Scenarios", RED_ACCENT)
scenarios = [
    ("Sync missing / not working", "Check: is the Sync button hidden?\nIs the device managed? Check sync\nhealth reports in Apps Admin Center."),
    ("Blocked download on\nunmanaged device", "Expected if limited access is set.\nVerify policy in Access control page.\nChanges take up to 24 hours."),
    ("User's OneDrive read-only", "Storage quota exceeded.\nIncrease quota or ask user\nto free space."),
    ("Departed user's files needed", "Check retention period setting.\nAccess as manager or SPO admin.\n93-day deleted state = admin restore."),
    ("External sharing not working", "OneDrive level ≤ SPO level.\nCheck both org level and site level.\nVerify user hasn't been restricted."),
]
for i, (issue, resolution) in enumerate(scenarios):
    y = Inches(1.35 + i * 1.1)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rounded_rect(sl, Inches(0.6), y, Inches(12), Inches(1), bg, MID_GRAY)
    add_text_box(sl, Inches(0.9), y + Inches(0.1), Inches(3.5), Inches(0.8),
                 issue, font_size=14, color=RED_ACCENT, bold=True)
    add_text_box(sl, Inches(4.8), y + Inches(0.05), Inches(7.5), Inches(0.9),
                 resolution, font_size=13, color=DARK_GRAY)
add_speaker_notes(sl,
    "Five common support scenarios that SharePoint admins encounter with OneDrive. "
    "Walk through each one and discuss how participants would investigate in a real environment. "
    "Sync issues are very common — point participants to the sync health reports.")

# ================================================================
# SLIDE 21 – Section Divider: Lab
# ================================================================
section_divider(prs, "Section 6", "Lab 9: OneDrive Administration", accent=GREEN)

# ================================================================
# SLIDE 22 – Lab preview
# ================================================================
sl = new_slide(prs, "Lab 9: Hands-On Exercises", GREEN)
exercises = [
    ("Task 1", "Document tenant settings", "Record current Sharing, Sync, Storage, Retention values"),
    ("Task 2", "Observe access control", "Review unmanaged device policy settings (read-only)"),
    ("Task 3", "Internal sharing test", "Share a OneDrive file with another participant"),
    ("Task 4", "Verify sharing behavior", "Confirm recipient can access the shared file"),
    ("Task 5", "Complete M09 worksheet", "Document observations in your participant pack"),
]
for i, (task, title, desc) in enumerate(exercises):
    y = Inches(1.4 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(1.2), Inches(0.5), GREEN)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = task
    p.font.size = Pt(13); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(2.3), y, Inches(3.5), Inches(0.5),
                 title, font_size=17, color=DARK_GRAY, bold=True)
    add_text_box(sl, Inches(6.2), y, Inches(6.5), Inches(0.5),
                 desc, font_size=14, color=DARK_GRAY)
# Time
add_rounded_rect(sl, Inches(8.5), Inches(6.3), Inches(4), Inches(0.6), GREEN)
add_text_box(sl, Inches(8.7), Inches(6.35), Inches(3.6), Inches(0.5),
             "⏱️  Estimated time: 25–35 min", font_size=15, color=WHITE, bold=True)
add_speaker_notes(sl,
    "Lab is observation-heavy by design — OneDrive policies are tenant-wide and we don't want "
    "participants changing settings in a shared tenant. The sharing test (tasks 3-4) is safe "
    "because it only creates an internal share between participants.")

# ================================================================
# SLIDE 23 – Validation checklist
# ================================================================
sl = new_slide(prs, "Lab 9: Validation Checklist", GREEN)
checks = [
    "Documented current Sharing level, Sync settings, Storage limit, and Retention value",
    "Reviewed Access control page — can describe unmanaged device policy options",
    "Successfully shared a OneDrive file internally with another participant",
    "Confirmed the recipient could access and open the shared file",
    "Completed M09 worksheet in participant pack",
    "Can explain the 30-day default + 93-day deleted state lifecycle",
]
for i, item in enumerate(checks):
    y = Inches(1.4 + i * 0.95)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.45), Inches(0.45), GREEN)
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = "✓"
    p.font.size = Pt(16); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.5), y, Inches(11), Inches(0.55),
                 item, font_size=15, color=DARK_GRAY)
add_speaker_notes(sl,
    "Six validation checkpoints. Tasks 1-2 are observation/documentation. "
    "Tasks 3-4 are hands-on sharing verification. Task 5 is worksheet completion. "
    "Checkpoint 6 is a knowledge check about the deleted-user lifecycle.")

# ================================================================
# SLIDE 24 – Common issues / troubleshooting
# ================================================================
sl = new_slide(prs, "Lab 9: Common Issues & Troubleshooting", RED_ACCENT)
issues = [
    ("Can't find Sync settings", "Ensure you're in the\nSharePoint admin center", "Go to Settings → Sync\n(not OneDrive admin)"),
    ("Sharing link doesn't work", "Recipient may be\nexternal or unlicensed", "Verify sharing level allows\nthe link type used"),
    ("Storage limit shows 0", "Tenant default may\nnot be set", "Settings → Storage limit →\nset default or per-user"),
    ("Access control page empty", "Requires SharePoint\nadmin or higher role", "Confirm role assignment\nin M365 admin center"),
]
for i, (issue, symptom, fix) in enumerate(issues):
    y = Inches(1.5 + i * 1.35)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rounded_rect(sl, Inches(0.6), y, Inches(12), Inches(1.2), bg, MID_GRAY)
    add_text_box(sl, Inches(0.9), y + Inches(0.1), Inches(3.2), Inches(0.4),
                 issue, font_size=14, color=RED_ACCENT, bold=True)
    add_text_box(sl, Inches(0.9), y + Inches(0.55), Inches(3.2), Inches(0.5),
                 symptom, font_size=12, color=DARK_GRAY)
    add_text_box(sl, Inches(4.5), y + Inches(0.1), Inches(7.8), Inches(1),
                 fix, font_size=13, color=DARK_GRAY)
add_speaker_notes(sl,
    "Common lab issues. The most frequent confusion is finding the right settings page — "
    "everything is in the SharePoint admin center, not a separate OneDrive admin portal. "
    "Sharing issues usually relate to link type vs sharing level mismatch.")

# ================================================================
# SLIDE 25 – Key takeaways
# ================================================================
sl = new_slide(prs, "Key Takeaways")
takeaways = [
    "OneDrive for Business = personal site collection on SharePoint Online",
    "All key OneDrive policies are managed via the SharePoint admin center",
    "Six setting areas: Sharing, Sync, Storage, Retention, Access control, Notifications",
    "Shortcuts > Sync button — Microsoft-recommended for file access",
    "Unmanaged device controls use Entra Conditional Access (24h propagation)",
    "Deleted-user OneDrive: default 30 days retention + 93 days deleted state",
    "Purview retention policies and eDiscovery holds can override deletion",
]
for i, item in enumerate(takeaways):
    y = Inches(1.4 + i * 0.8)
    add_shape_rect(sl, Inches(0.8), y + Inches(0.15), Inches(0.08), Inches(0.35), ACCENT_BLUE)
    add_text_box(sl, Inches(1.1), y, Inches(11.5), Inches(0.65),
                 item, font_size=16, color=DARK_GRAY)
add_speaker_notes(sl,
    "Seven key takeaways. Core message: OneDrive IS SharePoint — the admin surface is the same. "
    "The deleted-user lifecycle is a common exam/support topic. "
    "Shortcuts vs Sync is a practical recommendation admins should push to their org.")

# ================================================================
# SLIDE 26 – Knowledge check
# ================================================================
sl = new_slide(prs, "Knowledge Check")
questions = [
    ("Q1", "Where do OneDrive admin settings live?"),
    ("Q2", "What is the recommended way for users to access shared folders — Shortcuts or Sync?"),
    ("Q3", "What happens if you reduce the storage limit below a user's current usage?"),
    ("Q4", "How many days is the default OneDrive retention for deleted users?"),
    ("Q5", "After the retention period, how long does the 'deleted state' window last?"),
]
colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ORANGE, GREEN]
for i, (qnum, text) in enumerate(questions):
    y = Inches(1.4 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.7), Inches(0.55), colors[i])
    tf = badge.text_frame
    p = tf.paragraphs[0]; p.text = qnum
    p.font.size = Pt(16); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.8), y + Inches(0.05), Inches(10.5), Inches(0.55),
                 text, font_size=17, color=DARK_GRAY)
add_speaker_notes(sl,
    "Answers: Q1 — SharePoint admin center (Sharing, Settings, Access control pages). "
    "Q2 — Shortcuts (Add shortcut to OneDrive) — more performant, follows the user across devices. "
    "Q3 — The user's OneDrive becomes read-only until they reduce usage or admin increases quota. "
    "Q4 — 30 days (configurable 30–3650). "
    "Q5 — 93 days — only a SharePoint Administrator can restore during this window.")

# ================================================================
# SLIDE 27 – PowerShell quick reference
# ================================================================
sl = new_slide(prs, "PowerShell Quick Reference: OneDrive Admin")
commands = [
    ("Hide Sync button", "Set-SPOTenant -HideSyncButtonOnTeamSite $true"),
    ("Set retention (days)", "Set-SPOTenant -OrphanedPersonalSitesRetentionPeriod 365"),
    ("Set default storage (MB)", "Set-SPOTenant -OneDriveStorageQuota 5242880"),
    ("Set per-user storage", "Set-SPOSite -Identity https://tenant-my.sharepoint.com/\n    personal/user_domain_com -StorageQuotaWarningLevel\n    4194304 -StorageQuota 5242880"),
    ("Block unmanaged devices", "Set-SPOTenant -ConditionalAccessPolicy AllowLimitedAccess"),
]
for i, (label, cmd) in enumerate(commands):
    y = Inches(1.35 + i * 1.15)
    add_text_box(sl, Inches(0.8), y, Inches(3.5), Inches(0.4),
                 label, font_size=14, color=ACCENT_BLUE, bold=True)
    add_rounded_rect(sl, Inches(4.5), y, Inches(8.2), Inches(0.95),
                     RGBColor(0x1E, 0x1E, 0x2E))
    add_text_box(sl, Inches(4.7), y + Inches(0.1), Inches(7.8), Inches(0.75),
                 cmd, font_size=12, color=RGBColor(0xCE, 0xD4, 0xDA), font_name=FONT_CODE)
add_speaker_notes(sl,
    "PowerShell reference for common OneDrive admin tasks. "
    "All use Set-SPOTenant or Set-SPOSite from the SharePoint Online Management Shell. "
    "These are informational — participants should not run these in the shared training tenant.")

# ================================================================
# SLIDE 28 – Thank you / Next Module
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl, DARK_BG)
add_shape_rect(sl, Inches(1), Inches(3.1), Inches(1.2), Inches(0.06), ACCENT_TEAL)
add_text_box(sl, Inches(1), Inches(2.2), Inches(11), Inches(0.7),
             "End of Module 9", font_size=20, color=ACCENT_TEAL)
add_text_box(sl, Inches(1), Inches(3.35), Inches(11), Inches(1),
             "OneDrive for Business\nAdministration",
             font_size=36, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(4.8), Inches(10), Inches(0.6),
             "Up Next  →  Module 10: Automating SharePoint Administration",
             font_size=18, color=MID_GRAY)
add_footer_bar(sl)
add_speaker_notes(sl,
    "Module 9 complete. Next up: Module 10 dives into PowerShell and automation "
    "for SharePoint administration — connecting to SPO, scripting common tasks, "
    "and the admin toolkit.")

# ── Save ───────────────────────────────────────────────────────
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "Module-09-Slides.pptx")
prs.save(out_path)
print(f"✅ Saved {len(prs.slides)}-slide presentation → {out_path}")
