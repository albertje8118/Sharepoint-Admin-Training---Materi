"""
Generate Module 5 PPTX — Managing Metadata and the Term Store
Modern, engaging design matching Modules 1-4 style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG     = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEC, 0xF9)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT    = RGBColor(0x24, 0x24, 0x24)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT  = RGBColor(0xD1, 0x34, 0x38)
YELLOW_WARN = RGBColor(0xF2, 0xC8, 0x11)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ═══════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════
def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6), icon="▸", font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon} {bullet}"; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.name = font_name; p.space_after = spacing
    return txBox

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)

def add_footer_bar(slide, slide_num, total, module_label="Module 5"):
    bar_top = SLIDE_HEIGHT - Inches(0.45)
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_WIDTH, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), bar_top + Inches(0.05), Inches(7), Inches(0.35),
                 f"{module_label}  |  Managing Metadata and the Term Store",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), bar_top + Inches(0.05),
                 Inches(1.2), Inches(0.35), f"{slide_num} / {total}",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)

def section_divider(slide, section_title, subtitle="", icon=""):
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_TEAL)
    if icon:
        add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                     icon, font_size=60, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                 section_title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

TOTAL_SLIDES = 28
slide_counter = [0]

def new_slide(notes=""):
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if notes: add_speaker_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome to Module 5 — Managing Metadata and the Term Store. "
              "Metadata is the backbone of information architecture in SharePoint. "
              "In this module we learn why consistent tagging beats folders, "
              "how the term store works, and how to apply managed metadata to libraries.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
             "MODERN SHAREPOINT ONLINE FOR ADMINISTRATORS  ·  DAY 2",
             font_size=14, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(2.3), Inches(11), Inches(1.5),
             "Module 5", font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(1),
             "Managing Metadata and the Term Store",
             font_size=32, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(2), Inches(5.2), Inches(9), Inches(0.6),
             "Classify consistently  ·  Govern centrally  ·  Find anything fast",
             font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(6.1), Inches(2.3), Inches(0.04), ACCENT_PURPLE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Why Metadata Matters (Overview / Motivation)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Metadata answers the question 'What is this document about?' without relying "
              "on the file name or folder path. Consistent metadata drives better search refiners, "
              "reliable views, and enforceable governance policies. "
              "Use the Northwind contracts example to make this concrete.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Why Metadata Matters", font_size=36, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Move beyond folders — let documents describe themselves",
             font_size=16, color=MID_GRAY)

cards = [
    ("🏷️", "Consistent", "Classification", "Same terms used across\nsites, libraries, and teams"),
    ("🔍", "Better", "Search Refiners", "Metadata powers the\nrefinement panel in search"),
    ("📊", "Reliable", "Views & Filters", "Group, sort, and filter\nby metadata columns"),
    ("🛡️", "Stronger", "Governance", "Enforce retention, DLP,\nand compliance via tags"),
]
for i, (icon, stat, label, desc) in enumerate(cards):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(s, x, Inches(1.9), Inches(2.8), Inches(3.6), WHITE)
    card.shadow.inherit = False
    add_text_box(s, x + Inches(0.3), Inches(2.1), Inches(2.2), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.8),
                 stat, font_size=44, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(0.5),
                 label, font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.6), Inches(4.1), Inches(1.6), Inches(0.04), ACCENT_TEAL)
    add_text_box(s, x + Inches(0.2), Inches(4.3), Inches(2.4), Inches(1.0),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Learning Outcomes
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Set clear expectations: by the end of this module learners can explain "
              "IA principles, navigate the term store, and apply managed metadata columns.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🎯  Learning Outcomes", font_size=32, bold=True, color=DARK_BG)

outcomes = [
    "Explain basic IA principles and why metadata improves findability and governance",
    "Describe the Term Store hierarchy: groups → term sets → terms",
    "Describe term store roles and how delegated term management works",
    "Apply managed metadata to a document library via a managed metadata column",
]
for i, outcome in enumerate(outcomes):
    y = Inches(1.6 + i * 1.15)
    badge = add_rounded_rect(s, Inches(0.8), y, Inches(0.55), Inches(0.55), ACCENT_BLUE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.5),
                 outcome, font_size=18, color=DARK_TEXT)

add_footer_bar(s, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Section Divider: Information Architecture
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: before diving into the term store, let's understand *why* "
              "information architecture matters.")
section_divider(s, "Information Architecture",
                "Why structure matters before you touch the term store", "🏗️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — IA Principles
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Users don't search for 'the correct file name' — they search for concepts: "
              "contract type, department, customer, status. Good IA makes content findable "
              "by concept rather than by tribal knowledge of folder paths.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Information Architecture Principles", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "How users actually find information — and how admins should structure it",
             font_size=16, color=MID_GRAY)

principles = [
    ("🧭", "Findability", "Users search by concepts\n(type, department, status),\nnot by file names",
     ACCENT_BLUE),
    ("📐", "Consistency", "Same classification terms\nacross all sites and teams\nreduces confusion",
     ACCENT_TEAL),
    ("🔗", "Connectedness", "Related content linked via\nshared metadata enables\ncross-site discovery",
     ACCENT_PURPLE),
    ("📏", "Scalability", "A taxonomy designed today\nshould grow with the org\nwithout re-work",
     ORANGE),
]
for i, (icon, title, desc, color) in enumerate(principles):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(2.05), Inches(2.3), Inches(0.6),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.7), Inches(2.3), Inches(0.5),
                 title, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(1.2),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Anti-patterns callout
callout = add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.0), RGBColor(0xFD, 0xE7, 0xE9))
add_text_box(s, Inches(1.1), Inches(5.6), Inches(11.1), Inches(0.8),
             "⚠️  Common Anti-Patterns:  Too many folders encoding tribal knowledge  ·  "
             "Free-text columns for critical classification (spelling variants!)  ·  "
             "No governance = no consistency",
             font_size=14, color=RED_ACCENT)

add_footer_bar(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Section Divider: Metadata Types
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: now let's compare the two main metadata approaches.")
section_divider(s, "Metadata Types",
                "Managed metadata vs site columns — when to choose what", "🏷️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Site Columns vs Managed Metadata (Side-by-Side)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Site columns are great for simple, stable lists. Managed metadata columns "
              "are ideal when you need organization-wide controlled vocabulary, synonyms, "
              "hierarchy, and re-use across many site collections. "
              "Key decision: do you need governance and hierarchy? Choose managed metadata.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Site Columns vs Managed Metadata", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Right tool for the right job — both are valuable",
             font_size=16, color=MID_GRAY)

# LEFT card — Site Columns
left_card = add_rounded_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.6), WHITE)
left_card.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.6),
             "📋  Site Columns", font_size=24, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
site_col_items = [
    "Reusable field definition (date, number, choice, yes/no)",
    "Scope: site or content type",
    "Great for small, stable option lists",
    "No built-in synonyms or hierarchy",
    "Users type or select from a fixed list",
]
add_bullet_frame(s, Inches(1.2), Inches(2.7), Inches(5.0), Inches(3.4),
                 site_col_items, font_size=15, color=DARK_TEXT, icon="▸")

# RIGHT card — Managed Metadata
right_card = add_rounded_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.6), WHITE)
right_card.shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.6),
             "🏷️  Managed Metadata", font_size=24, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
mm_items = [
    "Terms from a centrally managed term store",
    "Scope: tenant-wide (or local term set)",
    "Ideal for organization-wide classification",
    "Supports synonyms, hierarchy, and multilingual labels",
    "Users pick from a type-ahead controlled picker",
]
add_bullet_frame(s, Inches(7.3), Inches(2.7), Inches(5.0), Inches(3.4),
                 mm_items, font_size=15, color=DARK_TEXT, icon="▸")

# Decision hint at bottom
hint_box = add_rounded_rect(s, Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.5), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(6.65), Inches(9.7), Inches(0.4),
             "💡 Rule of thumb:  Need governance, hierarchy, or reuse across sites?  → Managed Metadata",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Terminology (Microsoft Definitions)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("These are the official Microsoft terms. Make sure learners can distinguish "
              "taxonomy vs folksonomy, managed terms vs enterprise keywords, open vs closed term sets.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📖  Key Terminology", font_size=32, bold=True, color=DARK_BG)

terms = [
    ("Taxonomy", "Formal, hierarchical classification system\n(controlled, structured, top-down)", ACCENT_BLUE),
    ("Folksonomy", "Informal, user-driven tagging\n(unstructured, bottom-up, like tag clouds)", ACCENT_TEAL),
    ("Managed Terms", "Pre-defined terms organized in\nhierarchical term sets by admins", ACCENT_PURPLE),
    ("Enterprise Keywords", "Free-form words/phrases added by\nusers — can later be promoted to managed terms", ORANGE),
    ("Open Term Set", "Users can add new terms\nwhen tagging items", GREEN),
    ("Closed Term Set", "Only authorized users can add terms\n(stricter governance)", RED_ACCENT),
]
for i, (term, desc, color) in enumerate(terms):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.5 + row * 2.6)
    card = add_rounded_rect(s, x, y, Inches(3.8), Inches(2.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(3.8), Inches(0.06), color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.5),
                 term, font_size=20, bold=True, color=color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.2),
                 desc, font_size=14, color=DARK_TEXT)

add_footer_bar(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Benefits of Managed Metadata
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Per Microsoft docs: consistent metadata improves search, navigation, "
              "content discoverability, and enables metadata-driven navigation in lists and libraries.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, GREEN)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "✅  Benefits of Managed Metadata", font_size=32, bold=True, color=DARK_BG)

benefits = [
    ("🎯", "Consistent Use", "Control which terms users can apply.\nSame terms across all sites = reliable governance.",
     "Term sets enforce uniformity"),
    ("🔍", "Improved Discoverability", "Search refinement panel lets users filter\nresults by metadata facets.",
     "Metadata powers refiners in search"),
    ("🧭", "Metadata Navigation", "Site admins build navigation and views\nbased on metadata terms.",
     "Filter lists/libraries by metadata pivots"),
    ("⚡", "Flexibility", "Supports range from strict taxonomy\nto open folksonomy — your choice.",
     "Open or closed term sets per need"),
]
for i, (icon, title, desc, detail) in enumerate(benefits):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 2.6)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(2.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(2.2), GREEN)
    add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.5),
                 title, font_size=20, bold=True, color=GREEN)
    add_text_box(s, x + Inches(1.0), y + Inches(0.7), Inches(4.5), Inches(1.0),
                 desc, font_size=14, color=DARK_TEXT)
    add_text_box(s, x + Inches(1.0), y + Inches(1.7), Inches(4.5), Inches(0.4),
                 detail, font_size=12, color=MID_GRAY, font_name="Segoe UI")

add_footer_bar(s, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Section Divider: The Term Store
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: now let's look at the heart of managed metadata — the Term Store.")
section_divider(s, "The Term Store",
                "Hierarchy, roles, and delegated management", "🗂️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Term Store Hierarchy (Visual)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("The term store has a 3-level hierarchy: Term Group → Term Set → Terms. "
              "Groups provide security boundaries (who can manage what). "
              "Term sets can be global (tenant-wide) or local (site-scoped). "
              "Terms can have synonyms, translations, and custom properties.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Term Store Hierarchy", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Three levels — each with a distinct purpose",
             font_size=16, color=MID_GRAY)

# Level 1: Term Group
card1 = add_rounded_rect(s, Inches(0.8), Inches(1.9), Inches(3.6), Inches(4.0), WHITE)
card1.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.9), Inches(3.6), Inches(0.1), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(2.2), Inches(3.2), Inches(0.6),
             "📦  Term Group", font_size=24, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1.0), Inches(2.9), Inches(3.2), Inches(0.4),
             "Security boundary", font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
group_details = [
    "Container for term sets",
    "Controls who can manage",
    "Only term store admins create groups",
    "E.g. 'NW-Pxx-TermGroup'",
]
add_bullet_frame(s, Inches(1.2), Inches(3.5), Inches(3.0), Inches(2.0),
                 group_details, font_size=13, color=DARK_TEXT, icon="•")

# Arrow 1
add_text_box(s, Inches(4.5), Inches(3.6), Inches(0.6), Inches(0.6),
             "→", font_size=36, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Level 2: Term Set
card2 = add_rounded_rect(s, Inches(5.1), Inches(1.9), Inches(3.6), Inches(4.0), WHITE)
card2.shadow.inherit = False
add_shape_rect(s, Inches(5.1), Inches(1.9), Inches(3.6), Inches(0.1), ACCENT_TEAL)
add_text_box(s, Inches(5.3), Inches(2.2), Inches(3.2), Inches(0.6),
             "📋  Term Set", font_size=24, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(5.3), Inches(2.9), Inches(3.2), Inches(0.4),
             "Group of related terms", font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
tset_details = [
    "Global or local scope",
    "Open or closed for submissions",
    "Has owner, contact, stakeholders",
    "E.g. 'NW-Pxx-ContractType'",
]
add_bullet_frame(s, Inches(5.5), Inches(3.5), Inches(3.0), Inches(2.0),
                 tset_details, font_size=13, color=DARK_TEXT, icon="•")

# Arrow 2
add_text_box(s, Inches(8.8), Inches(3.6), Inches(0.6), Inches(0.6),
             "→", font_size=36, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Level 3: Terms
card3 = add_rounded_rect(s, Inches(9.4), Inches(1.9), Inches(3.6), Inches(4.0), WHITE)
card3.shadow.inherit = False
add_shape_rect(s, Inches(9.4), Inches(1.9), Inches(3.6), Inches(0.1), ACCENT_PURPLE)
add_text_box(s, Inches(9.6), Inches(2.2), Inches(3.2), Inches(0.6),
             "🏷️  Terms", font_size=24, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(9.6), Inches(2.9), Inches(3.2), Inches(0.4),
             "Individual labels", font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
term_details = [
    "Unique ID + text labels",
    "Supports synonyms",
    "Multilingual labels possible",
    "E.g. NDA, MSA, SOW, Renewal",
]
add_bullet_frame(s, Inches(9.8), Inches(3.5), Inches(3.0), Inches(2.0),
                 term_details, font_size=13, color=DARK_TEXT, icon="•")

# Key insight at bottom
insight_box = add_rounded_rect(s, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.6), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(6.35), Inches(9.7), Inches(0.5),
             "💡  Access the term store:  SharePoint admin center → Content services → Term store",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Global vs Local Term Sets
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Clarify the scope difference: global term sets are tenant-wide, local term sets "
              "are scoped to a site. In our shared training tenant we prefer global term sets "
              "within participant-specific groups, but fall back to local if needed.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Global vs Local Term Sets", font_size=32, bold=True, color=DARK_BG)

# LEFT — Global
gl = add_rounded_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.2), WHITE)
gl.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.6),
             "🌐  Global Term Sets", font_size=24, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
global_items = [
    "Available across ALL sites in the tenant",
    "Created in the Term Store admin center",
    "Ideal for org-wide classification (departments, regions)",
    "Requires Term Store Admin or Contributor role",
    "Our approach: NW-Pxx-TermGroup in global scope",
]
add_bullet_frame(s, Inches(1.2), Inches(2.4), Inches(5.0), Inches(3.0),
                 global_items, font_size=15, color=DARK_TEXT, icon="▸")

# RIGHT — Local
lo = add_rounded_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.2), WHITE)
lo.shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(0.08), ORANGE)
add_text_box(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.6),
             "📍  Local Term Sets", font_size=24, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
local_items = [
    "Scoped to a single site (site collection)",
    "Created when adding a MM column to a list/library",
    "Only visible within that site collection",
    "No special admin role required (site owner can create)",
    "Our fallback if term group creation is restricted",
]
add_bullet_frame(s, Inches(7.3), Inches(2.4), Inches(5.0), Inches(3.0),
                 local_items, font_size=15, color=DARK_TEXT, icon="▸")

add_footer_bar(s, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Delegated Management Roles
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Three roles in the term store: Term Store Admin has full control, "
              "Group Manager can manage term sets within their group and assign contributors, "
              "Contributor can create/edit terms and term sets. "
              "Important: Owner/Contact/Stakeholders labels are informational only — they don't grant permissions.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Delegated Term Management Roles", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Three levels of control — delegate wisely",
             font_size=16, color=MID_GRAY)

roles = [
    ("👑", "Term Store Admin", ACCENT_BLUE,
     [
         "Create / delete term groups",
         "Add or remove other admins",
         "Assign group managers & contributors",
         "Change working languages",
         "All actions of lower roles",
     ]),
    ("🛡️", "Group Manager", ACCENT_TEAL,
     [
         "Manage term sets within their group",
         "Add or remove contributors",
         "All actions of contributor role",
         "Cannot create new groups",
         "",
     ]),
    ("✏️", "Contributor", ACCENT_PURPLE,
     [
         "Create / edit term sets and terms",
         "Within assigned group only",
         "Cannot manage roles or groups",
         "Lowest management privilege",
         "",
     ]),
]
for i, (icon, title, color, items) in enumerate(roles):
    x = Inches(0.6 + i * 4.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(3.85), Inches(4.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(3.85), Inches(0.08), color)
    add_text_box(s, x + Inches(0.2), Inches(2.0), Inches(3.45), Inches(0.6),
                 f"{icon}  {title}", font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    filtered = [b for b in items if b]
    add_bullet_frame(s, x + Inches(0.3), Inches(2.7), Inches(3.25), Inches(3.0),
                     filtered, font_size=14, color=DARK_TEXT, icon="•")

# Warning callout
warn = add_rounded_rect(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6), RGBColor(0xFD, 0xE7, 0xE9))
add_text_box(s, Inches(1.1), Inches(6.35), Inches(11.1), Inches(0.5),
             "⚠️  Note:  Owner, Contact, and Stakeholders labels on term sets are informational only — "
             "they do NOT grant term store permissions.",
             font_size=14, color=RED_ACCENT)

add_footer_bar(s, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — How to Assign Roles (Admin Steps)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the admin center steps for assigning each role. "
              "All roles are managed from SharePoint admin center → Content services → Term store.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Assigning Term Store Roles", font_size=32, bold=True, color=DARK_BG)

steps_data = [
    ("Add Term Store Admin", ACCENT_BLUE, [
        "1. Go to Term store page",
        "2. Select the taxonomy in tree view",
        "3. Under Admins → click Edit",
        "4. Enter names/email → Save",
    ]),
    ("Add Group Manager", ACCENT_TEAL, [
        "1. Select the target term group",
        "2. Go to People page",
        "3. Under Group Managers → Edit",
        "4. Enter names/email → Save",
    ]),
    ("Add Contributor", ACCENT_PURPLE, [
        "1. Select the target term group",
        "2. Go to People page",
        "3. Under Contributors → Edit",
        "4. Enter names/email → Save",
    ]),
]
for i, (title, color, steps) in enumerate(steps_data):
    x = Inches(0.6 + i * 4.15)
    card = add_rounded_rect(s, x, Inches(1.4), Inches(3.85), Inches(3.6), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.4), Inches(3.85), Inches(0.08), color)
    add_text_box(s, x + Inches(0.2), Inches(1.6), Inches(3.45), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(s, x + Inches(0.4), Inches(2.2), Inches(3.0), Inches(2.5),
                     steps, font_size=14, color=DARK_TEXT, icon="")

# Path reminder
path_box = add_rounded_rect(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.5), LIGHT_BLUE)
add_text_box(s, Inches(1.1), Inches(5.35), Inches(11.1), Inches(0.4),
             "📍  Path:  SharePoint admin center → Content services → Term store",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — What Site Users Can Do
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Not just admins — regular site members interact with metadata too. "
              "They can tag items, add enterprise keywords, use metadata navigation, "
              "and even create local term sets if they're site owners.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "What Can Site Users Do with Metadata?", font_size=30, bold=True, color=DARK_BG)

user_actions = [
    ("🏷️", "Tag Items", "Update managed metadata columns\n(if term set is open or allows fill-in)"),
    ("🔑", "Add Keywords", "Add enterprise keywords when\nthe keywords column is enabled"),
    ("🧭", "Navigate by Metadata", "Use metadata navigation to filter\nand browse list/library items"),
    ("🔍", "Refine Search", "Use managed terms to refine\nsearch results via the panel"),
    ("📋", "Create Local Sets", "Site owners can create local term\nsets when adding MM columns"),
    ("📝", "Contribute to Open Sets", "Add new terms to open term sets\nwhen tagging content"),
]
for i, (icon, title, desc) in enumerate(user_actions):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.5 + row * 2.6)
    card = add_rounded_rect(s, x, y, Inches(3.8), Inches(2.2), WHITE)
    card.shadow.inherit = False
    add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(1.0), y + Inches(0.2), Inches(2.5), Inches(0.4),
                 title, font_size=18, bold=True, color=ACCENT_BLUE)
    add_text_box(s, x + Inches(1.0), y + Inches(0.7), Inches(2.5), Inches(1.2),
                 desc, font_size=13, color=DARK_TEXT)

add_footer_bar(s, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Section Divider: Applying Metadata
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: now let's put theory into practice — creating and applying metadata.")
section_divider(s, "Applying Metadata",
                "From term store to library columns — hands-on", "⚙️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Creating a Term Group (Step by Step)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Step-by-step guide to creating a term group. "
              "Must be a Term Store Admin to do this. If restricted in a shared tenant, "
              "participants fall back to local term sets.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Creating a Term Group", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Prerequisite: you must be a Term Store Admin",
             font_size=16, color=MID_GRAY)

create_steps = [
    ("1", "Navigate", "SharePoint admin center →\nContent services → Term store", ACCENT_BLUE),
    ("2", "Add Group", "Click 'Add term group'\nin the right pane", ACCENT_TEAL),
    ("3", "Name It", "Enter group name\n(e.g. NW-Pxx-TermGroup) → Enter", ACCENT_PURPLE),
    ("4", "Configure", "Add description, assign\ngroup managers & contributors", ORANGE),
]
for i, (num, title, desc, color) in enumerate(create_steps):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    badge = add_rounded_rect(s, x + Inches(1.05), Inches(2.0), Inches(0.7), Inches(0.7), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(6)
    add_text_box(s, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.5), Inches(2.5), Inches(1.2),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Creating a Term Set & Adding Terms
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("After the group exists, create a term set and populate it with terms. "
              "Demonstrate open vs closed submission policy and the tagging toggle.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Creating a Term Set & Adding Terms", font_size=32, bold=True, color=DARK_BG)

# Left: Term Set creation
left = add_rounded_rect(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(4.8), WHITE)
left.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.4), Inches(5.6), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(1.0), Inches(1.6), Inches(5.2), Inches(0.5),
             "📋  Create Term Set", font_size=22, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
ts_steps = [
    "1. Expand the group → click 'Add term set'",
    "2. Type name (e.g. NW-Pxx-ContractType) → Enter",
    "3. General tab: set Owner, Contact, Stakeholders",
    "4. Usage settings: choose Open or Closed",
    "5. Enable 'Available for tagging'",
]
add_bullet_frame(s, Inches(1.2), Inches(2.3), Inches(5.0), Inches(3.5),
                 ts_steps, font_size=15, color=DARK_TEXT, icon="")

# Right: Adding Terms
right = add_rounded_rect(s, Inches(6.9), Inches(1.4), Inches(5.6), Inches(4.8), WHITE)
right.shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.4), Inches(5.6), Inches(0.08), ACCENT_PURPLE)
add_text_box(s, Inches(7.1), Inches(1.6), Inches(5.2), Inches(0.5),
             "🏷️  Add Terms", font_size=22, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
term_steps = [
    "1. Select the term set in tree view",
    "2. Click 'Add term'",
    "3. Type term name (e.g. NDA) → Enter",
    "4. Optional: add synonyms, translations",
    "5. Repeat for MSA, SOW, Renewal",
]
add_bullet_frame(s, Inches(7.3), Inches(2.3), Inches(5.0), Inches(3.5),
                 term_steps, font_size=15, color=DARK_TEXT, icon="")

add_footer_bar(s, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Adding a Managed Metadata Column
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("The bridge between term store and library: the Managed Metadata column. "
              "When users add or edit documents, they pick from the type-ahead term picker.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Adding a Managed Metadata Column", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Connect the term store to your library",
             font_size=16, color=MID_GRAY)

col_steps = [
    ("1", "Open Library Settings", "Go to NW-Pxx-Contracts →\nSettings → Library settings →\n'Create column'", ACCENT_BLUE),
    ("2", "Choose Type", "Select 'Managed Metadata'\nas the column type → name it\n'Contract Type'", ACCENT_TEAL),
    ("3", "Connect Term Set", "Browse or search for your\nterm set (NW-Pxx-ContractType)\nand select it", ACCENT_PURPLE),
    ("4", "Tag Documents", "Edit document properties →\nuse type-ahead picker to\nselect terms (NDA, MSA…)", ORANGE),
]
for i, (num, title, desc, color) in enumerate(col_steps):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.8), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    badge = add_rounded_rect(s, x + Inches(1.05), Inches(2.0), Inches(0.7), Inches(0.7), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(6)
    add_text_box(s, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.5), Inches(2.5), Inches(1.5),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Tip at bottom
tip = add_rounded_rect(s, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.6), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(6.05), Inches(9.7), Inches(0.5),
             "💡  Enterprise Keywords column can also be added — it allows free-form tagging alongside managed terms",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Section Divider: Shared Tenant Safety
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: in our shared training tenant, we need to be extra careful.")
section_divider(s, "Shared-Tenant Safety",
                "Keeping the training environment clean", "🤝")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Shared Tenant Rules
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("In a shared training tenant, every participant touches the same term store. "
              "Following these rules prevents cross-impact between participants.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ORANGE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🤝  Shared-Tenant Safety Rules", font_size=32, bold=True, color=DARK_BG)

rules = [
    ("✅", "Use NW-Pxx- prefix", "Always name your term groups, sets,\nand columns with your participant ID",
     "E.g. NW-P03-TermGroup, NW-P03-ContractType", GREEN),
    ("🚫", "Never edit others' taxonomy", "Don't modify, rename, or delete\nterm groups/sets from other participants",
     "Treat other participants' metadata as read-only", RED_ACCENT),
    ("🔄", "Fallback to local", "If you can't create a global term group,\nuse a local term set instead",
     "Create the term set when adding the MM column", ORANGE),
    ("🧹", "Clean up after labs", "Delete your test term groups/sets\nwhen instructed during cleanup",
     "Keeps the tenant tidy for future sessions", ACCENT_PURPLE),
]
for i, (icon, title, desc, detail, color) in enumerate(rules):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 2.6)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(2.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(2.2), color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.9), y + Inches(0.2), Inches(4.6), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(s, x + Inches(0.9), y + Inches(0.7), Inches(4.6), Inches(0.8),
                 desc, font_size=14, color=DARK_TEXT)
    add_text_box(s, x + Inches(0.9), y + Inches(1.6), Inches(4.6), Inches(0.4),
                 detail, font_size=12, color=MID_GRAY)

add_footer_bar(s, 21, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Northwind Contracts Scenario
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("The Northwind scenario ties everything together: participants create a taxonomy "
              "for contract types and apply it to the contracts library created in Module 4.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🏢  Scenario: Northwind Contracts Taxonomy", font_size=30, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Applying metadata to the contracts library from Module 4",
             font_size=16, color=MID_GRAY)

# Taxonomy tree visualization
tree_card = add_rounded_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5), WHITE)
tree_card.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.5),
             "📦 Taxonomy Structure", font_size=20, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

tree_lines = [
    "📦  NW-Pxx-TermGroup",
    "    └── 📋  NW-Pxx-ContractType",
    "            ├── 🏷️  NDA (Non-Disclosure Agreement)",
    "            ├── 🏷️  MSA (Master Service Agreement)",
    "            ├── 🏷️  SOW (Statement of Work)",
    "            └── 🏷️  Renewal",
]
add_bullet_frame(s, Inches(1.2), Inches(2.7), Inches(5.0), Inches(3.2),
                 tree_lines, font_size=15, color=DARK_TEXT, icon="",
                 font_name="Cascadia Code")

# Right side — how it's used
use_card = add_rounded_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.5), WHITE)
use_card.shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.5),
             "⚙️ How It's Applied", font_size=20, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

apply_items = [
    "Library: NW-Pxx-Contracts",
    "New column: 'Contract Type' (Managed Metadata)",
    "Connected to: NW-Pxx-ContractType term set",
    "Users pick terms via type-ahead picker",
    "Views: group contracts by type",
    "Search: refine results by contract type",
]
add_bullet_frame(s, Inches(7.3), Inches(2.7), Inches(5.0), Inches(3.2),
                 apply_items, font_size=15, color=DARK_TEXT, icon="▸")

add_footer_bar(s, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Section Divider: Lab & Validation
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: now for the hands-on lab and what to validate afterward.")
section_divider(s, "Lab & Validation",
                "Hands-on practice and validation checks", "🔬")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Lab Preview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Preview of Lab 5: participants will create a term group, term set, "
              "add terms, then add a managed metadata column to the contracts library "
              "and tag documents. Reinforces the library created in Module 4.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🔬  Lab 5 Preview", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Hands-on: Create taxonomy + apply managed metadata",
             font_size=16, color=MID_GRAY)

lab_tasks = [
    ("1", "Create Term Group", "Navigate to the Term Store\nand create NW-Pxx-TermGroup\n(or use local fallback)", ACCENT_BLUE),
    ("2", "Create Term Set", "Within your group, create\nNW-Pxx-ContractType\nwith Closed submission policy", ACCENT_TEAL),
    ("3", "Add Terms", "Add NDA, MSA, SOW,\nand Renewal as terms\nto your term set", ACCENT_PURPLE),
    ("4", "Add MM Column", "In NW-Pxx-Contracts library,\ncreate 'Contract Type'\nmanaged metadata column", ORANGE),
]
for i, (num, title, desc, color) in enumerate(lab_tasks):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    badge = add_rounded_rect(s, x + Inches(1.05), Inches(2.0), Inches(0.7), Inches(0.7), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(6)
    add_text_box(s, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.5), Inches(2.5), Inches(1.2),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Bonus step
bonus = add_rounded_rect(s, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.8), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(5.45), Inches(9.7), Inches(0.7),
             "📝  Bonus:  Tag uploaded contract documents with the new column and create a view "
             "grouped by Contract Type to see metadata-driven organization in action.",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 24, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Validation & Troubleshooting
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Common issues and how to resolve them. "
              "These are the top 5 things that trip up learners during the lab.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, RED_ACCENT)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🛠️  Validation & Troubleshooting", font_size=32, bold=True, color=DARK_BG)

issues = [
    ("Can't create term group", "You need Term Store Admin role. "
     "Fallback: create a local term set when adding the column.",
     "Check: ask trainer to verify your role", RED_ACCENT),
    ("Term set not visible in column picker", "Ensure 'Available for tagging' is enabled "
     "in the term set's Usage Settings tab.",
     "Check: Term store → select set → Usage settings", ORANGE),
    ("Users adding unexpected free-text terms", "Your term set is Open. "
     "Change to Closed if you need strict control.",
     "Check: Term set → Usage settings → Submission policy", ACCENT_BLUE),
    ("Column shows GUID instead of term name", "Replication delay or column mapping issue. "
     "Wait a few minutes and refresh. Re-check the term set binding.",
     "Check: library settings → column → term set mapping", ACCENT_PURPLE),
    ("Terms not appearing for other users", "Check that terms have 'Available for tagging' enabled "
     "and the user has at least read access to the site.",
     "Check: Term store → select term → Usage settings", ACCENT_TEAL),
]
for i, (problem, solution, check, color) in enumerate(issues):
    y = Inches(1.3 + i * 1.15)
    card = add_rounded_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, Inches(0.6), y, Inches(0.08), Inches(1.0), color)
    add_text_box(s, Inches(0.9), y + Inches(0.05), Inches(3.2), Inches(0.4),
                 problem, font_size=14, bold=True, color=color)
    add_text_box(s, Inches(4.2), y + Inches(0.05), Inches(4.8), Inches(0.9),
                 solution, font_size=12, color=DARK_TEXT)
    add_text_box(s, Inches(9.2), y + Inches(0.05), Inches(3.3), Inches(0.9),
                 check, font_size=11, color=MID_GRAY)

add_footer_bar(s, 25, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Summary / Key Takeaways
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Recap the module: IA principles, term store hierarchy, delegated roles, "
              "and practical application of managed metadata columns.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, GREEN)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📌  Key Takeaways", font_size=32, bold=True, color=DARK_BG)

takeaways = [
    ("🏗️", "Information Architecture", "Good IA makes content findable by concept, not by folder path.\n"
     "Metadata > folders for classification at scale."),
    ("🏷️", "Managed Metadata", "Use managed metadata for org-wide controlled vocabulary.\n"
     "Supports synonyms, hierarchy, and type-ahead picking."),
    ("🗂️", "Term Store Hierarchy", "Term Group → Term Set → Terms. Groups = security boundaries.\n"
     "Global sets for tenant-wide use, local for site-scoped."),
    ("👥", "Delegated Roles", "Term Store Admin > Group Manager > Contributor.\n"
     "Owner/Contact/Stakeholders are labels, not permissions."),
    ("⚙️", "Practical Application", "Add a Managed Metadata column to connect term store to library.\n"
     "Users tag documents via type-ahead; views and search benefit."),
]
for i, (icon, title, desc) in enumerate(takeaways):
    y = Inches(1.3 + i * 1.15)
    card = add_rounded_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), WHITE)
    card.shadow.inherit = False
    add_text_box(s, Inches(0.8), y + Inches(0.1), Inches(0.6), Inches(0.5),
                 icon, font_size=26, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), y + Inches(0.1), Inches(2.5), Inches(0.4),
                 title, font_size=17, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(4.2), y + Inches(0.1), Inches(8.3), Inches(0.8),
                 desc, font_size=14, color=DARK_TEXT)

add_footer_bar(s, 26, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Knowledge Check
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("5 questions to check understanding. Encourage discussion rather than quick answers.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "❓  Knowledge Check", font_size=32, bold=True, color=DARK_BG)

questions = [
    'What is the difference between a "term set" and a "term group"?',
    "When would you prefer a managed metadata column over a standard Choice column?",
    "Who can create a new term group in the term store?",
    "What does 'delegated term management' mean in practice?",
    "What shared-tenant behaviors should you avoid when working with the term store?",
]
for i, q in enumerate(questions):
    y = Inches(1.4 + i * 1.1)
    qcard = add_rounded_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.9), WHITE)
    qcard.shadow.inherit = False
    badge = add_rounded_rect(s, Inches(0.8), y + Inches(0.15), Inches(0.55), Inches(0.55), ACCENT_PURPLE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, Inches(1.6), y + Inches(0.15), Inches(10.8), Inches(0.6),
                 q, font_size=16, color=DARK_TEXT)

add_footer_bar(s, 27, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Thank You / Next Module
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Module 5 complete. Next up: Module 6 — Configuring the Search Experience.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "✅  Module 5 Complete!", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
             "Managing Metadata and the Term Store",
             font_size=24, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(4.4), Inches(2.3), Inches(0.04), ACCENT_PURPLE)

add_text_box(s, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
             "Up Next  →  Module 6: Configuring the Search Experience",
             font_size=20, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(2), Inches(5.8), Inches(9), Inches(0.6),
             "🔍 Discover how metadata powers search refiners and discovery",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(__file__), "Module-05-Slides.pptx")
prs.save(out_path)
print(f"✅  Saved {TOTAL_SLIDES}-slide presentation → {out_path}")
