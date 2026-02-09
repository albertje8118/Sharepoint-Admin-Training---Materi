"""
Generate Module 3 PPTX — Working with Site Collections (Modern Sites)
Modern, engaging design matching Module 1 & 2 style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG     = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEC, 0xF9)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT    = RGBColor(0x24, 0x24, 0x24)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT  = RGBColor(0xD1, 0x34, 0x38)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ═══════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════
def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6), icon="▸", font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon} {bullet}"; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.name = font_name; p.space_after = spacing
    return txBox

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)

def add_footer_bar(slide, slide_num, total, module_label="Module 3"):
    bar_top = SLIDE_HEIGHT - Inches(0.45)
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_WIDTH, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), bar_top + Inches(0.05), Inches(7), Inches(0.35),
                 f"{module_label}  |  Working with Site Collections (Modern Sites)",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), bar_top + Inches(0.05),
                 Inches(1.2), Inches(0.35), f"{slide_num} / {total}",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)

def section_divider(slide, section_title, subtitle="", icon=""):
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_TEAL)
    if icon:
        add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                     icon, font_size=60, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                 section_title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

TOTAL_SLIDES = 28
slide_counter = [0]

def new_slide(notes=""):
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if notes: add_speaker_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome to Module 3. This is where participants get hands-on with their own practice sites.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "MODERN SHAREPOINT ONLINE FOR ADMINISTRATORS", font_size=18, bold=True,
             color=ACCENT_TEAL, font_name="Segoe UI Semibold")
add_text_box(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
             "Working with\nSite Collections", font_size=48, bold=True,
             color=WHITE, font_name="Segoe UI Light")
add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Day 1  ·  Tenant Foundations & Site Management", font_size=18,
             color=RGBColor(0xAA, 0xAA, 0xAA))
badge = add_rounded_rect(s, Inches(1), Inches(4.7), Inches(3.2), Inches(0.55), ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(4.75), Inches(2.8), Inches(0.45),
             "MODULE 3", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "🏗️  Create, manage, and administer modern SharePoint sites",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88))

add_footer_bar(s, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Module Agenda
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the module agenda. This is the most hands-on module so far.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "📋  Module 3 Agenda", font_size=30, bold=True, color=DARK_BG)

agenda = [
    ("1", "SharePoint Sites Overview — Types, when to use each", ACCENT_BLUE),
    ("2", "M365 Groups & Teams Integration — Behind the scenes", ACCENT_TEAL),
    ("3", "Hub Sites — Organising your site architecture", ACCENT_PURPLE),
    ("4", "Creating Sites — Admin center & PowerShell", ACCENT_BLUE),
    ("5", "Site Admin Operations — Membership, access, recycle bin", ACCENT_TEAL),
    ("6", "Storage Management — Tenant vs site-level quotas", ACCENT_PURPLE),
    ("7", "Site Lifecycle — Delete, restore, and retention", ORANGE),
    ("8", "PowerShell Basics — Connect and inspect sites", GREEN),
    ("9", "Lab 3 Preview & Knowledge Check", RGBColor(0xD1, 0x34, 0x38)),
]

y_pos = Inches(1.2)
for num, text, color in agenda:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.45), Inches(0.45))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    tf = circ.text_frame; p = tf.paragraphs[0]
    p.text = num; p.font.size = Pt(14); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Segoe UI"
    add_text_box(s, Inches(1.45), y_pos + Inches(0.02), Inches(10), Inches(0.4),
                 text, font_size=16, color=DARK_TEXT)
    y_pos += Inches(0.58)

add_footer_bar(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Module Objectives
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Review learning objectives. This module has the most lab tasks.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "🎯  Module Objectives", font_size=30, bold=True, color=DARK_BG)

objectives = [
    "Create a modern SharePoint site from the SharePoint admin center",
    "Identify where site ownership, admins, and settings are managed",
    "Perform day-to-day site admin tasks (membership, access requests, recycle bin)",
    "Explain and observe how storage limits work at tenant and site level",
    "Delete and restore a site safely using a test site",
    "Connect to SharePoint Online using PowerShell and retrieve site properties",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(5),
                 objectives, font_size=18, color=DARK_TEXT, icon="✅")

add_footer_bar(s, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Section Divider: What are SharePoint Sites?
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to the overview of SharePoint site types.")
section_divider(s, "SharePoint Sites Overview",
                "Understanding site types, architecture, and when to use each", "🏗️")
add_footer_bar(s, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What is a SharePoint Site?
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("A site is the top-level manageable entity in SharePoint. "
              "It has its own URL, owners, permissions, storage, and sharing config. "
              "Formerly called 'site collections' — now just 'sites' in modern terminology.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "What is a SharePoint Site?", font_size=28, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.7),
             "A site (formerly 'site collection') is the top-level manageable unit in SharePoint Online. "
             "It's the container for pages, document libraries, lists, and sub-resources.",
             font_size=16, color=MID_GRAY)

site_props = [
    ("🔗 Unique URL", "Each site gets its own URL\nunder your tenant domain", ACCENT_BLUE),
    ("👷 Owner / Admins", "Site-scoped administration\nand membership control", ACCENT_TEAL),
    ("🔐 Permissions", "SharePoint groups with\ninheritance model", ACCENT_PURPLE),
    ("📦 Storage", "Consumes tenant pool or\nhas a manual quota", ORANGE),
    ("🌍 Sharing Config", "Site-level sharing posture\n(same or more restrictive than org)", GREEN),
]

x_pos = Inches(0.3)
for title, desc, color in site_props:
    card = add_rounded_rect(s, x_pos, Inches(2.3), Inches(2.4), Inches(2.5), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(2.3), Inches(2.4), Inches(0.06), color)
    add_text_box(s, x_pos + Inches(0.15), Inches(2.55), Inches(2.1), Inches(0.5),
                 title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x_pos + Inches(0.15), Inches(3.15), Inches(2.1), Inches(1.3),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(2.6)

add_text_box(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
             "💡 Many governance and security controls are scoped to the site level — this is why site management matters.",
             font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Site Types: Team vs Communication
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Two primary modern site types: Team Sites (collaboration, groups-backed) "
              "and Communication Sites (broadcasting, not group-backed by default).")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Modern Site Types", font_size=28, bold=True, color=DARK_BG)

# Team Site
ts_card = add_rounded_rect(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(4.5), NEAR_WHITE)
add_shape_rect(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.1), ACCENT_BLUE)
add_text_box(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
             "👥 Team Site", font_size=22, bold=True, color=ACCENT_BLUE)
add_text_box(s, Inches(0.8), Inches(2.2), Inches(5.2), Inches(0.4),
             "Purpose: Team collaboration", font_size=14, bold=True, color=DARK_TEXT)
team_features = [
    "Backed by a Microsoft 365 Group",
    "Can be connected to a Microsoft Team",
    "Shared mailbox, calendar, Planner, OneNote",
    "Members collaborate; membership via group",
    "Default: private (members only)",
    "Best for: project teams, departments",
]
add_bullet_frame(s, Inches(0.8), Inches(2.7), Inches(5.2), Inches(3),
                 team_features, font_size=13, color=DARK_TEXT, icon="•")

# Communication Site
cs_card = add_rounded_rect(s, Inches(7), Inches(1.3), Inches(5.8), Inches(4.5), NEAR_WHITE)
add_shape_rect(s, Inches(7), Inches(1.3), Inches(5.8), Inches(0.1), ACCENT_TEAL)
add_text_box(s, Inches(7.3), Inches(1.6), Inches(5.2), Inches(0.5),
             "📢 Communication Site", font_size=22, bold=True, color=ACCENT_TEAL)
add_text_box(s, Inches(7.3), Inches(2.2), Inches(5.2), Inches(0.4),
             "Purpose: Broadcast & publish", font_size=14, bold=True, color=DARK_TEXT)
comm_features = [
    "NOT backed by a Microsoft 365 Group (by default)",
    "Cannot natively connect to Teams",
    "Beautiful page layouts and templates",
    "Permissions via SharePoint groups (classic model)",
    "Default: public within org (everyone can read)",
    "Best for: intranets, news, announcements",
]
add_bullet_frame(s, Inches(7.3), Inches(2.7), Inches(5.2), Inches(3),
                 comm_features, font_size=13, color=DARK_TEXT, icon="•")

add_text_box(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.4),
             "🎯 Choose Team Site for collaboration, Communication Site for broadcasting.",
             font_size=15, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — M365 Groups & Teams Integration
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("When you create a Team Site, an M365 Group is created. "
              "When you create a Teams team, it creates an M365 Group + Team Site. "
              "This bidirectional relationship means managing one affects the other.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Microsoft 365 Groups & Teams — Behind the Scenes",
             font_size=26, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Team Sites and Teams are connected through Microsoft 365 Groups. Understanding this relationship is crucial.",
             font_size=16, color=MID_GRAY)

# Flow diagram
arrows_data = [
    ("Create a Team Site\n(SP Admin Center)", "→", "M365 Group\nis auto-created", "→",
     "Can optionally add\na Teams team later"),
    ("Create a Teams team\n(Teams Admin/User)", "→", "M365 Group +\nSP Team Site created", "→",
     "Site appears in\nSP Admin Center"),
]

y_base = Inches(2.0)
for row_idx, (step1, arr1, step2, arr2, step3) in enumerate(arrows_data):
    y = y_base + Inches(row_idx * 2.2)
    color = ACCENT_BLUE if row_idx == 0 else ACCENT_TEAL
    for col_idx, (text, w) in enumerate([(step1, 3.2), (arr1, 0.6), (step2, 3.2), (arr2, 0.6), (step3, 3.2)]):
        x = Inches(0.5) + sum([Inches(x) for x in [3.2, 0.6, 3.2, 0.6, 3.2][:col_idx]]) + Inches(col_idx * 0.15)
        if text in ("→",):
            add_text_box(s, x, y + Inches(0.3), Inches(w), Inches(0.5),
                         "→", font_size=28, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        else:
            card = add_rounded_rect(s, x, y, Inches(w), Inches(1.3), NEAR_WHITE)
            add_shape_rect(s, x, y, Inches(w), Inches(0.06), color)
            add_text_box(s, x + Inches(0.15), y + Inches(0.2), Inches(w - 0.3), Inches(1),
                         text, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Warning
warn = add_rounded_rect(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
                        RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(6.22), Inches(11), Inches(0.45),
             "⚠️ Deleting a Teams-connected SP site can break the Team. Deleting a Team deletes the SP site.",
             font_size=13, bold=True, color=ORANGE)

add_footer_bar(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Hub Sites
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Hub sites are the organising layer. They provide shared navigation, branding, "
              "and search scope across associated sites. They don't affect permissions.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Hub Sites — Organising Your Architecture", font_size=28, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Hub sites group related sites under a shared brand, navigation, and search experience — "
             "without changing permissions.",
             font_size=16, color=MID_GRAY)

hub_features = [
    ("🧭 Shared Navigation", "Associated sites inherit the hub's\nglobal navigation bar", ACCENT_BLUE),
    ("🎨 Common Branding", "Theme, logo, and header\napplied across all associated sites", ACCENT_TEAL),
    ("🔍 Scoped Search", "Search can be scoped to all\nsite content across the hub", ACCENT_PURPLE),
    ("📰 News Rollup", "Aggregate news posts from\nall associated sites in one view", ORANGE),
]

x_pos = Inches(0.4)
for title, desc, color in hub_features:
    card = add_rounded_rect(s, x_pos, Inches(2.0), Inches(3), Inches(2.5), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(2.0), Inches(3), Inches(0.08), color)
    add_text_box(s, x_pos + Inches(0.2), Inches(2.3), Inches(2.6), Inches(0.5),
                 title, font_size=15, bold=True, color=color)
    add_text_box(s, x_pos + Inches(0.2), Inches(2.9), Inches(2.6), Inches(1.3),
                 desc, font_size=13, color=MID_GRAY)
    x_pos += Inches(3.25)

key_points = [
    "Hub sites do NOT affect permissions — each associated site keeps its own permissions",
    "A site can associate (join) or dissociate (leave) a hub at any time",
    "Hub registration is done by SharePoint admins; association can be delegated",
]
add_bullet_frame(s, Inches(0.8), Inches(4.8), Inches(11), Inches(1.5),
                 key_points, font_size=14, color=DARK_TEXT, icon="📌")

add_footer_bar(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — When to Use Which Site Type (Decision Guide)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("A quick decision matrix to help admins guide site creation requests.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "When to Use Which? — Quick Decision Guide", font_size=28, bold=True, color=DARK_BG)

# Table
rows, cols = 5, 4
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.0))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.5)
tbl.columns[1].width = Inches(3)
tbl.columns[2].width = Inches(3)
tbl.columns[3].width = Inches(2.8)

headers = ["Scenario", "Team Site", "Communication Site", "Hub Site"]
data = [
    ["Project team collaboration", "✅ Best fit", "❌ Not ideal", "Associate to hub"],
    ["Company intranet / news", "❌ Not ideal", "✅ Best fit", "Register as hub"],
    ["Department portal", "✅ For team work", "✅ For broadcasting", "Associate to hub"],
    ["Organise 10+ related sites", "N/A", "N/A", "✅ Register hub"],
]

for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_BLUE

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = tbl.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.color.rgb = DARK_TEXT
            p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEAR_WHITE if i % 2 == 0 else WHITE

add_footer_bar(s, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Section Divider: Creating & Managing Sites
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to hands-on site creation and management.")
section_divider(s, "Creating & Managing Sites",
                "Admin center workflow and day-to-day operations", "⚙️")
add_footer_bar(s, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Creating a Site from Admin Center
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the admin center flow: Active sites → Create → Choose type → Configure.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Creating a Site — Admin Center Workflow", font_size=28, bold=True, color=DARK_BG)

steps = [
    ("1", "Open Active Sites", "SharePoint admin center >\nSites > Active sites", ACCENT_BLUE),
    ("2", "Click Create", "Select '+ Create' to start\nthe site provision wizard", ACCENT_TEAL),
    ("3", "Choose Site Type", "Team site or\nCommunication site", ACCENT_PURPLE),
    ("4", "Configure Details", "Name, URL, owner, language,\ntime zone, storage", ORANGE),
    ("5", "Site is Provisioned", "Site appears in Active sites\nwithin seconds", GREEN),
]

x_pos = Inches(0.3)
for num, title, desc, color in steps:
    card = add_rounded_rect(s, x_pos, Inches(1.5), Inches(2.35), Inches(3), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(1.5), Inches(2.35), Inches(0.06), color)
    nbadge = s.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(0.8), Inches(1.7),
                                 Inches(0.6), Inches(0.6))
    nbadge.fill.solid(); nbadge.fill.fore_color.rgb = color; nbadge.line.fill.background()
    tf = nbadge.text_frame; p = tf.paragraphs[0]
    p.text = num; p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Segoe UI"
    add_text_box(s, x_pos + Inches(0.1), Inches(2.5), Inches(2.15), Inches(0.4),
                 title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x_pos + Inches(0.1), Inches(3.0), Inches(2.15), Inches(1.2),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(2.55)

tip = add_rounded_rect(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.05), Inches(11), Inches(0.6),
             "💡 UI can vary between tenants and over time. Learn to find the Create entry point, "
             "not memorize exact clicks.",
             font_size=14, color=ACCENT_BLUE)

add_footer_bar(s, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Site Details Panel
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("The site details panel provides General, Membership, and Settings tabs. "
              "This is where admins check and manage site-level properties.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Site Details Panel — What to Look For", font_size=28, bold=True, color=DARK_BG)

panels = [
    ("📋 General", [
        "Site URL and name",
        "Primary owner / admin",
        "Storage usage (current / limit)",
        "Hub association status",
        "Last activity date",
    ], ACCENT_BLUE),
    ("👥 Membership", [
        "Site admins (collection admins)",
        "Owners, Members, Visitors",
        "Microsoft 365 Group members (if applicable)",
        "Quick add/remove access",
    ], ACCENT_TEAL),
    ("⚙️ Settings", [
        "External sharing configuration",
        "Default sharing link type",
        "Storage limit (if manual mode)",
        "Conditional Access policy (if set)",
    ], ACCENT_PURPLE),
]

x_pos = Inches(0.4)
for title, items, color in panels:
    card = add_rounded_rect(s, x_pos, Inches(1.3), Inches(3.9), Inches(4.5), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(1.3), Inches(3.9), Inches(0.1), color)
    add_text_box(s, x_pos + Inches(0.3), Inches(1.6), Inches(3.3), Inches(0.5),
                 title, font_size=18, bold=True, color=color)
    add_bullet_frame(s, x_pos + Inches(0.3), Inches(2.2), Inches(3.3), Inches(3.3),
                     items, font_size=13, color=DARK_TEXT, icon="•")
    x_pos += Inches(4.15)

add_text_box(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.4),
             "📌 Always confirm you're looking at the correct site before making changes!",
             font_size=14, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Day-to-Day Site Admin Operations
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Common quick-win tasks: membership changes, access requests, recycle bin recovery. "
              "These are the 'small but urgent' issues admins handle daily.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Day-to-Day Site Admin Operations", font_size=28, bold=True, color=DARK_BG)

ops = [
    ("👥 Membership / Site Admins", "Add or remove additional site admins for break-glass access. "
     "Common when site owners leave or need a backup admin.", ACCENT_BLUE),
    ("🙋 Access Requests", "Verify whether access requests are enabled and who receives them. "
     "Route requests to the right person to avoid bottlenecks.", ACCENT_TEAL),
    ("🗑️ Recycle Bin Recovery", "Restore deleted files or pages when users make mistakes. "
     "Two-stage recycle bin: site-level → site collection-level.", ACCENT_PURPLE),
    ("🔍 Monthly Access Review", "Review who has access on a regular cadence (monthly recommended). "
     "Check Owners/Members/Visitors + Activity signals.", ORANGE),
]

y_pos = Inches(1.3)
for title, desc, color in ops:
    card = add_rounded_rect(s, Inches(0.5), y_pos, Inches(12.3), Inches(1.1), NEAR_WHITE)
    add_shape_rect(s, Inches(0.5), y_pos, Inches(0.12), Inches(1.1), color)
    add_text_box(s, Inches(0.9), y_pos + Inches(0.08), Inches(5), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(0.9), y_pos + Inches(0.5), Inches(11.5), Inches(0.55),
                 desc, font_size=13, color=MID_GRAY)
    y_pos += Inches(1.25)

add_footer_bar(s, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Section Divider: Storage Management
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to storage management.")
section_divider(s, "Storage Management",
                "Understanding tenant pool vs per-site quotas", "📦")
add_footer_bar(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Storage: Tenant vs Site Level
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Two models: Automatic (pooled) vs Manual (per-site limits). "
              "Most tenants default to Automatic. Changing this is Trainer-only in our lab.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Storage Controls: Two Models", font_size=28, bold=True, color=DARK_BG)

# Two columns
auto_card = add_rounded_rect(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(3.8), NEAR_WHITE)
add_shape_rect(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.1), ACCENT_BLUE)
add_text_box(s, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
             "🔄 Automatic (Pooled)", font_size=20, bold=True, color=ACCENT_BLUE)
auto_pts = [
    "Default for most tenants",
    "Storage shared across all sites",
    "SharePoint manages allocation",
    "No per-site limits to configure",
    "Simple but less granular control",
]
add_bullet_frame(s, Inches(0.8), Inches(2.2), Inches(5.2), Inches(2.5),
                 auto_pts, font_size=14, color=DARK_TEXT, icon="•")

man_card = add_rounded_rect(s, Inches(7), Inches(1.3), Inches(5.8), Inches(3.8), NEAR_WHITE)
add_shape_rect(s, Inches(7), Inches(1.3), Inches(5.8), Inches(0.1), ACCENT_TEAL)
add_text_box(s, Inches(7.3), Inches(1.6), Inches(5.2), Inches(0.5),
             "⚙️ Manual (Per-Site Limits)", font_size=20, bold=True, color=ACCENT_TEAL)
man_pts = [
    "Admin sets max GB per site",
    "Storage warnings configurable",
    "More granular governance",
    "Requires ongoing management",
    "Useful for large organisations",
]
add_bullet_frame(s, Inches(7.3), Inches(2.2), Inches(5.2), Inches(2.5),
                 man_pts, font_size=14, color=DARK_TEXT, icon="•")

# Storage formula
formula = add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.55), Inches(11), Inches(0.7),
             "📏 Total tenant storage = 1 TB base + 10 GB × number of licensed users  |  "
             "Max per site = 25 TB",
             font_size=15, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Section Divider: Site Lifecycle
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to delete and restore.")
section_divider(s, "Site Lifecycle",
                "Delete, restore, and the safety model", "♻️")
add_footer_bar(s, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Delete and Restore
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Deleting a site removes access. It goes to Deleted sites for ~93 days. "
              "After that, it's permanently deleted. Emphasize using a test site for the drill.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Site Delete & Restore — Safety Model", font_size=28, bold=True, color=DARK_BG)

lifecycle_steps = [
    ("1", "Active Site", "Site is live and\naccessible to users", GREEN),
    ("→", "", "", MID_GRAY),
    ("2", "Delete", "Admin deletes from\nActive sites list", RED_ACCENT),
    ("→", "", "", MID_GRAY),
    ("3", "Deleted Sites", "Recoverable for ~93 days\nfrom Deleted sites", ORANGE),
    ("→", "", "", MID_GRAY),
    ("4", "Restore / Purge", "Restore to Active sites\nor permanently delete", ACCENT_BLUE),
]

x_pos = Inches(0.3)
for num, title, desc, color in lifecycle_steps:
    if num == "→":
        add_text_box(s, x_pos, Inches(2.8), Inches(0.5), Inches(0.5),
                     "→", font_size=28, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(0.55)
    else:
        card = add_rounded_rect(s, x_pos, Inches(1.8), Inches(2.5), Inches(2.7), NEAR_WHITE)
        add_shape_rect(s, x_pos, Inches(1.8), Inches(2.5), Inches(0.06), color)
        nbadge = s.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(0.85), Inches(1.95),
                                     Inches(0.65), Inches(0.65))
        nbadge.fill.solid(); nbadge.fill.fore_color.rgb = color; nbadge.line.fill.background()
        tf = nbadge.text_frame; p = tf.paragraphs[0]
        p.text = num; p.font.size = Pt(18); p.font.bold = True
        p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Segoe UI"
        add_text_box(s, x_pos + Inches(0.1), Inches(2.8), Inches(2.3), Inches(0.4),
                     title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(s, x_pos + Inches(0.1), Inches(3.3), Inches(2.3), Inches(1),
                     desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(2.7)

# Warnings
warnings = [
    "🚫 Do NOT delete the organization's root site",
    "🚫 Do NOT delete any site you do not own",
    "✅ Use a dedicated NW-Pxx-RestoreTest site for the delete/restore drill",
]
add_bullet_frame(s, Inches(0.8), Inches(5.0), Inches(11), Inches(1.5),
                 warnings, font_size=15, color=DARK_TEXT, icon="")

add_footer_bar(s, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Recycle Bin Deep Dive
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Two-stage recycle bin: First stage (site level, user-accessible) → "
              "Second stage (site collection level, admin-accessible). Items auto-purge after 93 days total.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Recycle Bin — Two-Stage Recovery", font_size=28, bold=True, color=DARK_BG)

# Stage 1
s1_card = add_rounded_rect(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(3), NEAR_WHITE)
add_shape_rect(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.1), GREEN)
add_text_box(s, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5),
             "🗑️ First Stage (Site Recycle Bin)", font_size=18, bold=True, color=GREEN)
s1_pts = [
    "Accessible by site users and admins",
    "Deleted files, pages, list items",
    "Items stay for 93 days (total for both stages)",
    "Users can self-service restore",
]
add_bullet_frame(s, Inches(0.8), Inches(2.2), Inches(5), Inches(2),
                 s1_pts, font_size=14, color=DARK_TEXT, icon="•")

# Arrow
add_text_box(s, Inches(6.1), Inches(2.5), Inches(0.8), Inches(0.5),
             "→", font_size=36, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Stage 2
s2_card = add_rounded_rect(s, Inches(7), Inches(1.3), Inches(5.5), Inches(3), NEAR_WHITE)
add_shape_rect(s, Inches(7), Inches(1.3), Inches(5.5), Inches(0.1), ORANGE)
add_text_box(s, Inches(7.3), Inches(1.6), Inches(5), Inches(0.5),
             "🗑️ Second Stage (Site Collection)", font_size=18, bold=True, color=ORANGE)
s2_pts = [
    "Only accessible by site collection admins",
    "Items deleted from first-stage recycle bin",
    "Last chance before permanent deletion",
    "Admin recovery only (not self-service)",
]
add_bullet_frame(s, Inches(7.3), Inches(2.2), Inches(5), Inches(2),
                 s2_pts, font_size=14, color=DARK_TEXT, icon="•")

# Timeline
timeline = add_rounded_rect(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(4.85), Inches(11), Inches(0.6),
             "⏱️ Total retention: 93 days across both stages. After 93 days, items are permanently deleted.",
             font_size=15, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Section Divider: PowerShell
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to PowerShell basics for site inspection.")
section_divider(s, "PowerShell for Site Admin",
                "Connect, inspect, and manage sites via command line", "💻")
add_footer_bar(s, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — PowerShell: Connect & Inspect
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Show the essential workflow: Install module → Connect → Get sites. "
              "Emphasize read-only commands first, target only own sites.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "PowerShell Basics — Connect & Inspect", font_size=28, bold=True, color=DARK_BG)

# Code blocks
code_blocks = [
    ("Step 1: Connect to SharePoint Online",
     "Connect-SPOService -Url https://contoso-admin.sharepoint.com",
     "Uses your admin credentials to connect to the SP admin service.", ACCENT_BLUE),
    ("Step 2: List/Get Site Properties",
     "Get-SPOSite -Identity https://contoso.sharepoint.com/sites/NW-P01-ProjectSite",
     "Retrieves details for a specific site: URL, owner, storage, template.", ACCENT_TEAL),
    ("Step 3: Get All Sites (use with caution)",
     "Get-SPOSite -Limit All | Select-Object Url, Owner, StorageUsageCurrent",
     "Lists all sites. In shared tenant, just look — don't change.", ACCENT_PURPLE),
]

y_pos = Inches(1.2)
for title, code, desc, color in code_blocks:
    add_text_box(s, Inches(0.8), y_pos, Inches(10), Inches(0.35),
                 title, font_size=14, bold=True, color=color)
    # code box
    code_bg = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.35), Inches(11.5), Inches(0.55),
                               RGBColor(0x28, 0x28, 0x28))
    add_text_box(s, Inches(1), y_pos + Inches(0.38), Inches(11), Inches(0.5),
                 code, font_size=13, color=RGBColor(0x7E, 0xD3, 0x21), font_name="Cascadia Code")
    add_text_box(s, Inches(0.8), y_pos + Inches(0.95), Inches(11), Inches(0.3),
                 desc, font_size=12, color=MID_GRAY)
    y_pos += Inches(1.4)

# Safety note
safety = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.1), Inches(11.5), Inches(0.7),
                          RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, Inches(1), y_pos + Inches(0.15), Inches(11), Inches(0.6),
             "⚠️ In this course: read-only commands only. Target only YOUR NW-Pxx-... sites. "
             "Never run Set/Remove on other participants' sites.",
             font_size=13, bold=True, color=ORANGE)

add_footer_bar(s, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Additional PowerShell Cmdlets
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Quick reference of key cmdlets for site lifecycle management.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Key PowerShell Cmdlets — Quick Reference", font_size=28, bold=True, color=DARK_BG)

# Table
rows2, cols2 = 7, 3
tbl2_shape = s.shapes.add_table(rows2, cols2, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5))
tbl2 = tbl2_shape.table
tbl2.columns[0].width = Inches(3.5)
tbl2.columns[1].width = Inches(5)
tbl2.columns[2].width = Inches(3.8)

cmd_headers = ["Cmdlet", "Purpose", "Notes"]
cmd_data = [
    ["Connect-SPOService", "Connect to SharePoint admin service", "Required before any other cmdlet"],
    ["Get-SPOSite", "Retrieve site properties", "Read-only; safe to run"],
    ["New-SPOSite", "Create a new site (classic)", "Modern sites: use admin center"],
    ["Set-SPOSite", "Modify site properties", "⚠️ Use carefully in shared tenant"],
    ["Remove-SPOSite", "Delete a site", "⚠️ Lab: only NW-Pxx-RestoreTest"],
    ["Restore-SPODeletedSite", "Restore from Deleted sites", "Restores within 93-day window"],
]

for j, h in enumerate(cmd_headers):
    cell = tbl2.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_BLUE

for i, row_data in enumerate(cmd_data):
    for j, val in enumerate(row_data):
        cell = tbl2.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12); p.font.color.rgb = DARK_TEXT
            p.font.name = "Segoe UI"
            if j == 0: p.font.name = "Cascadia Code"
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEAR_WHITE if i % 2 == 0 else WHITE

add_footer_bar(s, 21, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Site Governance Best Practices
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Share governance best practices that connect modules 1-3 together.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Site Governance Best Practices", font_size=28, bold=True, color=DARK_BG)

practices = [
    ("📛 Naming Convention", "Enforce consistent naming (e.g., NW-Dept-Purpose) for discoverability and admin efficiency"),
    ("👤 Ownership Policy", "Every site MUST have at least 2 owners. Orphan sites are a governance risk"),
    ("📦 Storage Monitoring", "Review storage trends monthly. Set alerts before sites hit limits"),
    ("🔐 Sharing Posture", "Set site-level sharing tighter than org default for sensitive content"),
    ("🗑️ Lifecycle Policy", "Define when inactive sites should be archived or deleted. Use activity signals"),
    ("📋 Regular Access Reviews", "Review site membership monthly. Remove stale access. Document findings"),
]

y_pos = Inches(1.2)
for title, desc in practices:
    add_text_box(s, Inches(0.8), y_pos, Inches(4), Inches(0.35),
                 title, font_size=15, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(0.8), y_pos + Inches(0.35), Inches(11.5), Inches(0.35),
                 desc, font_size=13, color=MID_GRAY)
    y_pos += Inches(0.8)

add_footer_bar(s, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Scenario Rules Reminder
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Remind participants about shared-tenant rules for this module's lab.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🏢 Lab Scenario & Shared Tenant Rules", font_size=28, bold=True, color=DARK_BG)

rules = [
    "Scenario: Project Northwind Intranet Modernization",
    "Create your own practice sites using NW-Pxx-... naming only",
    "Persistent site: NW-Pxx-ProjectSite (keep this for later modules)",
    "Disposable site: NW-Pxx-RestoreTest (for delete/restore drill)",
    "Tenant-wide changes (e.g., storage mode) are TRAINER-ONLY",
    "PowerShell: read-only commands; target only YOUR sites",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(3.5),
                 rules, font_size=17, color=DARK_TEXT, icon="📋")

warn = add_rounded_rect(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.05), Inches(11), Inches(0.6),
             "🔒 NW-Pxx-ProjectSite is your foundation for Modules 4–8. Don't delete it!",
             font_size=14, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 23, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Lab 3 Preview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the lab tasks. This is the busiest lab so far — 7 tasks.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🔬 Lab 3 Preview — Managing Site Collections", font_size=26, bold=True, color=DARK_BG)

lab_tasks = [
    ("Task 1", "Create NW-Pxx-ProjectSite", "Your persistent practice site for the rest of the course"),
    ("Task 2", "Explore site details panel", "Capture owner, membership, activity, storage observations"),
    ("Task 3", "Check access request settings", "Observe who receives requests; document findings"),
    ("Task 4", "Monthly access review drill", "Review site membership and document gaps"),
    ("Task 5", "Recycle bin restore drill", "Delete a test file then restore it from recycle bin"),
    ("Task 6", "Delete/restore test site", "Create NW-Pxx-RestoreTest → delete → restore"),
    ("Task 7", "PowerShell connection", "Connect-SPOService + Get-SPOSite on your site"),
]

y_pos = Inches(1.1)
for task_id, task_title, task_desc in lab_tasks:
    card = add_rounded_rect(s, Inches(0.5), y_pos, Inches(12.3), Inches(0.68), NEAR_WHITE)
    tbadge = add_rounded_rect(s, Inches(0.7), y_pos + Inches(0.1), Inches(1), Inches(0.35), ACCENT_BLUE)
    add_text_box(s, Inches(0.72), y_pos + Inches(0.12), Inches(0.95), Inches(0.3),
                 task_id, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.9), y_pos + Inches(0.05), Inches(4), Inches(0.3),
                 task_title, font_size=14, bold=True, color=DARK_TEXT)
    add_text_box(s, Inches(6), y_pos + Inches(0.08), Inches(6.5), Inches(0.5),
                 task_desc, font_size=12, color=MID_GRAY)
    y_pos += Inches(0.73)

val_box = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.1), Inches(11.5), Inches(0.5), LIGHT_BLUE)
add_text_box(s, Inches(1), y_pos + Inches(0.12), Inches(11), Inches(0.45),
             "📸 Capture screenshots: site details, access review notes, recycle bin restore, PowerShell output.",
             font_size=13, color=ACCENT_BLUE)

add_footer_bar(s, 24, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Summary
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Recap key messages from Module 3.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "📝 Module 3 Summary", font_size=30, bold=True, color=DARK_BG)

summary = [
    "SharePoint sites are the primary manageable unit — each has URL, owners, permissions, storage",
    "Team Sites = collaboration (M365 Group-backed); Communication Sites = broadcasting",
    "Hub Sites organise sites with shared navigation, branding, and search — without changing permissions",
    "Sites are created from the admin center or PowerShell; always follow naming conventions",
    "Day-to-day ops: membership, access requests, recycle bin (two-stage recovery)",
    "Storage: tenant pool (automatic) vs per-site limits (manual) — know which mode you're in",
    "Delete/restore has a 93-day safety window — always use a test site for drills",
    "PowerShell: Connect-SPOService → Get-SPOSite — start with read-only",
]
add_bullet_frame(s, Inches(0.8), Inches(1.2), Inches(11), Inches(5),
                 summary, font_size=16, color=DARK_TEXT, icon="✅")

add_footer_bar(s, 25, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Day 1 Wrap-up
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("This is the last module of Day 1. Recap what was covered across all 3 modules.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
             "Day 1 Complete!", font_size=42, bold=True,
             color=WHITE, font_name="Segoe UI Light", alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(2.2), Inches(11), Inches(0.5),
             "Tenant Foundations & Site Management", font_size=20,
             color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

day1_modules = [
    ("Module 1", "Microsoft 365 & SharePoint Online overview, admin centers, service limits", ACCENT_BLUE),
    ("Module 2", "Identity, access, external sharing, guest access, Conditional Access", ACCENT_TEAL),
    ("Module 3", "Site types, creation, management, storage, lifecycle, PowerShell", ACCENT_PURPLE),
]

y_pos = Inches(3.2)
for mod, desc, color in day1_modules:
    card = add_rounded_rect(s, Inches(2), y_pos, Inches(9), Inches(0.8), RGBColor(0x2B, 0x2B, 0x3F))
    add_shape_rect(s, Inches(2), y_pos, Inches(0.1), Inches(0.8), color)
    add_text_box(s, Inches(2.3), y_pos + Inches(0.1), Inches(2), Inches(0.5),
                 mod, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(4.5), y_pos + Inches(0.15), Inches(6), Inches(0.5),
                 desc, font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB))
    y_pos += Inches(0.95)

add_text_box(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "🚀 Tomorrow: Permissions, Metadata, Search & Customization (Day 2)",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 26, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Knowledge Check
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Discussion slide. Keep answers short; prioritize scope and safety.")
add_solid_bg(s, WHITE)
add_top_bar(s, ACCENT_PURPLE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🧠 Knowledge Check", font_size=30, bold=True, color=DARK_BG)

questions = [
    ("Q1", "Why do we isolate work to NW-Pxx-... sites in this course?",
     "Prevent collisions; avoid impacting other participants in shared tenant"),
    ("Q2", "When can you set a per-site storage limit?",
     "Only when the tenant storage management mode is set to Manual"),
    ("Q3", "Why use a dedicated NW-Pxx-RestoreTest site for the delete drill?",
     "Avoid damaging the persistent practice site; avoid impacting others"),
    ("Q4", "What must happen before running Get-SPOSite?",
     "You must first connect via Connect-SPOService with the admin URL"),
]

y_pos = Inches(1.2)
for q_id, question, answer in questions:
    qbadge = add_rounded_rect(s, Inches(0.8), y_pos, Inches(0.6), Inches(0.5), ACCENT_PURPLE)
    add_text_box(s, Inches(0.82), y_pos + Inches(0.03), Inches(0.55), Inches(0.4),
                 q_id, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.6), y_pos, Inches(10.5), Inches(0.5),
                 question, font_size=16, color=DARK_TEXT)
    add_text_box(s, Inches(1.6), y_pos + Inches(0.5), Inches(10.5), Inches(0.4),
                 f"→ {answer}", font_size=13, color=RGBColor(0x99, 0x99, 0x99))
    y_pos += Inches(1.15)

add_text_box(s, Inches(0.8), y_pos + Inches(0.3), Inches(11.5), Inches(0.5),
             "💬 Discuss with your neighbour — then we'll share answers.", font_size=16,
             bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 27, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Thank You / End of Day 1
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Thank participants for Day 1. Preview Day 2 topics.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Thank You!", font_size=52, bold=True,
             color=WHITE, font_name="Segoe UI Light", alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "End of Day 1 — See you tomorrow for Day 2!", font_size=20,
             color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
             "Day 2: Permissions · Metadata · Search · Customization", font_size=16,
             color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 28, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
TOTAL_SLIDES = slide_counter[0]
output_path = os.path.join(os.path.dirname(__file__), "Module-03-Slides.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
