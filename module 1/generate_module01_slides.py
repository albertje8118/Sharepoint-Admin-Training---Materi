"""
Generate Module 1 PPTX — Introduction to Microsoft 365 and SharePoint Online
Modern, engaging design with gradient-style colours, icons via Unicode, and rich content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette (modern Microsoft-inspired) ──────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG     = RGBColor(0x1B, 0x1B, 0x2F)  # deep navy
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)  # Microsoft blue
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)  # teal accent
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6) # purple accent
LIGHT_BLUE  = RGBColor(0xDE, 0xEC, 0xF9)  # soft blue bg
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT    = RGBColor(0x24, 0x24, 0x24)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT  = RGBColor(0xD1, 0x34, 0x38)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ═══════════════════════════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════════════════════════

def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Set corner radius
    shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6), icon="▸", font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon} {bullet}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox

def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)


def add_footer_bar(slide, slide_num, total, module_label="Module 1"):
    bar_top = SLIDE_HEIGHT - Inches(0.45)
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_WIDTH, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), bar_top + Inches(0.05), Inches(6), Inches(0.35),
                 f"{module_label}  |  Introduction to Microsoft 365 & SharePoint Online",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), bar_top + Inches(0.05),
                 Inches(1.2), Inches(0.35),
                 f"{slide_num} / {total}", font_size=10,
                 color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)


def section_divider(slide, section_title, subtitle="", icon=""):
    add_solid_bg(slide, DARK_BG)
    # accent stripe
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_TEAL)
    if icon:
        add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                     icon, font_size=60, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                 section_title, font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB),
                     alignment=PP_ALIGN.CENTER)


TOTAL_SLIDES = 24  # We'll update this after building all slides
slide_counter = [0]  # mutable counter

def new_slide(notes=""):
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    if notes:
        add_speaker_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide (cover)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome participants. Introduce yourself and set the stage for 3 days of hands-on SharePoint Online administration.")
add_solid_bg(s, DARK_BG)
# accent rectangles
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

# Title
add_text_box(s, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "MODERN SHAREPOINT ONLINE", font_size=20, bold=True,
             color=ACCENT_TEAL, font_name="Segoe UI Semibold")
add_text_box(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
             "for Administrators", font_size=48, bold=True,
             color=WHITE, font_name="Segoe UI Light")
add_text_box(s, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
             "3-Day Instructor-Led Training  ·  2026 Aligned", font_size=18,
             color=RGBColor(0xAA, 0xAA, 0xAA))

# Module badge
badge = add_rounded_rect(s, Inches(1), Inches(4.5), Inches(3.2), Inches(0.55), ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(4.55), Inches(2.8), Inches(0.45),
             "MODULE 1", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(5.3), Inches(11), Inches(0.6),
             "Introduction to Microsoft 365 & SharePoint Online",
             font_size=24, bold=True, color=WHITE)

# Bottom info
add_text_box(s, Inches(1), Inches(6.4), Inches(5), Inches(0.4),
             "Day 1  ·  Tenant Foundations & Site Management",
             font_size=14, color=RGBColor(0x88, 0x88, 0x88))

add_footer_bar(s, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda / What We'll Cover Today
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the agenda so participants know what to expect from this module.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "📋  Module 1 Agenda", font_size=30, bold=True, color=DARK_BG)

agenda_items = [
    ("1", "SharePoint Online Overview — What it is & why it matters", ACCENT_BLUE),
    ("2", "SharePoint, OneDrive & Teams — How they work together", ACCENT_TEAL),
    ("3", "Microsoft 365 Licensing — Plans that include SharePoint", ACCENT_PURPLE),
    ("4", "Microsoft 365 Admin Mental Model — Architecture overview", ACCENT_BLUE),
    ("5", "Admin Centers — M365 & SharePoint admin portals", ACCENT_TEAL),
    ("6", "SPO vs SP Server — Conceptual comparison", ACCENT_PURPLE),
    ("7", "Service Limits & Boundaries — How to approach them", ACCENT_BLUE),
    ("8", "Lab 1 Preview & Knowledge Check", ORANGE),
]

y_pos = Inches(1.3)
for num, text, color in agenda_items:
    # number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    # number text
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Segoe UI"
    # item text
    add_text_box(s, Inches(1.5), y_pos + Inches(0.05), Inches(10), Inches(0.45),
                 text, font_size=17, color=DARK_TEXT)
    y_pos += Inches(0.65)

add_footer_bar(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Module Objectives
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Review learning objectives. These align with what participants will be assessed on.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "🎯  Module Objectives", font_size=30, bold=True, color=DARK_BG)

objectives = [
    "Describe Microsoft 365 service architecture at a practical admin level",
    "Explain the role of SharePoint Online in Microsoft 365 collaboration and content services",
    "Compare SharePoint Online and SharePoint Server conceptually",
    "Navigate Microsoft 365 admin center and SharePoint admin center for baseline tenant checks",
    "Locate and interpret service limits, quotas, and boundaries using official Microsoft documentation",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(5),
                 objectives, font_size=18, color=DARK_TEXT, icon="✅")

add_footer_bar(s, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Section Divider: SharePoint Online Overview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to the SharePoint Online overview section.")
section_divider(s, "SharePoint Online Overview", "What it is, what it does, and why admins should care", "🌐")
add_footer_bar(s, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What is SharePoint Online?
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("SharePoint Online is the content and collaboration backbone of Microsoft 365. "
              "Emphasize that it's not just a file store — it powers intranets, content services, and integrated workflows.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "What is SharePoint Online?", font_size=30, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "SharePoint Online is the content & collaboration platform in Microsoft 365 that helps "
             "organizations share and manage content, knowledge, and applications.",
             font_size=17, color=MID_GRAY)

# Cards row
cards = [
    ("🏗️ Sites", "Team sites, Communication\nsites, Hub sites for\ncollaboration & publishing", ACCENT_BLUE),
    ("📁 Document Libraries", "Controlled storage with\nversioning, metadata,\nco-authoring & sharing", ACCENT_TEAL),
    ("🔍 Content Services", "Microsoft Search, content\ntypes, managed metadata,\nretention & compliance", ACCENT_PURPLE),
    ("⚡ Integration", "Powers Teams files,\nOneDrive storage,\nPower Platform connectors", ORANGE),
]
x_pos = Inches(0.5)
for title, body, color in cards:
    card = add_rounded_rect(s, x_pos, Inches(2.5), Inches(2.9), Inches(3.2), NEAR_WHITE)
    # top accent bar
    add_shape_rect(s, x_pos, Inches(2.5), Inches(2.9), Inches(0.08), color)
    # title
    add_text_box(s, x_pos + Inches(0.2), Inches(2.75), Inches(2.5), Inches(0.6),
                 title, font_size=17, bold=True, color=color)
    # body
    add_text_box(s, x_pos + Inches(0.2), Inches(3.5), Inches(2.5), Inches(2),
                 body, font_size=14, color=MID_GRAY)
    x_pos += Inches(3.15)

add_footer_bar(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SharePoint Admin Responsibilities
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Highlight the key areas a SharePoint admin manages day-to-day.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Common SharePoint Admin Responsibilities", font_size=28, bold=True, color=DARK_BG)

responsibilities = [
    ("🔐", "Tenant Policies", "Sharing, access control, external collaboration settings"),
    ("🏠", "Site Lifecycle", "Creation, ownership, storage quotas, deletion & restore"),
    ("🏷️", "Information Architecture", "Metadata, term store, content types (covered later)"),
    ("📊", "Operational Health", "Service health dashboard, message center, admin notifications"),
    ("⚖️", "Governance & Compliance", "Purview policies, audit visibility, retention alignment"),
]

y_pos = Inches(1.3)
for icon, title, desc in responsibilities:
    # icon bg
    icon_bg = add_rounded_rect(s, Inches(0.8), y_pos, Inches(0.65), Inches(0.65), LIGHT_BLUE)
    add_text_box(s, Inches(0.82), y_pos + Inches(0.05), Inches(0.6), Inches(0.55),
                 icon, font_size=22, alignment=PP_ALIGN.CENTER, color=DARK_TEXT)
    add_text_box(s, Inches(1.7), y_pos, Inches(4), Inches(0.4),
                 title, font_size=18, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(1.7), y_pos + Inches(0.35), Inches(10), Inches(0.35),
                 desc, font_size=14, color=MID_GRAY)
    y_pos += Inches(0.9)

add_footer_bar(s, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Section Divider: SharePoint + OneDrive + Teams
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to the integration section.")
section_divider(s, "SharePoint + OneDrive + Teams",
                "Understanding the relationship between Microsoft 365's collaboration pillars", "🔗")
add_footer_bar(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — The Collaboration Triangle
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Key insight: SharePoint is the content layer; Teams is the collaboration UX; OneDrive is personal storage "
              "built on SharePoint technology. Files uploaded in Teams channels go to SharePoint; files in private chat go to OneDrive.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "The Microsoft 365 Collaboration Triangle", font_size=28, bold=True, color=DARK_BG)

# Three pillars
pillars = [
    ("SharePoint Online", "Content platform\n& intranet engine", "📄 Shared document libraries\n🌐 Team & communication sites\n🔎 Enterprise search\n🏷️ Metadata & governance", ACCENT_BLUE),
    ("Microsoft Teams", "Collaboration hub\n& user experience", "💬 Chat & channels\n📹 Meetings & calls\n📋 Tabs & apps\n📁 Files tab → SharePoint", ACCENT_PURPLE),
    ("OneDrive for Business", "Personal cloud storage\npowered by SharePoint", "👤 Individual file storage\n🔄 Sync to desktop\n↗️ Easy sharing via links\n💬 Private chat files land here", ACCENT_TEAL),
]

x_pos = Inches(0.4)
for title, subtitle, features, color in pillars:
    card = add_rounded_rect(s, x_pos, Inches(1.4), Inches(3.9), Inches(4.8), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(1.4), Inches(3.9), Inches(0.1), color)
    add_text_box(s, x_pos + Inches(0.2), Inches(1.7), Inches(3.5), Inches(0.5),
                 title, font_size=20, bold=True, color=color)
    add_text_box(s, x_pos + Inches(0.2), Inches(2.3), Inches(3.5), Inches(0.5),
                 subtitle, font_size=13, color=MID_GRAY)
    add_text_box(s, x_pos + Inches(0.2), Inches(3.0), Inches(3.5), Inches(3),
                 features, font_size=14, color=DARK_TEXT)
    x_pos += Inches(4.2)

# Bottom connector text
add_text_box(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "🔑 Key: Files shared in Teams channels are stored in SharePoint. Files in private chats are stored in the sender's OneDrive.",
             font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — How Teams Creates SharePoint Sites
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("When a new Team is created, Microsoft 365 automatically provisions: an M365 Group, "
              "a SharePoint Team site, an Exchange shared mailbox, a OneNote notebook, and Planner. "
              "This means every Team already has a SharePoint site behind it.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "What Happens When You Create a Team?", font_size=28, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.6),
             "Every new Microsoft Team automatically provisions these resources:",
             font_size=16, color=MID_GRAY)

resources = [
    ("👥", "Microsoft 365 Group", "Unified identity & membership", ACCENT_BLUE),
    ("📄", "SharePoint Team Site", "Document library for channel files", ACCENT_TEAL),
    ("📧", "Exchange Shared Mailbox", "Group conversations & calendar", ACCENT_PURPLE),
    ("📓", "OneNote Notebook", "Shared notebook for the team", ORANGE),
    ("✅", "Planner Board", "Task management for the group", GREEN),
]

x_pos = Inches(0.3)
for icon, title, desc, color in resources:
    card = add_rounded_rect(s, x_pos, Inches(2.2), Inches(2.35), Inches(2.8), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(2.2), Inches(2.35), Inches(0.06), color)
    add_text_box(s, x_pos + Inches(0.1), Inches(2.45), Inches(2.15), Inches(0.5),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER, color=DARK_TEXT)
    add_text_box(s, x_pos + Inches(0.1), Inches(3.1), Inches(2.15), Inches(0.5),
                 title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x_pos + Inches(0.1), Inches(3.65), Inches(2.15), Inches(0.6),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(2.55)

# callout
callout = add_rounded_rect(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.45), Inches(11), Inches(0.6),
             "💡 Admin Insight: Many sites in your SharePoint admin center are Teams-connected. "
             "Deleting a site can impact Teams!",
             font_size=14, bold=False, color=ACCENT_BLUE)

add_footer_bar(s, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Section Divider: Licensing
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to licensing section.")
section_divider(s, "Microsoft 365 Licensing", "Which plans include SharePoint Online, OneDrive, and Teams?", "🪪")
add_footer_bar(s, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Licensing Overview Table
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Emphasize that SharePoint, OneDrive, and Teams are included in most M365/O365 plans. "
              "Storage is pooled at the tenant level: 1 TB base + 10 GB per license. "
              "Frontline (F1/F3) plans have limited storage. Standalone SharePoint plans exist too.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Microsoft 365 Plans — SharePoint Inclusion", font_size=28, bold=True, color=DARK_BG)

# Table
from pptx.util import Inches, Pt
rows, cols = 7, 5
tbl_shape = s.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12.3), Inches(4.5))
tbl = tbl_shape.table

headers = ["Plan Family", "SharePoint", "OneDrive", "Teams", "Desktop Apps"]
data = [
    ["Business Basic",   "✅", "✅", "✅", "❌"],
    ["Business Standard","✅", "✅", "✅", "✅"],
    ["Business Premium", "✅", "✅", "✅", "✅"],
    ["Enterprise E3",    "✅", "✅", "✅", "✅"],
    ["Enterprise E5",    "✅", "✅", "✅", "✅"],
    ["Frontline F1/F3",  "✅ (limited)", "✅ (2 GB)", "✅", "❌ / ✅*"],
]

# style header
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE

# data rows
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = tbl.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEAR_WHITE if i % 2 == 0 else WHITE

# Storage note
add_text_box(s, Inches(0.8), Inches(5.9), Inches(11), Inches(0.8),
             "📦 Tenant Storage: 1 TB base + 10 GB per licensed user  |  Max per site: 25 TB  |  Up to 2 million sites per org",
             font_size=14, bold=True, color=ACCENT_TEAL)

add_footer_bar(s, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Standalone Plans & Add-ons
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Highlight standalone options for orgs that don't need full M365. "
              "SharePoint Plan 1 & 2, OneDrive Plan 1 & 2 are available. "
              "Compliance features (retention, DLP, eDiscovery) require E3/E5 or Purview add-ons.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Standalone Plans & Key Add-ons", font_size=28, bold=True, color=DARK_BG)

standalone_items = [
    ("SharePoint Online Plan 1", "Basic sites, document libraries, external sharing"),
    ("SharePoint Online Plan 2", "Plan 1 + advanced search, enterprise features"),
    ("OneDrive for Business Plan 1", "Personal storage (1 TB per user)"),
    ("OneDrive for Business Plan 2", "Plan 1 + unlimited storage, advanced compliance"),
    ("Microsoft Purview Add-ons", "Retention, DLP, eDiscovery — if not in your base plan"),
    ("Power Platform Add-ons", "Additional AI Builder credits, premium connectors"),
]

y_pos = Inches(1.3)
for title, desc in standalone_items:
    add_shape_rect(s, Inches(0.8), y_pos, Inches(0.08), Inches(0.6), ACCENT_TEAL)
    add_text_box(s, Inches(1.1), y_pos, Inches(5), Inches(0.35),
                 title, font_size=16, bold=True, color=DARK_TEXT)
    add_text_box(s, Inches(1.1), y_pos + Inches(0.3), Inches(10), Inches(0.35),
                 desc, font_size=13, color=MID_GRAY)
    y_pos += Inches(0.8)

# Tip box
tip = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.1), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), y_pos + Inches(0.15), Inches(11), Inches(0.6),
             "💡 Tip: You can combine Enterprise, Business, and standalone plans within a single tenant.",
             font_size=14, color=ACCENT_BLUE)

add_footer_bar(s, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Section Divider: M365 Architecture
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to architecture overview.")
section_divider(s, "Microsoft 365 Architecture", "The admin mental model for SharePoint professionals", "🧩")
add_footer_bar(s, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — M365 Architecture: Admin Mental Model
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Emphasize that many 'SharePoint issues' are actually identity/policy issues. "
              "SharePoint doesn't exist in isolation — it depends on Entra ID for auth, Purview for compliance, etc.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Microsoft 365: The Admin Mental Model", font_size=28, bold=True, color=DARK_BG)

layers = [
    ("🔑 Identity Layer", "Microsoft Entra ID", "Authentication, authorization, Conditional Access, device policies",
     ACCENT_BLUE, Inches(1.3)),
    ("⚙️ Workloads Layer", "SharePoint · OneDrive · Teams · Exchange · Purview · Search", 
     "The services your users interact with — each has its own admin surface",
     ACCENT_TEAL, Inches(2.7)),
    ("🖥️ Admin Surfaces", "M365 Admin Center + Workload Admin Centers",
     "Central portal for tenant-wide config; workload portals for service-specific settings",
     ACCENT_PURPLE, Inches(4.1)),
    ("🤖 Automation Layer", "Microsoft Graph API + PowerShell Modules",
     "Programmatic access for reporting, bulk operations, and integration",
     ORANGE, Inches(5.5)),
]

for icon_title, comp, desc, color, y in layers:
    bar = add_rounded_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.1), NEAR_WHITE)
    add_shape_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.1), color)
    add_text_box(s, Inches(0.9), y + Inches(0.05), Inches(5), Inches(0.4),
                 icon_title, font_size=17, bold=True, color=color)
    add_text_box(s, Inches(0.9), y + Inches(0.4), Inches(5), Inches(0.4),
                 comp, font_size=13, bold=True, color=DARK_TEXT)
    add_text_box(s, Inches(6.5), y + Inches(0.15), Inches(6), Inches(0.8),
                 desc, font_size=13, color=MID_GRAY)

add_footer_bar(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Why This Matters to SharePoint Admins
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Cross-service dependencies: a user can't access a SharePoint site if Conditional Access blocks them, "
              "even though SharePoint config looks fine. Train admins to think across services.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Why This Matters to SharePoint Admins", font_size=28, bold=True, color=DARK_BG)

cross_deps = [
    "Identity & access policies (Entra ID, Conditional Access) control who can reach SharePoint",
    "Compliance controls (Microsoft Purview) govern retention, DLP, and sensitivity labels",
    "Search & content experiences are integrated across all Microsoft 365 workloads",
    "Teams file storage IS SharePoint — admin changes affect both",
    "OneDrive inherits SharePoint sharing policies (OneDrive ≤ SharePoint)",
    "Power Platform connectors rely on SharePoint for content sources",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(4.5),
                 cross_deps, font_size=17, color=DARK_TEXT, icon="⚠️")

# Warning callout
warn = add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.7),
                        RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(5.55), Inches(11), Inches(0.6),
             "🔎 Many 'SharePoint issues' are actually identity or policy issues originating outside SharePoint.",
             font_size=15, bold=True, color=ORANGE)

add_footer_bar(s, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — M365 Admin Center: What SPO Admins Use It For
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Build the 'service health first' habit. Check service health before troubleshooting. "
              "Review message center to understand upcoming changes.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Microsoft 365 Admin Center — For SPO Admins", font_size=28, bold=True, color=DARK_BG)

tasks = [
    ("🔍 Service Health", "Is it you or is it the service? Check here FIRST.", ACCENT_BLUE),
    ("📢 Message Center", "Upcoming changes that may impact SharePoint governance.", ACCENT_TEAL),
    ("👤 Roles & Licenses", "Confirm admin roles and user licensing assignments.", ACCENT_PURPLE),
    ("📧 Tenant Context", "Verify tenant name, domains, and subscription info.", ORANGE),
]

y_pos = Inches(1.3)
for title, desc, color in tasks:
    card = add_rounded_rect(s, Inches(0.8), y_pos, Inches(11.5), Inches(1), NEAR_WHITE)
    add_shape_rect(s, Inches(0.8), y_pos, Inches(0.1), Inches(1), color)
    add_text_box(s, Inches(1.2), y_pos + Inches(0.1), Inches(5), Inches(0.4),
                 title, font_size=18, bold=True, color=color)
    add_text_box(s, Inches(1.2), y_pos + Inches(0.5), Inches(10.5), Inches(0.4),
                 desc, font_size=15, color=MID_GRAY)
    y_pos += Inches(1.2)

# Good practice callout
gp = add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.55), Inches(11), Inches(0.4),
             "✅ Good Practice", font_size=16, bold=True, color=ACCENT_BLUE)
add_text_box(s, Inches(1), Inches(5.9), Inches(11), Inches(0.5),
             "1) Service health first → 2) Message center next → 3) Then investigate your config",
             font_size=15, color=DARK_TEXT)

add_footer_bar(s, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — SharePoint Admin Center Overview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("SharePoint admin center focuses on sites, policies, and storage. "
              "UI labels can change — teach finding categories, not memorizing clicks.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "SharePoint Admin Center (Modern)", font_size=28, bold=True, color=DARK_BG)

areas = [
    ("Sites Management", "Create, view, manage all site collections at scale"),
    ("Sharing & Access Policies", "Tenant-wide sharing defaults, guest access, link types"),
    ("Storage & Usage", "Monitor storage consumption, set quotas, view usage reports"),
    ("Settings & Org-wide Config", "Default site creation, notifications, API access"),
    ("Content Services", "Term store, content type hub, search schema (linked)"),
]

y_pos = Inches(1.3)
for i, (title, desc) in enumerate(areas):
    color = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ORANGE, GREEN][i]
    num_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(0.5), Inches(0.5))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI"

    add_text_box(s, Inches(1.5), y_pos, Inches(4), Inches(0.35),
                 title, font_size=17, bold=True, color=DARK_TEXT)
    add_text_box(s, Inches(1.5), y_pos + Inches(0.32), Inches(10), Inches(0.3),
                 desc, font_size=14, color=MID_GRAY)
    y_pos += Inches(0.8)

# Navigation principle
nav = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.2), Inches(11.5), Inches(0.8), LIGHT_BLUE)
add_text_box(s, Inches(1), y_pos + Inches(0.25), Inches(11), Inches(0.7),
             "📌 Navigation principle: Understand which settings are tenant-wide vs site-specific, "
             "and which are 'policy' vs 'operational configuration'.",
             font_size=14, color=ACCENT_BLUE)

add_footer_bar(s, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Section Divider: SPO vs SP Server
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to comparison section.")
section_divider(s, "SPO vs SharePoint Server",
                "A conceptual comparison for modern admins", "⚖️")
add_footer_bar(s, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — SPO vs SP Server Comparison
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Keep this conceptual. Avoid deep feature checklists. "
              "Focus on the operating model difference: cloud-operated vs self-operated.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "SharePoint Online vs SharePoint Server", font_size=28, bold=True, color=DARK_BG)

# Comparison table
rows2, cols2 = 5, 3
tbl2_shape = s.shapes.add_table(rows2, cols2, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.0))
tbl2 = tbl2_shape.table
tbl2.columns[0].width = Inches(3)
tbl2.columns[1].width = Inches(4.65)
tbl2.columns[2].width = Inches(4.65)

comp_headers = ["Dimension", "SharePoint Online", "SharePoint Server"]
comp_data = [
    ["Deployment & Ownership", "Microsoft operates the service;\nyou control config & governance", "Your infrastructure, patching,\nupgrades, capacity planning"],
    ["Change Cadence", "Continuous service updates;\nfocus on governance & adoption", "Changes follow your maintenance\nand upgrade schedule"],
    ["Integration Model", "Designed for M365 integration\n(Entra ID, Purview, Teams, Graph)", "Often requires custom design\nand on-prem dependencies"],
    ["Customization", "Client-side solutions, SPFx,\nAPI-based (Graph/REST)", "Farm solutions, server-side\ncustomization patterns"],
]

for j, h in enumerate(comp_headers):
    cell = tbl2.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE

for i, row_data in enumerate(comp_data):
    for j, val in enumerate(row_data):
        cell = tbl2.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Segoe UI"
        cell.fill.solid()
        cell.fill.fore_color.rgb = NEAR_WHITE if i % 2 == 0 else WHITE

# Takeaway
add_text_box(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.6),
             "🎯 Admin Takeaway: SPO administration = policy + governance + lifecycle + integration (not server ops)",
             font_size=16, bold=True, color=ACCENT_TEAL)

add_footer_bar(s, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Section Divider: Service Limits
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to limits section.")
section_divider(s, "Service Limits & Boundaries",
                "How to approach platform constraints safely", "📏")
add_footer_bar(s, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Service Limits: The Right Approach
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Do NOT memorize numbers. Limits change. Always verify in official docs. "
              "Document which limit matters to YOUR scenario.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Service Limits, Quotas & Boundaries", font_size=28, bold=True, color=DARK_BG)

# Definitions
defs = [
    ("Limits / Boundaries", "Hard platform constraints — what the service CAN support", "🚧"),
    ("Quotas", "Configurable allocations (e.g., storage per site)", "📊"),
    ("Recommendations", "Guidance for performance and manageability", "💡"),
]

y_pos = Inches(1.3)
for title, desc, icon in defs:
    add_text_box(s, Inches(0.8), y_pos, Inches(0.5), Inches(0.5),
                 icon, font_size=24, alignment=PP_ALIGN.CENTER, color=DARK_TEXT)
    add_text_box(s, Inches(1.5), y_pos, Inches(4), Inches(0.35),
                 title, font_size=17, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(1.5), y_pos + Inches(0.32), Inches(10), Inches(0.3),
                 desc, font_size=14, color=MID_GRAY)
    y_pos += Inches(0.8)

# Key numbers (current as of docs)
add_text_box(s, Inches(0.8), y_pos + Inches(0.1), Inches(10), Inches(0.4),
             "Key Numbers (verify in docs — these change!):", font_size=16, bold=True, color=DARK_TEXT)

numbers = [
    "Max 25 TB per site collection",
    "Up to 2 million sites per organization", 
    "30 million items per list/library (but 5,000 list view threshold)",
    "Tenant storage = 1 TB + 10 GB per licensed user",
]
add_bullet_frame(s, Inches(0.8), y_pos + Inches(0.5), Inches(11), Inches(2.5),
                 numbers, font_size=15, color=DARK_TEXT, icon="📌")

# Safe approach
safe = add_rounded_rect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.85), Inches(11), Inches(0.6),
             "✅ Safe approach: Don't memorize → Verify in Microsoft Learn → Document scenario-specific limits",
             font_size=14, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 21, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Lab Scenario & Shared Tenant Rules
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Set shared-tenant rules. Emphasize participant IDs and collision avoidance.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🏢 Lab Scenario & Shared Tenant Rules", font_size=28, bold=True, color=DARK_BG)

rules = [
    "Scenario: Project Northwind Intranet Modernization",
    "One tenant shared by: Trainer + 10 Admin Participants (P01–P10)",
    "Use your Participant ID for all site/group/resource naming",
    "Avoid tenant-wide changes unless marked as Trainer-only",
    "Module 1 labs focus on orientation and verification — no destructive changes",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(3.5),
                 rules, font_size=17, color=DARK_TEXT, icon="📋")

# Warning
warn2 = add_rounded_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.7),
                         RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(4.55), Inches(11), Inches(0.6),
             "⚠️ Shared tenant = prevent collisions. Most hands-on changes happen in participant-isolated practice sites (later modules).",
             font_size=14, bold=True, color=ORANGE)

add_footer_bar(s, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Lab 1 Preview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Explain what learners should capture in validation checkpoints during the lab.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🔬 Lab 1 Preview — Exploring the Microsoft 365 Environment", font_size=26, bold=True, color=DARK_BG)

lab_tasks = [
    ("Task 1", "Access Microsoft 365 Admin Center", "Sign in, confirm your role and tenant context"),
    ("Task 2", "Review SharePoint Online tenant settings", "Navigate SharePoint admin center, explore settings areas"),
    ("Task 3", "Check Service Health", "Find current service status for SharePoint & OneDrive"),
    ("Task 4", "Review Message Center", "Identify recent announcements affecting SharePoint"),
]

y_pos = Inches(1.3)
for task_id, task_title, task_desc in lab_tasks:
    card = add_rounded_rect(s, Inches(0.8), y_pos, Inches(11.5), Inches(0.95), NEAR_WHITE)
    # task badge
    badge2 = add_rounded_rect(s, Inches(1), y_pos + Inches(0.15), Inches(1.2), Inches(0.4), ACCENT_BLUE)
    add_text_box(s, Inches(1.05), y_pos + Inches(0.17), Inches(1.1), Inches(0.35),
                 task_id, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.4), y_pos + Inches(0.1), Inches(9), Inches(0.35),
                 task_title, font_size=17, bold=True, color=DARK_TEXT)
    add_text_box(s, Inches(2.4), y_pos + Inches(0.48), Inches(9), Inches(0.35),
                 task_desc, font_size=13, color=MID_GRAY)
    y_pos += Inches(1.15)

# Validation tip
val_box = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.2), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), y_pos + Inches(0.25), Inches(11), Inches(0.6),
             "📸 Capture screenshots of: your admin role, service health status, and at least one Message Center post.",
             font_size=14, color=ACCENT_BLUE)

add_footer_bar(s, 23, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Knowledge Check
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Use this as a discussion slide. Encourage short answers. Gauge baseline understanding.")
add_solid_bg(s, WHITE)
add_top_bar(s, ACCENT_PURPLE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🧠 Knowledge Check", font_size=30, bold=True, color=DARK_BG)

questions = [
    ("Q1", "Which Microsoft 365 component is responsible for identity and access control?",
     "Microsoft Entra ID"),
    ("Q2", "Why should a SharePoint admin care about Message Center?",
     "It announces changes that can affect governance, features, and UX"),
    ("Q3", "Name two conceptual differences between SharePoint Online and SharePoint Server.",
     "Cloud vs self-operated; continuous vs scheduled; integration model differs"),
    ("Q4", "What is the safe approach to service limits in documentation?",
     "Verify current values in official docs; avoid memorized numbers"),
]

y_pos = Inches(1.2)
for q_id, question, answer in questions:
    # Q badge
    q_badge = add_rounded_rect(s, Inches(0.8), y_pos, Inches(0.6), Inches(0.5), ACCENT_PURPLE)
    add_text_box(s, Inches(0.82), y_pos + Inches(0.03), Inches(0.55), Inches(0.4),
                 q_id, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.6), y_pos, Inches(10.5), Inches(0.5),
                 question, font_size=16, bold=False, color=DARK_TEXT)
    # Answer (subtle)
    add_text_box(s, Inches(1.6), y_pos + Inches(0.5), Inches(10.5), Inches(0.4),
                 f"→ {answer}", font_size=13, color=RGBColor(0x99, 0x99, 0x99))
    y_pos += Inches(1.15)

# Discussion prompt
add_text_box(s, Inches(0.8), y_pos + Inches(0.3), Inches(11.5), Inches(0.5),
             "💬 Discuss with your neighbour — then we'll share answers.", font_size=16,
             bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 24, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# Update slide count and save
# ═══════════════════════════════════════════════════════════════════════════
TOTAL_SLIDES = slide_counter[0]
# Note: footer text already references TOTAL_SLIDES but it was set to 24 initially.
# Since we built exactly 24 slides, the count is correct.

output_path = os.path.join(os.path.dirname(__file__), "Module-01-Slides.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
