from __future__ import annotations

from dataclasses import dataclass
from datetime import date
from pathlib import Path

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas


COURSE_TITLE = "Modern SharePoint Online for Administrators (3-Day, 2026 aligned)"
SCENARIO_TITLE = "Project Northwind Intranet Modernization"
TODAY = date(2026, 2, 9)
TRAINER_ID = "TRAINER"


@dataclass(frozen=True)
class Participant:
    participant_id: str  # e.g. P01

    @property
    def pxx(self) -> str:
        return self.participant_id

    @property
    def nw_prefix(self) -> str:
        return f"NW-{self.pxx}-"

    @property
    def project_site_name(self) -> str:
        return f"NW-{self.pxx}-ProjectSite"

    @property
    def contracts_library_name(self) -> str:
        return f"NW-{self.pxx}-Contracts"

    @property
    def term_group_name(self) -> str:
        return f"NW-{self.pxx}-TermGroup"

    @property
    def termset_contract_type(self) -> str:
        return f"NW-{self.pxx}-ContractType"


def _set_doc_defaults(document: Document) -> None:
    style = document.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)


def build_docx(participant: Participant, out_path: Path) -> None:
    document = Document()
    _set_doc_defaults(document)

    document.add_heading(f"Participant Pack — {participant.pxx}", level=1)
    document.add_paragraph(COURSE_TITLE)
    document.add_paragraph(f"Scenario: {SCENARIO_TITLE}")
    document.add_paragraph(f"Date: {TODAY.isoformat()}")

    document.add_heading("Your lab naming", level=2)
    document.add_paragraph(f"Participant ID: {participant.pxx}")
    document.add_paragraph(f"Primary practice site (persistent): {participant.project_site_name}")
    document.add_paragraph(f"Module 4 library artifact: {participant.contracts_library_name}")

    document.add_heading("Shared-tenant safety rules (quick)", level=2)
    rules = [
        "Only change content/settings inside your own NW-Pxx artifacts.",
        "Treat tenant-wide settings as Trainer-only.",
        "Do not invite external guests unless the trainer explicitly asks.",
        "If you break inheritance, keep unique permissions minimal (library boundary first).",
    ]
    for rule in rules:
        document.add_paragraph(rule, style="List Bullet")

    document.add_heading("Quick checklist", level=2)
    checklist = [
        "Confirm you can access SharePoint admin center (read-only verification is fine).",
        "Confirm your NW-Pxx-ProjectSite URL is recorded in the tracker.",
        "Module 4: Create NW-Pxx-Contracts library and folders 01-Drafts / 02-InReview / 03-Final.",
        "Module 4: Break inheritance at library, and at one folder only (03-Final).",
        "Validate access using Check Permissions and Manage access.",
        "Module 5: Create (or use fallback) NW-Pxx-ContractType term set + add managed metadata column.",
        "Module 6: Upload search seed docs + validate you can find them via search.",
        "Module 7: Create NW-Pxx-AppRequests list + apply column and view formatting JSON.",
        "Module 8: Upload FAKE compliance doc + apply sensitivity/retention labels (if published) + complete worksheet.",
        "Module 9: Observe OneDrive tenant settings + complete worksheet; do internal-only sharing test.",
        "Module 10: Run PowerShell reporting (SPO) scoped to NW-Pxx + export CSV; optional Graph connect if trainer approves.",
        "Module 11: Complete ops worksheets (incident triage, lifecycle, external sharing governance) + provide a change request summary.",
        "Module 12 (optional): Build an AppRequests approval flow + customize the list form (Power Apps) + record governance notes.",
    ]
    for item in checklist:
        document.add_paragraph(item, style="List Number")

    document.add_heading("Notes", level=2)
    document.add_paragraph(
        "Use this page to jot down any URLs, group names, and screenshots you want to keep for the next modules."
    )

    document.save(out_path)


def build_docx_trainer(out_path: Path, participants: list[Participant]) -> None:
    document = Document()
    _set_doc_defaults(document)

    document.add_heading("Trainer Pack — Northwind Shared Tenant", level=1)
    document.add_paragraph(COURSE_TITLE)
    document.add_paragraph(f"Scenario: {SCENARIO_TITLE}")
    document.add_paragraph(f"Date: {TODAY.isoformat()}")

    document.add_heading("Class roster (IDs)", level=2)
    roster_line = ", ".join(p.pxx for p in participants)
    document.add_paragraph(roster_line)

    document.add_heading("Shared-tenant ground rules", level=2)
    for rule in [
        "Participants only modify their own NW-Pxx artifacts.",
        "Tenant-wide settings and policy changes are trainer-led only.",
        "Guest invitations: keep to a minimum; use trainer-prepared accounts if possible.",
        "Term store is tenant-wide: participants must only create NW-Pxx-prefixed term groups/sets.",
    ]:
        document.add_paragraph(rule, style="List Bullet")

    document.add_heading("Module 4 quick runbook (Permissions)", level=2)
    for step in [
        "Confirm each participant has NW-Pxx-ProjectSite and can access as owner.",
        "Remind: break inheritance at library, and at one folder only (03-Final) to keep scopes low.",
        "If anyone locks themselves out: re-add as Full Control at library level (site owner/admin).",
        "Keep sharing link drills internal-only unless you explicitly approve external tests.",
    ]:
        document.add_paragraph(step, style="List Number")

    document.add_heading("Module 5 quick runbook (Term store + metadata)", level=2)
    for step in [
        "Confirm who has Term store admin / Group manager permissions.",
        "Enforce NW-Pxx naming in term store (NW-Pxx-TermGroup, NW-Pxx-ContractType).",
        "If participants cannot create term groups: instruct them to use the local term set fallback.",
        "Spot-check: term set is Available for tagging; column created in NW-Pxx-Contracts; documents tagged.",
    ]:
        document.add_paragraph(step, style="List Number")

    document.add_heading("Modules 6–11 prep reminders (seed content)", level=2)
    for note in [
        "Module 6: ensure indexing time is accounted for; use unique phrases in seed docs.",
        "Module 6: bookmarks are immediate after publishing; acronyms can take up to a day (plan accordingly).",
            "Module 7: list formatting exercises are participant-safe; Apps/API access are trainer-led only.",
            "Module 7: do not approve API access requests in a shared training tenant.",
        "Module 8: use FAKE content only; avoid real personal data.",
        "Module 9: OneDrive settings are tenant-wide; treat as trainer-led if changed.",
        "Module 10: PowerShell is powerful—keep participants on read-only reporting and ensure outputs are NW-Pxx scoped.",
        "Module 11: Operations is largely observe/document; use worksheets to drive safe change requests.",
    ]:
        document.add_paragraph(note, style="List Bullet")

    document.add_heading("Module 10 quick runbook (PowerShell/automation)", level=2)
    for step in [
        "Confirm participants know their tenant admin URL and NW-Pxx site URL.",
        "Emphasize read-only reporting first (Get-SPOSite, Get-SPOUser) + export to CSV.",
        "Reinforce scoping: filter/export only NW-Pxx artifacts; no tenant-wide changes.",
        "If Graph PowerShell is used: keep scopes minimal and avoid write operations unless explicitly assigned.",
    ]:
        document.add_paragraph(step, style="List Number")

    document.add_heading("Module 11 quick runbook (Operations at scale)", level=2)
    for step in [
        "Lead with Service health + Message center checks before deep troubleshooting.",
        "For access issues: use Check Permissions and validate identity vs link type.",
        "For lifecycle topics: keep delete/restore actions trainer-led; participants document guardrails.",
        "For external sharing governance: review link types/defaults; keep policy changes trainer-led.",
        "Collect outputs: worksheets + any NW-Pxx-scoped CSV reports.",
    ]:
        document.add_paragraph(step, style="List Number")

    document.save(out_path)


def build_pdf(participant: Participant, out_path: Path) -> None:
    c = canvas.Canvas(str(out_path), pagesize=LETTER)
    width, height = LETTER

    x = 72
    y = height - 72
    line_h = 14

    def draw(line: str, bold: bool = False) -> None:
        nonlocal y
        c.setFont("Helvetica-Bold" if bold else "Helvetica", 11 if not bold else 12)
        c.drawString(x, y, line)
        y -= line_h

    draw(f"Participant Pack — {participant.pxx}", bold=True)
    draw(COURSE_TITLE)
    draw(f"Scenario: {SCENARIO_TITLE}")
    draw(f"Date: {TODAY.isoformat()}")
    y -= line_h

    draw("Your lab naming:", bold=True)
    draw(f"- Participant ID: {participant.pxx}")
    draw(f"- Practice site: {participant.project_site_name}")
    draw(f"- Contracts library: {participant.contracts_library_name}")
    y -= line_h

    draw("Safety rules (shared tenant):", bold=True)
    for rule in [
        "Only work inside your own NW-Pxx artifacts.",
        "Tenant-wide settings are Trainer-only.",
        "External guest invites only if trainer approves.",
        "Minimize unique permissions; use library boundary first.",
    ]:
        draw(f"- {rule}")

    c.showPage()
    c.save()


def build_pdf_trainer(out_path: Path, participants: list[Participant]) -> None:
    c = canvas.Canvas(str(out_path), pagesize=LETTER)
    width, height = LETTER

    x = 72
    y = height - 72
    line_h = 14

    def draw(line: str, bold: bool = False) -> None:
        nonlocal y
        c.setFont("Helvetica-Bold" if bold else "Helvetica", 11 if not bold else 12)
        c.drawString(x, y, line)
        y -= line_h

    draw("Trainer Pack — Northwind Shared Tenant", bold=True)
    draw(COURSE_TITLE)
    draw(f"Scenario: {SCENARIO_TITLE}")
    draw(f"Date: {TODAY.isoformat()}")
    y -= line_h

    draw("Roster:", bold=True)
    draw(", ".join(p.pxx for p in participants))
    y -= line_h

    draw("Ground rules:", bold=True)
    for rule in [
        "Participants only modify their own NW-Pxx artifacts.",
        "Tenant-wide policy changes are trainer-led only.",
        "Term store: NW-Pxx naming only; no cross-editing.",
    ]:
        draw(f"- {rule}")

    c.showPage()
    c.save()


def build_pptx(participant: Participant, out_path: Path) -> None:
    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Northwind — Participant {participant.pxx}"
    slide.placeholders[1].text = COURSE_TITLE

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Your key artifacts"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for line in [
        f"Practice site: {participant.project_site_name}",
        f"Module 4 library: {participant.contracts_library_name}",
        "Folders: 01-Drafts / 02-InReview / 03-Final",
    ]:
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        p.text = line

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Safety rules (shared tenant)"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    bullets = [
        "Work only inside your own NW-Pxx content.",
        "No tenant-wide changes unless trainer-led.",
        "Keep unique permissions minimal (library boundary first).",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b

    # Footer textbox (simple)
    for s in prs.slides:
        box = s.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.3))
        box.text_frame.text = f"{SCENARIO_TITLE} | {TODAY.isoformat()} | {participant.pxx}"

    prs.save(out_path)


def build_pptx_trainer(out_path: Path, participants: list[Participant]) -> None:
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Trainer Pack — Northwind Shared Tenant"
    slide.placeholders[1].text = COURSE_TITLE

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Ground rules"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    bullets = [
        "Participants only modify their own NW-Pxx artifacts.",
        "Tenant-wide settings/policies are trainer-led only.",
        "Term store is tenant-wide: NW-Pxx naming only.",
    ]
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Roster and checkpoints"
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "Roster: " + ", ".join(p.pxx for p in participants)
    for item in [
        "M4: Library boundary + one folder unique permissions",
        "M5: Term group/set created (or local fallback) + managed metadata column",
        "M6–M9: use seed content; keep policy changes trainer-only",
        "M10: PowerShell exports only; NW-Pxx scoped",
        "M11: Ops worksheets completed; change requests documented",
        "M12 (optional): Approval flow + Power Apps custom form; governance notes captured",
    ]:
        p = tf.add_paragraph()
        p.text = item

    for s in prs.slides:
        box = s.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.3))
        box.text_frame.text = f"{SCENARIO_TITLE} | {TODAY.isoformat()} | TRAINER"

    prs.save(out_path)


def build_xlsx(participant: Participant, out_path: Path) -> None:
    wb = Workbook()

    ws = wb.active
    ws.title = "Tracker"

    header_font = Font(bold=True)
    ws["A1"].value = "Participant ID"
    ws["B1"].value = participant.pxx
    ws["A2"].value = "Practice site name"
    ws["B2"].value = participant.project_site_name
    ws["A3"].value = "Contracts library"
    ws["B3"].value = participant.contracts_library_name

    for cell in ("A1", "A2", "A3"):
        ws[cell].font = header_font

    ws["A5"].value = "Module"
    ws["B5"].value = "Task"
    ws["C5"].value = "Done (Y/N)"
    ws["D5"].value = "Notes"
    for col in ("A5", "B5", "C5", "D5"):
        ws[col].font = header_font
        ws[col].alignment = Alignment(vertical="center")

    tasks = [
        ("M3", "Confirm NW-Pxx-ProjectSite URL", "", ""),
        ("M4", "Create NW-Pxx-Contracts library", "", ""),
        ("M4", "Break inheritance at library", "", ""),
        ("M4", "Break inheritance at one folder only (03-Final)", "", ""),
        ("M4", "Validate access: Check Permissions", "", ""),
        ("M4", "Remove any ad-hoc sharing links", "", ""),
        ("M5", "Create NW-Pxx-ContractType term set (or local fallback)", "", ""),
        ("M5", "Add NW-Pxx-ContractType column + tag 3 docs", "", ""),
        ("M6", "Upload seed docs to NW-Pxx-Contracts", "", ""),
        ("M6", "Validate search finds seed docs", "", ""),
        ("M7", "Create NW-Pxx-AppRequests list + apply formatting", "", ""),
        ("M8", "Apply sensitivity/retention labels to fake doc", "", ""),
        ("M9", "Complete OneDrive settings observation worksheet", "", ""),
        ("M10", "Export NW-Pxx sites + users reports (PowerShell)", "", ""),
        ("M11", "Complete ops worksheets + change request summary", "", ""),
        ("M12", "Build approval flow + customize form (Power Platform)", "", ""),
    ]

    start_row = 6
    for i, row in enumerate(tasks):
        r = start_row + i
        ws[f"A{r}"].value = row[0]
        ws[f"B{r}"].value = row[1]
        ws[f"C{r}"].value = row[2]
        ws[f"D{r}"].value = row[3]

    ws.column_dimensions["A"].width = 10
    ws.column_dimensions["B"].width = 50
    ws.column_dimensions["C"].width = 12
    ws.column_dimensions["D"].width = 40

    ws.freeze_panes = "A6"

    ws2 = wb.create_sheet("Links")
    ws2["A1"].value = "Reference"
    ws2["B1"].value = "URL"
    ws2["A1"].font = header_font
    ws2["B1"].font = header_font

    links = [
        ("Manage permission scopes", "https://learn.microsoft.com/en-us/sharepoint/manage-permission-scope"),
        ("Modern sharing & permissions", "https://learn.microsoft.com/en-us/sharepoint/modern-experience-sharing-permissions"),
        ("Shareable links", "https://learn.microsoft.com/en-us/sharepoint/shareable-links-anyone-specific-people-organization"),
        ("Troubleshoot Access Denied", "https://learn.microsoft.com/en-us/troubleshoot/sharepoint/administration/access-denied-or-need-permission-error-sharepoint-online-or-onedrive-for-business"),
    ]
    for i, (name, url) in enumerate(links, start=2):
        ws2[f"A{i}"].value = name
        ws2[f"B{i}"].value = url

    ws2.column_dimensions["A"].width = 28
    ws2.column_dimensions["B"].width = 90

    wb.save(out_path)


def build_xlsx_trainer(out_path: Path, participants: list[Participant]) -> None:
    wb = Workbook()

    header_font = Font(bold=True)

    ws = wb.active
    ws.title = "Roster"
    ws["A1"].value = "Participant ID"
    ws["B1"].value = "Practice site name"
    ws["C1"].value = "Notes"
    for cell in ("A1", "B1", "C1"):
        ws[cell].font = header_font

    for i, p in enumerate(participants, start=2):
        ws[f"A{i}"].value = p.pxx
        ws[f"B{i}"].value = p.project_site_name
        ws[f"C{i}"].value = ""

    ws.column_dimensions["A"].width = 14
    ws.column_dimensions["B"].width = 26
    ws.column_dimensions["C"].width = 50

    ws2 = wb.create_sheet("Module Checks")
    ws2["A1"].value = "Module"
    ws2["B1"].value = "Checkpoint"
    ws2["C1"].value = "Status"
    for cell in ("A1", "B1", "C1"):
        ws2[cell].font = header_font

    checks = [
        ("M4", "NW-Pxx-Contracts exists; inheritance broken at library", ""),
        ("M4", "Only one folder has unique permissions (03-Final)", ""),
        ("M5", "NW-Pxx-ContractType term set exists (or local fallback)", ""),
        ("M5", "Managed metadata column added + 3 docs tagged", ""),
        ("M6", "Seed docs uploaded + query works (bookmark optional)", ""),
        ("M7", "App governance artifact uploaded", ""),
        ("M8", "Fake compliance content available", ""),
        ("M9", "OneDrive test file available", ""),
        ("M10", "PowerShell reporting outputs collected (CSV)", ""),
        ("M11", "Ops worksheets collected; governance change requests documented", ""),
        ("M12", "Approval flow + Power Apps form customization completed (optional)", ""),
    ]
    for i, (m, chk, st) in enumerate(checks, start=2):
        ws2[f"A{i}"].value = m
        ws2[f"B{i}"].value = chk
        ws2[f"C{i}"].value = st

    ws2.column_dimensions["A"].width = 8
    ws2.column_dimensions["B"].width = 62
    ws2.column_dimensions["C"].width = 14

    ws3 = wb.create_sheet("Issue Log")
    ws3["A1"].value = "Time"
    ws3["B1"].value = "Participant"
    ws3["C1"].value = "Issue"
    ws3["D1"].value = "Resolution"
    for cell in ("A1", "B1", "C1", "D1"):
        ws3[cell].font = header_font

    ws3.column_dimensions["A"].width = 18
    ws3.column_dimensions["B"].width = 12
    ws3.column_dimensions["C"].width = 45
    ws3.column_dimensions["D"].width = 45

    wb.save(out_path)


def build_txt_templates_trainer(out_dir: Path, participants: list[Participant]) -> None:
    templates_dir = out_dir / "TXT-Templates"
    templates_dir.mkdir(parents=True, exist_ok=True)

    header = (
        "Trainer Pack — Templates\n"
        f"Course: {COURSE_TITLE}\n"
        f"Scenario: {SCENARIO_TITLE}\n"
        f"Date: {TODAY.isoformat()}\n\n"
    )

    (templates_dir / "Trainer-Announcements.txt").write_text(
        header
        + "Opening script (2–3 minutes)\n"
        + "- Welcome to the Northwind shared-tenant lab environment.\n"
        + "- Your Participant ID is P01–P10. Always prefix artifacts with NW-Pxx.\n"
        + "- Tenant-wide settings are trainer-led only.\n"
        + "- If you get stuck, capture the symptom and where you checked.\n",
        encoding="utf-8",
    )

    # Rename to reflect that this runbook spans multiple modules.
    # Remove the old filename (from earlier generations) to prevent confusion.
    (templates_dir / "Trainer-Module-4-5-Runbook.txt").unlink(missing_ok=True)

    # Remove the previous multi-module runbook name (from earlier generations)
    (templates_dir / "Trainer-Modules-04-09-Runbook.txt").unlink(missing_ok=True)

    # Remove the previous runbook name (from earlier generations)
    (templates_dir / "Trainer-Modules-04-10-Runbook.txt").unlink(missing_ok=True)

    # Remove the previous runbook name (from earlier generations)
    (templates_dir / "Trainer-Modules-04-11-Runbook.txt").unlink(missing_ok=True)

    (templates_dir / "Trainer-Modules-04-12-Runbook.txt").write_text(
        header
        + "Module 4 — Permissions (trainer runbook)\n"
        + "- Confirm each participant has NW-Pxx-ProjectSite.\n"
        + "- Emphasize: break inheritance at library; 1 folder only unique perms.\n"
        + "- Keep sharing drills internal-only unless you explicitly approve.\n\n"
        + "Module 5 — Term store + managed metadata (trainer runbook)\n"
        + "- Term store is tenant-wide: enforce NW-Pxx naming and isolation.\n"
        + "- If participants cannot create term groups, direct them to local fallback.\n"
        + "- Spot-check tagging: 3 docs tagged with ContractType terms.\n"
        + "\nModule 6 — Search (trainer runbook)\n"
        + "- Expect indexing delay; seed docs include participant-specific phrases.\n"
        + "- Bookmarks: visible immediately after publish; enforce NW-Pxx-only keywords.\n"
        + "- Acronyms: can take up to a day after publish; treat as concept if needed.\n"
        + "\nModule 7 — Apps and customization (trainer runbook)\n"
        + "- Participant exercise is list/view formatting (JSON) only; no app deployment required.\n"
        + "- Treat Apps and API access pages as trainer-led read-only tours in shared tenant.\n"
        + "- Do not approve API access requests during training.\n"
        + "\nModule 8 — Purview (trainer runbook)\n"
        + "- Publish training labels ahead of time if possible; plan for replication delays.\n"
        + "- Keep participant tasks scoped to their own NW-Pxx site content (apply labels only).\n"
        + "- eDiscovery: default delivery is demo-only; if hands-on, assign permissions and limit scope.\n"
        + "- DLP: do not enable blocking actions in a shared training tenant; use simulation-only if demonstrated.\n"
        + "\nModule 9 — OneDrive administration (trainer runbook)\n"
        + "- OneDrive settings are managed in SharePoint admin center and are tenant-wide: treat as trainer-led for changes.\n"
        + "- Participants should do read-only verification + internal-only OneDrive sharing test (no guests).\n"
        + "- If demonstrating unmanaged devices controls: note it can take up to 24 hours to apply and won't affect already signed-in users.\n"
        + "- If demonstrating sync governance: reinforce OneDrive shortcuts vs Sync, and clarify that hiding Sync blocks new syncs only.\n"
        + "\nModule 10 — PowerShell and automation (trainer runbook)\n"
        + "- Emphasize read-only reporting first; participants export CSV outputs for review.\n"
        + "- Treat tenant-wide changes as trainer-led only in a shared tenant.\n"
        + "- If using Microsoft Graph PowerShell: confirm scopes and consent behavior; avoid write operations unless explicitly assigned.\n"
        + "\nModule 11 — Operations at scale (trainer runbook)\n"
        + "- Lead with Service health and Message center: check for incidents/advisories and planned changes.\n"
        + "- Access issues: validate identity vs link type and use Check Permissions before changing anything.\n"
        + "- Lifecycle: keep delete/restore actions trainer-led; participants document guardrails and restore expectations.\n"
        + "- External sharing: review link types/defaults; keep org-level and site-level changes trainer-led.\n"
        + "\nModule 12 (Optional) — Power Automate + Power Apps (trainer runbook)\n"
        + "- Keep participant scope to their NW-Pxx-AppRequests list only.\n"
        + "- Approver should be self or trainer-designated to avoid spamming.\n"
        + "- Prefer 'When an item is created' trigger to avoid update loops.\n"
        + "- Environments and DLP policy changes are trainer-led demos unless explicitly assigned.\n",
        encoding="utf-8",
    )

    (templates_dir / "Trainer-Roster.txt").write_text(
        header + "Roster\n" + "\n".join(f"- {p.pxx}: {p.project_site_name}" for p in participants) + "\n",
        encoding="utf-8",
    )


def build_txt_templates(participant: Participant, out_dir: Path) -> None:
    templates_dir = out_dir / "TXT-Templates"
    templates_dir.mkdir(parents=True, exist_ok=True)

    common_header = (
        f"Northwind Training Scenario — {participant.pxx}\n"
        f"Course: {COURSE_TITLE}\n"
        f"Scenario: {SCENARIO_TITLE}\n"
        f"Date: {TODAY.isoformat()}\n"
        "\n"
        "Shared-tenant note: Use only your own NW-Pxx artifacts.\n"
        "\n"
    )

    files: dict[str, str] = {}

    # Module 4 — Contracts content (paste into Word docs)
    files["M04-Contract-Draft-001.txt"] = (
        common_header
        + "Module 4 (Permissions) — Document Template\n"
        + "Intended location: Library NW-Pxx-Contracts / Folder 01-Drafts\n\n"
        + "Document title: Contract Draft 001 (Training)\n"
        + "Status: Draft\n"
        + "Owner: Legal\n"
        + "\n"
        + "Purpose\n"
        + "This draft contract is used for permissions and collaboration drills.\n\n"
        + "Northwind ↔ Fabrikam Collaboration (Training Only)\n"
        + "- Counterparty: Fabrikam (external partner in story; do not invite guests unless trainer approves)\n"
        + "- Scope: Pilot collaboration for intranet modernization\n\n"
        + "Key clauses (draft)\n"
        + "1) Confidentiality: Internal discussion only; not final.\n"
        + "2) Deliverables: Discovery workshop, IA proposal, phased rollout plan.\n"
        + "3) Timeline: 4–6 weeks (draft estimate).\n"
        + "4) Notes: Replace placeholders before using in real life.\n\n"
        + "Revision notes\n"
        + "- Add comments here during editing exercises.\n"
    )

    files["M04-Contract-Final-001.txt"] = (
        common_header
        + "Module 4 (Permissions) — Document Template\n"
        + "Intended location: Library NW-Pxx-Contracts / Folder 03-Final\n\n"
        + "Document title: Contract Final 001 (Training)\n"
        + "Status: Final\n"
        + "Approved by: (training)\n"
        + "\n"
        + "Summary\n"
        + "This is the FINAL version used for folder-level permission drills (editors read-only).\n\n"
        + "Agreement Summary (Training Only)\n"
        + "- Parties: Northwind Traders and Fabrikam\n"
        + "- Effective date: TBD\n"
        + "- Access: Broad read access in Final folder; controlled editing\n\n"
        + "Sign-off\n"
        + "Northwind: ____________________\n"
        + "Fabrikam:  ____________________\n"
    )

    files["M04-Contract-Changelog.txt"] = (
        common_header
        + "Module 4 (Permissions) — Change Log Template\n"
        + "Intended location: Library NW-Pxx-Contracts / Folder 02-InReview\n\n"
        + "Change Log\n"
        + "- Date | Version | Author | Summary\n"
        + f"- {TODAY.isoformat()} | v0.1 | {participant.pxx} | Initial training draft\n"
    )

    # Module 5 — Metadata
    files["M05-TermSet-ContractType.txt"] = (
        common_header
        + "Module 5 (Metadata) — Term Set Design\n"
        + f"Term group: {participant.term_group_name}\n"
        + f"Term set:  {participant.termset_contract_type}\n\n"
        + "Baseline terms\n"
        + "- NDA\n"
        + "- MSA\n"
        + "- SOW\n"
        + "- Renewal\n\n"
        + "Notes\n"
        + "- Recommended: keep the term set Closed for controlled tagging in training.\n"
    )

    files["M05-Metadata-Tagging-Worksheet.txt"] = (
        common_header
        + "Module 5 (Metadata) — Tagging Worksheet\n"
        + f"Site: {participant.project_site_name}\n"
        + f"Library: {participant.contracts_library_name}\n"
        + "Column: NW-Pxx-ContractType (managed metadata)\n\n"
        + "Tag at least 3 documents:\n"
        + "1) Document: ____________________ | ContractType: ________\n"
        + "2) Document: ____________________ | ContractType: ________\n"
        + "3) Document: ____________________ | ContractType: ________\n\n"
        + "What did you observe?\n"
        + "- Type-ahead suggestions: Yes/No\n"
        + "- Fill-in choices allowed: Yes/No/Not sure\n"
    )

    # Module 6 — Search (seed content and queries)
    unique_alpha = f"Northwind Search Drill Alpha - {participant.pxx}"
    unique_beta = f"Northwind Search Drill Beta - {participant.pxx}"

    files["M06-Search-Seed-Document-A.txt"] = (
        common_header
        + "Module 6 (Search) — Seed Document A\n"
        + "Paste into a new document and upload to NW-Pxx-Contracts (or a lab library).\n\n"
        + "Unique phrases to search for:\n"
        + f"- {unique_alpha}\n"
        + "- ContractTypeRef: NDA\n"
        + "- ProjectCodename: HARBORLIGHT\n\n"
        + "Body\n"
        + "This document exists to validate indexing and search behavior in training.\n"
    )

    files["M06-Search-Seed-Document-B.txt"] = (
        common_header
        + "Module 6 (Search) — Seed Document B\n"
        + "Unique phrases to search for:\n"
        + f"- {unique_beta}\n"
        + "- ContractTypeRef: SOW\n"
        + "- ProjectCodename: RIVERSTONE\n\n"
        + "Body\n"
        + "Use this as a second document to compare search results and filters.\n"
    )

    files["M06-Search-Test-Queries.txt"] = (
        common_header
        + "Module 6 (Search) — Test Queries\n"
        + "Use these queries when validating search across SharePoint/Microsoft Search.\n\n"
        + "Queries\n"
        + f"- \"{unique_alpha}\"\n"
        + f"- \"{unique_beta}\"\n"
        + "- HARBORLIGHT\n"
        + "- RIVERSTONE\n"
        + "- ContractTypeRef\n"
        + "\nOptional (advanced / may vary by experience)\n"
        + "- filetype:docx\n"
    )

    files["M06-Search-Answers-Config.txt"] = (
        common_header
        + "Module 6 (Search) — Microsoft Search Answers (Config Worksheet)\n"
            + "Use this worksheet if you have Search admin/editor permissions.\n\n"
            + "Permissions and access\n"
            + "- To create/edit Bookmarks and Acronyms, you need one of these Microsoft 365 roles:\n"
            + "  - Search admin, or\n"
            + "  - Search editor\n"
            + "- These roles are assigned by a Global admin.\n\n"
        + "Bookmark\n"
        + f"- Title: NW-{participant.pxx} Project Site\n"
        + "- URL: (paste your NW-Pxx-ProjectSite home page URL)\n"
        + "- Keywords (training-only; include Pxx):\n"
        + f"  - NW {participant.pxx} ProjectSite\n"
        + f"  - Northwind {participant.pxx} site\n"
        + f"  - NW-{participant.pxx} Project Site\n\n"
        + "Acronym\n"
        + f"- Acronym: NW{participant.pxx} (no spaces)\n"
        + f"- Stands for: Northwind Participant {participant.pxx}\n"
        + "- Source URL: (paste your NW-Pxx-ProjectSite URL)\n\n"
        + "Operational notes\n"
        + "- Bookmarks: visible immediately after publishing.\n"
        + "- Acronyms: can take up to a day after publishing.\n"
    )

    # Module 7 — Apps/customization (paperwork-style content, safe to upload)
    files["M07-App-Request-Form.txt"] = (
        common_header
        + "Module 7 (Apps) — App Request Form (Training)\n"
        + "Purpose: A simple document you can upload and use in app governance discussions.\n\n"
        + "Request summary\n"
        + "- Requestor department: ________\n"
        + "- Business need: Add a lightweight page component for contracts dashboard\n"
        + "- Data sensitivity: Internal\n"
        + "- Target site: NW-Pxx-ProjectSite\n\n"
        + "Approval checklist\n"
        + "- Security review completed: Y/N\n"
        + "- Licensing confirmed: Y/N\n"
        + "- Rollback plan documented: Y/N\n"
    )

    files["M07-AppRequests-Formatting-JSON.txt"] = (
        common_header
        + "Module 7 (Apps/Customization) — AppRequests Formatting (JSON)\n"
        + "Use this template to paste JSON into SharePoint list formatting panels.\n"
        + "List: NW-Pxx-AppRequests\n\n"
        + "What this JSON is (and is not)\n"
        + "- This is SharePoint list formatting (declarative customization).\n"
        + "- It changes how fields/rows are displayed (icons, colors, layout).\n"
        + "- It does NOT change the underlying list data.\n\n"
        + "Permissions\n"
        + "- You typically need Site Owner / Full Control on NW-Pxx-ProjectSite to format columns/views.\n\n"
        + "No-JSON option (no code)\n"
        + "- In many tenants, the formatting pane includes a visual/rules editor that generates JSON for you.\n"
        + "- If you see options like Rules / Conditional formatting / Style, use those first.\n"
        + "- If you do not see those options, your tenant/UI may require Advanced mode (JSON).\n\n"
        + "A) Column formatting — Status column (Format this column)\n"
        + "What it does (summary)\n"
        + "- Uses the Status value to pick a severity style and an icon.\n"
        + "- Shows the icon + the Status text in the cell.\n\n"
        + "Without typing JSON (if available)\n"
        + "1) Status column menu → Column settings → Format this column\n"
        + "2) Choose Rules/Conditional formatting\n"
        + "3) Create rules: Approved=success, In review=warning, Submitted=neutral, Rejected=blocked\n"
        + "4) (Optional) Pick icons per rule → Preview → Save\n\n"
        + "Paste and Save:\n"
        + "{\n"
        + "  \"$schema\": \"https://developer.microsoft.com/json-schemas/sp/v2/column-formatting.schema.json\",\n"
        + "  \"elmType\": \"div\",\n"
        + "  \"attributes\": {\n"
        + "    \"class\": \"=if(@currentField == 'Approved', 'sp-field-severity--good', if(@currentField == 'In review', 'sp-field-severity--warning', if(@currentField == 'Submitted', 'sp-field-severity--low', 'sp-field-severity--blocked'))) + ' ms-fontColor-neutralSecondary'\"\n"
        + "  },\n"
        + "  \"children\": [\n"
        + "    {\n"
        + "      \"elmType\": \"span\",\n"
        + "      \"style\": { \"display\": \"inline-block\", \"padding\": \"0 4px\" },\n"
        + "      \"attributes\": {\n"
        + "        \"iconName\": \"=if(@currentField == 'Approved', 'CheckMark', if(@currentField == 'In review', 'Error', if(@currentField == 'Submitted', 'Forward', 'ErrorBadge')))\"\n"
        + "      }\n"
        + "    },\n"
        + "    {\n"
        + "      \"elmType\": \"span\",\n"
        + "      \"txtContent\": \"@currentField\"\n"
        + "    }\n"
        + "  ]\n"
        + "}\n\n"
        + "B) View formatting — current view (Format current view)\n"
        + "What it does (summary)\n"
        + "- Applies an extra row class to alternating rows (zebra striping) using @rowIndex.\n\n"
        + "Without typing JSON (if available)\n"
        + "1) View dropdown → Format current view\n"
        + "2) Look for Alternating rows / Row styling / a preset with zebra shading\n"
        + "3) Save\n\n"
        + "Paste and Save:\n"
        + "{\n"
        + "  \"$schema\": \"https://developer.microsoft.com/json-schemas/sp/view-formatting.schema.json\",\n"
        + "  \"additionalRowClass\": \"=if(@rowIndex % 2 == 0,'ms-bgColor-themeLighter ms-bgColor-themeLight--hover','')\"\n"
        + "}\n"
    )

    # Module 8 — Purview/compliance (safe fake content)
    files["M08-Confidential-Statement-FAKE.txt"] = (
        common_header
        + "Module 8 (Purview) — Confidential Statement (FAKE)\n"
        + "Use this for sensitivity/retention demonstrations if your lab instructions require content.\n"
        + "Do not use real personal data.\n\n"
        + "CONFIDENTIAL — NORTHWIND (TRAINING ONLY)\n"
        + "This document is intended to be classified as internal/confidential in training.\n\n"
        + "Project: Intranet Modernization\n"
        + "Document type: Decision memo\n"
        + "Audience: Northwind leadership team\n"
    )

    files["M08-DLP-Test-Data-FAKE.txt"] = (
        common_header
        + "Module 8 (Purview) — DLP Test Data (FAKE)\n"
        + "WARNING: Only use if the trainer explicitly enables a DLP test exercise.\n"
        + "All data below is fake and for pattern-testing only.\n\n"
        + "Fake identifiers (training)\n"
        + "- Employee ID: NW-EMP-000123\n"
        + "- Invoice ID:  NW-INV-2026-000045\n"
        + "- Account (fake): ACCT-0000-0000\n\n"
        + "Notes\n"
        + "- If your tenant has DLP policies, this content might be detected depending on configuration.\n"
    )

    files["M08-Labeling-Worksheet.txt"] = (
        common_header
        + "Module 8 (Purview) — Labeling Worksheet\n"
        + f"Site: {participant.project_site_name}\n"
        + f"Library: {participant.contracts_library_name}\n"
        + "Target document: NW-Pxx-Confidential-Statement.docx\n\n"
        + "Sensitivity label\n"
        + "- Label name applied: ____________________\n"
        + "- Where applied (SharePoint/Office): ____________________\n\n"
        + "Retention label\n"
        + "- Label name applied: ____________________\n"
        + "- Where applied (SharePoint properties/other): ____________________\n\n"
        + "Notes (timing/replication)\n"
        + "- Did labels appear immediately? Yes/No\n"
        + "- If not, how long did it take? ____________________\n"
    )

    files["M08-eDiscovery-Case-Notes.txt"] = (
        common_header
        + "Module 8 (Purview) — eDiscovery Case Notes (Worksheet)\n"
        + "Default delivery: trainer demo. Hands-on only if trainer assigns permissions.\n\n"
        + "Case planning\n"
        + f"- Proposed case name: NW-{participant.pxx}-eDiscovery-Case\n"
        + f"- SharePoint site in scope: {participant.project_site_name} (only)\n"
        + "- Keywords to test: CONFIDENTIAL\n\n"
        + "Search results\n"
        + "- Did search find your document? Yes/No\n"
        + "- Any false positives? ____________________\n\n"
        + "Hold decision (trainer-approved only)\n"
        + "- Hold created? Yes/No\n"
        + "- Scope limited to NW-Pxx locations only? Yes/No\n"
    )

    # Module 9 — OneDrive admin (safe personal-work style content)
    files["M09-OneDrive-Settings-Observation.txt"] = (
        common_header
        + "Module 9 (OneDrive) — Settings Observation Worksheet\n"
        + "Purpose: Record current tenant settings (read-only verification).\n"
        + "Shared-tenant note: Do NOT change tenant-wide settings unless trainer instructs.\n\n"
        + "A) Where OneDrive settings live\n"
        + "- SharePoint admin center pages checked: Sharing / Settings / Access control\n"
        + "- Notes: ____________________\n\n"
        + "B) Sharing (organization level)\n"
        + "- External sharing level observed: ____________________\n"
        + "- Default link type (if shown): ____________________\n"
        + "- Default permissions (if shown): ____________________\n\n"
        + "C) Sync\n"
        + "- Sync button shown on OneDrive website? Yes/No/Not sure\n"
        + "- Any sync restrictions shown (domains/file types)? ____________________\n\n"
        + "D) Storage limit\n"
        + "- Default OneDrive storage limit (if shown): ____________________\n"
        + "- Notes/warnings observed: ____________________\n\n"
        + "E) Retention (deleted users)\n"
        + "- Days to retain files a deleted user's OneDrive: ________\n"
        + "- Notes: retention starts when user is deleted; restore scenarios are admin-led\n\n"
        + "F) Access control (unmanaged devices)\n"
        + "- Current unmanaged devices setting: Allow full / Limited web-only / Block\n"
        + "- Operational note: changes can take up to 24 hours and won't affect already signed-in users\n"
    )

    files["M09-OneDrive-Sharing-Test.txt"] = (
        common_header
        + "Module 9 (OneDrive) — Sharing Test Note\n"
        + "Paste into a new file and store in OneDrive (if your lab includes OneDrive sharing/sync).\n\n"
        + "Title: Personal Work Notes (Training)\n"
        + "- This file is used to observe OneDrive sharing behavior under policy.\n"
        + "- Keep it non-sensitive and training-only.\n"
    )

    # Module 10 — PowerShell / automation
    files["M10-PowerShell-Setup-Notes.txt"] = (
        common_header
        + "Module 10 (PowerShell) — Setup Notes\n"
        + "Goal: Record how you connected and what context you used.\n\n"
        + "SharePoint Online PowerShell\n"
        + "- Admin URL used: https://<tenant>-admin.sharepoint.com\n"
        + "- Connect-SPOService succeeded: Yes/No\n"
        + "- Any errors/warnings: ____________________\n\n"
        + "Microsoft Graph PowerShell (optional, trainer-approved)\n"
        + "- Scopes requested (example: User.Read.All): ____________________\n"
        + "- Admin consent prompt occurred: Yes/No/Not sure\n"
    )

    files["M10-Reporting-Worksheet.txt"] = (
        common_header
        + "Module 10 (PowerShell) — Reporting Worksheet\n"
        + "Purpose: Capture what you exported and validate safe scoping.\n\n"
        + "A) Site inventory export\n"
        + "- Output file: NW-Pxx-Sites.csv\n"
        + "- Filter used to scope to NW-Pxx: ____________________\n"
        + "- Count of sites in output: ________\n\n"
        + "B) Site users export\n"
        + "- Target site URL: https://<tenant>.sharepoint.com/sites/NW-Pxx-ProjectSite\n"
        + "- Output file: NW-Pxx-Users.csv\n"
        + "- Did output include your account? Yes/No\n\n"
        + "C) CSV-driven loop\n"
        + "- Targets.csv created with only your NW-Pxx site? Yes/No\n"
        + "- Output file: Targets-Users.csv\n\n"
        + "Safety check\n"
        + "- Did you avoid tenant-wide changes? Yes/No\n"
        + "- Notes: ____________________\n"
    )

    # Module 11 — Operations at scale (worksheets)
    files["M11-Incident-Triage-Worksheet.txt"] = (
        common_header
        + "Module 11 (Ops) — Incident Triage Worksheet\n"
        + "Goal: Capture the first 5–10 minutes of admin triage.\n\n"
        + "Service health\n"
        + "- Any active incident/advisory for SharePoint/OneDrive? Yes/No\n"
        + "- ID (if present): ____________________\n"
        + "- User impact summary: ____________________\n"
        + "- Last updated: ____________________\n\n"
        + "Message center\n"
        + "- Any recent messages affecting SharePoint/OneDrive? Yes/No\n"
        + "- Message title(s): ____________________\n\n"
        + "Decision\n"
        + "- Wait/communicate vs troubleshoot locally: ____________________\n"
    )

    files["M11-Site-Lifecycle-Checklist.txt"] = (
        common_header
        + "Module 11 (Ops) — Site Lifecycle Checklist\n"
        + "Purpose: Document guardrails and restore expectations (shared tenant safe).\n\n"
        + "Restore model\n"
        + "- Who can restore deleted sites? ____________________\n"
        + "- Deleted sites retention window (days): ________\n\n"
        + "Guardrails\n"
        + "- Do NOT delete root site (high impact): Confirmed\n"
        + "- Delete/restore drills are trainer-led unless assigned: Confirmed\n\n"
        + "Notes\n"
        + "- ____________________\n"
    )

    files["M11-External-Sharing-Governance-Worksheet.txt"] = (
        common_header
        + "Module 11 (Ops) — External Sharing Governance Worksheet\n"
        + "Purpose: Record sharing posture and defaults (read-only verification).\n\n"
        + "Organization-level\n"
        + "- External sharing level (SharePoint): ____________________\n"
        + "- External sharing level (OneDrive): ____________________\n"
        + "- Default link type: ____________________\n"
        + "- Anyone link expiration/permissions (if enabled): ____________________\n\n"
        + "Site-level (your NW-Pxx site)\n"
        + "- Can site be more permissive than org-level? No\n"
        + "- Any site-level restriction observed (do not change): ____________________\n\n"
        + "Change request draft (trainer-led approval required)\n"
        + "- Proposed change: ____________________\n"
        + "- Risk/impact: ____________________\n"
        + "- Rollback plan: ____________________\n"
    )

    files["M11-Change-Request-Template.txt"] = (
        common_header
        + "Module 11 (Ops) — Change Request Template (Training)\n"
        + "Purpose: Convert observations into a safe, reviewable change request.\n\n"
        + "Change summary\n"
        + "- Title: ____________________\n"
        + "- Requested by (Pxx/team): ____________________\n"
        + "- Date: ____________________\n\n"
        + "Scope\n"
        + "- Org-level or site-level? ____________________\n"
        + "- Target (site URL / policy name): ____________________\n\n"
        + "Current state (evidence)\n"
        + "- What did you observe? ____________________\n"
        + "- Where did you verify it (admin center page / report)? ____________________\n\n"
        + "Proposed change\n"
        + "- Setting/value to change: ____________________\n"
        + "- Desired value: ____________________\n\n"
        + "Risk and impact\n"
        + "- Who is impacted? ____________________\n"
        + "- Risk level (Low/Med/High): ________\n"
        + "- Notes: ____________________\n\n"
        + "Rollback plan\n"
        + "- How to revert: ____________________\n"
        + "- Validation after change: ____________________\n\n"
        + "Approvals (trainer-led in shared tenant)\n"
        + "- Approved by: ____________________\n"
        + "- Approved time: ____________________\n"
    )

    # Module 12 (optional) — Power Automate + Power Apps (worksheets)
    files["M12-Flow-Design-Worksheet.txt"] = (
        common_header
        + "Module 12 (Optional) — Power Automate Flow Design Worksheet\n"
        + "Goal: Document a simple approval flow connected to your NW-Pxx-AppRequests list.\n\n"
        + "Flow identity\n"
        + "- Proposed flow name: NW-Pxx-AppRequests-Approval\n"
        + "- Environment (if known): ____________________\n"
        + "- Owner: ____________________\n"
        + "- Backup owner: ____________________\n\n"
        + "Trigger\n"
        + "- Trigger used (created vs created/modified): ____________________\n"
        + "- Site address: https://<tenant>.sharepoint.com/sites/NW-Pxx-ProjectSite\n"
        + "- List name: NW-Pxx-AppRequests\n\n"
        + "Approval action\n"
        + "- Approval type (First to respond / Everyone must approve / etc.): ____________________\n"
        + "- Assigned to (who receives approval): ____________________\n"
        + "- Title format used: ____________________\n\n"
        + "Outcome handling\n"
        + "- Approved → Status value: Approved\n"
        + "- Rejected → Status value: Rejected\n"
        + "- Comments stored in Notes? Yes/No\n\n"
        + "Safety checks\n"
        + "- Approver is self or trainer-designated (avoid spamming): Yes/No\n"
        + "- Loop risk addressed (trigger choice/conditions): Yes/No\n"
    )

    files["M12-PowerApps-Form-Customization-Checklist.txt"] = (
        common_header
        + "Module 12 (Optional) — Power Apps Form Customization Checklist\n"
        + "Target list: NW-Pxx-AppRequests\n\n"
        + "Access\n"
        + "- Can open Integrate → Power Apps → Customize forms: Yes/No\n\n"
        + "Customization (choose at least one)\n"
        + "- Re-ordered fields to put Status/Owner near the top: Yes/No\n"
        + "- Added a label/header text for the form: Yes/No\n"
        + "- Added basic conditional visibility (optional): Yes/No\n\n"
        + "Publish and validation\n"
        + "- Saved successfully in Power Apps Studio: Yes/No\n"
        + "- Published to SharePoint successfully: Yes/No\n"
        + "- Verified form change in the list item UI: Yes/No\n\n"
        + "Notes\n"
        + "- ____________________\n"
    )

    files["M12-Governance-DLP-Notes.txt"] = (
        common_header
        + "Module 12 (Optional) — Governance Notes (Environments + DLP)\n"
        + "Purpose: Record what governance constraints apply to makers and connectors.\n\n"
        + "Environments\n"
        + "- Where are flows/apps being created for this lab? Default / Other: ____________________\n"
        + "- Who can create environments in this tenant? ____________________\n\n"
        + "DLP (Data Loss Prevention)\n"
        + "- Data policy name observed (if trainer shows): ____________________\n"
        + "- Key connectors allowed for productivity flows (example): SharePoint, Outlook, Approvals\n"
        + "- Any blocked connectors discussed: ____________________\n\n"
        + "Operational notes\n"
        + "- What happens when a flow violates DLP? ____________________\n"
        + "- Who approves exceptions/changes? ____________________\n"
    )

    for filename, content in files.items():
        (templates_dir / filename).write_text(content, encoding="utf-8")


def main() -> None:
    root = Path(__file__).resolve().parent

    participants = [Participant(f"P{i:02d}") for i in range(1, 11)]

    for participant in participants:
        out_dir = root / participant.pxx
        out_dir.mkdir(parents=True, exist_ok=True)

        base = f"{participant.pxx}-Training-Pack"

        build_docx(participant, out_dir / f"{base}.docx")
        build_pdf(participant, out_dir / f"{base}.pdf")
        build_pptx(participant, out_dir / f"{base}.pptx")
        build_xlsx(participant, out_dir / f"{base}.xlsx")
        build_txt_templates(participant, out_dir)

    trainer_dir = root / TRAINER_ID
    trainer_dir.mkdir(parents=True, exist_ok=True)
    build_docx_trainer(trainer_dir / "Trainer-Training-Pack.docx", participants)
    build_pdf_trainer(trainer_dir / "Trainer-Training-Pack.pdf", participants)
    build_pptx_trainer(trainer_dir / "Trainer-Training-Pack.pptx", participants)
    build_xlsx_trainer(trainer_dir / "Trainer-Training-Pack.xlsx", participants)
    build_txt_templates_trainer(trainer_dir, participants)


if __name__ == "__main__":
    main()
