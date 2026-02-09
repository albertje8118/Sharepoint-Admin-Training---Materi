"""
Generate Module 4 PPTX — Permissions and Collaboration Model
Modern, engaging design matching Modules 1-3 style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG     = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEC, 0xF9)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT    = RGBColor(0x24, 0x24, 0x24)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT  = RGBColor(0xD1, 0x34, 0x38)
YELLOW_WARN = RGBColor(0xF2, 0xC8, 0x11)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ═══════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════
def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6), icon="▸", font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon} {bullet}"; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.name = font_name; p.space_after = spacing
    return txBox

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)

def add_footer_bar(slide, slide_num, total, module_label="Module 4"):
    bar_top = SLIDE_HEIGHT - Inches(0.45)
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_WIDTH, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), bar_top + Inches(0.05), Inches(7), Inches(0.35),
                 f"{module_label}  |  Permissions and Collaboration Model",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), bar_top + Inches(0.05),
                 Inches(1.2), Inches(0.35), f"{slide_num} / {total}",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)

def section_divider(slide, section_title, subtitle="", icon=""):
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_TEAL)
    if icon:
        add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                     icon, font_size=60, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                 section_title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

TOTAL_SLIDES = 28
slide_counter = [0]

def new_slide(notes=""):
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if notes: add_speaker_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome to Module 4 — Permissions and Collaboration Model. This module covers the core admin skill: managing who can access what, using inheritance, groups, and sharing links. "
              "Permissions are the #1 topic in SharePoint support tickets, so this module is high-impact.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

# course label
add_text_box(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
             "MODERN SHAREPOINT ONLINE FOR ADMINISTRATORS  ·  DAY 2",
             font_size=14, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(2.3), Inches(11), Inches(1.5),
             "Module 4", font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(1),
             "Permissions and Collaboration Model",
             font_size=32, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

# tagline
add_text_box(s, Inches(2), Inches(5.2), Inches(9), Inches(0.6),
             "Control access  ·  Simplify governance  ·  Enable secure collaboration",
             font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(6.1), Inches(2.3), Inches(0.04), ACCENT_PURPLE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Why Permissions? (Overview / Motivation)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Permissions are the single most common support topic for SharePoint admins. "
              "Most 'Access Denied' tickets trace back to misconfigured permissions or broken inheritance. "
              "This overview sets the scene for why mastering permissions is essential.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Why Permissions Matter", font_size=36, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "The #1 admin support topic — understanding access is a core competency",
             font_size=16, color=MID_GRAY)

# 4 stat cards
cards = [
    ("🎫", "#1", "Support Topic", "Access-related tickets\ndominate helpdesk queues"),
    ("🔓", "50K", "Scope Limit", "Max unique permission scopes\nper document library"),
    ("⚡", "5K", "Recommended", "Keep unique scopes under\n5,000 for best performance"),
    ("👥", "3", "Default Groups", "Owners · Members · Visitors\ncreated with every site"),
]
for i, (icon, stat, label, desc) in enumerate(cards):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(s, x, Inches(1.9), Inches(2.8), Inches(3.6), WHITE)
    card.shadow.inherit = False
    add_text_box(s, x + Inches(0.3), Inches(2.1), Inches(2.2), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.8),
                 stat, font_size=44, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(0.5),
                 label, font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.6), Inches(4.1), Inches(1.6), Inches(0.04), ACCENT_TEAL)
    add_text_box(s, x + Inches(0.2), Inches(4.3), Inches(2.4), Inches(1.0),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Learning Outcomes
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Set clear expectations: by end of this module, learners can design and validate a permission model.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🎯  Learning Outcomes", font_size=32, bold=True, color=DARK_BG)

outcomes = [
    "Explain permission inheritance and permission scopes in practical admin terms",
    "Describe SharePoint groups vs Microsoft 365 group-connected permissions",
    "Design a permission model for a realistic collaboration scenario",
    "Implement and validate site/library/folder permissions safely in a shared tenant",
    "Explain how SharePoint and OneDrive sharing policies relate at the org and site level",
]
for i, outcome in enumerate(outcomes):
    y = Inches(1.5 + i * 1.05)
    badge = add_rounded_rect(s, Inches(0.8), y, Inches(0.55), Inches(0.55), ACCENT_BLUE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.5),
                 outcome, font_size=18, color=DARK_TEXT)

add_footer_bar(s, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Section Divider: Permission Fundamentals
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: let's begin with the building blocks of SharePoint permissions.")
section_divider(s, "Permission Fundamentals", "What every admin should know — and troubleshoot", "🔐")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Permission Building Blocks
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("These four concepts are the foundation. Principal = who. Permission Level = what they can do. "
              "Inheritance = where the rule comes from. Unique permissions = where the rule is overridden.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Four Building Blocks of SharePoint Permissions", font_size=30, bold=True, color=DARK_BG)

blocks = [
    ("👤", "Principal", "The identity you grant\npermissions to", "User, SharePoint Group,\nor Entra ID Group", ACCENT_BLUE),
    ("🔑", "Permission Level", "A named set of\npermissions", "Full Control · Edit\nRead · Contribute", ACCENT_TEAL),
    ("⬇️", "Inheritance", "Child inherits parent's\npermissions by default", "Site → Library →\nFolder → File", ACCENT_PURPLE),
    ("✂️", "Unique Permissions", "Break inheritance to\ncreate a new scope", "Object gets its own\nAccess Control List", ORANGE),
]
for i, (icon, title, desc, detail, color) in enumerate(blocks):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.6), Inches(2.9), Inches(4.5), WHITE)
    card.shadow.inherit = False
    # top color strip
    add_shape_rect(s, x, Inches(1.6), Inches(2.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(1.85), Inches(2.3), Inches(0.6),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.5), Inches(2.3), Inches(0.5),
                 title, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(1.0),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.6), Inches(4.2), Inches(1.7), Inches(0.03), LIGHT_GRAY)
    add_text_box(s, x + Inches(0.2), Inches(4.4), Inches(2.5), Inches(1.0),
                 detail, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Inheritance Flow (Visual)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Visualize the inheritance chain: Site → Library → Folder → File. "
              "Each level can inherit or break. When broken, a new permission scope is created.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Permission Inheritance Flow", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Permissions cascade downward — breaking creates a new scope",
             font_size=16, color=MID_GRAY)

# Chain of 4 levels
levels = [
    ("🌐", "Site", "Root permissions\n(Owners / Members / Visitors)", ACCENT_BLUE, True),
    ("📚", "Library", "Inherits from site\n(or unique if broken)", ACCENT_TEAL, True),
    ("📁", "Folder", "Inherits from library\n(break here for isolation)", ORANGE, False),
    ("📄", "File", "Inherits from folder\n(avoid per-file unique perms)", RED_ACCENT, False),
]
for i, (icon, label, desc, color, inherited) in enumerate(levels):
    x = Inches(0.8 + i * 3.15)
    y = Inches(2.2)
    card = add_rounded_rect(s, x, y, Inches(2.8), Inches(2.4), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(2.8), Inches(0.08), color)
    add_text_box(s, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.5),
                 f"{icon}  {label}", font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), y + Inches(0.9), Inches(2.4), Inches(1.0),
                 desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # status badge
    status_color = GREEN if inherited else ORANGE
    status_text = "✓ Inherits" if inherited else "✂ Can Break"
    badge = add_rounded_rect(s, x + Inches(0.6), y + Inches(1.9), Inches(1.6), Inches(0.35), status_color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = status_text; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"

    # Arrow between cards
    if i < 3:
        arrow_x = x + Inches(2.85)
        add_text_box(s, arrow_x, Inches(3.0), Inches(0.3), Inches(0.5),
                     "→", font_size=28, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Key insight box at bottom
insight_box = add_rounded_rect(s, Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.0), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(5.3), Inches(9.7), Inches(0.8),
             "💡 Best Practice:  Keep unique permission scopes under 5,000 per library (recommended limit). "
             "The hard limit is 50,000 unique ACLs per document library. Share folders, not individual files.",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Why Unique Permissions Are Risky
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("This slide hammers home the operational cost of breaking inheritance too often. "
              "Reference Microsoft guidance on permission scopes and best practice.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, RED_ACCENT)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "⚠️  Why Unique Permissions Are Risky", font_size=32, bold=True, color=RED_ACCENT)

risks = [
    ("🔍", "Harder to Audit", "Each unique scope means another set of ACLs to review.\n"
     "Large-scale audits become exponentially complex."),
    ("🛠️", "Harder to Troubleshoot", "\"Who has access?\" becomes a detective game\nwhen scopes are scattered across folders and files."),
    ("⚡", "Performance Impact", "Exceeding 5,000 unique scopes increases SQL\nround trips, degrading list view performance."),
    ("❌", "Easier to Misconfigure", "Ad-hoc per-file permissions are often forgotten,\nleading to accidental data exposure or access loss."),
]
for i, (icon, title, desc) in enumerate(risks):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 2.6)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(2.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(2.2), RED_ACCENT)
    add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.5),
                 title, font_size=20, bold=True, color=RED_ACCENT)
    add_text_box(s, x + Inches(1.0), y + Inches(0.8), Inches(4.5), Inches(1.2),
                 desc, font_size=14, color=DARK_TEXT)

add_footer_bar(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Best Practices for Managing Scopes
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Microsoft official recommendations for managing permission scopes efficiently.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, GREEN)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "✅  Best Practices: Minimizing Permission Scopes", font_size=30, bold=True, color=DARK_BG)

practices = [
    ("1", "Leverage inheritance", "Let children inherit from parent — don't break unless necessary"),
    ("2", "Share folders, not files", "One shared folder = 1 scope; 10K individual files = 10K scopes"),
    ("3", "Use groups, not individuals", "SharePoint groups + Entra ID groups keep ACLs clean"),
    ("4", "Share large folders early", "Folders with >100K items cannot break inheritance later"),
    ("5", "Regular audits", "Review and clean up unique scopes and stale sharing links"),
    ("6", "Design boundaries first", "Decide site vs library vs folder scope before content arrives"),
]
for i, (num, title, desc) in enumerate(practices):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 1.7)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(1.4), WHITE)
    card.shadow.inherit = False
    badge = add_rounded_rect(s, x + Inches(0.2), y + Inches(0.35), Inches(0.55), Inches(0.55), GREEN)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, x + Inches(1.0), y + Inches(0.15), Inches(4.5), Inches(0.5),
                 title, font_size=18, bold=True, color=DARK_TEXT)
    add_text_box(s, x + Inches(1.0), y + Inches(0.65), Inches(4.5), Inches(0.6),
                 desc, font_size=13, color=MID_GRAY)

add_footer_bar(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Section Divider: Groups & Site Types
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: Now let's see how groups and site types interact with permissions.")
section_divider(s, "Groups and Site Types", "SharePoint Groups · M365 Groups · Teams Connections", "👥")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Default SharePoint Groups
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Every SharePoint site gets three default groups. This is the foundation of the permission model. "
              "Teach learners to always use groups instead of direct user assignments.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Default SharePoint Groups", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Three groups created automatically with every site",
             font_size=16, color=MID_GRAY)

groups = [
    ("👑", "Owners", "Full Control", [
        "Manage site settings & permissions",
        "Add/remove members from all groups",
        "Create subsites and manage navigation",
        "Delete the site",
    ], ACCENT_BLUE),
    ("✏️", "Members", "Edit", [
        "Add/edit/delete content in libraries & lists",
        "Create document libraries and lists",
        "Cannot manage site settings",
        "Cannot change permissions",
    ], ACCENT_TEAL),
    ("👁️", "Visitors", "Read", [
        "View content only",
        "Cannot add, edit, or delete items",
        "Cannot change site settings",
        "Ideal for broad read-only access",
    ], ACCENT_PURPLE),
]
for i, (icon, name, level, perms, color) in enumerate(groups):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(3.8), Inches(4.5), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(3.8), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(2.05), Inches(3.2), Inches(0.5),
                 f"{icon}  {name}", font_size=24, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # permission level badge
    badge = add_rounded_rect(s, x + Inches(1.0), Inches(2.7), Inches(1.8), Inches(0.4), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = level; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    # permissions list
    add_bullet_frame(s, x + Inches(0.3), Inches(3.3), Inches(3.3), Inches(2.8),
                     perms, font_size=13, color=DARK_TEXT, icon="•")

add_footer_bar(s, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SharePoint Groups vs M365 Group-Connected Permissions
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("This is a key conceptual distinction. Communication sites use SharePoint groups. "
              "Team sites often have M365 group backing, which means membership is managed from M365/Entra.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "SharePoint Groups vs M365 Group-Connected", font_size=30, bold=True, color=DARK_BG)

# Two-column comparison
# Left column - SP Groups
add_rounded_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.8), WHITE).shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.5),
             "📋  SharePoint Groups", font_size=22, bold=True, color=ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(2.2), Inches(5.2), Inches(0.4),
             "Communication Sites & Classic Sites", font_size=14, bold=True, color=MID_GRAY)
sp_items = [
    "Site-scoped: exist only within the site",
    "Managed in SharePoint site settings",
    "3 defaults: Owners, Members, Visitors",
    "Can add users or Entra security groups",
    "Full admin control over membership",
    "Best for: intranet portals, publishing sites",
]
add_bullet_frame(s, Inches(1.0), Inches(2.7), Inches(5.0), Inches(3.2),
                 sp_items, font_size=14, color=DARK_TEXT, icon="▸")

# Right column - M365 Groups
add_rounded_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.8), WHITE).shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
             "🔗  M365 Group-Connected", font_size=22, bold=True, color=ACCENT_TEAL)
add_text_box(s, Inches(7.1), Inches(2.2), Inches(5.2), Inches(0.4),
             "Team Sites & Teams-Connected Sites", font_size=14, bold=True, color=MID_GRAY)
m365_items = [
    "M365 Group owners → site Owners",
    "M365 Group members → site Members",
    "Visitors group still SharePoint-only",
    "Membership managed via M365 or Teams",
    "Entra ID backs the group identity",
    "Best for: project teams, departments",
]
add_bullet_frame(s, Inches(7.1), Inches(2.7), Inches(5.0), Inches(3.2),
                 m365_items, font_size=14, color=DARK_TEXT, icon="▸")

# Tip box
tip = add_rounded_rect(s, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.55), LIGHT_BLUE)
add_text_box(s, Inches(2.7), Inches(6.55), Inches(7.9), Inches(0.45),
             "💡 Tip:  Teams-connected sites may show permissions as read-only in SharePoint — manage via Teams instead.",
             font_size=13, color=DARK_TEXT)

add_footer_bar(s, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Site Type → Permission Model Decision
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Quick decision reference: which permission model applies based on site type.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Site Type → Permission Model", font_size=32, bold=True, color=DARK_BG)

# Table header
header_y = Inches(1.5)
cols = [Inches(0.8), Inches(3.8), Inches(6.5), Inches(9.5)]
col_widths = [Inches(3.0), Inches(2.7), Inches(3.0), Inches(3.0)]
headers = ["Site Type", "Group Model", "Managed In", "Use Case"]
add_shape_rect(s, Inches(0.8), header_y, Inches(11.7), Inches(0.55), ACCENT_BLUE)
for j, h in enumerate(headers):
    add_text_box(s, cols[j], header_y + Inches(0.05), col_widths[j], Inches(0.45),
                 h, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

rows_data = [
    ("Communication Site", "SharePoint Groups", "SharePoint Site Settings", "Intranet · Portal · News"),
    ("Team Site (no group)", "SharePoint Groups", "SharePoint Site Settings", "Classic · Standalone"),
    ("Team Site (M365 Group)", "M365 Group + SP Visitors", "M365 Admin / Entra ID", "Projects · Departments"),
    ("Teams-Connected Site", "Teams Owners/Members", "Microsoft Teams", "Team Channels · Collab"),
    ("Teams Channel Site", "Channel membership", "Microsoft Teams", "Private/Shared Channels"),
]
for i, (site_type, group_model, managed_in, use_case) in enumerate(rows_data):
    y = header_y + Inches(0.55 + i * 0.85)
    bg_color = WHITE if i % 2 == 0 else NEAR_WHITE
    add_shape_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.85), bg_color, LIGHT_GRAY)
    vals = [site_type, group_model, managed_in, use_case]
    for j, v in enumerate(vals):
        weight = True if j == 0 else False
        add_text_box(s, cols[j], y + Inches(0.15), col_widths[j], Inches(0.55),
                     v, font_size=13, bold=weight, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Section Divider: Sharing Links
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: Sharing links are the other side of the access coin. "
              "Users create them daily — admins need to understand and govern them.")
section_divider(s, "Sharing Links vs Permissions", "Collaboration actions vs durable access design", "🔗")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Permissions vs Sharing Links (Conceptual)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Key distinction: Permissions = durable governance design. Sharing links = collaboration actions. "
              "In incidents, 'mystery access' often comes from a forgotten sharing link.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Permissions vs Sharing Links", font_size=32, bold=True, color=DARK_BG)

# Left: Permissions
add_rounded_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.2), WHITE).shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.5),
             "🛡️  Permissions (Governance)", font_size=22, bold=True, color=ACCENT_BLUE)
perm_bullets = [
    "Durable access design at site/library/folder level",
    "Applied via groups for manageability",
    "Controlled by site owners and admins",
    "Auditable through site settings",
    "Persist until explicitly removed",
]
add_bullet_frame(s, Inches(1.0), Inches(2.4), Inches(5.0), Inches(3.0),
                 perm_bullets, font_size=14, color=DARK_TEXT, icon="▸")

# Right: Sharing Links
add_rounded_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.2), WHITE).shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(0.08), ORANGE)
add_text_box(s, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.5),
             "🔗  Sharing Links (Collaboration)", font_size=22, bold=True, color=ORANGE)
link_bullets = [
    "Ad-hoc access at file/folder level",
    "Created by users as they collaborate",
    "Can be forwarded (People in org / Anyone)",
    "May bypass group governance if unchecked",
    "Can be revoked or set to expire",
]
add_bullet_frame(s, Inches(7.1), Inches(2.4), Inches(5.0), Inches(3.0),
                 link_bullets, font_size=14, color=DARK_TEXT, icon="▸")

# Bottom insight
insight = add_rounded_rect(s, Inches(2.0), Inches(6.0), Inches(9.3), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(2.3), Inches(6.1), Inches(8.7), Inches(0.5),
             "🔍 Admin Question:  When a user reports access issues, always ask — "
             "\"Is access from a group or a link?\"  Check both sources.",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Three Types of Sharing Links
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Deep-dive into the three primary sharing link types. "
              "Each has different authentication, auditability, and risk characteristics.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Three Types of Sharing Links", font_size=32, bold=True, color=DARK_BG)

link_types = [
    ("🌐", "Anyone", "No authentication\nrequired", [
        "Works for anyone with the link",
        "Cannot audit who accessed",
        "Can set expiration & view-only",
        "May be disabled by policy",
    ], RED_ACCENT, "HIGH RISK"),
    ("🏢", "People in Your\nOrganization", "Internal users\nauthenticate", [
        "Works for all org members",
        "Forwarded links still work",
        "Users must sign in",
        "Good for broad internal sharing",
    ], ORANGE, "MEDIUM RISK"),
    ("👤", "Specific People", "Named recipients\nauthenticate", [
        "Only specified people can access",
        "Requires authentication",
        "Fully auditable",
        "Works for internal & external",
    ], GREEN, "LOW RISK"),
]
for i, (icon, name, auth, features, color, risk) in enumerate(link_types):
    x = Inches(0.6 + i * 4.2)
    card = add_rounded_rect(s, x, Inches(1.5), Inches(3.9), Inches(5.0), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.5), Inches(3.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(1.75), Inches(3.3), Inches(0.5),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.3), Inches(3.3), Inches(0.7),
                 name, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Risk badge
    badge = add_rounded_rect(s, x + Inches(1.0), Inches(3.05), Inches(1.9), Inches(0.35), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = risk; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    # Auth info
    add_text_box(s, x + Inches(0.3), Inches(3.55), Inches(3.3), Inches(0.6),
                 auth, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.6), Inches(4.2), Inches(2.7), Inches(0.03), LIGHT_GRAY)
    # Features
    add_bullet_frame(s, x + Inches(0.3), Inches(4.35), Inches(3.3), Inches(2.0),
                     features, font_size=12, color=DARK_TEXT, icon="•")

add_footer_bar(s, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Sharing Link Settings (Admin Controls)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Show the admin controls available in SharePoint admin center for managing sharing links. "
              "Default link type, expiration, and permissions can all be configured.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Admin Controls for Sharing Links", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "SharePoint admin center → Policies → Sharing",
             font_size=16, color=MID_GRAY)

controls = [
    ("🔧", "Default Link Type", "Set org-wide default:\nSpecific People (most restrictive)\nPeople in org (balanced)\nAnyone (least restrictive)", ACCENT_BLUE),
    ("⏱️", "Link Expiration", "Set max days for Anyone links\nExisting links keep their expiration\nif new setting is longer", ACCENT_TEAL),
    ("🔒", "Link Permissions", "View only or View + Edit\nSet separately for files & folders\nAnyone links can be restricted", ACCENT_PURPLE),
    ("🌐", "Site-Level Override", "Per-site sharing can be MORE\nrestrictive than org, never more\npermissive", ORANGE),
]
for i, (icon, title, desc, color) in enumerate(controls):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 2.5)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(2.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(2.2), color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.5),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(1.0), y + Inches(0.2), Inches(4.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color)
    add_text_box(s, x + Inches(1.0), y + Inches(0.8), Inches(4.5), Inches(1.2),
                 desc, font_size=13, color=DARK_TEXT)

add_footer_bar(s, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — SharePoint vs OneDrive Sharing Policy
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Important relationship: OneDrive sharing is controlled alongside SharePoint. "
              "OneDrive can be equal or MORE restrictive than SharePoint, but never more permissive.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "SharePoint vs OneDrive Sharing Policy", font_size=30, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "The sharing hierarchy: Org → SharePoint → OneDrive → Site",
             font_size=16, color=MID_GRAY)

# Flow diagram: 4 levels
levels_data = [
    ("🏛️", "Organization\nLevel", "Sets maximum\npermissiveness", ACCENT_BLUE),
    ("📡", "SharePoint\nAdmin", "Equal or more\nrestrictive", ACCENT_TEAL),
    ("☁️", "OneDrive\nAdmin", "Equal or more\nrestrictive than SP", ACCENT_PURPLE),
    ("🌐", "Individual\nSite", "Most restrictive\n(per-site override)", ORANGE),
]
for i, (icon, label, desc, color) in enumerate(levels_data):
    x = Inches(0.6 + i * 3.2)
    card = add_rounded_rect(s, x, Inches(2.0), Inches(2.8), Inches(2.6), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(2.0), Inches(2.8), Inches(0.08), color)
    add_text_box(s, x + Inches(0.2), Inches(2.25), Inches(2.4), Inches(0.5),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.7),
                 label, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.55), Inches(2.4), Inches(0.7),
                 desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < 3:
        arrow_x = x + Inches(2.85)
        add_text_box(s, arrow_x, Inches(2.9), Inches(0.3), Inches(0.5),
                     "→", font_size=24, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Key rules at bottom
rules_box = add_rounded_rect(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.6), LIGHT_BLUE)
rules = [
    "🔹 A site can NEVER be more permissive than the org-level setting",
    "🔹 OneDrive sharing can be equal or more restrictive than SharePoint — not more permissive",
    "🔹 Changing org-level setting affects existing sites only if they were using the old maximum",
    "🔹 External sharing level: Anyone > New & Existing Guests > Existing Guests > Only Org",
]
add_bullet_frame(s, Inches(1.1), Inches(5.1), Inches(11.0), Inches(1.4),
                 rules, font_size=13, color=DARK_TEXT, icon="")

add_footer_bar(s, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Section Divider: The Northwind Scenario
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: Apply the theory to a realistic scenario — Northwind Contracts workflow.")
section_divider(s, "Northwind Scenario", "Designing a permission model for the Contracts workflow", "📋")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Scenario: Contracts Workflow
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Paint the scenario: Northwind has a Contracts workflow with Drafts (editors collaborate) "
              "and Finals (broad read, restricted edit). Goal: isolate sensitive content in a dedicated library.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📋  Scenario: Northwind Contracts Workflow", font_size=30, bold=True, color=DARK_BG)

# Workflow visual: 3 stages
stages = [
    ("📝", "Drafts Library", "Editors collaborate on\ncontract drafts", [
        "Restricted to contract editors",
        "Members can edit & upload",
        "Owners have full control",
    ], ACCENT_BLUE, "Broken Inheritance"),
    ("✅", "Finals Folder", "Approved contracts\nfor broad read access", [
        "Broad read access (Visitors)",
        "Only Owners can edit",
        "One folder with unique perms",
    ], GREEN, "Unique Permissions"),
    ("🔒", "Confidential Folder", "Sensitive contracts\n(board-level only)", [
        "Board members only",
        "Strict need-to-know basis",
        "Unique permissions within library",
    ], RED_ACCENT, "Unique Permissions"),
]
for i, (icon, title, desc, points, color, perms_status) in enumerate(stages):
    x = Inches(0.6 + i * 4.2)
    card = add_rounded_rect(s, x, Inches(1.4), Inches(3.9), Inches(4.8), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.4), Inches(3.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(1.65), Inches(3.3), Inches(0.5),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.2), Inches(3.3), Inches(0.5),
                 title, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.7), Inches(3.3), Inches(0.6),
                 desc, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    # Status badge
    badge = add_rounded_rect(s, x + Inches(0.7), Inches(3.35), Inches(2.5), Inches(0.35), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = perms_status; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    # Bullet points
    add_bullet_frame(s, x + Inches(0.3), Inches(3.9), Inches(3.3), Inches(2.0),
                     points, font_size=12, color=DARK_TEXT, icon="•")

add_footer_bar(s, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Permission Design Principles
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Four principles to guide permission design decisions. These apply beyond the Northwind scenario.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Permission Design Principles", font_size=32, bold=True, color=DARK_BG)

principles = [
    ("🎯", "Broad at the Top,\nRestricted Below", "Keep site-level access simple.\nIsolate sensitive content in\ndedicated libraries or folders.", ACCENT_BLUE),
    ("👥", "Groups Over\nIndividuals", "Always use SharePoint groups or\nEntra ID groups. Never assign\npermissions to individual users.", ACCENT_TEAL),
    ("✂️", "Minimize\nUnique Scopes", "Every break costs manageability.\nAim for library-level isolation,\nnot file-by-file permissions.", ORANGE),
    ("📊", "Document\nthe Model", "Record who has access and why.\nReview quarterly at minimum.\nClean up stale access.", ACCENT_PURPLE),
]
for i, (icon, title, desc, color) in enumerate(principles):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.5), Inches(2.9), Inches(4.5), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.5), Inches(2.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(1.75), Inches(2.3), Inches(0.5),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(2.35), Inches(2.5), Inches(0.8),
                 title, font_size=17, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.6), Inches(3.2), Inches(1.7), Inches(0.03), LIGHT_GRAY)
    add_text_box(s, x + Inches(0.2), Inches(3.4), Inches(2.5), Inches(2.0),
                 desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Section Divider: Troubleshooting
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: Now let's cover the practical troubleshooting skills needed for access issues.")
section_divider(s, "Troubleshooting Access Issues", "The admin mindset for resolving 'Access Denied'", "🔍")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Troubleshooting Workflow (Step-by-step)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Five-step troubleshooting flow for access issues. This is a practical admin skill "
              "that learners will use in real environments.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Access Denied Troubleshooting Workflow", font_size=30, bold=True, color=DARK_BG)

steps = [
    ("1", "Identify the Boundary", "Where is the resource?\nSite → Library → Folder → File", ACCENT_BLUE),
    ("2", "Check Group\nMembership", "Is the user in Owners,\nMembers, or Visitors?", ACCENT_TEAL),
    ("3", "Check Direct\nPermissions", "Any per-user or per-group\nassignment at this level?", ACCENT_PURPLE),
    ("4", "Check Sharing Links", "Use 'Manage Access' panel\nto see active sharing links", ORANGE),
    ("5", "Use 'Check\nPermissions'", "Site Settings → Check Permissions\nfor definitive effective access", GREEN),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.4 + i * 2.6)
    card = add_rounded_rect(s, x, Inches(1.5), Inches(2.35), Inches(4.0), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.5), Inches(2.35), Inches(0.08), color)
    # Number badge
    badge = add_rounded_rect(s, x + Inches(0.75), Inches(1.75), Inches(0.65), Inches(0.65), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(6)
    add_text_box(s, x + Inches(0.15), Inches(2.6), Inches(2.05), Inches(0.7),
                 title, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.15), Inches(3.35), Inches(2.05), Inches(1.5),
                 desc, font_size=12, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < 4:
        arrow_x = x + Inches(2.4)
        add_text_box(s, arrow_x, Inches(3.0), Inches(0.2), Inches(0.5),
                     "→", font_size=20, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Pro-tip
tip = add_rounded_rect(s, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(5.9), Inches(9.7), Inches(0.5),
             "💡 Pro Tip:  Encourage learners to take screenshots of each check. This builds evidence for "
             "escalation and helps document the resolution.",
             font_size=13, color=DARK_TEXT)

add_footer_bar(s, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Top 5 Common Permission Failures
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Real-world failure scenarios that admins encounter frequently. "
              "Each has a specific resolution path.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, RED_ACCENT)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "⚠️  Top 5 Common Permission Failures", font_size=30, bold=True, color=RED_ACCENT)

failures = [
    ("1", "Broken Inheritance Forgotten", "Permissions changed at library/folder level but not documented. "
     "New content inherits old scope.", ACCENT_BLUE),
    ("2", "Sharing Link Sprawl", "Users create 'Anyone' links freely. Access cannot be tracked or revoked "
     "without manual cleanup.", ORANGE),
    ("3", "M365 Group vs SP Group Confusion", "Admin changes SharePoint group membership but actual access is "
     "driven by M365 group. Change has no effect.", ACCENT_TEAL),
    ("4", "Direct User Assignment", "Permissions granted to individual users instead of groups. "
     "Offboarding misses these direct grants.", RED_ACCENT),
    ("5", "Site vs File-Level Mismatch", "User has site access but file has unique permissions that exclude them. "
     "Or: user lacks site access but has a sharing link.", ACCENT_PURPLE),
]
for i, (num, title, desc, color) in enumerate(failures):
    y = Inches(1.3 + i * 1.15)
    card = add_rounded_rect(s, Inches(0.8), y, Inches(11.7), Inches(1.0), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, Inches(0.8), y, Inches(0.08), Inches(1.0), color)
    badge = add_rounded_rect(s, Inches(1.1), y + Inches(0.18), Inches(0.55), Inches(0.55), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, Inches(1.9), y + Inches(0.05), Inches(3.0), Inches(0.45),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(1.9), y + Inches(0.5), Inches(10.3), Inches(0.45),
                 desc, font_size=12, color=DARK_TEXT)

add_footer_bar(s, 23, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Section Divider: Lab Preview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: Lab preview — learners will implement the Northwind permission model.")
section_divider(s, "Lab Preview", "Designing a Permission Model — Northwind Contracts", "🧪")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Lab 04 Preview (Detailed)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Lab preview with step-by-step tasks. Emphasize shared-tenant safety rules. "
              "Each participant works inside their own NW-Pxx-ProjectSite.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_TEAL)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🧪  Lab 04: Designing a Permission Model", font_size=30, bold=True, color=DARK_BG)

# Lab steps
lab_steps = [
    ("1", "Navigate to NW-Pxx-ProjectSite", "Open your assigned site in SharePoint"),
    ("2", "Create NW-Pxx-Contracts library", "New document library for contracts workflow"),
    ("3", "Break inheritance at library level", "Stop inheriting from site, set editor-only access"),
    ("4", "Create 'Finals' folder", "Add folder for approved contracts with unique permissions"),
    ("5", "Break inheritance at folder level", "Set broad read access; restrict edit to Owners only"),
    ("6", "Validate with Check Permissions", "Confirm effective access for test users"),
]
for i, (num, title, desc) in enumerate(lab_steps):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.3 + row * 1.55)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(1.3), WHITE)
    card.shadow.inherit = False
    badge = add_rounded_rect(s, x + Inches(0.15), y + Inches(0.28), Inches(0.55), Inches(0.55), ACCENT_TEAL)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, x + Inches(0.9), y + Inches(0.15), Inches(4.7), Inches(0.45),
                 title, font_size=16, bold=True, color=DARK_TEXT)
    add_text_box(s, x + Inches(0.9), y + Inches(0.6), Inches(4.7), Inches(0.5),
                 desc, font_size=13, color=MID_GRAY)

# Safety reminder
safety = add_rounded_rect(s, Inches(2.0), Inches(5.9), Inches(9.3), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(2.3), Inches(6.0), Inches(8.7), Inches(0.5),
             "⚠️ Shared Tenant Rule:  Work ONLY inside your NW-Pxx-ProjectSite. "
             "Do not modify tenant-wide settings or other participants' sites.",
             font_size=13, color=DARK_TEXT)

add_footer_bar(s, 25, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Summary: Key Takeaways
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Wrap up with key takeaways from the module.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📌  Module 4 Summary", font_size=32, bold=True, color=DARK_BG)

takeaways = [
    ("✅", "Inheritance is your friend", "Use it by default; break only when necessary for isolation"),
    ("✅", "Groups over individuals", "SharePoint groups or Entra groups — never assign to individual users"),
    ("✅", "Know your sharing links", "Understand the risk profile: Anyone > People in Org > Specific People"),
    ("✅", "OneDrive follows SharePoint", "OneDrive sharing can be equal or more restrictive, never more permissive"),
    ("✅", "Troubleshoot systematically", "Boundary → Group membership → Direct perms → Sharing links → Check Permissions"),
    ("✅", "Document and review", "Record your permission model; audit quarterly to prevent scope sprawl"),
]
for i, (icon, title, desc) in enumerate(takeaways):
    y = Inches(1.3 + i * 0.9)
    add_rounded_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.75), WHITE).shadow.inherit = False
    add_text_box(s, Inches(1.0), y + Inches(0.1), Inches(0.5), Inches(0.5),
                 icon, font_size=20, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), y + Inches(0.07), Inches(3.5), Inches(0.4),
                 title, font_size=16, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(5.2), y + Inches(0.07), Inches(7.0), Inches(0.55),
                 desc, font_size=14, color=DARK_TEXT)

add_footer_bar(s, 26, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Knowledge Check
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Interactive knowledge check to reinforce module content. "
              "Have learners discuss in pairs before sharing answers.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "❓  Knowledge Check", font_size=32, bold=True, color=DARK_BG)

questions = [
    "Why is breaking inheritance repeatedly considered risky in SharePoint?",
    "On a group-connected team site, what drives SharePoint permissions?",
    "When would you prefer SharePoint groups over individual user permissions?",
    "What's the difference between 'People in your org' and 'Specific people' links?",
    "What is the FIRST place you check for 'Access Denied' — and why?",
]
for i, q in enumerate(questions):
    y = Inches(1.3 + i * 1.05)
    card = add_rounded_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.85), WHITE)
    card.shadow.inherit = False
    badge = add_rounded_rect(s, Inches(1.1), y + Inches(0.12), Inches(0.55), Inches(0.55), ACCENT_PURPLE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, Inches(1.9), y + Inches(0.15), Inches(10.3), Inches(0.55),
                 q, font_size=16, color=DARK_TEXT)

add_footer_bar(s, 27, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Thank You / End
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("End of Module 4. Transition to Module 5: Managing Metadata and the Term Store.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "Module 4 Complete ✓", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
             "Permissions and Collaboration Model",
             font_size=24, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.04), ACCENT_PURPLE)

add_text_box(s, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
             "Next → Module 5: Managing Metadata and the Term Store",
             font_size=18, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

# References
add_text_box(s, Inches(2), Inches(5.8), Inches(9), Inches(1.2),
             "References:\n"
             "• Manage Permission Scopes: learn.microsoft.com/sharepoint/manage-permission-scope\n"
             "• Sharing & Permissions (Modern): learn.microsoft.com/sharepoint/modern-experience-sharing-permissions\n"
             "• Manage Sharing Settings: learn.microsoft.com/sharepoint/turn-external-sharing-on-or-off",
             font_size=11, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Module-04-Slides.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {slide_counter[0]}")
