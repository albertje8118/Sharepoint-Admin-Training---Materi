#!/usr/bin/env python3
"""
Module 7 – Apps and Customization in SharePoint Online
Generates a 28-slide PPTX with the established design system.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, math

# ── Design tokens ──────────────────────────────────────────────
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
DARK_BG   = RGBColor(0x1B, 0x1B, 0x2F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE   = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL   = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
ORANGE        = RGBColor(0xFF, 0x8C, 0x00)
GREEN         = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT    = RGBColor(0xD1, 0x34, 0x38)

FONT_BODY = "Segoe UI"
FONT_CODE = "Cascadia Code"
FOOTER_TEXT = "Module 7 | Apps and Customization in SharePoint Online"

# ── Helper functions ───────────────────────────────────────────
def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_BODY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tb

def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=16, color=DARK_GRAY, bold_first=False,
                     spacing=Pt(10), icon=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = f"{icon} " if icon else "• "
        p.text = prefix + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = spacing
        if bold_first and i == 0:
            p.font.bold = True
    return tb

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE):
    add_shape_rect(slide, 0, 0, SLIDE_W, Inches(0.06), color)

def add_footer_bar(slide, label=FOOTER_TEXT):
    add_shape_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), Inches(7.08), Inches(8), Inches(0.35),
                 label, font_size=10, color=WHITE)

def section_divider(prs, title, subtitle="", accent=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(1), Inches(3.1), Inches(1.2), Inches(0.06), accent)
    add_text_box(slide, Inches(1), Inches(3.35), Inches(11), Inches(1),
                 title, font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(4.3), Inches(10), Inches(0.7),
                     subtitle, font_size=18, color=MID_GRAY)
    add_footer_bar(slide)
    return slide

def new_slide(prs, title, accent=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    add_top_bar(slide, accent)
    add_shape_rect(slide, Inches(0.6), Inches(0.45), Inches(0.08), Inches(0.55), accent)
    add_text_box(slide, Inches(0.85), Inches(0.4), Inches(11), Inches(0.65),
                 title, font_size=28, color=DARK_GRAY, bold=True)
    add_footer_bar(slide)
    return slide

# ── Presentation setup ─────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ================================================================
# SLIDE 1 – Title
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl, DARK_BG)
add_shape_rect(sl, 0, Inches(2.6), SLIDE_W, Inches(2.6), RGBColor(0x22, 0x22, 0x3A))
add_shape_rect(sl, Inches(1), Inches(3.55), Inches(1.5), Inches(0.07), ACCENT_BLUE)
add_text_box(sl, Inches(1), Inches(2.75), Inches(11), Inches(0.7),
             "Module 7", font_size=22, color=ACCENT_TEAL, bold=True)
add_text_box(sl, Inches(1), Inches(3.7), Inches(11), Inches(1),
             "Apps and Customization\nin SharePoint Online", font_size=40, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(5.0), Inches(10), Inches(0.5),
             "Scenario: Project Northwind Intranet Modernization  ·  Day 2 of 3",
             font_size=16, color=MID_GRAY)
add_footer_bar(sl)
add_speaker_notes(sl,
    "Module 7 covers the customization spectrum from out-of-box to SPFx, "
    "declarative JSON formatting, app governance, and API access. "
    "This is the last module of Day 2.")

# ================================================================
# SLIDE 2 – Why admins care
# ================================================================
sl = new_slide(prs, "Why Admins Care About Customization")
cards = [
    ("🚀", "Productivity", "Customized views and apps\ncan dramatically improve\nuser adoption."),
    ("🔒", "Security", "Every customization adds a\npotential attack surface.\nGovern wisely."),
    ("🛠️", "Supportability", "Custom code must be\nmaintainable over time.\nPrefer low-code first."),
    ("📋", "Governance", "Who can deploy what?\nTenant-wide vs site-scoped\ndecisions matter."),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(sl, x, Inches(1.6), Inches(2.8), Inches(4.6), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.3), Inches(1.8), Inches(2.2), Inches(0.7),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(2.5), Inches(2.2), Inches(0.5),
                 title, font_size=20, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.25), Inches(3.2), Inches(2.3), Inches(2.8),
                 desc, font_size=15, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Customization is a double-edged sword. It improves productivity but adds risk. "
    "Administrators must balance flexibility with governance, security, and long-term supportability.")

# ================================================================
# SLIDE 3 – Learning outcomes
# ================================================================
sl = new_slide(prs, "Learning Outcomes")
outcomes = [
    "Explain modern customization models in SharePoint Online",
    "Customize a list safely using JSON-based column and view formatting",
    "Describe how SPFx solutions are deployed and governed",
    "Describe what API access is for and why it matters",
]
for i, item in enumerate(outcomes):
    y = Inches(1.6 + i * 1.3)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.55), Inches(0.55), ACCENT_BLUE)
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.55),
                 item, font_size=20, color=DARK_GRAY)
add_speaker_notes(sl,
    "Four learning outcomes covering the full customization spectrum, "
    "from JSON formatting to SPFx governance and API access management.")

# ================================================================
# SLIDE 4 – Section divider: Customization Models
# ================================================================
section_divider(prs, "Section 1", "Customization Models in SharePoint Online")

# ================================================================
# SLIDE 5 – The customization spectrum
# ================================================================
sl = new_slide(prs, "The Customization Spectrum")
tiers = [
    ("Out-of-Box Configuration", "Lowest risk  •  No code", "Settings, web parts,\npages, permissions", GREEN, "⚙️"),
    ("Declarative (JSON Formatting)", "Low risk  •  JSON only", "Column & view formatting,\nlist form config", ACCENT_BLUE, "📝"),
    ("SharePoint Framework (SPFx)", "Highest flexibility  •  Code", "Web parts, extensions,\nTeams + Viva", ACCENT_PURPLE, "🧩"),
]
for i, (title, risk, desc, color, icon) in enumerate(tiers):
    y = Inches(1.5 + i * 1.8)
    # Arrow / tier bar
    add_shape_rect(sl, Inches(0.8), y, Inches(0.12), Inches(1.4), color)
    add_text_box(sl, Inches(1.2), y, Inches(1), Inches(0.6),
                 icon, font_size=32, color=DARK_GRAY)
    add_text_box(sl, Inches(2.0), y, Inches(4), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(sl, Inches(2.0), y + Inches(0.5), Inches(4), Inches(0.4),
                 risk, font_size=14, color=DARK_GRAY)
    add_text_box(sl, Inches(7), y, Inches(5), Inches(1.2),
                 desc, font_size=15, color=DARK_GRAY)
# Arrow label
add_text_box(sl, Inches(0.3), Inches(1.2), Inches(0.6), Inches(0.4),
             "▲", font_size=14, color=GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(0.15), Inches(5.7), Inches(0.9), Inches(0.4),
             "Flexibility ▼", font_size=10, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Present the spectrum from safest to most flexible. "
    "Admin recommendation: prefer out-of-box first, then JSON formatting, "
    "and only use SPFx when the other options aren't sufficient.")

# ================================================================
# SLIDE 6 – Out-of-box configuration
# ================================================================
sl = new_slide(prs, "Tier 1: Out-of-Box Configuration", GREEN)
items = [
    "List / library settings (columns, views, content types)",
    "Modern pages and web parts (no code needed)",
    "Permissions, sharing, and governance controls",
    "Hub sites, navigation, site templates",
    "Microsoft Lists templates (Issue Tracker, etc.)",
]
add_bullet_frame(sl, Inches(0.8), Inches(1.5), Inches(7), Inches(4.5),
                 items, font_size=18, icon="✅")
add_rounded_rect(sl, Inches(8.5), Inches(1.5), Inches(4.2), Inches(2),
                 RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text_box(sl, Inches(8.8), Inches(1.7), Inches(3.6), Inches(0.5),
             "🎯  Admin Take-Away", font_size=18, color=GREEN, bold=True)
add_text_box(sl, Inches(8.8), Inches(2.3), Inches(3.6), Inches(1),
             "Always prefer out-of-box features first.\nThey are easiest to support and\nrequire no code deployment.",
             font_size=15, color=DARK_GRAY)
add_speaker_notes(sl,
    "Out-of-box configuration = zero deployment risk. Covers most standard business needs. "
    "Modern pages + web parts give rich layout without custom code. "
    "Lists templates provide pre-built structures.")

# ================================================================
# SLIDE 7 – Declarative customization overview
# ================================================================
sl = new_slide(prs, "Tier 2: Declarative Customization (JSON Formatting)")
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
             "JSON-based formatting changes how lists and libraries are displayed — not the data.",
             font_size=18, color=DARK_GRAY)
types = [
    ("Column Formatting", "Customize how a single\nfield is rendered in a view", "📊"),
    ("View Formatting", "Customize how rows/cards\nare rendered in the view", "📋"),
    ("List Form Configuration", "Customize the item\nform layout (header/body/footer)", "📄"),
]
for i, (title, desc, icon) in enumerate(types):
    x = Inches(0.6 + i * 4.1)
    card = add_rounded_rect(sl, x, Inches(2.5), Inches(3.8), Inches(3.2), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.3), Inches(2.7), Inches(3.2), Inches(0.7),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.4), Inches(3.2), Inches(0.5),
                 title, font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(4.1), Inches(3.2), Inches(1.2),
                 desc, font_size=15, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Three types of declarative customization. Column formatting targets individual fields. "
    "View formatting targets the entire row/card layout. List form configuration targets the item form. "
    "All use JSON — no compiled code, no deployment package needed.")

# ================================================================
# SLIDE 8 – Column formatting: how to access
# ================================================================
sl = new_slide(prs, "Column Formatting: How It Works")
steps = [
    ("1", "Click column header → Column settings → Format this column"),
    ("2", "Choose a Conditional Formatting rule (designer) – or switch to Advanced mode"),
    ("3", "Paste or write JSON in the Advanced editor"),
    ("4", "Click Preview to verify the rendering"),
    ("5", "Click Save — formatting applies to all users in that view"),
]
for i, (num, text) in enumerate(steps):
    y = Inches(1.5 + i * 1.0)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.5), Inches(0.5), ACCENT_BLUE)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.6), y + Inches(0.05), Inches(10), Inches(0.5),
                 text, font_size=17, color=DARK_GRAY)
# Key point callout
add_rounded_rect(sl, Inches(8.5), Inches(5.5), Inches(4.2), Inches(1),
                 RGBColor(0xE3, 0xF2, 0xFD), ACCENT_BLUE)
add_text_box(sl, Inches(8.7), Inches(5.6), Inches(3.8), Inches(0.8),
             "💡 Formatting changes rendering only\n— the underlying data is unchanged.",
             font_size=14, color=ACCENT_BLUE, bold=True)
add_speaker_notes(sl,
    "Walk through the 5-step process. Emphasize that the designer mode (conditional formatting) "
    "is the easiest entry point. Advanced mode gives full JSON control. "
    "Key point: formatting never modifies data, only visual rendering.")

# ================================================================
# SLIDE 9 – Column formatting JSON example
# ================================================================
sl = new_slide(prs, "Column Formatting: JSON Example (Status Field)")
# Code block
code_bg = add_rounded_rect(sl, Inches(0.6), Inches(1.5), Inches(7), Inches(4.8),
                           RGBColor(0x1E, 0x1E, 0x2E))
code_text = (
    '{\n'
    '  "$schema": "https://developer.microsoft.com/\n'
    '    json-schemas/sp/v2/column-formatting.schema.json",\n'
    '  "elmType": "div",\n'
    '  "style": {\n'
    '    "padding": "4px 8px",\n'
    '    "border-radius": "4px",\n'
    '    "background-color": {\n'
    '      "operator": "?",\n'
    '      "operands": [\n'
    '        { "operator": "==",\n'
    '          "operands": ["[$Status]","Active"] },\n'
    '        "#107C10",   // green\n'
    '        "#D13438"    // red\n'
    '      ]\n'
    '    }\n'
    '  },\n'
    '  "txtContent": "@currentField"\n'
    '}'
)
add_text_box(sl, Inches(0.9), Inches(1.7), Inches(6.4), Inches(4.5),
             code_text, font_size=13, color=RGBColor(0xCE, 0xD4, 0xDA),
             font_name=FONT_CODE)
# Explanation
add_text_box(sl, Inches(8.2), Inches(1.5), Inches(4.5), Inches(0.5),
             "What this does:", font_size=18, color=ACCENT_BLUE, bold=True)
bullets = [
    "Renders Status field as a colored badge",
    "Green (#107C10) if Status = 'Active'",
    "Red (#D13438) otherwise",
    "Uses the v2 column formatting schema",
    "No code deployment needed",
]
add_bullet_frame(sl, Inches(8.2), Inches(2.2), Inches(4.5), Inches(3.5),
                 bullets, font_size=15, icon="→")
add_speaker_notes(sl,
    "Show a real JSON example. The ternary operator (?) checks the Status field value. "
    "This is a great teaching moment: simple conditional logic, no SPFx required. "
    "Samples available at github.com/SharePoint/sp-dev-column-formatting. "
    "Schema URL: developer.microsoft.com/json-schemas/sp/v2/column-formatting.schema.json")

# ================================================================
# SLIDE 10 – View formatting
# ================================================================
sl = new_slide(prs, "View Formatting: Customizing Rows & Cards")
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Format current view → Choose layout → Customize rendering",
             font_size=18, color=DARK_GRAY, bold=True)
layouts = [
    ("List / Compact List", "Default row-based layout.\nApply alternating row colors\nor conditional row styles.", ACCENT_BLUE),
    ("Gallery", "Card-based layout.\nGreat for visual content\nlike projects or requests.", ACCENT_TEAL),
    ("Board", "Kanban-style view.\nGroup items by a column\n(e.g., Status or Priority).", ACCENT_PURPLE),
]
for i, (name, desc, color) in enumerate(layouts):
    x = Inches(0.6 + i * 4.1)
    card = add_rounded_rect(sl, x, Inches(2.4), Inches(3.8), Inches(3.6), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(2.4), Inches(3.8), Inches(0.07), color)
    add_text_box(sl, x + Inches(0.3), Inches(2.7), Inches(3.2), Inches(0.5),
                 name, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.4), Inches(3.2), Inches(2.2),
                 desc, font_size=15, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "Three layout options for view formatting. List is the default table view. "
    "Gallery renders cards — great for visual items. Board gives a Kanban experience. "
    "All three can be customized via the designer or advanced JSON mode.")

# ================================================================
# SLIDE 11 – When formatting is 'enough'
# ================================================================
sl = new_slide(prs, "When Is Formatting Enough?")
# Left: formatting is enough
add_rounded_rect(sl, Inches(0.6), Inches(1.5), Inches(5.5), Inches(4.8),
                 RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text_box(sl, Inches(0.9), Inches(1.7), Inches(5), Inches(0.5),
             "✅ Use JSON Formatting When…", font_size=18, color=GREEN, bold=True)
left_items = [
    "Highlighting status fields with colors/icons",
    "Improving scannability of list data",
    "Showing progress bars or conditional badges",
    "Creating card-based project views",
    "Adding mailto: or URL action links",
]
add_bullet_frame(sl, Inches(0.9), Inches(2.4), Inches(5), Inches(3.5),
                 left_items, font_size=16, icon="•", color=DARK_GRAY)
# Right: need SPFx
add_rounded_rect(sl, Inches(6.8), Inches(1.5), Inches(5.9), Inches(4.8),
                 RGBColor(0xFD, 0xE8, 0xE8), RED_ACCENT)
add_text_box(sl, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.5),
             "⚠️ Consider SPFx When…", font_size=18, color=RED_ACCENT, bold=True)
right_items = [
    "You need custom web parts or extensions",
    "Complex business logic or external API calls",
    "Fully custom UI beyond list/library rendering",
    "Integration with Teams tabs or Viva Connections",
    "Rich interactive dashboards",
]
add_bullet_frame(sl, Inches(7.1), Inches(2.4), Inches(5.3), Inches(3.5),
                 right_items, font_size=16, icon="•", color=DARK_GRAY)
add_speaker_notes(sl,
    "Help participants draw the line. JSON formatting is remarkably powerful for 'app-like' "
    "experiences without deployment overhead. SPFx is for truly custom UI/logic scenarios. "
    "Most admin teams should exhaust formatting options before requesting SPFx development.")

# ================================================================
# SLIDE 12 – Section Divider: SPFx & App Governance
# ================================================================
section_divider(prs, "Section 2", "SPFx Solutions & App Governance", accent=ACCENT_PURPLE)

# ================================================================
# SLIDE 13 – SPFx overview
# ================================================================
sl = new_slide(prs, "SharePoint Framework (SPFx) Overview", ACCENT_PURPLE)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "The modern extensibility model for SharePoint Online, Teams, and Viva Connections",
             font_size=18, color=DARK_GRAY)
facts = [
    ("📦", "Package Format", ".sppkg file uploaded\nto App Catalog"),
    ("🌐", "Runs Client-Side", "Executes in the browser\nin the user's context"),
    ("🔑", "Same Permissions", "Always runs with the\ncurrent user's permissions"),
    ("☁️", "Code Hosting", "Bundle hosted on CDN,\nAzure, or SharePoint"),
]
for i, (icon, title, desc) in enumerate(facts):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(sl, x, Inches(2.4), Inches(2.8), Inches(3.5), LIGHT_GRAY, MID_GRAY)
    add_text_box(sl, x + Inches(0.3), Inches(2.6), Inches(2.2), Inches(0.6),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.2), Inches(2.2), Inches(0.5),
                 title, font_size=17, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.3), Inches(3.9), Inches(2.2), Inches(1.5),
                 desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_speaker_notes(sl,
    "SPFx key facts: packages are .sppkg files. Code runs in the browser as the current user. "
    "The code bundle can be hosted anywhere — Office 365 CDN, Azure blob storage, or SPO itself. "
    "SPFx is the primary replacement for retired SharePoint Add-ins.")

# ================================================================
# SLIDE 14 – SPFx anatomy diagram
# ================================================================
sl = new_slide(prs, "Anatomy of an SPFx Solution", ACCENT_PURPLE)
# Left box: Package
add_rounded_rect(sl, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.2),
                 RGBColor(0xEE, 0xE8, 0xF5), ACCENT_PURPLE)
add_text_box(sl, Inches(1.1), Inches(2.0), Inches(4), Inches(0.5),
             "📦  .sppkg Package", font_size=20, color=ACCENT_PURPLE, bold=True)
pkg_items = [
    "Component manifest (metadata)",
    "URL pointing to bundle location",
    "Permissions declarations",
    "Feature definitions (optional)",
]
add_bullet_frame(sl, Inches(1.1), Inches(2.7), Inches(4), Inches(3),
                 pkg_items, font_size=15, icon="→")
# Arrow
add_text_box(sl, Inches(5.5), Inches(3.3), Inches(1.5), Inches(0.6),
             "→  deploys to →", font_size=14, color=ACCENT_PURPLE, bold=True,
             alignment=PP_ALIGN.CENTER)
# Right box: App Catalog
add_rounded_rect(sl, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.2),
                 RGBColor(0xEE, 0xE8, 0xF5), ACCENT_PURPLE)
add_text_box(sl, Inches(7.5), Inches(2.0), Inches(4.7), Inches(0.5),
             "🏪  App Catalog", font_size=20, color=ACCENT_PURPLE, bold=True)
cat_items = [
    "Tenant App Catalog — available org-wide",
    "Site Collection App Catalog — scoped to one site",
    "Admin reviews trust dialog before approval",
    "Enable / Disable solutions at any time",
    "Monitor installation count across sites",
]
add_bullet_frame(sl, Inches(7.5), Inches(2.7), Inches(4.7), Inches(3),
                 cat_items, font_size=15, icon="→")
add_speaker_notes(sl,
    "Two-part structure: the .sppkg contains manifest + CDN URL; the app catalog stores and governs it. "
    "Tenant app catalog for org-wide; site collection app catalog for isolated deployments. "
    "Admin always gets a trust dialog before approving. Can disable or remove at any time.")

# ================================================================
# SLIDE 15 – App Catalog: Tenant vs Site Collection
# ================================================================
sl = new_slide(prs, "Tenant vs Site Collection App Catalog")
# Headers
add_shape_rect(sl, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.65), ACCENT_BLUE)
add_text_box(sl, Inches(0.8), Inches(1.55), Inches(5.4), Inches(0.55),
             "Tenant App Catalog", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape_rect(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.65), ACCENT_TEAL)
add_text_box(sl, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.55),
             "Site Collection App Catalog", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
# Tenant items
t_items = [
    "Managed by SharePoint admin",
    "Solutions available org-wide",
    "Tenant-wide deployment option",
    "Central governance point",
    "One per tenant (+ one per geo in multi-geo)",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.5), Inches(5.4), Inches(3.5),
                 t_items, font_size=16, icon="•")
# Site col items
s_items = [
    "Managed by site collection admin",
    "Solutions scoped to that site only",
    "No 'Make available to all sites'",
    "Useful for isolated testing",
    "Multiple per tenant possible",
]
add_bullet_frame(sl, Inches(7.0), Inches(2.5), Inches(5.4), Inches(3.5),
                 s_items, font_size=16, icon="•")
add_speaker_notes(sl,
    "Tenant app catalog is the primary governance point. Site collection app catalogs allow "
    "isolated deployment for specific sites. In multi-geo tenants, each geo location gets its own catalog. "
    "For training: tenant-wide deployment is trainer-only in shared environments.")

# ================================================================
# SLIDE 16 – Deploying SPFx: governance flow
# ================================================================
sl = new_slide(prs, "SPFx Deployment Governance Flow", ACCENT_PURPLE)
flow_steps = [
    ("1", "Developer builds & bundles", "gulp bundle --ship\ngulp package-solution --ship", GREEN),
    ("2", "Admin uploads .sppkg", "Apps for SharePoint library\nin App Catalog site", ACCENT_BLUE),
    ("3", "Trust dialog review", "Full trust client-side code?\nCheck CDN domain origin", ORANGE),
    ("4", "Enable & deploy scope", "Tenant-wide vs per-site\ninstall decision", ACCENT_PURPLE),
    ("5", "Monitor & maintain", "Track installations, update\nor disable as needed", ACCENT_TEAL),
]
for i, (num, title, desc, color) in enumerate(flow_steps):
    y = Inches(1.4 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.5), Inches(0.5), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.6), y, Inches(4), Inches(0.5),
                 title, font_size=17, color=color, bold=True)
    add_text_box(sl, Inches(6), y, Inches(6.5), Inches(0.7),
                 desc, font_size=14, color=DARK_GRAY)
add_speaker_notes(sl,
    "Five-step governance flow from development to production. "
    "Emphasize step 3: the trust dialog is the admin's key checkpoint. "
    "Step 5: admins can disable a solution immediately across all sites if needed.")

# ================================================================
# SLIDE 17 – 2026 Alignment: Add-ins Retirement
# ================================================================
sl = new_slide(prs, "2026 Alignment: SharePoint Add-Ins Retirement", ORANGE)
# Warning box
add_rounded_rect(sl, Inches(0.6), Inches(1.5), Inches(12), Inches(1.4),
                 RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text_box(sl, Inches(0.9), Inches(1.65), Inches(11.4), Inches(1),
             "⚠️  Microsoft has announced the retirement of SharePoint Add-ins for SharePoint Online.\n"
             "The SPFx is the primary replacement technology. Add-ins should be treated as legacy.",
             font_size=16, color=DARK_GRAY)
add_text_box(sl, Inches(0.8), Inches(3.3), Inches(11), Inches(0.5),
             "What This Means for Admins:", font_size=20, color=ORANGE, bold=True)
items = [
    "New customization projects → SPFx or declarative (JSON)",
    "Existing add-ins → plan migration timeline",
    "App catalog supports both .sppkg (SPFx) and .app (legacy) — for now",
    "Training focus: SPFx + declarative customization going forward",
    "Add-in model retirement does NOT affect SPFx",
]
add_bullet_frame(sl, Inches(0.8), Inches(4.0), Inches(11), Inches(2.5),
                 items, font_size=17, icon="→")
add_speaker_notes(sl,
    "Important 2026 context. The add-in model is being retired for SharePoint Online. "
    "SPFx continues to be fully supported and invested in. "
    "From Microsoft docs: 'The SharePoint add-in model deprecation does not impact SPFx.' "
    "Train admins to focus on SPFx and declarative approaches.")

# ================================================================
# SLIDE 18 – Section Divider: API Access & Permissions
# ================================================================
section_divider(prs, "Section 3", "API Access & Permissions Governance", accent=ACCENT_TEAL)

# ================================================================
# SLIDE 19 – API access: why it exists
# ================================================================
sl = new_slide(prs, "API Access: Why It Exists", ACCENT_TEAL)
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "SPFx solutions and custom scripts can request permissions to\n"
             "Microsoft Entra ID-secured APIs. Admins manage these via API access.",
             font_size=18, color=DARK_GRAY)
# Diagram: Solution → requests → API access → Entra ID → API
boxes_data = [
    ("SPFx Solution", ACCENT_PURPLE),
    ("requests\npermission", None),
    ("API Access\n(Admin Center)", ACCENT_TEAL),
    ("approves via", None),
    ("Microsoft\nEntra ID", ACCENT_BLUE),
]
x_positions = [Inches(0.5), Inches(2.9), Inches(4.3), Inches(7.5), Inches(9)]
for i, (label, color) in enumerate(boxes_data):
    x = x_positions[i]
    if color:
        w = Inches(2.2) if i != 4 else Inches(2.5)
        box = add_rounded_rect(sl, x, Inches(3.5), w, Inches(1.4), color)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        add_text_box(sl, x, Inches(3.8), Inches(1.5), Inches(0.6),
                     label, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
# Key point
add_rounded_rect(sl, Inches(0.6), Inches(5.5), Inches(12), Inches(1),
                 RGBColor(0xE0, 0xF7, 0xF1), ACCENT_TEAL)
add_text_box(sl, Inches(0.9), Inches(5.6), Inches(11.4), Inches(0.8),
             "🔐  API access is a governance & security surface — not a routine click-through.\n"
             "Every approval grants the solution access to APIs on behalf of all tenant users.",
             font_size=15, color=DARK_GRAY)
add_speaker_notes(sl,
    "API access is where solutions request Entra ID-secured permissions. "
    "The admin center page (SharePoint admin center → API access) shows pending and approved requests. "
    "Each approval grants application-level consent — affects the entire tenant.")

# ================================================================
# SLIDE 20 – Roles and approvals
# ================================================================
sl = new_slide(prs, "API Access: Roles and Approvals", ACCENT_TEAL)
# Table-like layout
headers = ["API Scope", "Required Role", "Notes"]
rows = [
    ["Third-party APIs", "Application Administrator", "Sufficient for most external APIs"],
    ["Microsoft Graph", "Global Administrator", "Highest privilege required"],
    ["Other Microsoft APIs", "Global Administrator", "e.g., Outlook, Teams APIs"],
    ["Custom line-of-business", "Application Administrator", "Your org's own APIs via Entra ID"],
]
# Header row
for j, h in enumerate(headers):
    x = Inches(0.6 + j * 4)
    w = Inches(3.8) if j < 2 else Inches(4.5)
    add_shape_rect(sl, x, Inches(1.5), w, Inches(0.6), ACCENT_TEAL)
    add_text_box(sl, x + Inches(0.15), Inches(1.53), w - Inches(0.3), Inches(0.55),
                 h, font_size=15, color=WHITE, bold=True)
# Data rows
for i, row in enumerate(rows):
    y = Inches(2.15 + i * 0.9)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        x = Inches(0.6 + j * 4)
        w = Inches(3.8) if j < 2 else Inches(4.5)
        add_shape_rect(sl, x, y, w, Inches(0.85), bg, MID_GRAY)
        add_text_box(sl, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.65),
                     cell, font_size=14, color=DARK_GRAY)
# Warning
add_rounded_rect(sl, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
                 RGBColor(0xFF, 0xF3, 0xE0), ORANGE)
add_text_box(sl, Inches(0.9), Inches(5.9), Inches(11.4), Inches(0.6),
             "⚠️  Approvals affect the entire tenant. Review each request carefully before approving.",
             font_size=15, color=DARK_GRAY, bold=True)
add_speaker_notes(sl,
    "Key governance table. Third-party APIs need Application Administrator. "
    "Microsoft Graph and other Microsoft APIs require Global Administrator. "
    "Emphasize that each approval is tenant-wide — not scoped to a single site or user.")

# ================================================================
# SLIDE 21 – Permissions and Governance: Scope Matters
# ================================================================
sl = new_slide(prs, "Governance: Scope Matters")
add_text_box(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             '"Don\'t surprise the tenant" — the #1 rule for app deployment',
             font_size=20, color=ACCENT_BLUE, bold=True)
# Two columns
col_data = [
    ("Site-Scoped Deployment", GREEN,
     ["Solution available in one site only",
      "Lower risk, easier to roll back",
      "Good for testing and pilots",
      "Site collection admin can manage"]),
    ("Tenant-Wide Deployment", RED_ACCENT,
     ["Solution available everywhere immediately",
      "Higher risk — affects all users",
      "Requires SharePoint admin approval",
      "Cannot be 'un-deployed' per site"]),
]
for i, (title, color, items) in enumerate(col_data):
    x = Inches(0.6 + i * 6.3)
    add_rounded_rect(sl, x, Inches(2.3), Inches(5.8), Inches(4.0), LIGHT_GRAY, color)
    add_shape_rect(sl, x, Inches(2.3), Inches(5.8), Inches(0.07), color)
    add_text_box(sl, x + Inches(0.3), Inches(2.5), Inches(5.2), Inches(0.5),
                 title, font_size=18, color=color, bold=True)
    add_bullet_frame(sl, x + Inches(0.3), Inches(3.2), Inches(5.2), Inches(2.8),
                     items, font_size=15, icon="•")
add_speaker_notes(sl,
    "Scope is the most important governance consideration. Tenant-wide deployment affects all users "
    "immediately. In training environments, tenant-wide deployment should be trainer-only. "
    "In production, always test in a site collection app catalog first.")

# ================================================================
# SLIDE 22 – Section Divider: Lab
# ================================================================
section_divider(prs, "Section 4", "Lab 7: Apps and Customization", accent=ORANGE)

# ================================================================
# SLIDE 23 – Lab preview
# ================================================================
sl = new_slide(prs, "Lab 7: Hands-On Exercises", ORANGE)
exercises = [
    ("Task 1", "Create list NW-Pxx-AppRequests", "Build a list to track app/customization requests"),
    ("Task 2", "Apply column formatting (Status)", "Use conditional formatting to color-code the Status field"),
    ("Task 3", "Apply view formatting (row shading)", "Add alternating row colors for improved readability"),
    ("Task 4", "Trainer-led tour: Apps page", "Explore More features → Apps in SharePoint admin center"),
    ("Task 5", "Trainer-led tour: API access", "Review the API access page and pending requests"),
]
for i, (task, title, desc) in enumerate(exercises):
    y = Inches(1.4 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(1.2), Inches(0.5), ORANGE)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = task
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(2.3), y, Inches(4.5), Inches(0.5),
                 title, font_size=17, color=DARK_GRAY, bold=True)
    add_text_box(sl, Inches(7.2), y, Inches(5.5), Inches(0.5),
                 desc, font_size=14, color=DARK_GRAY)
# Time estimate
add_rounded_rect(sl, Inches(8.5), Inches(6.2), Inches(4), Inches(0.6), ORANGE)
add_text_box(sl, Inches(8.7), Inches(6.25), Inches(3.6), Inches(0.5),
             "⏱️  Estimated time: 30–40 min", font_size=15, color=WHITE, bold=True)
add_speaker_notes(sl,
    "Tasks 1-3 are hands-on for participants. Tasks 4-5 are trainer-led demonstrations "
    "in the SharePoint admin center. Keep the app catalog and API access tours brief — "
    "participants observe while the trainer navigates.")

# ================================================================
# SLIDE 24 – Validation checklist
# ================================================================
sl = new_slide(prs, "Lab 7: Validation Checklist", GREEN)
checks = [
    "NW-Pxx-AppRequests list created with Title, Status, RequestedBy, Description columns",
    "Column formatting applied — Status shows colored badges (green/red/yellow)",
    "View formatting applied — alternating row colors visible in All Items view",
    "Trainer demo: Apps page shown in SharePoint admin center → More features",
    "Trainer demo: API access page reviewed — participants can explain its purpose",
]
for i, item in enumerate(checks):
    y = Inches(1.5 + i * 1.05)
    # Checkbox icon
    add_rounded_rect(sl, Inches(0.8), y, Inches(0.45), Inches(0.45), GREEN)
    tf = sl.shapes[-1].text_frame
    p = tf.paragraphs[0]
    p.text = "✓"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.5), y, Inches(11), Inches(0.6),
                 item, font_size=16, color=DARK_GRAY)
add_speaker_notes(sl,
    "Walk through each checkpoint. Verify formatting works in participants' browsers. "
    "For trainer-led items, confirm participants can articulate what the Apps page "
    "and API access page are used for.")

# ================================================================
# SLIDE 25 – Common Issues / Troubleshooting
# ================================================================
sl = new_slide(prs, "Common Issues & Troubleshooting", RED_ACCENT)
issues = [
    ("JSON syntax error", "Preview shows error", "Check for missing commas, brackets;\nuse a JSON validator or VS Code"),
    ("Formatting not appearing", "Column looks normal", "Verify you saved; check you're\nviewing the correct view"),
    ("Apps page not visible", "Admin center navigation", "Go to More features → Apps;\nrequires SharePoint admin role"),
    ("API access empty", "No pending requests", "Normal if no solutions have\nrequested API permissions yet"),
]
for i, (issue, symptom, fix) in enumerate(issues):
    y = Inches(1.5 + i * 1.35)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rounded_rect(sl, Inches(0.6), y, Inches(12), Inches(1.2), bg, MID_GRAY)
    add_text_box(sl, Inches(0.9), y + Inches(0.1), Inches(3), Inches(0.4),
                 issue, font_size=15, color=RED_ACCENT, bold=True)
    add_text_box(sl, Inches(0.9), y + Inches(0.55), Inches(3), Inches(0.5),
                 symptom, font_size=13, color=DARK_GRAY)
    add_text_box(sl, Inches(4.5), y + Inches(0.1), Inches(7.8), Inches(1),
                 fix, font_size=14, color=DARK_GRAY)
add_speaker_notes(sl,
    "Most common lab issue is JSON syntax errors — remind participants to use the "
    "designer (conditional formatting mode) first before switching to advanced JSON. "
    "Apps page requires admin role; participants with viewer-only access won't see it.")

# ================================================================
# SLIDE 26 – Key takeaways
# ================================================================
sl = new_slide(prs, "Key Takeaways")
takeaways = [
    "Prefer OOB first → then JSON formatting → then SPFx (customization spectrum)",
    "JSON formatting changes rendering only — never modifies list data",
    "Column formatting = single field; View formatting = rows/cards/boards",
    "SPFx solutions are .sppkg files deployed via the App Catalog",
    "Tenant-wide deployment needs careful governance — don't surprise the tenant",
    "API access approvals are tenant-wide and require appropriate admin roles",
    "SharePoint Add-ins are being retired — focus on SPFx + declarative",
]
for i, item in enumerate(takeaways):
    y = Inches(1.4 + i * 0.8)
    add_shape_rect(sl, Inches(0.8), y + Inches(0.15), Inches(0.08), Inches(0.35), ACCENT_BLUE)
    add_text_box(sl, Inches(1.1), y, Inches(11.5), Inches(0.65),
                 item, font_size=16, color=DARK_GRAY)
add_speaker_notes(sl,
    "Seven key takeaways. Reinforce the customization spectrum and governance principles. "
    "The add-in retirement note is important for 2026-aligned training.")

# ================================================================
# SLIDE 27 – Knowledge check
# ================================================================
sl = new_slide(prs, "Knowledge Check")
questions = [
    ("Q1", "What are the three tiers of the customization spectrum?"),
    ("Q2", "Does column formatting change the underlying list data? Why or why not?"),
    ("Q3", "What file format is used to deploy SPFx solutions to the App Catalog?"),
    ("Q4", "What admin role is needed to approve Microsoft Graph API access requests?"),
    ("Q5", "Why is tenant-wide deployment considered higher risk than site-scoped?"),
]
colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ORANGE, GREEN]
for i, (qnum, text) in enumerate(questions):
    y = Inches(1.4 + i * 1.1)
    badge = add_rounded_rect(sl, Inches(0.8), y, Inches(0.7), Inches(0.55), colors[i])
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = qnum
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_text_box(sl, Inches(1.8), y + Inches(0.05), Inches(10.5), Inches(0.55),
                 text, font_size=17, color=DARK_GRAY)
add_speaker_notes(sl,
    "Answers: Q1 — OOB Configuration, Declarative (JSON Formatting), SPFx. "
    "Q2 — No, formatting only changes rendering/display, not the stored data. "
    "Q3 — .sppkg files. "
    "Q4 — Global Administrator. "
    "Q5 — Because it immediately makes the solution available to all sites/users, "
    "and cannot be un-deployed per-site; affects the entire tenant.")

# ================================================================
# SLIDE 28 – Thank you / Next Module
# ================================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(sl, DARK_BG)
add_shape_rect(sl, Inches(1), Inches(3.1), Inches(1.2), Inches(0.06), ACCENT_TEAL)
add_text_box(sl, Inches(1), Inches(2.2), Inches(11), Inches(0.7),
             "End of Module 7", font_size=20, color=ACCENT_TEAL)
add_text_box(sl, Inches(1), Inches(3.35), Inches(11), Inches(1),
             "Apps and Customization\nin SharePoint Online", font_size=36, color=WHITE, bold=True)
add_text_box(sl, Inches(1), Inches(4.8), Inches(10), Inches(0.6),
             "Up Next  →  Module 8: Compliance and Governance with Microsoft Purview",
             font_size=18, color=MID_GRAY)
add_text_box(sl, Inches(1), Inches(5.6), Inches(10), Inches(0.5),
             "🎉  Day 2 Complete!  Take a well-deserved break.",
             font_size=16, color=ORANGE)
add_footer_bar(sl)
add_speaker_notes(sl,
    "This concludes Module 7 and Day 2 of the training. "
    "Day 3 begins with Module 8: Compliance and Governance with Microsoft Purview. "
    "Remind participants to save their lab work and take any screenshots needed.")

# ── Save ───────────────────────────────────────────────────────
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "Module-07-Slides.pptx")
prs.save(out_path)
print(f"✅ Saved {len(prs.slides)}-slide presentation → {out_path}")
