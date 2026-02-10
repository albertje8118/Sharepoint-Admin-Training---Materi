"""
Generate Module 6 PPTX — Search in SharePoint Online and Microsoft Search
Modern, engaging design matching Modules 1-5 style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG     = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEC, 0xF9)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT    = RGBColor(0x24, 0x24, 0x24)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT  = RGBColor(0xD1, 0x34, 0x38)
YELLOW_WARN = RGBColor(0xF2, 0xC8, 0x11)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ═══════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════
def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6), icon="▸", font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon} {bullet}"; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.name = font_name; p.space_after = spacing
    return txBox

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)

def add_footer_bar(slide, slide_num, total, module_label="Module 6"):
    bar_top = SLIDE_HEIGHT - Inches(0.45)
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_WIDTH, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), bar_top + Inches(0.05), Inches(7), Inches(0.35),
                 f"{module_label}  |  Search in SharePoint Online & Microsoft Search",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), bar_top + Inches(0.05),
                 Inches(1.2), Inches(0.35), f"{slide_num} / {total}",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)

def section_divider(slide, section_title, subtitle="", icon=""):
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_TEAL)
    if icon:
        add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                     icon, font_size=60, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                 section_title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

TOTAL_SLIDES = 28
slide_counter = [0]

def new_slide(notes=""):
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if notes: add_speaker_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome to Module 6 — Search in SharePoint Online and Microsoft Search. "
              "Search is only as good as the content quality, metadata, and permissions behind it. "
              "In this module we cover security trimming, Microsoft Search answers "
              "(Bookmarks & Acronyms), search schema basics, and how to troubleshoot.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
             "MODERN SHAREPOINT ONLINE FOR ADMINISTRATORS  ·  DAY 2",
             font_size=14, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(2.3), Inches(11), Inches(1.5),
             "Module 6", font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(1),
             "Search in SharePoint Online\n& Microsoft Search",
             font_size=32, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(2), Inches(5.4), Inches(9), Inches(0.6),
             "Find anything  ·  Curate answers  ·  Troubleshoot like a pro",
             font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(6.1), Inches(2.3), Inches(0.04), ACCENT_PURPLE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Why Admins Care About Search (Overview / Motivation)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Information discovery equals productivity. If users can't find content, "
              "it might as well not exist. Search quality depends on four pillars: "
              "permissions, content quality, metadata, and governance.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Why Admins Care About Search", font_size=36, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "If users can't find it, it doesn't exist — search is a core admin responsibility",
             font_size=16, color=MID_GRAY)

cards = [
    ("🔐", "Permissions", "Foundation", "Security trimming ensures\nusers only see what they\nhave access to"),
    ("📄", "Content", "Quality", "Well-named files with\nconsistent metadata are\nmore discoverable"),
    ("🏷️", "Metadata", "Drives Refiners", "Managed metadata from\nModule 5 powers search\nrefinement panels"),
    ("🛡️", "Governance", "Controls", "Admin-curated answers,\nschema settings, and\nindexing policies"),
]
for i, (icon, stat, label, desc) in enumerate(cards):
    x = Inches(0.8 + i * 3.1)
    card = add_rounded_rect(s, x, Inches(1.9), Inches(2.8), Inches(3.6), WHITE)
    card.shadow.inherit = False
    add_text_box(s, x + Inches(0.3), Inches(2.1), Inches(2.2), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.8),
                 stat, font_size=44, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(3.5), Inches(2.2), Inches(0.5),
                 label, font_size=16, bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_shape_rect(s, x + Inches(0.6), Inches(4.1), Inches(1.6), Inches(0.04), ACCENT_TEAL)
    add_text_box(s, x + Inches(0.2), Inches(4.3), Inches(2.4), Inches(1.0),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Learning Outcomes
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Four learning objectives covering the admin's view of search.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🎯  Learning Outcomes", font_size=32, bold=True, color=DARK_BG)

outcomes = [
    "Explain how SharePoint Search and Microsoft Search relate, and what security trimming means",
    "Describe Microsoft Search answers (Bookmarks & Acronyms) and who manages them",
    "Validate whether a site/library is searchable and request reindexing when appropriate",
    "Use basic query syntax (phrases, AND/OR/NOT, property restrictions) to troubleshoot",
]
for i, outcome in enumerate(outcomes):
    y = Inches(1.6 + i * 1.15)
    badge = add_rounded_rect(s, Inches(0.8), y, Inches(0.55), Inches(0.55), ACCENT_BLUE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.5),
                 outcome, font_size=18, color=DARK_TEXT)

add_footer_bar(s, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Section Divider: Search Concepts
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: let's start with the core search concepts every admin needs.")
section_divider(s, "Search Concepts",
                "Security trimming, indexing, and how search really works", "🔍")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — How Search Works (High-Level Flow)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("The search pipeline: content is crawled from lists and libraries, "
              "site columns are mapped to managed properties in the index, "
              "user queries are matched against the index, and results are security-trimmed "
              "before being returned. This is the conceptual foundation.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "How Search Works — The Pipeline", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "From content to results — four stages",
             font_size=16, color=MID_GRAY)

stages = [
    ("📚", "Crawl", "Content & metadata\ndiscovered from lists,\nlibraries, and sites", ACCENT_BLUE),
    ("🗂️", "Index", "Crawled properties mapped\nto managed properties;\nstored in search index", ACCENT_TEAL),
    ("🔎", "Query", "User enters search query;\nmatched against managed\nproperties in the index", ACCENT_PURPLE),
    ("🔐", "Trim & Return", "Results filtered by user\npermissions (security\ntrimming) and displayed", ORANGE),
]
for i, (icon, title, desc, color) in enumerate(stages):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.4), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(2.05), Inches(2.3), Inches(0.6),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.3), Inches(2.7), Inches(2.3), Inches(0.5),
                 title, font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.3), Inches(2.5), Inches(1.5),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    # Arrow between cards
    if i < 3:
        arrow_x = x + Inches(2.95)
        add_text_box(s, arrow_x, Inches(3.1), Inches(0.3), Inches(0.5),
                     "→", font_size=28, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Insight
insight = add_rounded_rect(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(5.65), Inches(9.7), Inches(0.6),
             "💡  Key Insight:  Search only finds what's in the index, and only shows results "
             "the user has permission to see. Both content quality and permissions matter.",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Security Trimming (The #1 Rule)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Security trimming is non-negotiable. Search NEVER overrides permissions. "
              "If a document is missing from search results for a user, the first thing "
              "to check is: does that user have access to the item?")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, RED_ACCENT)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🔐  Security Trimming — The #1 Rule", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Search never overrides permissions — this is by design, non-negotiable",
             font_size=16, color=MID_GRAY)

# Large central visual
center_card = add_rounded_rect(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.0), WHITE)
center_card.shadow.inherit = False
add_shape_rect(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.08), RED_ACCENT)

# Left scenario
add_text_box(s, Inches(2.0), Inches(2.1), Inches(4.5), Inches(0.5),
             "👤 User A (has access)", font_size=20, bold=True, color=GREEN)
add_bullet_frame(s, Inches(2.0), Inches(2.7), Inches(4.5), Inches(1.8),
                 ["Searches for 'Northwind Contract'",
                  "Has Read access to the library",
                  "✅ Document appears in results",
                  "Search refiners show metadata"],
                 font_size=14, color=DARK_TEXT, icon="")

# Divider
add_shape_rect(s, Inches(6.6), Inches(2.1), Inches(0.04), Inches(2.4), LIGHT_GRAY)

# Right scenario
add_text_box(s, Inches(7.0), Inches(2.1), Inches(4.5), Inches(0.5),
             "👤 User B (no access)", font_size=20, bold=True, color=RED_ACCENT)
add_bullet_frame(s, Inches(7.0), Inches(2.7), Inches(4.5), Inches(1.8),
                 ["Searches for 'Northwind Contract'",
                  "Has NO access to the library",
                  "❌ Document does NOT appear",
                  "As if the document doesn't exist"],
                 font_size=14, color=DARK_TEXT, icon="")

# Takeaway boxes
boxes = [
    ("🔑", "Admin Takeaway", "'Missing' search results? Check permissions first!", ACCENT_BLUE),
    ("⚡", "Performance", "Security trimming happens at query time, per user", ACCENT_TEAL),
    ("🛡️", "Compliance", "No search backdoor — data classification is enforced", ACCENT_PURPLE),
]
for i, (icon, title, desc, color) in enumerate(boxes):
    x = Inches(0.8 + i * 4.15)
    box = add_rounded_rect(s, x, Inches(5.2), Inches(3.85), Inches(1.0), WHITE)
    box.shadow.inherit = False
    add_shape_rect(s, x, Inches(5.2), Inches(3.85), Inches(0.06), color)
    add_text_box(s, x + Inches(0.2), Inches(5.35), Inches(0.5), Inches(0.4),
                 icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.7), Inches(5.35), Inches(2.9), Inches(0.35),
                 title, font_size=14, bold=True, color=color)
    add_text_box(s, x + Inches(0.7), Inches(5.7), Inches(2.9), Inches(0.4),
                 desc, font_size=12, color=DARK_TEXT)

add_footer_bar(s, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Search Entry Points
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Users search from multiple entry points: the SharePoint search box, "
              "Microsoft 365 app bar, Teams, Outlook, and more. The underlying engine "
              "is increasingly unified (Microsoft Search), but admin surfaces differ.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Search Entry Points — Where Users Search", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Multiple experiences, one underlying engine — Microsoft Search",
             font_size=16, color=MID_GRAY)

entry_points = [
    ("🌐", "SharePoint\nSearch Box", "Site-scoped or hub-scoped\nsearch within SharePoint\nOnline sites", ACCENT_BLUE),
    ("📱", "Microsoft 365\nApp Bar", "Organization-wide search\nacross M365 from the\ntop navigation bar", ACCENT_TEAL),
    ("💬", "Microsoft\nTeams", "Search files, messages,\nand people directly\nfrom Teams", ACCENT_PURPLE),
    ("📧", "Outlook &\nOther Apps", "Find files and people\nfrom Outlook, Word,\nand other M365 apps", ORANGE),
]
for i, (icon, title, desc, color) in enumerate(entry_points):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.5), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(2.05), Inches(2.3), Inches(0.6),
                 icon, font_size=40, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.8),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.6), Inches(2.5), Inches(1.2),
                 desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Admin note
note = add_rounded_rect(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1.1), Inches(5.75), Inches(11.1), Inches(0.6),
             "💡  Admin Note:  Some config is site/library-scoped (safe for participant labs). "
             "Some config is organization-level (use NW-Pxx naming and trainer governance).",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — What Admins Can Tune
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Four big buckets of admin control over search experience.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "What Admins Can Tune", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Four levers to improve the search experience",
             font_size=16, color=MID_GRAY)

buckets = [
    ("📄", "Content Quality", ACCENT_BLUE,
     ["Good file names & descriptions",
      "Consistent metadata (Module 5)",
      "Clean information architecture",
      "Avoid duplicate / outdated content"]),
    ("👁️", "Search Visibility", ACCENT_TEAL,
     ["Site/library search settings",
      "Control what appears in results",
      "'Allow items to appear in search'",
      "Site-level search visibility toggle"]),
    ("📌", "Curated Answers", ACCENT_PURPLE,
     ["Bookmarks (promoted links)",
      "Acronyms (definitions)",
      "Org-wide answers at top of results",
      "Managed in Search & intelligence"]),
    ("⚙️", "Search Schema", ORANGE,
     ["Crawled → managed property mapping",
      "Custom refiners for navigation",
      "Advanced: requires reindex",
      "Treat as trainer-led in this course"]),
]
for i, (icon, title, color, items) in enumerate(buckets):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(4.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    add_text_box(s, x + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.6),
                 f"{icon}  {title}", font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(s, x + Inches(0.3), Inches(2.7), Inches(2.3), Inches(3.0),
                     items, font_size=13, color=DARK_TEXT, icon="•")

add_footer_bar(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Section Divider: Microsoft Search Answers
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: now let's look at the admin-curated answers that surface "
              "at the top of search results.")
section_divider(s, "Microsoft Search Answers",
                "Bookmarks, Acronyms, and curated content", "📌")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Microsoft Search Admin Surface
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("The admin entry point for Microsoft Search answers is in the "
              "Microsoft 365 admin center under Settings → Search & intelligence. "
              "Two key roles: Search admin and Search editor.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Microsoft Search Admin Surface", font_size=32, bold=True, color=DARK_BG)

# Path card
path_card = add_rounded_rect(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.2), WHITE)
path_card.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(1.6), Inches(10.9), Inches(0.4),
             "📍  Microsoft 365 admin center → Settings → Search & intelligence",
             font_size=18, bold=True, color=ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(2.1), Inches(10.9), Inches(0.4),
             "Also accessible via:  Copilot → Search  (in updated tenants, June 2025+)",
             font_size=14, color=MID_GRAY)

# Roles side by side
role_data = [
    ("🔧", "Search Admin", ACCENT_BLUE,
     ["Full access to Search & intelligence",
      "Create/manage all answer types",
      "Manage search schema settings",
      "Assign Search editor role",
      "View search analytics"]),
    ("✏️", "Search Editor", ACCENT_TEAL,
     ["Create & manage Bookmarks",
      "Create & manage Acronyms",
      "Create & manage Q&As",
      "Cannot manage schema or roles",
      "Ideal for content stewards"]),
]
for i, (icon, title, color, items) in enumerate(role_data):
    x = Inches(0.8 + i * 6.2)
    card = add_rounded_rect(s, x, Inches(3.0), Inches(5.8), Inches(3.5), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(3.0), Inches(5.8), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(3.2), Inches(5.2), Inches(0.5),
                 f"{icon}  {title}", font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(s, x + Inches(0.5), Inches(3.8), Inches(4.8), Inches(2.5),
                     items, font_size=14, color=DARK_TEXT, icon="▸")

add_footer_bar(s, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Bookmarks
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Bookmarks are admin-curated links triggered by keywords. "
              "They appear at the top of search results immediately after publishing. "
              "In our shared tenant, use NW-Pxx prefixed keywords to avoid conflicts.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📌  Bookmarks — Curated Links", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Promoted answers triggered by keywords — visible immediately",
             font_size=16, color=MID_GRAY)

# Left: What they are
left = add_rounded_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.8), WHITE)
left.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.5),
             "What Are Bookmarks?", font_size=20, bold=True, color=ACCENT_BLUE)
bm_items = [
    "Admin-curated links triggered by keywords",
    "Appear at TOP of search results",
    "Available immediately after publishing",
    "Can be Draft, Published, or Scheduled",
    "Customizable: title, URL, description, audience",
]
add_bullet_frame(s, Inches(1.2), Inches(2.6), Inches(5.0), Inches(1.8),
                 bm_items, font_size=14, color=DARK_TEXT, icon="▸")

# Right: How to create
right = add_rounded_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.8), WHITE)
right.shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.5),
             "How to Create", font_size=20, bold=True, color=ACCENT_TEAL)
create_items = [
    "1. Go to M365 admin center → Search & intelligence",
    "2. Select Bookmarks → Add bookmark",
    "3. Enter title, URL, description",
    "4. Add trigger keywords (use NW-Pxx-…)",
    "5. Publish → available immediately",
]
add_bullet_frame(s, Inches(7.3), Inches(2.6), Inches(5.0), Inches(1.8),
                 create_items, font_size=14, color=DARK_TEXT, icon="")

# Shared tenant warning
warn = add_rounded_rect(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.7), RGBColor(0xFD, 0xE7, 0xE9))
add_text_box(s, Inches(1.1), Inches(5.05), Inches(11.1), Inches(0.6),
             "⚠️  Shared Tenant:  Only use training-specific keywords with NW-Pxx prefix. "
             "Never use generic terms like 'HR', 'IT', or 'Benefits' — they affect the whole organization.",
             font_size=14, color=RED_ACCENT)

add_footer_bar(s, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Acronyms
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Acronyms can be admin-curated or system-curated. Admin acronyms "
              "can be Draft or Published. Published acronyms take up to a day to appear. "
              "System-curated acronyms are discovered from emails and documents automatically.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📝  Acronyms — Definitions at a Glance", font_size=32, bold=True, color=DARK_BG)

# Two types side by side
types = [
    ("👤", "Admin-Curated", ACCENT_BLUE,
     ["Manually created by Search admin/editor",
      "Set to Draft or Published state",
      "Published → available within ~24 hours",
      "Can also be Excluded (blocked)",
      "Bulk import via CSV supported"]),
    ("🤖", "System-Curated", ACCENT_TEAL,
     ["Automatically discovered by Microsoft Search",
      "Mined from emails, documents, public data",
      "No admin action required",
      "Admins can exclude unwanted acronyms",
      "Supplements admin-curated definitions"]),
]
for i, (icon, title, color, items) in enumerate(types):
    x = Inches(0.8 + i * 6.2)
    card = add_rounded_rect(s, x, Inches(1.4), Inches(5.8), Inches(3.8), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.4), Inches(5.8), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(1.6), Inches(5.2), Inches(0.5),
                 f"{icon}  {title}", font_size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_bullet_frame(s, x + Inches(0.5), Inches(2.2), Inches(4.8), Inches(2.8),
                     items, font_size=14, color=DARK_TEXT, icon="▸")

# Timing note
timing = add_rounded_rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1.1), Inches(5.65), Inches(11.1), Inches(0.6),
             "⏱️  Important:  Published acronyms take up to a day to appear in search results. "
             "Plan accordingly during labs — Bookmarks (immediate) are easier to validate.",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Q&A Answers (Status Note)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Q&A answers were historically part of Microsoft Search. With the retirement "
              "of Microsoft Search in Bing (March 2025), some answer types may not be available. "
              "We teach Q&As as a concept but focus labs on Bookmarks and Acronyms.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ORANGE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "❓  Q&A Answers — Status Note (2026)", font_size=30, bold=True, color=DARK_BG)

# Status card
status = add_rounded_rect(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.2), WHITE)
status.shadow.inherit = False
add_shape_rect(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.08), ORANGE)

add_text_box(s, Inches(2.0), Inches(1.8), Inches(9.3), Inches(0.5),
             "📋  What You Need to Know", font_size=22, bold=True, color=ORANGE)

qna_items = [
    ("📖", "Concept", "Q&As are question-answer pairs curated by admins,\nhistorically shown at the top of search results"),
    ("🔄", "Status", "Microsoft Search in Bing retired (March 2025). Q&A availability\nvaries by tenant — check your admin center"),
    ("🎯", "Our Approach", "Teach Q&As as a concept for exam/knowledge purposes.\nDesign hands-on labs around Bookmarks + Acronyms"),
    ("✅", "If Available", "If your tenant still shows Q&As in Search & intelligence,\nthey work similarly to Bookmarks (keyword-triggered)"),
]
for i, (icon, label, desc) in enumerate(qna_items):
    y = Inches(2.5 + i * 0.75)
    add_text_box(s, Inches(2.0), y, Inches(0.5), Inches(0.4),
                 icon, font_size=22, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.6), y, Inches(1.5), Inches(0.4),
                 label, font_size=16, bold=True, color=ORANGE)
    add_text_box(s, Inches(4.3), y, Inches(7.0), Inches(0.7),
                 desc, font_size=13, color=DARK_TEXT)

add_footer_bar(s, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Answer Types Comparison
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Quick comparison table of the three answer types: Bookmarks vs Acronyms vs Q&As.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Answer Types Comparison", font_size=32, bold=True, color=DARK_BG)

# Table header
header_y = Inches(1.4)
cols = [
    (Inches(0.8), Inches(2.5), "Feature"),
    (Inches(3.3), Inches(3.3), "📌 Bookmarks"),
    (Inches(6.6), Inches(3.3), "📝 Acronyms"),
    (Inches(9.9), Inches(2.9), "❓ Q&As"),
]
for x, w, label in cols:
    add_shape_rect(s, x, header_y, w, Inches(0.5), ACCENT_BLUE)
    add_text_box(s, x + Inches(0.1), header_y + Inches(0.05), w - Inches(0.2), Inches(0.4),
                 label, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

rows = [
    ("Purpose", "Promoted links", "Definitions", "FAQ answers"),
    ("Trigger", "Keyword match", "Acronym query", "Keyword match"),
    ("Created by", "Search admin/editor", "Admin or system", "Search admin/editor"),
    ("Availability", "Immediate", "Up to 24 hours", "Varies by tenant"),
    ("States", "Draft / Published /\nScheduled", "Draft / Published /\nExcluded", "Draft / Published"),
    ("Lab focus", "✅ Primary", "✅ Secondary", "⚪ Concept only"),
]
for ri, (feature, bm, acr, qa) in enumerate(rows):
    y = Inches(1.95 + ri * 0.65)
    bg = WHITE if ri % 2 == 0 else NEAR_WHITE
    vals = [feature, bm, acr, qa]
    for ci, (x, w, _) in enumerate(cols):
        add_shape_rect(s, x, y, w, Inches(0.6), bg, border_color=LIGHT_GRAY)
        fsize = 12 if ci > 0 else 13
        fbold = ci == 0
        add_text_box(s, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.5),
                     vals[ci], font_size=fsize, bold=fbold, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Section Divider: Search Schema
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: now let's look under the hood — crawled vs managed properties.")
section_divider(s, "Search Schema & Indexing",
                "Crawled properties, managed properties, and reindexing", "⚙️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Crawled vs Managed Properties
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Crawled properties are discovered during crawl — they're raw. "
              "Managed properties are what's kept in the index and can be queried. "
              "Only managed properties are searchable. Schema changes require reindex.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Crawled vs Managed Properties", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "The bridge between content and searchable index",
             font_size=16, color=MID_GRAY)

# LEFT — Crawled
lc = add_rounded_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.2), WHITE)
lc.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.5),
             "📥  Crawled Properties", font_size=22, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
crawled_items = [
    "Discovered automatically during content crawl",
    "Raw metadata from documents, lists, libraries",
    "Examples: ows_Title, ows_Author, ows_Created",
    "Not directly searchable by users",
    "Must be mapped to managed properties",
]
add_bullet_frame(s, Inches(1.2), Inches(2.6), Inches(5.0), Inches(3.0),
                 crawled_items, font_size=14, color=DARK_TEXT, icon="▸")

# Arrow
add_text_box(s, Inches(6.4), Inches(3.5), Inches(0.5), Inches(0.6),
             "→", font_size=36, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

# RIGHT — Managed
rc = add_rounded_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.2), WHITE)
rc.shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.8), Inches(5.6), Inches(0.08), ACCENT_TEAL)
add_text_box(s, Inches(7.1), Inches(2.0), Inches(5.2), Inches(0.5),
             "📤  Managed Properties", font_size=22, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
managed_items = [
    "Kept in the search index (queryable/retrievable)",
    "Users search against managed properties",
    "Settings: queryable, searchable, retrievable, refinable, sortable",
    "Built-in (e.g. Author, Title) or custom (RefinableString00…)",
    "Changes require reindexing the affected content",
]
add_bullet_frame(s, Inches(7.3), Inches(2.6), Inches(5.0), Inches(3.0),
                 managed_items, font_size=14, color=DARK_TEXT, icon="▸")

# Warning
warn = add_rounded_rect(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6), RGBColor(0xFD, 0xE7, 0xE9))
add_text_box(s, Inches(1.1), Inches(6.35), Inches(11.1), Inches(0.5),
             "⚠️  Warning:  Changing managed property mappings can affect other M365 experiences. "
             "In this course, treat schema changes as trainer-led unless explicitly assigned.",
             font_size=14, color=RED_ACCENT)

add_footer_bar(s, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Built-in Managed Properties (Quick Reference)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("SharePoint comes with hundreds of pre-mapped managed properties. "
              "For custom needs, use the RefinableStringXX or RefinableDateXX properties "
              "and rename them via alias.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Built-in Managed Properties (Quick Reference)", font_size=30, bold=True, color=DARK_BG)

props = [
    ("Author", "Document author", "✅", "✅", "✅", "✅"),
    ("Title", "Document title", "✅", "✅", "✅", "✅"),
    ("FileType", "File extension", "✅", "✅", "✅", "✅"),
    ("Created", "Date created", "✅", "—", "✅", "✅"),
    ("Path", "Document URL", "✅", "—", "✅", "—"),
    ("RefinableString00…199", "Custom (alias)", "✅", "—", "✅", "✅"),
]

# Table header
header_y = Inches(1.3)
col_defs = [
    (Inches(0.8), Inches(3.0), "Property"),
    (Inches(3.8), Inches(2.5), "Description"),
    (Inches(6.3), Inches(1.4), "Queryable"),
    (Inches(7.7), Inches(1.4), "Searchable"),
    (Inches(9.1), Inches(1.5), "Retrievable"),
    (Inches(10.6), Inches(1.5), "Refinable"),
]
for x, w, label in col_defs:
    add_shape_rect(s, x, header_y, w, Inches(0.5), ACCENT_BLUE)
    add_text_box(s, x + Inches(0.05), header_y + Inches(0.05), w - Inches(0.1), Inches(0.4),
                 label, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for ri, (prop, desc, q, s_, ret, ref) in enumerate(props):
    y = Inches(1.85 + ri * 0.58)
    bg = WHITE if ri % 2 == 0 else NEAR_WHITE
    vals = [prop, desc, q, s_, ret, ref]
    for ci, (x, w, _) in enumerate(col_defs):
        add_shape_rect(s, x, y, w, Inches(0.53), bg, border_color=LIGHT_GRAY)
        fb = ci == 0
        fc = ACCENT_BLUE if ci == 0 else DARK_TEXT
        add_text_box(s, x + Inches(0.05), y + Inches(0.05), w - Inches(0.1), Inches(0.43),
                     vals[ci], font_size=12, bold=fb, color=fc, alignment=PP_ALIGN.CENTER,
                     font_name="Cascadia Code" if ci == 0 else "Segoe UI")

# Tip
tip = add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1.1), Inches(5.55), Inches(11.1), Inches(0.6),
             "💡  Tip:  To create a custom refiner, map a crawled property to an unused RefinableStringXX, "
             "set an alias, then reindex. Microsoft recommends tenant-level mapping for consistency.",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Reindexing (When and Why)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Reindexing is needed after schema changes or search visibility changes. "
              "Caution: reindexing can create heavy load. Only reindex when necessary.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Reindexing — When and Why", font_size=32, bold=True, color=DARK_BG)

# When to Reindex
when_card = add_rounded_rect(s, Inches(0.8), Inches(1.3), Inches(5.6), Inches(3.2), WHITE)
when_card.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(1.3), Inches(5.6), Inches(0.08), GREEN)
add_text_box(s, Inches(1.0), Inches(1.5), Inches(5.2), Inches(0.5),
             "✅  When to Reindex", font_size=20, bold=True, color=GREEN)
when_items = [
    "After changing managed property mappings",
    "After changing search visibility settings",
    "After adding/modifying site columns used in search",
    "When content 'should be there' but isn't showing up",
]
add_bullet_frame(s, Inches(1.2), Inches(2.1), Inches(5.0), Inches(2.0),
                 when_items, font_size=14, color=DARK_TEXT, icon="▸")

# When NOT to
not_card = add_rounded_rect(s, Inches(6.9), Inches(1.3), Inches(5.6), Inches(3.2), WHITE)
not_card.shadow.inherit = False
add_shape_rect(s, Inches(6.9), Inches(1.3), Inches(5.6), Inches(0.08), RED_ACCENT)
add_text_box(s, Inches(7.1), Inches(1.5), Inches(5.2), Inches(0.5),
             "❌  When NOT to Reindex", font_size=20, bold=True, color=RED_ACCENT)
not_items = [
    "Routinely or 'just in case' (causes load)",
    "After simply uploading new content (auto-crawled)",
    "When the issue is permissions (security trimming)",
    "On large sites without good reason",
]
add_bullet_frame(s, Inches(7.3), Inches(2.1), Inches(5.0), Inches(2.0),
                 not_items, font_size=14, color=DARK_TEXT, icon="▸")

# How to Reindex steps
how_card = add_rounded_rect(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.5), WHITE)
how_card.shadow.inherit = False
add_shape_rect(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.08), ACCENT_BLUE)
add_text_box(s, Inches(1.0), Inches(4.95), Inches(11.3), Inches(0.4),
             "🔄  How to Request Reindex (Library)", font_size=18, bold=True, color=ACCENT_BLUE)

reindex_steps = [
    "1. Go to the library → Settings → Library settings",
    "2. Under General Settings → Advanced settings",
    "3. Scroll to 'Reindex Document Library'",
    "4. Click the button → content is re-crawled at next scheduled crawl",
]
add_bullet_frame(s, Inches(1.2), Inches(5.4), Inches(11.0), Inches(0.8),
                 reindex_steps, font_size=13, color=DARK_TEXT, icon="", spacing=Pt(2))

add_footer_bar(s, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Section Divider: Query Syntax
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: knowing basic query syntax helps admins troubleshoot search issues.")
section_divider(s, "Query Syntax for Admins",
                "KQL basics for troubleshooting search results", "⌨️")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — KQL Query Basics
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("KQL (Keyword Query Language) is the query language behind SharePoint search. "
              "Admins don't need to be experts, but knowing basics helps troubleshoot "
              "'why isn't this showing up?' scenarios.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "KQL Query Basics for Admins", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Keyword Query Language — your troubleshooting Swiss Army knife",
             font_size=16, color=MID_GRAY)

queries = [
    ("Phrase Search", '"Northwind Contract Alpha"', "Exact phrase match\n(quotes required)", ACCENT_BLUE),
    ("Boolean AND", "Alpha AND Harborlight", "Both terms must appear\n(AND must be uppercase)", ACCENT_TEAL),
    ("Boolean OR", "Contract OR Agreement", "Either term matches\n(OR must be uppercase)", ACCENT_PURPLE),
    ("Exclusion", "Alpha -Beta", "Exclude results containing\n'Beta' from results", ORANGE),
    ("Author Filter", 'author:"Jane Smith"', "Restrict by document\nauthor property", GREEN),
    ("File Type", "filetype:docx", "Restrict by file extension\n(docx, pdf, xlsx, etc.)", RED_ACCENT),
]
for i, (title, syntax, desc, color) in enumerate(queries):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.8 + row * 2.3)
    card = add_rounded_rect(s, x, y, Inches(3.8), Inches(2.0), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(3.8), Inches(0.06), color)
    add_text_box(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.35),
                 title, font_size=16, bold=True, color=color)
    # Syntax in code font
    code_box = add_rounded_rect(s, x + Inches(0.2), y + Inches(0.55), Inches(3.4), Inches(0.45),
                                 RGBColor(0xF0, 0xF0, 0xF0))
    add_text_box(s, x + Inches(0.3), y + Inches(0.58), Inches(3.2), Inches(0.4),
                 syntax, font_size=13, color=DARK_TEXT, font_name="Cascadia Code")
    add_text_box(s, x + Inches(0.2), y + Inches(1.1), Inches(3.4), Inches(0.7),
                 desc, font_size=12, color=MID_GRAY)

add_footer_bar(s, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Troubleshooting Flow
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("When a document doesn't appear in search results, follow this troubleshooting flow: "
              "1) Check permissions, 2) Check search visibility, 3) Request reindex if needed, "
              "4) Allow time for indexing.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🛠️  Search Troubleshooting Flow", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "\"I uploaded a document but it doesn't show in search\" — follow these steps",
             font_size=16, color=MID_GRAY)

flow_steps = [
    ("1", "Check Permissions", "Does the user have access\nto the item? (Security\ntrimming is #1 cause)", ACCENT_BLUE),
    ("2", "Check Visibility", "Is the site/library allowed\nto appear in search?\n(Search visibility setting)", ACCENT_TEAL),
    ("3", "Request Reindex", "If settings changed recently,\nrequest a reindex of\nthe library or site", ACCENT_PURPLE),
    ("4", "Wait & Validate", "Allow time for crawling\n(minutes to hours).\nThen search again.", ORANGE),
]
for i, (num, title, desc, color) in enumerate(flow_steps):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.4), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    badge = add_rounded_rect(s, x + Inches(1.05), Inches(2.0), Inches(0.7), Inches(0.7), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(6)
    add_text_box(s, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.5), Inches(2.5), Inches(1.2),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    if i < 3:
        arrow_x = x + Inches(2.95)
        add_text_box(s, arrow_x, Inches(3.2), Inches(0.3), Inches(0.5),
                     "→", font_size=28, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Still not working
still = add_rounded_rect(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.7), RGBColor(0xFD, 0xE7, 0xE9))
add_text_box(s, Inches(1.8), Inches(5.65), Inches(9.7), Inches(0.6),
             "🔴  Still not showing?  Check if numeric-only content (not indexed in Excel). "
             "Verify the managed property mapping. Escalate to Microsoft support if persistent.",
             font_size=14, color=RED_ACCENT)

add_footer_bar(s, 21, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Section Divider: Lab & Validation
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition: time for hands-on lab preview and validation checkpoints.")
section_divider(s, "Lab & Validation",
                "Hands-on search labs and validation checks", "🔬")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Lab Preview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Lab 6 overview: upload seed docs, validate search results, "
              "reindex a library, create a Bookmark (immediate), and create an Acronym (delayed).")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🔬  Lab 6 Preview", font_size=32, bold=True, color=DARK_BG)
add_text_box(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
             "Hands-on: validate search, reindex, and curate answers",
             font_size=16, color=MID_GRAY)

lab_tasks = [
    ("1", "Upload Seed Docs", "Upload seed documents to your\nNW-Pxx library and verify\nthey appear in search", ACCENT_BLUE),
    ("2", "Validate Search", "Use phrase search and KQL\nto locate your documents.\nTest security trimming.", ACCENT_TEAL),
    ("3", "Reindex Library", "Request a reindex of your\nlibrary and observe the\neffect on search results", ACCENT_PURPLE),
    ("4", "Create Bookmark", "Create an NW-Pxx Bookmark\nin Search & intelligence.\nValidate immediate availability", ORANGE),
]
for i, (num, title, desc, color) in enumerate(lab_tasks):
    x = Inches(0.6 + i * 3.15)
    card = add_rounded_rect(s, x, Inches(1.8), Inches(2.9), Inches(3.2), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, Inches(1.8), Inches(2.9), Inches(0.08), color)
    badge = add_rounded_rect(s, x + Inches(1.05), Inches(2.0), Inches(0.7), Inches(0.7), color)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(6)
    add_text_box(s, x + Inches(0.2), Inches(2.9), Inches(2.5), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.5), Inches(2.5), Inches(1.2),
                 desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# Bonus
bonus = add_rounded_rect(s, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.8), LIGHT_BLUE)
add_text_box(s, Inches(1.8), Inches(5.45), Inches(9.7), Inches(0.7),
             "📝  Bonus:  Create an Acronym (e.g. NW-Pxx-NDA = 'Non-Disclosure Agreement'). "
             "Note: published acronyms may take up to 24 hours to appear — validate next day.",
             font_size=14, color=DARK_TEXT)

add_footer_bar(s, 23, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Validation Checklist
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Pre- and post-lab validation checklist. These are the key checkpoints "
              "that confirm the lab objectives have been met.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, GREEN)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "✅  Validation Checklist", font_size=32, bold=True, color=DARK_BG)

checks = [
    ("Seed docs discoverable", "Search for your uploaded documents —\nthey should appear in results for your account", GREEN),
    ("Security trimming works", "Ask a neighbor to search for your docs —\nif they lack access, docs shouldn't appear", ACCENT_BLUE),
    ("Reindex completed", "After requesting reindex, wait for crawl,\nthen verify updated content appears", ACCENT_TEAL),
    ("Bookmark works", "Search your NW-Pxx keyword — the bookmark\nshould appear at the top of results immediately", ACCENT_PURPLE),
    ("Acronym created", "Published acronym visible (may take up to 24h).\nVerify Draft vs Published state in admin center", ORANGE),
    ("KQL queries work", 'Test: "exact phrase", author:"Name", filetype:docx\nVerify results match expectations', RED_ACCENT),
]
for i, (title, desc, color) in enumerate(checks):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.3 + row * 1.65)
    card = add_rounded_rect(s, x, y, Inches(5.8), Inches(1.4), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, x, y, Inches(0.08), Inches(1.4), color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.1), Inches(5.2), Inches(0.4),
                 f"☐  {title}", font_size=16, bold=True, color=color)
    add_text_box(s, x + Inches(0.3), y + Inches(0.55), Inches(5.2), Inches(0.7),
                 desc, font_size=13, color=DARK_TEXT)

add_footer_bar(s, 24, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Common Troubleshooting Issues
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Top 5 lab issues and their resolutions.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, RED_ACCENT)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "🛠️  Common Troubleshooting Issues", font_size=32, bold=True, color=DARK_BG)

issues = [
    ("Document not in search results", "Check: user permissions, search visibility settings, "
     "indexing delay. Try reindex if recently changed settings.",
     "Permissions → Visibility → Reindex → Wait", RED_ACCENT),
    ("Bookmark not appearing", "Verify the bookmark is Published (not Draft). "
     "Check keyword spelling. Bookmarks are immediate after publish.",
     "Admin center → Search & intelligence → Bookmarks", ORANGE),
    ("Acronym not showing after publish", "Published acronyms can take up to 24 hours. "
     "Verify it's in Published state, not Draft.",
     "Wait 24h. Check state in admin center.", ACCENT_BLUE),
    ("KQL query returns no results", "Check operator case (AND/OR/NOT must be uppercase). "
     "Verify the property is queryable (not all properties are).",
     "Test with simple phrase first, then add filters", ACCENT_PURPLE),
    ("Reindex seems to have no effect", "Reindex marks content for next crawl — this isn't instant. "
     "Wait 15-60 minutes. Avoid re-triggering.",
     "Be patient. Check again after 1 hour.", ACCENT_TEAL),
]
for i, (problem, solution, check, color) in enumerate(issues):
    y = Inches(1.3 + i * 1.15)
    card = add_rounded_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), WHITE)
    card.shadow.inherit = False
    add_shape_rect(s, Inches(0.6), y, Inches(0.08), Inches(1.0), color)
    add_text_box(s, Inches(0.9), y + Inches(0.05), Inches(3.2), Inches(0.4),
                 problem, font_size=14, bold=True, color=color)
    add_text_box(s, Inches(4.2), y + Inches(0.05), Inches(4.8), Inches(0.9),
                 solution, font_size=12, color=DARK_TEXT)
    add_text_box(s, Inches(9.2), y + Inches(0.05), Inches(3.3), Inches(0.9),
                 check, font_size=11, color=MID_GRAY)

add_footer_bar(s, 25, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Key Takeaways
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Recap the module: search pipeline, security trimming, answers, schema, "
              "reindexing, and KQL basics.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, GREEN)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "📌  Key Takeaways", font_size=32, bold=True, color=DARK_BG)

takeaways = [
    ("🔐", "Security Trimming", "Search never overrides permissions. 'Missing' results?\nCheck access first — this is the #1 troubleshooting step."),
    ("📌", "Curated Answers", "Bookmarks (immediate) and Acronyms (up to 24h) help users\nfind the right content faster. Use NW-Pxx keywords in shared tenants."),
    ("⚙️", "Search Schema", "Crawled properties → managed properties → index.\nSchema changes require reindex. Treat as trainer-led."),
    ("🔄", "Reindexing", "Reindex after schema/visibility changes, not routinely.\nCauses load — be deliberate and patient."),
    ("⌨️", "KQL Queries", 'Phrase search with quotes, AND/OR/NOT (uppercase),\nproperty restrictions (author:, filetype:) for troubleshooting.'),
]
for i, (icon, title, desc) in enumerate(takeaways):
    y = Inches(1.3 + i * 1.15)
    card = add_rounded_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.0), WHITE)
    card.shadow.inherit = False
    add_text_box(s, Inches(0.8), y + Inches(0.1), Inches(0.6), Inches(0.5),
                 icon, font_size=26, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), y + Inches(0.1), Inches(2.5), Inches(0.4),
                 title, font_size=17, bold=True, color=ACCENT_BLUE)
    add_text_box(s, Inches(4.2), y + Inches(0.1), Inches(8.3), Inches(0.8),
                 desc, font_size=14, color=DARK_TEXT)

add_footer_bar(s, 26, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Knowledge Check
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("5 discussion questions to check understanding.")
add_solid_bg(s, NEAR_WHITE)
add_top_bar(s, ACCENT_PURPLE)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "❓  Knowledge Check", font_size=32, bold=True, color=DARK_BG)

questions = [
    "What is security trimming and why is it the #1 rule for search?",
    "What is the difference between Bookmarks and Acronyms in Microsoft Search?",
    "When should you request a reindex of a library — and when should you NOT?",
    "What are crawled properties vs managed properties, and how do they relate?",
    "Write a KQL query to find all DOCX files authored by 'Jane Smith' about 'Northwind'.",
]
for i, q in enumerate(questions):
    y = Inches(1.4 + i * 1.1)
    qcard = add_rounded_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.9), WHITE)
    qcard.shadow.inherit = False
    badge = add_rounded_rect(s, Inches(0.8), y + Inches(0.15), Inches(0.55), Inches(0.55), ACCENT_PURPLE)
    tf = badge.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Segoe UI"
    tf.paragraphs[0].space_before = Pt(4)
    add_text_box(s, Inches(1.6), y + Inches(0.15), Inches(10.8), Inches(0.6),
                 q, font_size=16, color=DARK_TEXT)

add_footer_bar(s, 27, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Thank You / Next Module
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Module 6 complete. Next up: Module 7 — Apps and Customization.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "✅  Module 6 Complete!", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
             "Search in SharePoint Online & Microsoft Search",
             font_size=24, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

add_shape_rect(s, Inches(5.5), Inches(4.4), Inches(2.3), Inches(0.04), ACCENT_PURPLE)

add_text_box(s, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
             "Up Next  →  Module 7: Apps and Customization",
             font_size=20, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(2), Inches(5.8), Inches(9), Inches(0.6),
             "🧩 Governance, deployment, and custom solutions in SharePoint",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(__file__), "Module-06-Slides.pptx")
prs.save(out_path)
print(f"✅  Saved {TOTAL_SLIDES}-slide presentation → {out_path}")
