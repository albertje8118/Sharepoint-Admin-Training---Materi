"""
Generate Module 2 PPTX — Identity, Access, and External Sharing
Modern, engaging design matching Module 1 style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xF5, 0xF5, 0xF5)
DARK_BG     = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)
ACCENT_PURPLE = RGBColor(0x6B, 0x69, 0xD6)
LIGHT_BLUE  = RGBColor(0xDE, 0xEC, 0xF9)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY    = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT    = RGBColor(0x24, 0x24, 0x24)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)
GREEN       = RGBColor(0x10, 0x7C, 0x10)
RED_ACCENT  = RGBColor(0xD1, 0x34, 0x38)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ═══════════════════════════════════════════════════════════════════════════
# Helpers (same as Module 1)
# ═══════════════════════════════════════════════════════════════════════════
def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color; shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6), icon="▸", font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{icon} {bullet}"; p.font.size = Pt(font_size)
        p.font.color.rgb = color; p.font.name = font_name; p.space_after = spacing
    return txBox

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, height, color)

def add_footer_bar(slide, slide_num, total, module_label="Module 2"):
    bar_top = SLIDE_HEIGHT - Inches(0.45)
    add_shape_rect(slide, Inches(0), bar_top, SLIDE_WIDTH, Inches(0.45), DARK_BG)
    add_text_box(slide, Inches(0.5), bar_top + Inches(0.05), Inches(7), Inches(0.35),
                 f"{module_label}  |  Identity, Access & External Sharing",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text_box(slide, SLIDE_WIDTH - Inches(1.5), bar_top + Inches(0.05),
                 Inches(1.2), Inches(0.35), f"{slide_num} / {total}",
                 font_size=10, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)

def section_divider(slide, section_title, subtitle="", icon=""):
    add_solid_bg(slide, DARK_BG)
    add_shape_rect(slide, Inches(0), Inches(3.2), Inches(2.5), Inches(0.06), ACCENT_TEAL)
    if icon:
        add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                     icon, font_size=60, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                 section_title, font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
                     subtitle, font_size=18, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

TOTAL_SLIDES = 26
slide_counter = [0]

def new_slide(notes=""):
    slide_counter[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if notes: add_speaker_notes(slide, notes)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Cover
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Welcome to Module 2. This module bridges identity fundamentals with SharePoint sharing controls.")
add_solid_bg(s, DARK_BG)
add_shape_rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12), ACCENT_BLUE)
add_shape_rect(s, Inches(0), Inches(0.12), SLIDE_WIDTH, Inches(0.04), ACCENT_TEAL)

add_text_box(s, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "MODERN SHAREPOINT ONLINE FOR ADMINISTRATORS", font_size=18, bold=True,
             color=ACCENT_TEAL, font_name="Segoe UI Semibold")
add_text_box(s, Inches(1), Inches(2.1), Inches(11), Inches(1.2),
             "Identity, Access &\nExternal Sharing", font_size=48, bold=True,
             color=WHITE, font_name="Segoe UI Light")
add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Day 1  ·  Tenant Foundations & Site Management", font_size=18,
             color=RGBColor(0xAA, 0xAA, 0xAA))
badge = add_rounded_rect(s, Inches(1), Inches(4.7), Inches(3.2), Inches(0.55), ACCENT_BLUE)
add_text_box(s, Inches(1.2), Inches(4.75), Inches(2.8), Inches(0.45),
             "MODULE 2", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "🔐  Securing collaboration through identity & policy layers",
             font_size=16, color=RGBColor(0x88, 0x88, 0x88))

add_footer_bar(s, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Module Agenda
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the topics. Emphasize that many 'SharePoint issues' are identity issues.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
             "📋  Module 2 Agenda", font_size=30, bold=True, color=DARK_BG)

agenda = [
    ("1", "Why Identity Matters — Background for SPO admins", ACCENT_BLUE),
    ("2", "Microsoft Entra ID Fundamentals — The identity backbone", ACCENT_TEAL),
    ("3", "Zero Trust & Least Privilege — Security principles", ACCENT_PURPLE),
    ("4", "Admin Roles vs SharePoint Permissions — Know the difference", ACCENT_BLUE),
    ("5", "External Sharing Model — Org-level & site-level controls", ACCENT_TEAL),
    ("6", "Guest Access (B2B) — Lifecycle & troubleshooting", ACCENT_PURPLE),
    ("7", "Conditional Access — Overview for SPO admins", ORANGE),
    ("8", "Lab 2 Preview & Knowledge Check", GREEN),
]

y_pos = Inches(1.3)
for num, text, color in agenda:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y_pos, Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    tf = circ.text_frame; p = tf.paragraphs[0]
    p.text = num; p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Segoe UI"
    add_text_box(s, Inches(1.5), y_pos + Inches(0.05), Inches(10), Inches(0.45),
                 text, font_size=17, color=DARK_TEXT)
    y_pos += Inches(0.65)

add_footer_bar(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Module Objectives
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Review learning objectives. These are what participants will be assessed on.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "🎯  Module Objectives", font_size=30, bold=True, color=DARK_BG)

objectives = [
    "Explain Microsoft Entra ID fundamentals relevant to SharePoint Online",
    "Distinguish admin roles (tenant) from SharePoint permissions (site)",
    "Describe guest access (B2B collaboration) and how it interacts with SharePoint sharing",
    "Identify and apply the correct external sharing control at the correct scope (org vs site)",
    "Explain Conditional Access at a high level and how it impacts SharePoint access",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(5),
                 objectives, font_size=18, color=DARK_TEXT, icon="✅")

add_footer_bar(s, 3, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Section Divider: Why Identity Matters
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to the background/overview section on identity.")
section_divider(s, "Why Identity Matters",
                "The foundation that every SharePoint admin must understand", "🔑")
add_footer_bar(s, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Identity Challenge for SharePoint Admins
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Set the business context. Modern collaboration means internal + external users, "
              "multiple devices, and compliance obligations. Identity is the perimeter.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "The Modern Identity Challenge", font_size=28, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.8),
             "In today's work environment, collaboration crosses organisational boundaries. "
             "Identity management is now the primary security perimeter — not the network firewall.",
             font_size=16, color=MID_GRAY)

challenges = [
    ("👤  Internal Users", "Full-time employees, contractors\nwith corporate accounts", ACCENT_BLUE),
    ("🤝  External Partners", "Vendors, clients, consultants\nneeding project access", ACCENT_TEAL),
    ("📱  Multiple Devices", "Corporate laptops, personal phones,\nshared kiosks", ACCENT_PURPLE),
    ("⚖️  Compliance", "Data protection, retention,\nregulatory requirements", ORANGE),
]

x_pos = Inches(0.4)
for title, desc, color in challenges:
    card = add_rounded_rect(s, x_pos, Inches(2.3), Inches(3), Inches(2.5), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(2.3), Inches(3), Inches(0.08), color)
    add_text_box(s, x_pos + Inches(0.2), Inches(2.6), Inches(2.6), Inches(0.5),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s, x_pos + Inches(0.2), Inches(3.2), Inches(2.6), Inches(1.2),
                 desc, font_size=14, color=MID_GRAY)
    x_pos += Inches(3.25)

# Key insight
insight = add_rounded_rect(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.25), Inches(11), Inches(0.7),
             "🔑 Key Insight: Many 'SharePoint access problems' are actually identity & policy "
             "problems originating in Microsoft Entra ID, not in SharePoint itself.",
             font_size=15, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The Access Decision Flow
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Show the logical flow: user authenticates via Entra ID → policies are evaluated "
              "(Conditional Access, MFA, device) → SharePoint checks tenant sharing → site sharing → permissions.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "How Access Decisions Are Made", font_size=28, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "When a user tries to access a SharePoint resource, multiple layers are evaluated:",
             font_size=16, color=MID_GRAY)

flow_steps = [
    ("1", "Authentication", "User signs in via\nMicrosoft Entra ID", ACCENT_BLUE),
    ("→", "", "", MID_GRAY),
    ("2", "Policy Evaluation", "Conditional Access,\nMFA, device compliance", ACCENT_PURPLE),
    ("→", "", "", MID_GRAY),
    ("3", "Tenant Sharing", "Org-level SharePoint\nsharing settings", ACCENT_TEAL),
    ("→", "", "", MID_GRAY),
    ("4", "Site Sharing", "Site-level sharing\nconfiguration", ORANGE),
    ("→", "", "", MID_GRAY),
    ("5", "Permissions", "SharePoint groups &\npermission levels", GREEN),
]

x_pos = Inches(0.3)
for num, title, desc, color in flow_steps:
    if num == "→":
        add_text_box(s, x_pos, Inches(2.8), Inches(0.4), Inches(0.5),
                     "→", font_size=28, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(0.5)
    else:
        card = add_rounded_rect(s, x_pos, Inches(2.0), Inches(2.1), Inches(2.8), NEAR_WHITE)
        add_shape_rect(s, x_pos, Inches(2.0), Inches(2.1), Inches(0.06), color)
        # number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(0.7), Inches(2.2),
                                   Inches(0.6), Inches(0.6))
        badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
        tf = badge.text_frame; p = tf.paragraphs[0]
        p.text = num; p.font.size = Pt(18); p.font.bold = True
        p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Segoe UI"
        add_text_box(s, x_pos + Inches(0.1), Inches(3.0), Inches(1.9), Inches(0.4),
                     title, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(s, x_pos + Inches(0.1), Inches(3.5), Inches(1.9), Inches(1),
                     desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(2.2)

# Blocked callout
add_text_box(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
             "🚫 If ANY layer blocks access, the user is denied — even if SharePoint permissions look correct.",
             font_size=15, bold=True, color=RED_ACCENT, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 6, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Zero Trust Principles
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Introduce Zero Trust as the guiding security philosophy. "
              "Microsoft 365 (and thus SharePoint) operates on these principles.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Zero Trust: The Security Foundation", font_size=28, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Microsoft 365 security is built on Zero Trust principles — never trust, always verify.",
             font_size=16, color=MID_GRAY)

zt_principles = [
    ("🔍", "Verify Explicitly", "Always authenticate and authorize\nbased on all available data:\nidentity, location, device, service,\ndata classification, anomalies.", ACCENT_BLUE),
    ("🔒", "Least Privilege Access", "Limit user access with just-in-time\nand just-enough-access (JIT/JEA),\nrisk-based adaptive policies,\nand data protection.", ACCENT_TEAL),
    ("💥", "Assume Breach", "Minimize blast radius and segment\naccess. Verify end-to-end encryption.\nUse analytics to detect threats,\nimprove defenses.", RED_ACCENT),
]

x_pos = Inches(0.5)
for icon, title, desc, color in zt_principles:
    card = add_rounded_rect(s, x_pos, Inches(2.0), Inches(3.9), Inches(3.8), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(2.0), Inches(3.9), Inches(0.1), color)
    add_text_box(s, x_pos + Inches(0.3), Inches(2.3), Inches(3.3), Inches(0.6),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER, color=DARK_TEXT)
    add_text_box(s, x_pos + Inches(0.3), Inches(2.9), Inches(3.3), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x_pos + Inches(0.3), Inches(3.5), Inches(3.3), Inches(2),
                 desc, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(4.15)

add_text_box(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.4),
             "💡 As a SharePoint admin, you operate inside this model — your policies should align with these principles.",
             font_size=14, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Section Divider: Microsoft Entra ID
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to Entra ID fundamentals.")
section_divider(s, "Microsoft Entra ID Fundamentals",
                "The identity backbone of Microsoft 365", "🪪")
add_footer_bar(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — What is Microsoft Entra ID?
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Entra ID (formerly Azure AD) is the cloud identity service. "
              "It handles authentication, directory, and policy enforcement for all M365 services.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Microsoft Entra ID — What SharePoint Admins Need to Know",
             font_size=26, bold=True, color=DARK_BG)

entra_features = [
    ("🔐 Authentication", "Users sign in via Entra ID — MFA, passwordless, SSO",
     "Every SharePoint access starts with an Entra ID sign-in", ACCENT_BLUE),
    ("📋 Directory", "Users, groups, guest accounts, app registrations",
     "This is where members and guests are managed", ACCENT_TEAL),
    ("🛡️ Policies", "Conditional Access, external collaboration settings",
     "Controls WHO can access WHAT, from WHERE, on WHICH device", ACCENT_PURPLE),
    ("👥 Groups", "Microsoft 365 Groups, Security Groups, Distribution Lists",
     "M365 Groups power Teams-connected SharePoint sites", ORANGE),
]

y_pos = Inches(1.3)
for title, desc, relevance, color in entra_features:
    card = add_rounded_rect(s, Inches(0.5), y_pos, Inches(12.3), Inches(1.05), NEAR_WHITE)
    add_shape_rect(s, Inches(0.5), y_pos, Inches(0.12), Inches(1.05), color)
    add_text_box(s, Inches(0.9), y_pos + Inches(0.05), Inches(4.5), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(0.9), y_pos + Inches(0.4), Inches(5), Inches(0.55),
                 desc, font_size=13, color=DARK_TEXT)
    add_text_box(s, Inches(6.8), y_pos + Inches(0.2), Inches(5.5), Inches(0.6),
                 f"SPO relevance: {relevance}", font_size=12, color=MID_GRAY)
    y_pos += Inches(1.2)

add_footer_bar(s, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Member vs Guest
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Members = internal users with full directory privileges. "
              "Guests = external B2B users with limited directory access. "
              "Guest behaviour is shaped by both Entra AND SharePoint settings.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Member vs Guest Users", font_size=28, bold=True, color=DARK_BG)

# Comparison columns
for col_data in [
    (Inches(0.5), "👤 Member (Internal)", [
        "Full user in your directory",
        "Corporate email & license",
        "Full directory browsing by default",
        "Full access to licensed services",
        "Managed by HR / IT provisioning",
    ], ACCENT_BLUE),
    (Inches(6.8), "🤝 Guest (External / B2B)", [
        "External user invited to your directory",
        "Uses their own email / identity",
        "Limited directory visibility",
        "Access scoped to invited resources",
        "Managed via Entra external settings",
    ], ACCENT_TEAL),
]:
    x, title, items, color = col_data
    card = add_rounded_rect(s, x, Inches(1.3), Inches(5.8), Inches(4.2), NEAR_WHITE)
    add_shape_rect(s, x, Inches(1.3), Inches(5.8), Inches(0.08), color)
    add_text_box(s, x + Inches(0.3), Inches(1.55), Inches(5.2), Inches(0.5),
                 title, font_size=18, bold=True, color=color)
    add_bullet_frame(s, x + Inches(0.3), Inches(2.2), Inches(5.2), Inches(3),
                     items, font_size=15, color=DARK_TEXT, icon="•")

callout = add_rounded_rect(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.85), Inches(11), Inches(0.6),
             "⚠️ Guest access depends on: Entra collaboration settings + SharePoint org sharing + Site sharing + Permissions",
             font_size=14, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Section Divider: Roles vs Permissions
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to the roles vs permissions section.")
section_divider(s, "Admin Roles vs Site Permissions",
                "Two different scopes — don't confuse them", "⚙️")
add_footer_bar(s, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Admin Roles (Tenant Level)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Admin roles grant access to admin centers and tenant-wide config. "
              "Being a SP Admin does NOT make you a site owner, and vice versa.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Tenant Admin Roles (Directory Level)", font_size=28, bold=True, color=DARK_BG)

roles = [
    ("Global Administrator", "Full access to all admin centers and settings — use sparingly", RED_ACCENT),
    ("SharePoint Administrator", "Manages SharePoint admin center, all sites, tenant policies", ACCENT_BLUE),
    ("Teams Administrator", "Manages Teams settings; affects Teams-connected SP sites", ACCENT_PURPLE),
    ("Exchange Administrator", "Manages Exchange; relevant for mail-enabled groups", ACCENT_TEAL),
    ("Compliance Administrator", "Manages Purview policies affecting SP content", ORANGE),
]

y_pos = Inches(1.3)
for role, desc, color in roles:
    add_shape_rect(s, Inches(0.8), y_pos, Inches(0.08), Inches(0.75), color)
    add_text_box(s, Inches(1.1), y_pos, Inches(5), Inches(0.35),
                 role, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(1.1), y_pos + Inches(0.35), Inches(11), Inches(0.35),
                 desc, font_size=13, color=MID_GRAY)
    y_pos += Inches(0.85)

warn = add_rounded_rect(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.7),
                        RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, Inches(1), Inches(5.75), Inches(11), Inches(0.6),
             "🔑 Principle: Use least privilege. Avoid Global Admin for day-to-day SharePoint work.",
             font_size=14, bold=True, color=ORANGE)

add_footer_bar(s, 12, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SharePoint Permissions (Site Level)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("SharePoint permissions control what users can do INSIDE a specific site. "
              "Default groups: Owners, Members, Visitors.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "SharePoint Permissions (Site-Scoped)", font_size=28, bold=True, color=DARK_BG)

perm_groups = [
    ("🔑 Site Owners", "Full control within the site\nManage permissions, settings, pages\nNOT necessarily a tenant admin",
     ACCENT_BLUE),
    ("✏️ Site Members", "Contribute content (add, edit, delete)\nCollaborate on documents\nCannot change site settings",
     ACCENT_TEAL),
    ("👁️ Site Visitors", "View/read content only\nCannot edit or upload\nIdeal for stakeholders/readers",
     ACCENT_PURPLE),
]

x_pos = Inches(0.5)
for title, desc, color in perm_groups:
    card = add_rounded_rect(s, x_pos, Inches(1.3), Inches(3.9), Inches(3.5), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(1.3), Inches(3.9), Inches(0.1), color)
    add_text_box(s, x_pos + Inches(0.3), Inches(1.6), Inches(3.3), Inches(0.5),
                 title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x_pos + Inches(0.3), Inches(2.3), Inches(3.3), Inches(2.2),
                 desc, font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(4.15)

# Key distinction
add_text_box(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.8),
             "⚡ Key Distinction:\n"
             "   Tenant admin role ≠ Site owner  |  Site owner ≠ Tenant admin  |  Assign only what's needed",
             font_size=16, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Section Divider: External Sharing
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to external sharing model.")
section_divider(s, "External Sharing Model",
                "The multi-layer control stack for secure collaboration", "🌍")
add_footer_bar(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — The External Sharing Control Stack
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Three layers: Entra collaboration settings → Org-level SharePoint sharing → Site-level sharing. "
              "Site can NEVER be more permissive than org level.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "The External Sharing Control Stack", font_size=28, bold=True, color=DARK_BG)

layers = [
    ("Layer 1", "Entra External Collaboration", "Who can invite guests?\nGuest invitation restrictions\nCross-tenant access policies", ACCENT_BLUE, "Broadest scope"),
    ("Layer 2", "Org-Level Sharing (SPO Admin Center)", "Tenant-wide baseline for SharePoint + OneDrive\nOneDrive ≤ SharePoint (never more permissive)\nAnyone / New+Existing / Existing / Only org", ACCENT_TEAL, "Tenant baseline"),
    ("Layer 3", "Site-Level Sharing", "Per-site override (same or MORE restrictive)\nScoped to individual site collections\nIdeal for extranet / project sites", ACCENT_PURPLE, "Most specific"),
]

y_pos = Inches(1.3)
for layer_id, title, desc, color, scope in layers:
    card = add_rounded_rect(s, Inches(0.5), y_pos, Inches(12.3), Inches(1.35), NEAR_WHITE)
    add_shape_rect(s, Inches(0.5), y_pos, Inches(0.12), Inches(1.35), color)
    # Layer badge
    lbadge = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.15), Inches(1.2), Inches(0.4), color)
    add_text_box(s, Inches(0.85), y_pos + Inches(0.17), Inches(1.1), Inches(0.35),
                 layer_id, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.2), y_pos + Inches(0.1), Inches(4.5), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(2.2), y_pos + Inches(0.5), Inches(5), Inches(0.8),
                 desc, font_size=12, color=MID_GRAY)
    add_text_box(s, Inches(10), y_pos + Inches(0.3), Inches(2.5), Inches(0.5),
                 f"📌 {scope}", font_size=12, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    y_pos += Inches(1.55)

# Rule
rule = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.1), Inches(11.5), Inches(0.7),
                        RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(s, Inches(1), y_pos + Inches(0.15), Inches(11), Inches(0.6),
             "⚠️ Golden Rule: A site-level sharing setting can NEVER be more permissive than the org-level setting.",
             font_size=14, bold=True, color=ORANGE)

add_footer_bar(s, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Sharing Levels Explained
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the 4 sharing levels from most to least permissive. "
              "Advise against 'Anyone' unless deliberately accepted.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Sharing Levels — Most to Least Permissive", font_size=28, bold=True, color=DARK_BG)

sharing_levels = [
    ("🌐 Anyone", "Anonymous links — no sign-in required\nRiskiest option; use with caution",
     RED_ACCENT, "MOST PERMISSIVE"),
    ("🤝 New & Existing Guests", "Guests must authenticate (sign-in)\nNew guests can be invited",
     ORANGE, "RECOMMENDED"),
    ("👤 Existing Guests Only", "Only guests already in the directory\nNo new invitations via sharing",
     ACCENT_TEAL, "RESTRICTIVE"),
    ("🏢 Only People in Your Org", "No external sharing at all\nInternal collaboration only",
     ACCENT_BLUE, "MOST RESTRICTIVE"),
]

y_pos = Inches(1.3)
for title, desc, color, label in sharing_levels:
    card = add_rounded_rect(s, Inches(0.5), y_pos, Inches(9.5), Inches(1.05), NEAR_WHITE)
    add_shape_rect(s, Inches(0.5), y_pos, Inches(0.12), Inches(1.05), color)
    add_text_box(s, Inches(0.9), y_pos + Inches(0.05), Inches(4), Inches(0.4),
                 title, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(0.9), y_pos + Inches(0.45), Inches(8.5), Inches(0.55),
                 desc, font_size=13, color=MID_GRAY)
    # label badge
    lbl = add_rounded_rect(s, Inches(10.3), y_pos + Inches(0.2), Inches(2.3), Inches(0.5), color)
    add_text_box(s, Inches(10.35), y_pos + Inches(0.22), Inches(2.2), Inches(0.45),
                 label, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    y_pos += Inches(1.25)

add_footer_bar(s, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — OneDrive Sharing Relationship
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("OneDrive sharing can never be more permissive than SharePoint sharing. "
              "Both are controlled from the SharePoint admin center.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "SharePoint ↔ OneDrive Sharing Relationship", font_size=28, bold=True, color=DARK_BG)

# Visual: SP >= OD
sp_card = add_rounded_rect(s, Inches(1), Inches(1.5), Inches(5), Inches(3.5), NEAR_WHITE)
add_shape_rect(s, Inches(1), Inches(1.5), Inches(5), Inches(0.1), ACCENT_BLUE)
add_text_box(s, Inches(1.3), Inches(1.8), Inches(4.4), Inches(0.5),
             "📄 SharePoint Online Sharing", font_size=20, bold=True, color=ACCENT_BLUE)
add_text_box(s, Inches(1.3), Inches(2.5), Inches(4.4), Inches(2),
             "• Sets the MAXIMUM permissiveness\n• Controls the sharing ceiling\n• All 4 levels available\n• Configured in: SP admin center → Policies → Sharing",
             font_size=14, color=DARK_TEXT)

od_card = add_rounded_rect(s, Inches(7), Inches(1.5), Inches(5), Inches(3.5), NEAR_WHITE)
add_shape_rect(s, Inches(7), Inches(1.5), Inches(5), Inches(0.1), ACCENT_TEAL)
add_text_box(s, Inches(7.3), Inches(1.8), Inches(4.4), Inches(0.5),
             "☁️ OneDrive Sharing", font_size=20, bold=True, color=ACCENT_TEAL)
add_text_box(s, Inches(7.3), Inches(2.5), Inches(4.4), Inches(2),
             "• CANNOT exceed SharePoint level\n• Same or more restrictive\n• Personal file sharing scope\n• Configured in: SP admin center → Policies → Sharing",
             font_size=14, color=DARK_TEXT)

# Arrow
add_text_box(s, Inches(5.5), Inches(2.8), Inches(2), Inches(0.6),
             "≥", font_size=48, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

rule2 = add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(5.55), Inches(11), Inches(0.6),
             "📌 Rule: SharePoint sharing level ≥ OneDrive sharing level ≥ Site-level sharing",
             font_size=15, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 17, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Section Divider: Guest Access (B2B)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to guest access lifecycle.")
section_divider(s, "Guest Access (B2B Collaboration)",
                "Lifecycle, common pitfalls, and troubleshooting", "🤝")
add_footer_bar(s, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Guest Lifecycle
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Guest lifecycle: invite → redeem → access. Invitation state matters! "
              "If a guest hasn't redeemed, they can't access resources.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Guest User Lifecycle", font_size=28, bold=True, color=DARK_BG)

lifecycle = [
    ("1", "Invite", "Admin or user invites\nan external user via\nemail or sharing link", ACCENT_BLUE),
    ("→", "", "", MID_GRAY),
    ("2", "Redeem", "Guest clicks the link\nand authenticates\nwith their identity", ACCENT_TEAL),
    ("→", "", "", MID_GRAY),
    ("3", "Access", "Guest accesses the\nshared resource based\non permissions granted", ACCENT_PURPLE),
    ("→", "", "", MID_GRAY),
    ("4", "Review / Remove", "Admin reviews guest\naccess periodically\nand removes if needed", ORANGE),
]

x_pos = Inches(0.3)
for num, title, desc, color in lifecycle:
    if num == "→":
        add_text_box(s, x_pos, Inches(2.8), Inches(0.5), Inches(0.5),
                     "→", font_size=32, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(0.6)
    else:
        card = add_rounded_rect(s, x_pos, Inches(1.8), Inches(2.6), Inches(3), NEAR_WHITE)
        add_shape_rect(s, x_pos, Inches(1.8), Inches(2.6), Inches(0.08), color)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(0.9), Inches(2.0),
                                   Inches(0.7), Inches(0.7))
        badge.fill.solid(); badge.fill.fore_color.rgb = color; badge.line.fill.background()
        tf = badge.text_frame; p = tf.paragraphs[0]
        p.text = num; p.font.size = Pt(20); p.font.bold = True
        p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Segoe UI"
        add_text_box(s, x_pos + Inches(0.1), Inches(2.9), Inches(2.4), Inches(0.4),
                     title, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(s, x_pos + Inches(0.1), Inches(3.4), Inches(2.4), Inches(1.2),
                     desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(2.75)

add_footer_bar(s, 19, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Common Guest Access Failures
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Walk through the top 5 causes of 'guest can't access' issues. "
              "This is the most common support scenario for SharePoint admins.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🔧 Top 5 Guest Access Failure Causes", font_size=28, bold=True, color=DARK_BG)

failures = [
    ("1", "Org-level sharing is too restrictive", "Tenant blocks guests entirely or limits to existing only", RED_ACCENT),
    ("2", "Site-level sharing is more restrictive", "Site overrides org baseline with tighter settings", ORANGE),
    ("3", "Guest hasn't redeemed the invitation", "Invitation pending — guest never clicked the link", ACCENT_BLUE),
    ("4", "M365 Group / Teams guest settings conflict", "Group-connected site respects Teams guest policies", ACCENT_TEAL),
    ("5", "Conditional Access blocks the guest", "Device compliance or location policy prevents access", ACCENT_PURPLE),
]

y_pos = Inches(1.3)
for num, cause, detail, color in failures:
    card = add_rounded_rect(s, Inches(0.5), y_pos, Inches(12.3), Inches(0.9), NEAR_WHITE)
    nbadge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y_pos + Inches(0.15),
                                 Inches(0.55), Inches(0.55))
    nbadge.fill.solid(); nbadge.fill.fore_color.rgb = color; nbadge.line.fill.background()
    tf = nbadge.text_frame; p = tf.paragraphs[0]
    p.text = num; p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; p.font.name = "Segoe UI"
    add_text_box(s, Inches(1.5), y_pos + Inches(0.05), Inches(5.5), Inches(0.4),
                 cause, font_size=16, bold=True, color=color)
    add_text_box(s, Inches(1.5), y_pos + Inches(0.45), Inches(10.5), Inches(0.4),
                 detail, font_size=13, color=MID_GRAY)
    y_pos += Inches(1.05)

add_text_box(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             "💡 Troubleshooting tip: Check layers from top (Entra) to bottom (site permissions).",
             font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 20, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Section Divider: Conditional Access
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Transition to Conditional Access overview.")
section_divider(s, "Conditional Access",
                "Identity-driven policy enforcement for SharePoint", "🛡️")
add_footer_bar(s, 21, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Conditional Access Overview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("CA is NOT a SharePoint feature — it's an Entra ID feature that affects SharePoint. "
              "Explain common signals: user risk, device compliance, location, auth strength.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "Conditional Access — How It Affects SharePoint", font_size=26, bold=True, color=DARK_BG)

add_text_box(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Conditional Access evaluates signals and enforces access decisions in real-time:",
             font_size=16, color=MID_GRAY)

ca_signals = [
    ("👤 User / Risk", "Who is signing in?\nUser risk level\nGroup membership", ACCENT_BLUE),
    ("📱 Device", "Is the device compliant?\nManaged vs unmanaged\nOS platform", ACCENT_TEAL),
    ("📍 Location", "Where is the request from?\nTrusted vs unknown network\nCountry/region", ACCENT_PURPLE),
    ("🔐 Auth Strength", "MFA required?\nPasswordless?\nPhishing-resistant?", ORANGE),
]

x_pos = Inches(0.4)
for title, desc, color in ca_signals:
    card = add_rounded_rect(s, x_pos, Inches(2.0), Inches(3), Inches(2.5), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(2.0), Inches(3), Inches(0.08), color)
    add_text_box(s, x_pos + Inches(0.2), Inches(2.3), Inches(2.6), Inches(0.5),
                 title, font_size=15, bold=True, color=color)
    add_text_box(s, x_pos + Inches(0.2), Inches(2.9), Inches(2.6), Inches(1.3),
                 desc, font_size=13, color=MID_GRAY)
    x_pos += Inches(3.25)

# Outcomes
add_text_box(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
             "Possible outcomes:", font_size=16, bold=True, color=DARK_TEXT)

outcomes = [
    ("✅ Allow access", GREEN),
    ("🔐 Require MFA", ACCENT_BLUE),
    ("📱 Require compliant device", ACCENT_TEAL),
    ("🚫 Block access", RED_ACCENT),
]
x_pos = Inches(0.5)
for text, color in outcomes:
    obadge = add_rounded_rect(s, x_pos, Inches(5.3), Inches(2.9), Inches(0.6), NEAR_WHITE)
    add_shape_rect(s, x_pos, Inches(5.3), Inches(2.9), Inches(0.06), color)
    add_text_box(s, x_pos + Inches(0.1), Inches(5.35), Inches(2.7), Inches(0.5),
                 text, font_size=14, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    x_pos += Inches(3.15)

add_footer_bar(s, 22, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Scenario Rules Reminder
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Remind participants about shared-tenant rules for this module's lab.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🏢 Lab Scenario & Shared Tenant Rules", font_size=28, bold=True, color=DARK_BG)

rules = [
    "Scenario: Project Northwind — Fabrikam is the external partner",
    "Tenant-wide policy changes are TRAINER-ONLY in this module",
    "Participant hands-on work is scoped to participant-isolated practice sites",
    "Use your Participant ID (P01–P10) with NW-Pxx-... naming",
    "Trainer has pre-provisioned a test guest user for the exercise",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(3.5),
                 rules, font_size=17, color=DARK_TEXT, icon="📋")

warn3 = add_rounded_rect(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), Inches(4.85), Inches(11), Inches(0.6),
             "🔒 Shared tenant = be careful. Org-level sharing changes affect ALL participants.",
             font_size=14, bold=True, color=ACCENT_BLUE)

add_footer_bar(s, 23, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Lab 2 Preview
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Explain lab tasks. Participants will observe org-level, configure site-level, "
              "verify guest status, and review admin roles vs site permissions.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🔬 Lab 2 Preview — Configuring Secure Access", font_size=26, bold=True, color=DARK_BG)

lab_tasks = [
    ("Task 1", "Review org-level sharing settings", "Observe (Trainer-only changes) — document current baseline"),
    ("Task 2", "Configure site-level sharing", "Set sharing on your NW-Pxx-ProjectSite to 'Existing guests only'"),
    ("Task 3", "Verify guest account status", "Check the trainer-provisioned guest in Entra ID — Pending vs Accepted"),
    ("Task 4", "Review admin roles vs site permissions", "Compare your SP Admin role access to your site owner permissions"),
]

y_pos = Inches(1.2)
for task_id, task_title, task_desc in lab_tasks:
    card = add_rounded_rect(s, Inches(0.8), y_pos, Inches(11.5), Inches(0.95), NEAR_WHITE)
    tbadge = add_rounded_rect(s, Inches(1), y_pos + Inches(0.15), Inches(1.2), Inches(0.4), ACCENT_BLUE)
    add_text_box(s, Inches(1.05), y_pos + Inches(0.17), Inches(1.1), Inches(0.35),
                 task_id, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2.4), y_pos + Inches(0.1), Inches(9), Inches(0.35),
                 task_title, font_size=17, bold=True, color=DARK_TEXT)
    add_text_box(s, Inches(2.4), y_pos + Inches(0.48), Inches(9), Inches(0.35),
                 task_desc, font_size=13, color=MID_GRAY)
    y_pos += Inches(1.15)

val_box = add_rounded_rect(s, Inches(0.8), y_pos + Inches(0.2), Inches(11.5), Inches(0.7), LIGHT_BLUE)
add_text_box(s, Inches(1), y_pos + Inches(0.25), Inches(11), Inches(0.6),
             "📸 Capture: org sharing level, site sharing config, guest invite status, role vs permission comparison.",
             font_size=14, color=ACCENT_BLUE)

add_footer_bar(s, 24, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Summary
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Recap the key messages before knowledge check.")
add_solid_bg(s, WHITE)
add_top_bar(s)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "📝 Module 2 Summary", font_size=30, bold=True, color=DARK_BG)

summary = [
    "Identity (Entra ID) is the foundation — many 'SharePoint problems' are identity problems",
    "Admin roles (tenant) ≠ Site permissions (resource) — use least privilege",
    "External sharing has 3 layers: Entra → Org-level → Site-level (each constrains the next)",
    "OneDrive sharing ≤ SharePoint sharing — always",
    "Guest lifecycle: Invite → Redeem → Access → Review — check invitation status when troubleshooting",
    "Conditional Access adds context-aware enforcement (MFA, device, location, risk)",
]
add_bullet_frame(s, Inches(0.8), Inches(1.3), Inches(11), Inches(4.5),
                 summary, font_size=17, color=DARK_TEXT, icon="✅")

add_footer_bar(s, 25, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Knowledge Check
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide("Use as discussion slide. Encourage short answers. Focus on scope and troubleshooting logic.")
add_solid_bg(s, WHITE)
add_top_bar(s, ACCENT_PURPLE)

add_text_box(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "🧠 Knowledge Check", font_size=30, bold=True, color=DARK_BG)

questions = [
    ("Q1", "What is the difference between a SharePoint Administrator role and a site owner?",
     "Admin role → admin center/tenant; Site owner → permissions within a specific site"),
    ("Q2", "Why can a site-level sharing setting never be more permissive than org-level?",
     "Org-level is the tenant baseline/ceiling; sites inherit that maximum"),
    ("Q3", "Name two common causes of guest access failures.",
     "Org sharing disabled; invitation not redeemed; CA blocks; site sharing too restrictive"),
    ("Q4", "What is Conditional Access trying to accomplish?",
     "Enforce access requirements based on identity and context (device, location, risk)"),
]

y_pos = Inches(1.2)
for q_id, question, answer in questions:
    qbadge = add_rounded_rect(s, Inches(0.8), y_pos, Inches(0.6), Inches(0.5), ACCENT_PURPLE)
    add_text_box(s, Inches(0.82), y_pos + Inches(0.03), Inches(0.55), Inches(0.4),
                 q_id, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.6), y_pos, Inches(10.5), Inches(0.5),
                 question, font_size=16, color=DARK_TEXT)
    add_text_box(s, Inches(1.6), y_pos + Inches(0.5), Inches(10.5), Inches(0.4),
                 f"→ {answer}", font_size=13, color=RGBColor(0x99, 0x99, 0x99))
    y_pos += Inches(1.15)

add_text_box(s, Inches(0.8), y_pos + Inches(0.3), Inches(11.5), Inches(0.5),
             "💬 Discuss with your neighbour — then we'll share answers.", font_size=16,
             bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

add_footer_bar(s, 26, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
TOTAL_SLIDES = slide_counter[0]
output_path = os.path.join(os.path.dirname(__file__), "Module-02-Slides.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
